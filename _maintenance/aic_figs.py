"""AI-Infra-Compute 讲义信息图（5 张）。复用 kb_draw + agent_figs helper。
嫁接：p6 后(五层栈)、p15 后(CPU vs GPU)、p17 后(精度阶梯)、p18 后(Roofline)、p44 后(带宽阶梯)。
AIC 无页码，图系统页脚匹配。独立预览：python3 aic_figs.py"""
import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "AI-Infra-Compute 讲义 · 算力硬件"
GINK = RGBColor(0x27, 0x50, 0x0A)

# ---------- p6 五层栈 ----------
def draw_aic_five_stack(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 全景总览  FIVE-LAYER STACK")
    title(slide, "五层栈：AI 工厂的分层地图")
    label(slide, 0.55, 1.58, 12.2, "每往上一层，瓶颈就换一次——瓶颈在哪一层，钱就该花在哪一层", SUB, 12, PP_ALIGN.LEFT, True)
    layers = [("数据中心", "供电与冷却", "拼「电」", 5.7, RGBColor(0x5A,0x24,0x12), RGBColor(0xF3,0xE2,0xD8)),
              ("集群", "成千上万卡组网", "拼网络与容错", 5.05, ODARK, CORAL),
              ("机柜", "NVL72 这类「机柜级计算机」", "拼互联带宽", 4.4, ODARK, CORAL),
              ("服务器", "8 卡一机", "拼装配与散热", 3.75, DTEAL, CARD),
              ("芯片", "GPU / ASIC", "拼算力", 3.1, DTEAL, CARD)]
    # 从下（芯片）往上（数据中心）画，宽度递增制造「金字塔倒过来」的工厂感
    for i, (zh, en, bott, y, ec, fill) in enumerate(layers):
        w = 6.4 + i*0.9; x = (13.333 - w)/2
        band(slide, x, y, w, 0.6, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, y+0.09, w-0.5, 0.45); p = tf.paragraphs[0]
        sr(p.add_run(), zh + "　", BODY_F, 13, ec if fill != CORAL else ODARK, bold=True)
        sr(p.add_run(), en, BODY_F, 10, SUB)
        label(slide, x+w+0.1, y+0.14, 2.2, "→ " + bott, ec if fill != CORAL else ODARK, 10, PP_ALIGN.LEFT, False, True)
    label(slide, 1.0, 2.72, 5.0, "▲ 芯片 → 数据中心，逐层向上", SUB, 10, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 6.4, 12.03, 0.44, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 6.44, 11.5, 0.38); p = tf.paragraphs[0]
    sr(p.add_run(), "类比 · 发电厂：", BODY_F, 11, DTEAL, bold=True)
    sr(p.add_run(), "客户买的不是发电机，是稳定出来的电——算力也一样，客户要的是 token，不是卡。", BODY_F, 11, INK)
    footer(slide, MOD, "第 1 章 · 全景总览")

# ---------- p15 CPU vs GPU ----------
def draw_aic_cpu_vs_gpu(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · GPU 解剖  LATENCY vs THROUGHPUT")
    title(slide, "延迟机器 vs 吞吐机器：两种性格")
    label(slide, 0.55, 1.58, 12.2, "深度学习 90%+ 是矩阵乘法——天生就是「一万件一样的事」", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 3.0, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.52, 5.6, "CPU · 延迟机器 —— 几位教授", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    for r in range(2):
        for c in range(3):
            band(slide, 0.95 + c*0.75, 3.05 + r*0.55, 0.62, 0.44, fill=DTEAL, edge=NAVY)
    label(slide, 3.35, 3.25, 3.0, "几十个复杂核心", DTEAL, 10.5, PP_ALIGN.LEFT, True)
    tf = tb(slide, 0.9, 4.35, 5.5, 1.0)
    para(tf, "「一件事尽快做完」：分支预测、乱序执行都为单任务延迟服务。", BODY_F, 11, SUB, first=True)
    para(tf, "什么难题都能解，但一次只能解几道。", BODY_F, 11, INK, bold=True)
    band(slide, 6.75, 2.4, 6.0, 3.0, fill=CORAL, edge=ORANGE)
    label(slide, 7.0, 2.52, 5.6, "GPU · 吞吐机器 —— 一万名小学生", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    for r in range(4):
        for c in range(12):
            band(slide, 7.1 + c*0.36, 3.02 + r*0.24, 0.28, 0.17, fill=ORANGE, edge=ODARK)
    tf2 = tb(slide, 7.05, 4.35, 5.5, 1.0)
    para(tf2, "「一万件一样的事一起做」：牺牲单任务速度，换总吞吐。", BODY_F, 11, SUB, first=True)
    para(tf2, "只会加减乘，但一声令下一万道口算题同时开工。", BODY_F, 11, INK, bold=True)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.67, 11.5, 1.1); p = tf3.paragraphs[0]
    sr(p.add_run(), "案例：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "同一个 70B 矩阵乘，顶级 CPU 数百 GFLOPS 量级，一张 B200 的 Tensor Core 到 PFLOPS 量级——差三个数量级。", BODY_F, 11.5, INK)
    para(tf3, "记住分工而不是优劣：调度 / 预处理 / 业务逻辑归 CPU，矩阵乘法归 GPU——所以每台 GPU 服务器里还有两颗不小的 CPU。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 2 章 · GPU 解剖")

# ---------- p17 精度阶梯 ----------
def draw_aic_precision_ladder(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 精度阶梯  PRECISION LADDER")
    title(slide, "精度阶梯：FP32 → FP4，数字每砍一半，算力翻一倍")
    label(slide, 0.55, 1.58, 12.2, "位宽减半 → 同样硅片与带宽能处理两倍的数 → 标称算力翻倍（代价是数值精细度）", SUB, 12, PP_ALIGN.LEFT, True)
    tiers = [("FP32", "4 字节", "1×", 0.16, RGBColor(0x5A,0x24,0x12), RGBColor(0xF3,0xE2,0xD8)),
             ("FP16 / BF16", "2 字节", "2×", 0.33, ODARK, CORAL),
             ("FP8", "1 字节", "4×", 0.62, DTEAL, CARD),
             ("FP4 / NVFP4", "半字节", "8×", 1.0, GINK, GBG)]
    x0 = 0.7; y = 2.5; h = 0.72; gap = 0.13; barmax = 7.5
    for i, (name, byt, mult, frac, tc, fill) in enumerate(tiers):
        yy = y + i*(h+gap)
        band(slide, x0, yy, 2.7, h, fill=fill, edge=tc)
        tf = tb(slide, x0+0.2, yy+0.08, 2.4, h-0.12); p = tf.paragraphs[0]
        sr(p.add_run(), name + "　", BODY_F, 13, tc, bold=True)
        sr(p.add_run(), byt, BODY_F, 10, SUB)
        band(slide, x0+2.9, yy+0.13, barmax*frac, h-0.26, fill=tc, edge=tc)
        vt = tb(slide, x0+2.95+barmax*frac, yy+0.06, 1.6, h-0.08); q = vt.paragraphs[0]
        sr(q.add_run(), "算力 " + mult, BODY_F, 13, tc, bold=True)
    label(slide, 0.7, 2.2, 6.0, "数字越窄（字节越少）→ 每秒能算的数越多 →", SUB, 10, PP_ALIGN.LEFT, True)
    band(slide, 0.7, 5.85, 11.93, 0.97, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 5.96, 11.4, 0.82); p = tf.paragraphs[0]
    sr(p.add_run(), "对客提醒：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "对比两张卡的算力必须对齐精度——「B 卡是 A 卡 5 倍」常有一半来自 FP8→FP4 的口径切换，而不是硅片进步。", BODY_F, 11.5, INK)
    para(tf, "所以训练用混合精度（主权重保精度），推理敢一路砍到 FP8 / FP4。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 2 章 · GPU 解剖")

# ---------- p18 Roofline ----------
def draw_aic_roofline(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · Roofline  BOUND CHECK")
    title(slide, "屋顶线：一眼判断卡在算力还是带宽")
    label(slide, 0.55, 1.58, 12.2, "强度低顶到「斜屋顶」= 带宽受限；强度高顶到「平屋顶」= 算力受限", SUB, 12, PP_ALIGN.LEFT, True)
    ox, oy, pw, ph = 1.4, 5.4, 5.6, 2.7
    line(slide, ox, oy-ph, ox, oy, color=SUB, w=1.3); line(slide, ox, oy, ox+pw, oy, color=SUB, w=1.3)
    label(slide, 0.45, oy-ph-0.05, 1.0, "性能", SUB, 10, PP_ALIGN.LEFT, False, True)
    label(slide, ox+pw-3.0, oy+0.08, 3.1, "计算强度 FLOPs/Byte →", SUB, 9.5, PP_ALIGN.RIGHT, True)
    # 屋顶线：斜坡 + 平顶
    knee = (ox + 0.4*pw, oy - 0.82*ph)
    line(slide, ox, oy, knee[0], knee[1], color=TEAL, w=2.6)
    line(slide, knee[0], knee[1], ox+pw, knee[1], color=ORANGE, w=2.6)
    line(slide, knee[0], oy, knee[0], knee[1], color=LINE, w=1.0, dashed=True)
    label(slide, ox+0.15, oy-0.55*ph-0.35, 2.2, "斜屋顶=带宽", DTEAL, 9.5, PP_ALIGN.LEFT, False, True)
    label(slide, knee[0]+0.3, knee[1]-0.32, 2.6, "平屋顶=算力", ODARK, 9.5, PP_ALIGN.LEFT, False, True)
    label(slide, knee[0]-0.5, oy+0.04, 1.6, "平衡点 ~300", SUB, 9, PP_ALIGN.CENTER, True)
    # 两个负载点
    dot(slide, ox+0.16*pw, oy-0.33*ph, 0.09, DTEAL)
    label(slide, ox+0.02*pw, oy-0.33*ph-0.34, 2.6, "decode 逐 token", DTEAL, 9.5, PP_ALIGN.LEFT, False, True)
    label(slide, ox+0.02*pw, oy-0.33*ph-0.1, 2.6, "低强度→带宽受限", DTEAL, 8.5, PP_ALIGN.LEFT, True)
    dot(slide, ox+0.72*pw, oy-0.82*ph, 0.09, ODARK)
    label(slide, ox+0.55*pw, oy-0.82*ph+0.08, 2.8, "大矩阵乘 / prefill", ODARK, 9.5, PP_ALIGN.LEFT, False, True)
    label(slide, ox+0.55*pw, oy-0.82*ph+0.3, 2.8, "高强度→算力受限", ODARK, 8.5, PP_ALIGN.LEFT, True)
    band(slide, 7.4, 2.5, 5.25, 2.85, fill=CARD, edge=TEAL)
    tf = tb(slide, 7.65, 2.63, 4.75, 2.65)
    para(tf, "类比 · 自助餐厅", BODY_F, 12.5, DTEAL, bold=True, first=True)
    para(tf, "厨师做菜快不快是算力，传菜通道宽不宽是带宽。大锅菜（prefill / 训练）拼厨师；一人一碗面（decode）厨师闲着，全堵在传菜通道。", BODY_F, 11, SUB)
    para(tf, "H100 级：算力 ~1000 TFLOPS、带宽 ~3.35TB/s，平衡点 ~300 FLOPs/Byte；decode 每字节只个位数运算——所以永远带宽受限。", BODY_F, 11, SUB)
    band(slide, 0.6, 5.7, 12.03, 1.12, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.82, 11.5, 0.95); p = tf3.paragraphs[0]
    sr(p.add_run(), "对客一句话：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "客户报「GPU 利用率低」，先用 Roofline 问——你的负载是算力型还是带宽型？带宽型算力利用率低是物理规律，不是配置错。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 2 章 · GPU 解剖")

# ---------- p44 带宽阶梯 + scale-up/out ----------
def draw_aic_bandwidth_ladder(slide):
    bg(slide)
    eyebrow(slide, "第 5 章 · Scale-up 互联  BANDWIDTH LADDER")
    title(slide, "带宽阶梯：把最话痨的切法，塞进最快的域")
    label(slide, 0.55, 1.58, 12.2, "每跨一级，慢一个数量级——并行策略的本质就是「按话痨程度选域」", SUB, 12, PP_ALIGN.LEFT, True)
    tiers = [("卡内 HBM", "~8 TB/s", "最快", 1.0, DTEAL, PALE, "scale-up 域"),
             ("机柜内 NVLink", "~1.8 TB/s", "快一个量级差", 0.55, TEAL, CARD, "scale-up 域"),
             ("跨机柜网络", "~100 GB/s（800Gb/s）", "再慢一个量级", 0.18, ORANGE, CORAL, "scale-out 域")]
    x0 = 0.7; y = 2.5; h = 0.8; gap = 0.16; barmax = 7.0
    for i, (name, bw, note, frac, tc, fill, dom) in enumerate(tiers):
        yy = y + i*(h+gap)
        band(slide, x0, yy, 3.0, h, fill=fill, edge=tc)
        tf = tb(slide, x0+0.2, yy+0.1, 2.7, h-0.15); p = tf.paragraphs[0]
        sr(p.add_run(), name, BODY_F, 12.5, tc, bold=True)
        para(tf, bw, BODY_F, 10.5, SUB)
        band(slide, x0+3.2, yy+0.16, barmax*frac, h-0.32, fill=tc, edge=tc)
        label(slide, x0+3.3, yy+0.24, barmax*frac, note if frac > 0.3 else "", WHITE if frac > 0.5 else tc, 10, PP_ALIGN.LEFT, False, True)
        label(slide, x0+3.25+barmax*frac+0.1, yy+0.24, 2.4, dom, tc, 10.5, PP_ALIGN.LEFT, False, True)
    band(slide, 0.7, 5.28, 5.95, 1.15, fill=PALE, edge=TEAL)
    tf = tb(slide, 0.95, 5.4, 5.5, 1.0)
    para(tf, "scale-up（机柜内）", BODY_F, 11.5, DTEAL, bold=True, first=True)
    para(tf, "TP / EP 一步一同步、通信极密——两个人抬桌子必须肩并肩，只能塞进最快的域。", BODY_F, 10.5, SUB)
    band(slide, 6.78, 5.28, 5.85, 1.15, fill=PALE, edge=ORANGE)
    tf2 = tb(slide, 7.03, 5.4, 5.4, 1.0)
    para(tf2, "scale-out（跨机柜）", BODY_F, 11.5, ODARK, bold=True, first=True)
    para(tf2, "DP 数据并行、定期汇总——各小组分楼层干活，走普通网络就够。", BODY_F, 10.5, SUB)
    label(slide, 0.7, 6.55, 12.0, "小结：scale-up 域的大小（几张卡能当一张用）直接决定能高效跑多大的模型——这就是机柜级产品存在的意义。", DTEAL, 10.5, PP_ALIGN.LEFT, False, True)
    footer(slide, MOD, "第 5 章 · Scale-up 互联")

# ---------- p9 训练 vs 推理需求曲线 ----------
def draw_aic_train_vs_infer(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 两条需求曲线  TRAIN vs SERVE")
    title(slide, "训练与推理：两条完全不同的需求曲线")
    label(slide, 0.55, 1.58, 12.2, "先问客户：你的负载是训练、微调还是推理？三者对硬件 / 网络 / 采购几乎是三套方案", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 3.0, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.52, 5.6, "训练 —— 造船（看总价）", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 3.1, 5.3, 0.56, "一次性重投入", "", "几千~几万卡同步跑几周~几月", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, zh_size=12.5, desc_size=9)
    tf = tb(slide, 0.9, 3.85, 5.5, 1.4)
    para(tf, "· 要求「同时、同地、超高互联」", BODY_F, 11, SUB, first=True)
    para(tf, "· 中断代价大", BODY_F, 11, ODARK, bold=True)
    para(tf, "· 船坞里集中所有工人干半年", BODY_F, 11, SUB)
    band(slide, 6.75, 2.4, 6.0, 3.0, fill=CORAL, edge=ORANGE)
    label(slide, 7.0, 2.52, 5.6, "推理 —— 跑航线（看每海里成本）", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.1, 5.3, 0.56, "持续账单", "", "跟用户流量走、可弹性伸缩", fill=ORANGE, edge=ODARK, zh_c=OINK, sub_c=OINK, zh_size=12.5, desc_size=9)
    tf2 = tb(slide, 7.05, 3.85, 5.5, 1.4)
    para(tf2, "· 可分区域部署", BODY_F, 11, SUB, first=True)
    para(tf2, "· 追求每 token 成本最低", BODY_F, 11, ODARK, bold=True)
    para(tf2, "· 船一下水，烧油钱按天算、按客流调船期", BODY_F, 11, SUB)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=GBG, edge=GREEN)
    tf3 = tb(slide, 0.9, 5.67, 11.5, 1.1); p = tf3.paragraphs[0]
    sr(p.add_run(), "数字：", BODY_F, 11.5, GINK, bold=True)
    sr(p.add_run(), "行业预测 2030 年推理占 AI 总算力约 75%——算力市场的主战场，正从「谁能训」转向「谁跑得省」。", BODY_F, 11.5, INK)
    para(tf3, "造船看总价，跑船看每海里成本——负载定错，硬件 / 网络 / 采购模式全错。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 1 章 · 全景总览")

# ---------- p18 GPU 内部组织 ----------
def draw_aic_gpu_internals(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · GPU 内部组织  WHERE FLOPS COME FROM")
    title(slide, "算力从哪来：SM、CUDA Core 与 Tensor Core")
    label(slide, 0.55, 1.58, 12.2, "AI 算力标称值几乎全来自 Tensor Core 的矩阵乘加——看规格表先看它和显存带宽两行", SUB, 12, PP_ALIGN.LEFT, True)
    # GPU → SM（车间，放大看一个）→ 工位/机床
    band(slide, 0.6, 2.4, 7.6, 3.0, fill=RGBColor(0xF7,0xFB,0xFC), edge=LINE)
    node(slide, 0.9, 3.35, 1.85, 1.15, "一张 GPU", "", "= 100+ 个 SM", fill=DTEAL, edge=NAVY, zh_c=WHITE, sub_c=PALE, zh_size=14, desc_size=10)
    arrow(slide, 2.8, 3.92, 3.35, 3.92, color=TEAL, w=2.0); label(slide, 2.75, 3.5, 0.9, "放大", TEAL, 9.5, PP_ALIGN.CENTER, True)
    band(slide, 3.45, 2.72, 4.5, 2.55, fill=CARD, edge=TEAL)
    label(slide, 3.45, 2.82, 4.5, "一个 SM（车间）", DTEAL, 12, PP_ALIGN.CENTER, False, True)
    node(slide, 3.7, 3.3, 4.0, 0.78, "CUDA Core · 手工工位", "", "标量 / 向量 · 啥都能干、单件慢", fill=WHITE, edge=SUB, zh_c=INK, sub_c=SUB, zh_size=12.5, desc_size=10)
    node(slide, 3.7, 4.25, 4.0, 0.9, "Tensor Core · 矩阵机床", "", "只干矩阵乘 · 一次冲 16×16 · AI 算力全靠它", fill=ORANGE, edge=ODARK, zh_c=OINK, sub_c=OINK, zh_size=12.5, desc_size=10)
    band(slide, 8.35, 2.4, 4.28, 3.0, fill=CORAL, edge=ORANGE)
    tf = tb(slide, 8.6, 2.53, 3.8, 2.8)
    para(tf, "案例 · 小 batch 跑不满", BODY_F, 12.5, ODARK, bold=True, first=True)
    para(tf, "机床一次冲一整块料，你只送来一根料条（batch=1 的小矩阵）——机床空转大半。", BODY_F, 11, SUB)
    para(tf, "这就是推理侧要拼批处理的硬件根源（串 LLM-Inference 第 3 章）。", BODY_F, 11, DTEAL, bold=True)
    para(tf, "看规格表：Tensor Core 算力 + 显存带宽两行是主角；CUDA Core 数量只是配角。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 2 章 · GPU 解剖")

# ---------- p29 HBM 堆叠 ----------
def draw_aic_hbm(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 显存与 HBM  WHY STACK")
    title(slide, "HBM：把内存立起来、焊在 GPU 旁边")
    label(slide, 0.55, 1.58, 12.2, "距离近 + 接口宽（数千根数据线），带宽做到普通内存的几十倍——「买算力」很大程度是买显存带宽", SUB, 12, PP_ALIGN.LEFT, True)
    # GPU die + 堆叠 DRAM
    node(slide, 1.2, 3.35, 2.0, 1.5, "GPU", "Die", "计算单元", fill=DTEAL, edge=NAVY, zh_c=WHITE, sub_c=PALE, en_c=PALE, zh_size=16, en_size=11, desc_size=9)
    for i in range(8):
        yy = 4.7 - i*0.19
        band(slide, 3.4, yy, 2.2, 0.16, fill=(PALE if i % 2 else CARD), edge=TEAL)
    label(slide, 3.4, 2.98, 2.2, "DRAM ×8–16 层", DTEAL, 11, PP_ALIGN.CENTER, False, True)
    label(slide, 3.3, 4.94, 2.4, "垂直堆叠 · TSV 硅通孔打通", SUB, 11, PP_ALIGN.CENTER, True)
    line(slide, 3.2, 4.1, 3.4, 4.1, color=SUB, w=1.2)
    label(slide, 1.2, 5.05, 2.0, "同一封装基板", SUB, 11, PP_ALIGN.CENTER, True)
    # 带宽对比
    band(slide, 6.1, 3.2, 6.5, 0.72, fill=GBG, edge=GREEN)
    label(slide, 6.3, 3.35, 3.0, "HBM（B200 级）", GINK, 12, PP_ALIGN.LEFT, False, True)
    label(slide, 9.2, 3.34, 3.2, "≈ 8 TB/s", GINK, 13.5, PP_ALIGN.LEFT, False, True)
    band(slide, 6.1, 4.1, 2.6, 0.72, fill=CORAL, edge=ORANGE)
    label(slide, 6.3, 4.26, 2.4, "服务器 DDR 内存", ODARK, 11, PP_ALIGN.LEFT, False, True)
    label(slide, 8.8, 4.25, 3.6, "几百 GB/s（差几十倍）", ODARK, 12, PP_ALIGN.LEFT, False, True)
    label(slide, 6.1, 5.06, 6.5, "类比：平房改高层公寓——同一地皮住进十倍的人，电梯井（TSV）直通每层，工厂隔壁通勤短。", SUB, 11, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 5.6, 12.03, 1.22, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.72, 11.5, 1.05); p = tf.paragraphs[0]
    sr(p.add_run(), "案例：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "decode 是带宽受限，所以推理体验基本由 HBM 带宽决定——H200 比 H100 只是显存从 80GB/3.35TB/s 升到 141GB/4.8TB/s，长上下文吞吐就明显拉开。", BODY_F, 11.5, INK)
    para(tf, "HBM 占 AI 芯片物料成本的相当大头，且长期供不应求。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 3 章 · 显存与 HBM")

# ---------- p31 训练显存账 ----------
def draw_aic_train_memory(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 训练显存账  16 BYTES/PARAM")
    title(slide, "训练显存账：一个参数要花 16 个字节")
    label(slide, 0.55, 1.58, 12.2, "混合精度 + Adam 标准配方下，每个参数 16 字节——权重只占小头", SUB, 12, PP_ALIGN.LEFT, True)
    # 16 字节拆条
    X0 = 0.7; y = 2.6; BW = 11.9; h = 0.95
    segs = [("权重 BF16", 2, TEAL, WHITE), ("梯度", 2, ORANGE, OINK),
            ("优化器：主权重 FP32", 4, RGBColor(0x9F,0xB3,0xC0), WHITE), ("动量", 4, RGBColor(0x9F,0xB3,0xC0), WHITE), ("方差", 4, RGBColor(0x9F,0xB3,0xC0), WHITE)]
    x = X0
    for name, b, c, tc in segs:
        w = BW*b/16
        band(slide, x, y, w, h, fill=c, edge=(DTEAL if c == TEAL else ODARK if c == ORANGE else SUB))
        t = tb(slide, x+0.05, y+0.16, w-0.1, 0.65); p = t.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        sr(p.add_run(), name, BODY_F, 10.5 if w > 1.3 else 9, tc, bold=True)
        p2 = t.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        sr(p2.add_run(), f"{b} 字节", BODY_F, 9.5, tc)
        x += w
    label(slide, X0, y-0.3, 6.0, "权重 2 + 梯度 2 + 优化器状态 12 = 16 字节／参数（另加激活值）", SUB, 10.5, PP_ALIGN.LEFT, False, True)
    label(slide, X0+BW*4/16, y+h+0.06, BW*12/16, "↑ 优化器状态（主权重 4 + 动量 4 + 方差 4）才是大头", ODARK, 10, PP_ALIGN.CENTER, False, True)
    band(slide, 0.6, 4.35, 6.0, 1.85, fill=CORAL, edge=ORANGE)
    tf = tb(slide, 0.85, 4.47, 5.5, 1.65)
    para(tf, "案例 · 7B 全参训练", BODY_F, 12.5, ODARK, bold=True, first=True)
    para(tf, "7B × 16 字节 ≈ 112GB——单张 80GB 卡装不下。这就是「7B 也要多卡训练」、「LoRA 只训小矩阵就能单卡跑」的算术根源。", BODY_F, 11, SUB)
    band(slide, 6.75, 4.35, 5.88, 1.85, fill=CARD, edge=TEAL)
    tf2 = tb(slide, 7.0, 4.47, 5.4, 1.65)
    para(tf2, "类比 · 装修预算", BODY_F, 12.5, DTEAL, bold=True, first=True)
    para(tf2, "家具本身（权重）只占小头，施工队工具材料（梯度 + 优化器）占大头——只按家具报价必然爆预算。", BODY_F, 11, SUB)
    para(tf2, "这也是 ZeRO / FSDP「把账本切开分摊」存在的理由（LLM-Training 第 7 章）。", BODY_F, 11, DTEAL, bold=True)
    band(slide, 0.6, 6.3, 12.03, 0.52, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 6.36, 11.5, 0.42); p = tf3.paragraphs[0]
    sr(p.add_run(), "心算：", BODY_F, 11, DTEAL, bold=True)
    sr(p.add_run(), "全参训练显存 ≈ 参数量 × 16 字节 + 激活。", BODY_F, 11, INK)
    footer(slide, MOD, "第 3 章 · 显存与 HBM")

FIGS = [draw_aic_five_stack, draw_aic_cpu_vs_gpu, draw_aic_precision_ladder,
        draw_aic_roofline, draw_aic_bandwidth_ladder,
        draw_aic_train_vs_infer, draw_aic_gpu_internals, draw_aic_hbm, draw_aic_train_memory]
NEW = [draw_aic_train_vs_infer, draw_aic_gpu_internals, draw_aic_hbm, draw_aic_train_memory]

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in FIGS:
        fn(prs.slides.add_slide(B))
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_aic_preview.pptx")
    prs.save(out); print(f"saved {len(FIGS)} ->", out)
