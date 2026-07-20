"""AI-Ops 讲义信息图（5 张）。复用 kb_draw + agent_figs helper。
嫁接：p21 后(在线vs离线)、p22 后(四级漏斗)、p30 后(静默退化)、p38 后(发布流水线)、p48 后(事故响应)。
AI-Ops 无页码，图系统页脚匹配。独立预览：python3 ops_figs.py"""
import sys
sys.path.insert(0, "/Users/lijiaxiang/project/myAILearning/_maintenance")
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "运维观测讲义 · AI-Ops"
GINK = RGBColor(0x27, 0x50, 0x0A)
REDBG = RGBColor(0xFB, 0xE9, 0xE7)
REDINK = RGBColor(0x9B, 0x2C, 0x2C)

# ---------- p21 在线 vs 离线 ----------
def draw_ops_online_offline(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 在线 vs 离线  ACCEPTANCE vs PATROL")
    title(slide, "两种评估的分工：验收 vs 运营")
    label(slide, 0.55, 1.58, 12.2, "共享同一套判分器与指标——离线用它验收，在线用它巡逻；定义不一致，两边分数就没法互相解释", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 3.0, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.52, 5.6, "离线评估 · 像「高考」", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 3.12, 5.3, 0.56, "回答：能不能上线？", "", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, zh_size=12)
    tf = tb(slide, 0.9, 3.85, 5.5, 1.4)
    para(tf, "· 固定评估集、发布前跑", BODY_F, 11, SUB, first=True)
    para(tf, "· 一次性、覆盖面广、有标准答案", BODY_F, 11, INK, bold=True)
    para(tf, "· 方法论深潜见 Evaluation 第 5 章", BODY_F, 11, SUB)
    band(slide, 6.75, 2.4, 6.0, 3.0, fill=CORAL, edge=ORANGE)
    label(slide, 7.0, 2.52, 5.6, "在线评估 · 像「体检」（本章）", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.12, 5.3, 0.56, "回答：线上还好吗？", "", "", fill=ORANGE, edge=ODARK, zh_c=OINK, zh_size=12)
    tf2 = tb(slide, 7.05, 3.85, 5.5, 1.4)
    para(tf2, "· 真实流量采样、持续跑", BODY_F, 11, SUB, first=True)
    para(tf2, "· 周期性、抓趋势、早发现", BODY_F, 11, INK, bold=True)
    para(tf2, "· 评估指标在这里变成了监控指标", BODY_F, 11, SUB)
    label(slide, 5.35, 4.55, 2.65, "同一套判分器", DTEAL, 10.5, PP_ALIGN.CENTER, False, True)
    band(slide, 0.6, 5.5, 12.03, 1.32, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.62, 11.5, 1.15); p = tf3.paragraphs[0]
    sr(p.add_run(), "2026 年的共识句：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "「生产团队把模型当移动靶」——离线过了不等于线上一直好。", BODY_F, 11.5, INK)
    para(tf3, "在线评估是移动靶时代的标配：离线验收 + 在线巡逻，同一套指标定义两头对齐。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 3 章 · 在线评估")

# ---------- p22 四级评估漏斗 ----------
def draw_ops_eval_funnel(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 评估漏斗  FOUR LEVELS")
    title(slide, "四级漏斗：越贵的手段，处理越少的样本")
    label(slide, 0.55, 1.58, 12.2, "从毫秒级检查到数据集学习——inline 全量、judge 采样、人工只看队列", SUB, 12, PP_ALIGN.LEFT, True)
    levels = [("inline 检查", "格式 / 护栏 · 同步毫秒级 · 全量", "唯一同步层，挡住格式错误与护栏违规", TEAL, DTEAL, CARD, 7.0),
              ("judge 打分", "采样 5–10% · 异步 · LLM 判官", "质量分与失败模式标签，成本可控的主力", TEAL, DTEAL, CARD, 5.4),
              ("人工评审", "低分 / 争议进队列", "判官拿不准的 + 高价值样本抽检", ORANGE, ODARK, CORAL, 3.9),
              ("离线学习", "失败进数据集", "回流到评估集，让下一轮更强", ORANGE, ODARK, CORAL, 2.5)]
    y = 2.45; h = 0.72; gap = 0.16; cx = 4.2
    for i, (name, meta, note, tc, ec, fill, w) in enumerate(levels):
        yy = y + i*(h+gap); x = cx - w/2
        band(slide, x, yy, w, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, yy+0.08, w-0.4, h-0.12); p = tf.paragraphs[0]
        sr(p.add_run(), name + "　", BODY_F, 12.5, tc, bold=True)
        sr(p.add_run(), meta, BODY_F, 9.5, SUB)
        nt = tb(slide, cx + w/2 + 0.25, yy+0.12, 12.5 - (cx + w/2 + 0.3), h-0.1)
        sr(nt.paragraphs[0].add_run(), note, BODY_F, 10.5, INK)
        if i < 3:
            arrow(slide, cx, yy+h, cx, yy+h+gap, color=SUB, w=1.6)
    band(slide, 0.6, 5.9, 12.03, 0.92, fill=GBG, edge=GREEN)
    tf = tb(slide, 0.9, 6.0, 11.5, 0.78); p = tf.paragraphs[0]
    sr(p.add_run(), "成本控制：", BODY_F, 11.5, GINK, bold=True)
    sr(p.add_run(), "采样 5–10% 真实流量、异步执行（先回复用户后打分，零应用时延）、判官用便宜档 + 定期人工校准。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 3 章 · 在线评估")

# ---------- p30 静默退化 ----------
def draw_ops_silent_degradation(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · 静默退化  THE GREEN LIE")
    title(slide, "为什么全绿的仪表盘会骗人")
    label(slide, 0.55, 1.58, 12.2, "200 OK 不代表答案 OK——模型「答了」，只是答错了", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 2.95, fill=GBG, edge=GREEN)
    label(slide, 0.85, 2.52, 5.6, "传统仪表盘看到的 —— 一切「健康」", GINK, 12.5, PP_ALIGN.LEFT, False, True)
    for i, (t, c) in enumerate([("可用性 99.9%", GINK), ("延迟正常", GINK), ("错误率 0.1%", GINK), ("——因为模型「答了」，只是答错了", REDINK)]):
        tf = tb(slide, 0.95, 3.12 + i*0.5, 5.4, 0.45); q = tf.paragraphs[0]
        sr(q.add_run(), "· " + t, BODY_F, 11.5, c, bold=(i == 3))
    band(slide, 6.75, 2.4, 6.0, 2.95, fill=REDBG, edge=REDINK)
    label(slide, 7.0, 2.52, 5.6, "质量仪表盘看到的 —— 正在退化", REDINK, 12.5, PP_ALIGN.LEFT, False, True)
    # 缓降小折线
    ox, oy, pw, ph = 7.35, 4.15, 2.0, 0.85
    line(slide, ox, oy-ph, ox, oy, color=SUB, w=1.0); line(slide, ox, oy, ox+pw, oy, color=SUB, w=1.0)
    pts = [(0.0, 0.85), (0.4, 0.7), (0.7, 0.5), (1.0, 0.35)]
    P = [(ox+a*pw, oy-b*ph) for a, b in pts]
    for j in range(len(P)-1):
        line(slide, P[j][0], P[j][1], P[j+1][0], P[j+1][1], color=REDINK, w=2.2)
    label(slide, 7.0, 3.0, 2.3, "judge 分 4.2→3.6（三周）", REDINK, 9.5, PP_ALIGN.LEFT, False, True)
    tf2 = tb(slide, 9.6, 3.15, 3.05, 2.0)
    para(tf2, "· 「无法回答」类回复占比翻倍", BODY_F, 10.5, INK, first=True)
    para(tf2, "· 退款类意图失败模式聚集", BODY_F, 10.5, INK)
    para(tf2, "· 缓降以「周」为单位，客诉爆发时已煮熟", BODY_F, 10.5, SUB)
    band(slide, 0.6, 5.5, 12.03, 1.32, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.62, 11.5, 1.15); p = tf3.paragraphs[0]
    sr(p.add_run(), "类比 · 温水煮青蛙：", BODY_F, 11.5, ODARK, bold=True)
    sr(p.add_run(), "没有质量遥测的 AI 应用，退化是以周为单位缓慢发生的。", BODY_F, 11.5, INK)
    para(tf3, "破法只有一个：让「质量」成为和「可用性」并列的一等监控对象——这正是在线评估存在的理由。", BODY_F, 11, DTEAL, bold=True)
    footer(slide, MOD, "第 4 章 · 漂移监测")

# ---------- p38 五步发布流水线 ----------
def draw_ops_release_pipeline(slide):
    bg(slide)
    eyebrow(slide, "第 5 章 · 发布流水线  FIVE STEPS")
    title(slide, "从改动到全量：五步流水线")
    label(slide, 0.55, 1.58, 12.2, "没有在线评估的金丝雀是裸奔——LLM 金丝雀必须在金丝雀流量上跑自动评估", SUB, 12, PP_ALIGN.LEFT, True)
    steps = [("① 版本注册", "工件入册打 tag"), ("② 评估门禁", "离线集必须过线"), ("③ 金丝雀", "5–10% 流量"),
             ("④ 观测窗", "质量分/错误率/延迟"), ("⑤ 晋升或回滚", "改指针，秒级")]
    W = 2.28; gap = 0.14; x0 = 0.6; y = 2.55; h = 1.15; xs = []
    for i, (zh, d) in enumerate(steps):
        x = x0 + i*(W+gap); xs.append(x)
        solid = (i == 1)
        node(slide, x, y, W, h, zh, "", d, fill=(TEAL if solid else CARD), edge=(DTEAL if solid else TEAL),
             zh_c=(WHITE if solid else DTEAL), sub_c=(PALE if solid else SUB), zh_size=13, desc_size=9.5)
        if i < 4:
            arrow(slide, x+W+0.01, y+h/2, x+W+gap-0.01, y+h/2, color=TEAL, w=1.6)
    band(slide, 0.6, 4.05, 5.95, 1.3, fill=CARD, edge=TEAL)
    tf = tb(slide, 0.85, 4.17, 5.5, 1.1)
    para(tf, "评估门禁 · 不过线不出门", BODY_F, 12, DTEAL, bold=True, first=True)
    para(tf, "· 自动：回归集通过率不得低于基线", BODY_F, 10.5, SUB)
    para(tf, "· 人工：高风险变更（人设 / 护栏）加专家评审", BODY_F, 10.5, SUB)
    band(slide, 6.68, 4.05, 5.95, 1.3, fill=CORAL, edge=ORANGE)
    tf2 = tb(slide, 6.93, 4.17, 5.5, 1.1)
    para(tf2, "与传统金丝雀的关键差异", BODY_F, 12, ODARK, bold=True, first=True)
    para(tf2, "· 传统金丝雀看健康指标（错误率/延迟）就够", BODY_F, 10.5, SUB)
    para(tf2, "· LLM 金丝雀必须在流量上跑自动评估——坏掉的是答案质量，不是 HTTP 码", BODY_F, 10.5, SUB)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.7, 11.5, 1.05); p = tf3.paragraphs[0]
    sr(p.add_run(), "要点：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "这条流水线把「评估」和「基线」串起来了——改个提示词也要走它，因为一个词就能改变线上行为。", BODY_F, 12, INK)
    footer(slide, MOD, "第 5 章 · 发布管理")

# ---------- p48 事故响应 ----------
def draw_ops_incident(slide):
    bg(slide)
    eyebrow(slide, "第 6 章 · 信号与急停  BILL & KILL-SWITCH")
    title(slide, "第一信号常是账单，第一动作是急停")
    label(slide, 0.55, 1.58, 12.2, "客户高管最容易被打动的一页：「你们的 AI 有没有刹车？」", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 3.0, fill=REDBG, edge=REDINK)
    label(slide, 0.85, 2.52, 5.6, "成本尖峰 —— 最早的警报", REDINK, 12.5, PP_ALIGN.LEFT, False, True)
    # token 暴涨小尖峰
    ox, oy, pw, ph = 1.0, 3.75, 3.2, 0.75
    line(slide, ox, oy, ox+pw, oy, color=SUB, w=1.0)
    sp = [(0.0, 0.08), (0.35, 0.1), (0.55, 0.12), (0.68, 0.9), (0.8, 0.95), (1.0, 0.92)]
    P = [(ox+a*pw, oy-b*ph) for a, b in sp]
    for j in range(len(P)-1):
        line(slide, P[j][0], P[j][1], P[j+1][0], P[j+1][1], color=REDINK, w=2.2)
    label(slide, 3.2, 2.95, 3.2, "token 暴涨", REDINK, 10, PP_ALIGN.LEFT, False, True)
    tf = tb(slide, 0.9, 4.15, 5.5, 1.2)
    para(tf, "· 死循环 / 被注入的 agent 先表现为 token 暴涨", BODY_F, 11, SUB, first=True)
    para(tf, "· 有写权限的 agent：等人发现时可能已执行了动作", BODY_F, 11, SUB)
    para(tf, "· 成本告警要分钟级粒度，不是日账单", BODY_F, 11, REDINK, bold=True)
    band(slide, 6.75, 2.4, 6.0, 3.0, fill=CARD, edge=TEAL)
    label(slide, 7.0, 2.52, 5.6, "急停开关 kill switch —— 第一动作", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.1, 5.3, 0.62, "一键急停（IM 机器人指令）", "", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, zh_size=12.5)
    tf2 = tb(slide, 7.05, 3.9, 5.5, 1.3)
    para(tf2, "· 停任务不砸状态：暂停新任务、保留执行现场", BODY_F, 11, SUB, first=True)
    para(tf2, "· 恢复要走发布门禁，不是「再点一下」", BODY_F, 11, SUB)
    para(tf2, "· 案例：agent 递归循环，十分钟烧掉一天预算", BODY_F, 11, ODARK, bold=True)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.7, 11.5, 1.05); p = tf3.paragraphs[0]
    sr(p.add_run(), "一句话：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "急停开关 + 分钟级成本告警，就是 AI 的「刹车」的工程形态——成本告警常比任何质量告警都先响。", BODY_F, 12, INK)
    footer(slide, MOD, "第 6 章 · 事故响应")

# ---------- p54 HITL 三分级 ----------
def draw_ops_hitl(slide):
    bg(slide)
    eyebrow(slide, "第 6 章 · HITL 分级  HUMAN IN/ON/OUT")
    title(slide, "人机边界三分级：按风险配人")
    label(slide, 0.55, 1.58, 12.2, "分级不是一次定终身——新场景从 in-the-loop 起步，评估分稳定后逐级放权", SUB, 12, PP_ALIGN.LEFT, True)
    levels = [("Human-in-the-Loop", "人开车，AI 副驾", "每个动作人工确认后执行", "高危写操作：退款 / 删除 / 对外发送", REDINK, RGBColor(0xFB,0xE9,0xE7)),
              ("Human-on-the-Loop", "人监督，AI 有界行动", "低危动作自动执行，人看仪表盘可随时接管", "客服回复、内部查询类", ODARK, CORAL),
              ("Human-out-of-the-Loop", "人审计，AI 按策略行动", "事后抽检 + 完整 trace 留证", "低风险高频批量任务", GINK, GBG)]
    W = 3.95; gap = 0.14; x0 = 0.6; y = 2.4; h = 2.75
    for i, (en, zh, how, use, tc, fill) in enumerate(levels):
        x = x0 + i*(W+gap)
        band(slide, x, y, W, h, fill=fill, edge=tc)
        tf = tb(slide, x+0.25, y+0.16, W-0.5, 0.72); p = tf.paragraphs[0]
        sr(p.add_run(), en, BODY_F, 13, tc, bold=True)
        para(tf, zh, BODY_F, 11.5, INK, bold=True)
        node(slide, x+0.3, y+1.0, W-0.6, 0.62, how, "", "", fill=WHITE, edge=tc, zh_c=tc, zh_size=11)
        ub = tb(slide, x+0.28, y+1.78, W-0.56, 0.8); q = ub.paragraphs[0]
        sr(q.add_run(), "适用：", BODY_F, 10.5, tc, bold=True); sr(q.add_run(), use, BODY_F, 10.5, SUB)
    label(slide, 0.6, 5.28, 12.0, "风险越高 →→ 越往左，人越靠前；风险越低 →→ 越往右，越自动化", SUB, 10.5, PP_ALIGN.CENTER, True)
    band(slide, 0.6, 5.62, 12.03, 1.2, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.74, 11.5, 1.0); p = tf.paragraphs[0]
    sr(p.add_run(), "「放权路线图」本身就是方案书的一节：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "新场景从 in-the-loop 起步，等在线评估分稳定、坏例可控后，再逐级放权到 on / out。", BODY_F, 11.5, INK)
    para(tf, "对客一句话：不是「敢不敢让 AI 自动干」，而是「用什么条件、分几步把手放开」。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 6 章 · 事故响应")

FIGS = [draw_ops_online_offline, draw_ops_eval_funnel, draw_ops_silent_degradation,
        draw_ops_release_pipeline, draw_ops_incident, draw_ops_hitl]
NEW = [draw_ops_hitl]

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in FIGS:
        fn(prs.slides.add_slide(B))
    out = "/Users/lijiaxiang/project/myAILearning/_maintenance/_ops_preview.pptx"
    prs.save(out); print(f"saved {len(FIGS)} ->", out)
