"""AI-Infra-Platform 讲义信息图（5 张）。复用 kb_draw + agent_figs helper。
嫁接：p7 后(四大职责)、p25 后(gang)、p35 后(MIG)、p44 后(checkpoint)、p62 后(推理vs训练)。
AIP 无页码，图系统页脚匹配。独立预览：python3 aip_figs.py"""
import sys
sys.path.insert(0, "/Users/lijiaxiang/project/myAILearning/_maintenance")
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "AI-Infra-Platform 讲义 · 集群平台"
GINK = RGBColor(0x27, 0x50, 0x0A)
REDBG = RGBColor(0xFB, 0xE9, 0xE7)
REDINK = RGBColor(0x9B, 0x2C, 0x2C)

# ---------- p7 平台四大职责 ----------
def draw_aip_four_duties(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 平台全景  FOUR DUTIES")
    title(slide, "平台层的四大职责：调度、隔离、容错、观测")
    label(slide, 0.55, 1.58, 12.2, "客户的任何平台需求，先归类到这四项里的哪一项，方案就有了骨架", SUB, 12, PP_ALIGN.LEFT, True)
    duties = [("调度", "第 3 章", "哪个作业上哪些卡、谁排队谁抢占——最贵的卡永远有活干", TEAL, DTEAL, CARD),
              ("隔离", "第 4 章", "把一张卡安全切给多个团队 / 任务，互不干扰", TEAL, DTEAL, CARD),
              ("容错", "第 5 章", "万卡天天坏卡坏网——自动检测-驱逐-续训，别让一颗螺丝钉停掉整条流水线", ORANGE, ODARK, CORAL),
              ("观测", "第 6 章", "把利用率 / 健康度 / 成本变成看得见的指标，否则优化全是盲人摸象", ORANGE, ODARK, CORAL)]
    W = 5.95; H = 1.5; gap = 0.13; x0 = 0.6; y0 = 2.4
    for i, (zh, ch, d, tc, ec, fill) in enumerate(duties):
        x = x0 + (i % 2)*(W+gap); y = y0 + (i//2)*(H+gap)
        band(slide, x, y, W, H, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, y+0.16, W-0.5, 0.55); p = tf.paragraphs[0]
        sr(p.add_run(), f"{i+1}. {zh}　", BODY_F, 15, tc, bold=True)
        sr(p.add_run(), "（" + ch + "）", BODY_F, 10.5, SUB)
        dt = tb(slide, x+0.28, y+0.74, W-0.56, 0.7); sr(dt.paragraphs[0].add_run(), d, BODY_F, 11, INK)
    band(slide, 0.6, 5.68, 12.03, 1.14, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.8, 11.5, 0.98); p = tf.paragraphs[0]
    sr(p.add_run(), "类比 · 机场塔台：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "给航班排跑道（调度）、给航司分廊桥不打架（隔离）、天气突变改航保安全（容错）、雷达全程盯着（观测）。", BODY_F, 11.5, INK)
    para(tf, "塔台不开飞机，但没它机场就瘫痪——平台不算 token，但没它上万张卡就是一盘散沙。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 1 章 · 平台全景")

# ---------- p25 gang scheduling ----------
def draw_aip_gang(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 作业调度  GANG SCHEDULING")
    title(slide, "成组调度：半个作业等于零个作业")
    label(slide, 0.55, 1.58, 12.2, "要么一次给齐、要么一张都不给——这是 AI 调度和普通调度的分水岭", SUB, 12, PP_ALIGN.LEFT, True)
    # 左：全到位 → 开划
    band(slide, 0.6, 2.4, 6.0, 2.4, fill=GBG, edge=GREEN)
    label(slide, 0.85, 2.52, 5.6, "64 / 64 到位 —— 一起做集合通信，跑！", GINK, 12, PP_ALIGN.LEFT, False, True)
    for r in range(4):
        for c in range(16):
            band(slide, 0.95 + c*0.33, 3.05 + r*0.28, 0.26, 0.2, fill=GREEN, edge=GINK)
    label(slide, 0.95, 4.35, 5.5, "64 张卡同时上船 → 龙舟开划", GINK, 10.5, PP_ALIGN.LEFT, True)
    # 右：半个 → 死锁
    band(slide, 6.75, 2.4, 6.0, 2.4, fill=REDBG, edge=REDINK)
    label(slide, 7.0, 2.52, 5.6, "各拿一半 —— 谁都启动不了，集体罚站", REDINK, 12, PP_ALIGN.LEFT, False, True)
    for r in range(4):
        for c in range(16):
            filled = (c < 8)
            band(slide, 7.1 + c*0.33, 3.05 + r*0.28, 0.26, 0.2, fill=(REDINK if filled else WHITE), edge=REDINK)
    label(slide, 7.1, 4.35, 5.6, "作业 A、B 各拿 48/96 → 都不释放、都不启动", REDINK, 10.5, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 5.0, 12.03, 1.82, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.12, 11.5, 1.65); p = tf.paragraphs[0]
    sr(p.add_run(), "类比 · 划龙舟：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "20 个桨手必须同时上船才能开划；先上去 12 个不但划不动，还白占着船（GPU），让别的队也用不了。", BODY_F, 11.5, INK)
    para(tf, "普通 Web 是「来一个客人开一桌」，训练是「凑不齐一桌就别开」——缺一张卡，其余 63 张全在空等。", BODY_F, 11, SUB)
    para(tf, "判断一个调度器能不能跑训练，第一问就是「支不支持 gang scheduling」。", BODY_F, 11, DTEAL, bold=True)
    footer(slide, MOD, "第 3 章 · 作业调度")

# ---------- p35 MIG ----------
def draw_aip_mig(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · GPU 切分  MIG")
    title(slide, "MIG：硬件级切分，一张卡变七个小 GPU")
    label(slide, 0.55, 1.58, 12.2, "在硬件层把一张卡切成最多 7 个物理隔离的实例——真·硬隔离", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 0.7, 3.05, 2.3, 1.2, "整卡", "A100 / H100 / B 系", "一张物理 GPU", fill=DTEAL, edge=NAVY, zh_c=WHITE, sub_c=PALE, en_c=PALE, zh_size=15, en_size=9.5, desc_size=9)
    arrow(slide, 3.05, 3.65, 3.75, 3.65, color=TEAL, w=2.2)
    label(slide, 3.0, 3.2, 0.9, "切 7", TEAL, 10, PP_ALIGN.CENTER, False, True)
    for i in range(7):
        x = 3.85 + i*1.24
        band(slide, x, 3.05, 1.12, 1.2, fill=CARD, edge=TEAL)
        tf = tb(slide, x+0.05, 3.2, 1.02, 0.95); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        sr(p.add_run(), f"实例{i+1}", BODY_F, 10.5, DTEAL, bold=True)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        sr(p2.add_run(), "独享资源", BODY_F, 8.5, SUB)
    label(slide, 3.85, 4.35, 8.7, "每个实例独享显存 / 算力 / 缓存通路——一个崩了、打满了，不影响其它", SUB, 10.5, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 4.85, 6.0, 1.35, fill=CARD, edge=TEAL)
    tf = tb(slide, 0.85, 4.97, 5.5, 1.2)
    para(tf, "优点 · 隔离最强、性能可预测", BODY_F, 11.5, DTEAL, bold=True, first=True)
    para(tf, "多租户生产的首选——尤其对外部多租户或线上服务。", BODY_F, 11, SUB)
    para(tf, "类比：大开间隔成独立公寓，一户失火不殃及邻居。", BODY_F, 10.5, SUB)
    band(slide, 6.75, 4.85, 5.88, 1.35, fill=CORAL, edge=ORANGE)
    tf2 = tb(slide, 7.0, 4.97, 5.4, 1.2)
    para(tf2, "缺点 · 不够灵活", BODY_F, 11.5, ODARK, bold=True, first=True)
    para(tf2, "切法档位固定（1/7、2/7…）；改配置要重置 GPU；只有特定卡型支持。", BODY_F, 11, SUB)
    para(tf2, "墙一旦砌好、户型就固定了。", BODY_F, 10.5, SUB)
    band(slide, 0.6, 6.32, 12.03, 0.5, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 6.37, 11.5, 0.42); p = tf3.paragraphs[0]
    sr(p.add_run(), "选型：", BODY_F, 11, DTEAL, bold=True)
    sr(p.add_run(), "要「安全隔离 + 性能可保证」（对外多租户 / 线上服务）→ MIG 最稳；要极致灵活 → 时间片 / MPS / HAMi。", BODY_F, 11, INK)
    footer(slide, MOD, "第 4 章 · GPU 切分与多租户")

# ---------- p44 checkpoint 两难 ----------
def draw_aip_checkpoint(slide):
    bg(slide)
    eyebrow(slide, "第 5 章 · 训练容错  CHECKPOINT")
    title(slide, "Checkpoint：存太勤浪费，存太疏重算")
    label(slide, 0.55, 1.58, 12.2, "checkpoint 频率是「拿存储换算力」的保险费——异步和分级让这份保险变便宜", SUB, 12, PP_ALIGN.LEFT, True)
    # 两难天平
    band(slide, 0.6, 2.4, 6.0, 2.5, fill=CORAL, edge=ORANGE)
    label(slide, 0.85, 2.52, 5.6, "两难：一根轴上的两端", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 3.15, 2.45, 0.95, "存太勤", "", "暂停写 TB 级、费吞吐", fill=WHITE, edge=ODARK, zh_c=ODARK, sub_c=SUB, zh_size=13, desc_size=9.5)
    node(slide, 3.85, 3.15, 2.45, 0.95, "存太疏", "", "一崩重算好几小时", fill=WHITE, edge=ODARK, zh_c=ODARK, sub_c=SUB, zh_size=13, desc_size=9.5)
    label(slide, 0.95, 4.25, 5.5, "写论文存档：每句 Ctrl-S 太耽误 / 几小时不存怕断电全丢", SUB, 10, PP_ALIGN.LEFT, True)
    # 破局
    band(slide, 6.75, 2.4, 6.0, 2.5, fill=GBG, edge=GREEN)
    label(slide, 7.0, 2.52, 5.6, "破局：异步 + 分级", GINK, 12.5, PP_ALIGN.LEFT, False, True)
    tf = tb(slide, 7.05, 3.05, 5.6, 1.8)
    para(tf, "· 异步 checkpoint：先把状态快速拷到内存 / 本地盘（暂停秒级），再后台慢慢刷到远端——解耦「暂停时间」与「落盘时间」", BODY_F, 10.5, SUB, first=True)
    para(tf, "· 分级 checkpoint：高频存内存、低频存远端——兼顾恢复速度与持久性", BODY_F, 10.5, SUB)
    para(tf, "· 像「后台自动静默保存」，不打断你写作", BODY_F, 10.5, GINK, bold=True)
    band(slide, 0.6, 5.1, 12.03, 1.1, fill=PALE, edge=LINE)
    tf2 = tb(slide, 0.9, 5.22, 11.5, 0.95); p = tf2.paragraphs[0]
    sr(p.add_run(), "根因：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "大模型训练是同步的——所有卡每一步都要对齐，一张卡崩，整个作业回滚到上一个 checkpoint。规模 × 时间 = 故障必然。", BODY_F, 11.5, INK)
    para(tf2, "这条直接串到 Compute 模块第 7 章的存储洪峰设计（TB 级并发写）。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 5 章 · 训练容错工程")

# ---------- p62 推理 vs 训练调度 ----------
def draw_aip_infer_vs_train(slide):
    bg(slide)
    eyebrow(slide, "第 7 章 · 推理服务平台化  INFER vs TRAIN")
    title(slide, "推理 vs 训练调度：几乎相反的两套要求")
    label(slide, 0.55, 1.58, 12.2, "同一个平台上，两类作业的调度策略截然不同——一套集群两种节奏", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 3.0, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.52, 5.6, "训练调度 · 批处理心态 —— 包场办婚宴", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    tf = tb(slide, 0.9, 3.1, 5.5, 2.1)
    para(tf, "· gang：一次要一大把卡", BODY_F, 11.5, INK, bold=True, first=True)
    para(tf, "· 跑完就释放", BODY_F, 11, SUB)
    para(tf, "· 要吞吐", BODY_F, 11, SUB)
    para(tf, "· 能等排队", BODY_F, 11, SUB)
    para(tf, "一次占满整个宴会厅、办完撤场。", BODY_F, 10.5, DTEAL, italic=True)
    band(slide, 6.75, 2.4, 6.0, 3.0, fill=CORAL, edge=ORANGE)
    label(slide, 7.0, 2.52, 5.6, "推理调度 · 服务心态 —— 开快餐店", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    tf2 = tb(slide, 7.05, 3.1, 5.5, 2.1)
    para(tf2, "· 常驻不走", BODY_F, 11.5, INK, bold=True, first=True)
    para(tf2, "· 延迟敏感（用户等着首字）", BODY_F, 11, SUB)
    para(tf2, "· 按流量弹性伸缩", BODY_F, 11, SUB)
    para(tf2, "· 多副本高可用", BODY_F, 11, SUB)
    para(tf2, "全天营业、随到随吃、高峰加窗口、低谷收人力。", BODY_F, 10.5, ODARK, italic=True)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.67, 11.5, 1.1); p = tf3.paragraphs[0]
    sr(p.add_run(), "混部的常见做法：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "推理占基线卡、保 SLA；训练 / 批处理作为低优先级填空，低谷用空闲卡，流量一涨就被抢占让位。", BODY_F, 11.5, INK)
    para(tf3, "设计推理平台先记住：它是「服务」不是「作业」——高可用、弹性、延迟 SLA 第一位，吞吐和成本在此前提下优化。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 7 章 · 推理服务平台化")

# ---------- p28 碎片 ----------
def draw_aip_fragmentation(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 碎片  FRAGMENTATION")
    title(slide, "碎片：卡是整的，需求是零散的")
    label(slide, 0.55, 1.58, 12.2, "利用率不只看「空闲卡多少」，更看「空闲卡能不能拼成需要的形状」", SUB, 12, PP_ALIGN.LEFT, True)
    # 左：碎片
    band(slide, 0.6, 2.4, 6.0, 3.0, fill=REDBG, edge=REDINK)
    label(slide, 0.85, 2.5, 5.6, "碎片：总量够、拓扑不够", REDINK, 12, PP_ALIGN.LEFT, False, True)
    frag = [[1,0,1,1],[0,1,1,0],[1,1,0,1]]  # 1=占用 0=空闲，空闲零散
    for r in range(3):
        for c in range(4):
            occ = frag[r][c]
            band(slide, 1.0+c*0.62, 3.0+r*0.5, 0.52, 0.4, fill=(RGBColor(0xE0,0xC8,0xC4) if occ else WHITE), edge=REDINK)
    label(slide, 3.7, 3.05, 2.8, "空闲卡东一个西一个", REDINK, 11, PP_ALIGN.LEFT, True)
    node(slide, 3.7, 3.85, 2.8, 0.78, "要「同机 8 卡」的大作业", "", "凑不出 → 排不进去", fill=WHITE, edge=REDINK, zh_c=REDINK, sub_c=SUB, zh_size=11.5, desc_size=10.5)
    label(slide, 0.9, 4.78, 5.6, "类比：电影院还剩 20 个座，但都是单座——5 口想坐一起就是买不到", SUB, 10.5, PP_ALIGN.LEFT, True)
    # 右：装箱
    band(slide, 6.75, 2.4, 6.0, 3.0, fill=GBG, edge=GREEN)
    label(slide, 7.0, 2.5, 5.6, "装箱 bin-packing：拼出整块", GINK, 12, PP_ALIGN.LEFT, False, True)
    pack = [[1,1,1,1],[1,1,0,0],[0,0,0,0]]
    for r in range(3):
        for c in range(4):
            occ = pack[r][c]
            fill = (GBG if occ else GREEN) if r < 2 else GREEN
            band(slide, 7.15+c*0.62, 3.0+r*0.5, 0.52, 0.4, fill=(RGBColor(0xCF,0xE3,0xD5) if occ else GREEN), edge=GINK)
    label(slide, 9.85, 3.05, 2.8, "小作业塞满同节点", GINK, 11, PP_ALIGN.LEFT, True)
    node(slide, 9.85, 3.85, 2.8, 0.78, "给大作业留整块空间", "", "KAI / Volcano 拓扑感知 + 装箱", fill=WHITE, edge=GREEN, zh_c=GINK, sub_c=SUB, zh_size=11.5, desc_size=10)
    label(slide, 7.05, 4.78, 5.6, "不撒胡椒面——碎片治理直接等于利用率", SUB, 10.5, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.68, 11.5, 1.1); p = tf.paragraphs[0]
    sr(p.add_run(), "一句话：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "总空闲卡够、却凑不出一个作业要的「同机 8 卡 / 同域 64 卡」——大作业照样排不进去，利用率虚高实低。", BODY_F, 12, INK)
    para(tf, "碎片是贵卡空转的隐形元凶；调度器把「装箱 + 拓扑感知」做进核心，就是在把利用率捞回来。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 3 章 · 作业调度")

# ---------- p29 四调度器分层 ----------
def draw_aip_schedulers(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 调度器图鉴  LAYERED, NOT EITHER-OR")
    title(slide, "四个调度器的分工：不是二选一，是分层叠用")
    label(slide, 0.55, 1.58, 12.2, "选型别问「哪个最好」，问「准入配额用什么、成组放置用什么」——两个位置各选一个", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 5.0, 2.35, 3.35, 0.6, "作业进来", "", "", fill=NAVY, edge=NAVY, zh_c=WHITE, zh_size=13)
    arrow(slide, 6.65, 2.97, 6.65, 3.28, color=SUB, w=1.7)
    band(slide, 0.7, 3.3, 11.93, 1.0, fill=CARD, edge=TEAL)
    label(slide, 0.95, 3.4, 4.0, "① 准入 / 配额层", DTEAL, 12, PP_ALIGN.LEFT, False, True)
    node(slide, 4.5, 3.5, 3.6, 0.62, "Kueue", "K8s 社区标配", "管谁能放多少作业、排队顺序", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=PALE, zh_size=12.5, en_size=9, desc_size=8.5)
    label(slide, 8.4, 3.5, 4.0, "本身不做复杂放置", SUB, 10, PP_ALIGN.LEFT, True)
    arrow(slide, 6.65, 4.32, 6.65, 4.6, color=SUB, w=1.7)
    band(slide, 0.7, 4.62, 11.93, 1.15, fill=CORAL, edge=ORANGE)
    label(slide, 0.95, 4.72, 4.0, "② 成组放置层（gang + 拓扑 + 装箱）", ODARK, 12, PP_ALIGN.LEFT, False, True)
    node(slide, 3.9, 5.05, 2.7, 0.6, "Volcano", "老牌批调度", "gang/公平/优先级·稳", fill=WHITE, edge=ODARK, zh_c=ODARK, sub_c=SUB, en_c=ODARK, zh_size=11.5, en_size=8.5, desc_size=8)
    node(slide, 6.75, 5.05, 2.7, 0.6, "KAI Scheduler", "NVIDIA 2025", "拓扑+装箱+DRA·为 GPU 生", fill=ORANGE, edge=ODARK, zh_c=OINK, sub_c=OINK, en_c=OINK, zh_size=11.5, en_size=8.5, desc_size=8)
    node(slide, 9.6, 5.05, 2.7, 0.6, "Slurm", "HPC 老将", "大集群/托管仍是主力", fill=WHITE, edge=ODARK, zh_c=ODARK, sub_c=SUB, en_c=ODARK, zh_size=11.5, en_size=8.5, desc_size=8)
    band(slide, 0.7, 5.95, 11.93, 0.87, fill=GBG, edge=GREEN)
    tf = tb(slide, 1.0, 6.04, 11.4, 0.72); p = tf.paragraphs[0]
    sr(p.add_run(), "2026 共识：", BODY_F, 11.5, GINK, bold=True)
    sr(p.add_run(), "生产里常是 Kueue 管准入配额 + Volcano/KAI 管 gang 与放置，两层配合而非互斥；小集群（<16 卡）单用 Kueue 或 KAI 就够。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 3 章 · 作业调度")

# ---------- p46 规模×时间=故障必然 ----------
def draw_aip_failure(slide):
    bg(slide)
    eyebrow(slide, "第 5 章 · 故障是常态  SCALE × TIME")
    title(slide, "规模 × 时间 = 故障必然")
    label(slide, 0.55, 1.58, 12.2, "不是「会不会坏」，而是「坏了能不能几分钟自动恢复」——容错能力直接决定有效算力", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 3.0, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.52, 5.6, "为什么必然：同步训练的脆弱", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    tf = tb(slide, 0.9, 3.05, 5.5, 1.1)
    para(tf, "所有卡每一步都要对齐——任何一张卡 / 一条链路出问题，整个作业就卡住或崩溃。", BODY_F, 11, SUB, first=True)
    para(tf, "卡越多、跑越久，「至少坏一个」的概率就越接近 100%。", BODY_F, 11, INK, bold=True)
    tfa = tb(slide, 0.9, 4.2, 5.6, 0.95)
    para(tfa, "类比：一万人乐团连奏一个月——任一乐手打喷嚏、断根弦，整场就得停。", BODY_F, 11, SUB, first=True)
    # 右：Llama 3 数字
    band(slide, 6.75, 2.4, 6.0, 3.0, fill=RGBColor(0xFB,0xE9,0xE7), edge=REDINK)
    label(slide, 7.0, 2.52, 5.6, "Llama 3 的真实账本（已发表论文）", REDINK, 12, PP_ALIGN.LEFT, False, True)
    stats = [("16k", "张 H100"), ("54", "天连续跑"), ("466", "次作业中断"), ("78%", "是硬件相关")]
    for i, (n, d) in enumerate(stats):
        x = 7.05 + (i % 2)*2.9; y = 3.05 + (i//2)*0.98
        band(slide, x, y, 2.7, 0.85, fill=WHITE, edge=REDINK)
        t = tb(slide, x+0.1, y+0.08, 2.5, 0.72); p = t.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        sr(p.add_run(), n, BODY_F, 20, REDINK, bold=True)
        p2 = t.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        sr(p2.add_run(), d, BODY_F, 10.5, SUB)
    label(slide, 7.0, 5.04, 5.6, "= 平均每天坏好几次；万卡尺度的物理规律。", REDINK, 11, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=PALE, edge=LINE)
    tf2 = tb(slide, 0.9, 5.68, 11.5, 1.1); p = tf2.paragraphs[0]
    sr(p.add_run(), "给客户建立正确预期：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "大规模训练必须有「随时补位、从中断处接着奏」的机制——这就是为什么第 5 章要讲 checkpoint 与自动续训。", BODY_F, 11.5, INK)
    para(tf2, "容错不是加分项，是有效算力的乘数：坏了 30 分钟才恢复 vs 3 分钟自动恢复，一个月下来有效算力差一大截。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 5 章 · 训练容错工程")

FIGS = [draw_aip_four_duties, draw_aip_gang, draw_aip_mig, draw_aip_checkpoint, draw_aip_infer_vs_train,
        draw_aip_fragmentation, draw_aip_schedulers, draw_aip_failure]
NEW = [draw_aip_fragmentation, draw_aip_schedulers, draw_aip_failure]

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in FIGS:
        fn(prs.slides.add_slide(B))
    out = "/Users/lijiaxiang/project/myAILearning/_maintenance/_aip_preview.pptx"
    prs.save(out); print(f"saved {len(FIGS)} ->", out)
