"""通用 zip 级插页嫁接：把新绘制的信息图页嫁接进讲义，原有部件字节一律不动。
不用 python-pptx round-trip 整包重存（那会保留/引入病根）。app.xml 计数中英文版通吃。

  insert_figures(src, dst, specs)   specs = [(after_page_1based, draw_fn, title), ...]
"""
import zipfile, re, os
from pptx import Presentation
from pptx.util import Inches

SLIDE_CT = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NEWRELS = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\r\n'
           '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
           '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" '
           'Target="../slideLayouts/slideLayout1.xml"/></Relationships>')
_TMP = "/Users/lijiaxiang/project/myAILearning/_maintenance/_kbins_new.pptx"

def _gen_slide_xmls(draw_fns):
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in draw_fns:
        fn(prs.slides.add_slide(B))
    prs.save(_TMP)
    out = []
    with zipfile.ZipFile(_TMP) as z:
        ns = sorted([n for n in z.namelist() if re.match(r'ppt/slides/slide\d+\.xml$', n)],
                    key=lambda s: int(re.search(r'slide(\d+)', s).group(1)))
        for n in ns:
            out.append(z.read(n))
    os.remove(_TMP)
    return out

def insert_figures(src, dst, specs):
    new_xmls = _gen_slide_xmls([fn for _, fn, _ in specs])
    with zipfile.ZipFile(src) as z:
        names = z.namelist()
        parts = {n: z.read(n) for n in names}
    ct = parts['[Content_Types].xml'].decode('utf-8')
    pres = parts['ppt/presentation.xml'].decode('utf-8')
    prels = parts['ppt/_rels/presentation.xml.rels'].decode('utf-8')
    app = parts['docProps/app.xml'].decode('utf-8')

    slidenums = [int(re.search(r'slide(\d+)\.xml', n).group(1)) for n in names if re.match(r'ppt/slides/slide\d+\.xml$', n)]
    base_num = max(slidenums)
    base_rid = max(int(m) for m in re.findall(r'Id="rId(\d+)"', prels))
    base_sldid = max(int(m) for m in re.findall(r'<p:sldId id="(\d+)"', pres))
    entries = re.findall(r'<p:sldId id="\d+" r:id="rId\d+"\s*/>', pres)
    assert len(entries) == len(slidenums), f"sldId {len(entries)} != slides {len(slidenums)}"

    assigned = []
    for i, ((after, fn, title), xml) in enumerate(zip(specs, new_xmls)):
        num, rid, sid = base_num+1+i, base_rid+1+i, base_sldid+1+i
        assigned.append((after, num, rid, sid, title))
        parts[f'ppt/slides/slide{num}.xml'] = xml
        parts[f'ppt/slides/_rels/slide{num}.xml.rels'] = NEWRELS.encode('utf-8')
    for after, num, rid, sid, title in sorted(assigned, key=lambda x: -x[0]):
        entries.insert(after, f'<p:sldId id="{sid}" r:id="rId{rid}"/>')
    pres = re.sub(r'(<p:sldIdLst>).*?(</p:sldIdLst>)', lambda m: m.group(1)+''.join(entries)+m.group(2), pres, flags=re.S)
    prels = prels.replace('</Relationships>',
        ''.join(f'<Relationship Id="rId{rid}" Type="{NS_R}/slide" Target="slides/slide{num}.xml"/>' for _, num, rid, _, _ in assigned) + '</Relationships>')
    ct = ct.replace('</Types>',
        ''.join(f'<Override PartName="/ppt/slides/slide{num}.xml" ContentType="{SLIDE_CT}"/>' for _, num, _, _, _ in assigned) + '</Types>')

    esc = lambda t: t.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    npages = len(entries)
    old_slides = int(re.search(r'<Slides>(\d+)</Slides>', app).group(1))
    new_lps = ''.join(f'<vt:lpstr>{esc(t)}</vt:lpstr>' for _, _, _, _, t in assigned)
    app = re.sub(r'<Slides>\d+</Slides>', f'<Slides>{npages}</Slides>', app)
    hp = re.search(r'<HeadingPairs>.*?</HeadingPairs>', app, re.S).group(0)
    app = app.replace(hp, hp.replace(f'<vt:i4>{old_slides}</vt:i4>', f'<vt:i4>{npages}</vt:i4>'))
    cur_tp = int(re.search(r'<TitlesOfParts><vt:vector size="(\d+)"', app).group(1))
    app = re.sub(r'(<TitlesOfParts><vt:vector size=")\d+(")', rf'\g<1>{cur_tp + len(assigned)}\g<2>', app)
    app = app.replace('</vt:vector></TitlesOfParts>', new_lps + '</vt:vector></TitlesOfParts>')

    parts['[Content_Types].xml'] = ct.encode('utf-8')
    parts['ppt/presentation.xml'] = pres.encode('utf-8')
    parts['ppt/_rels/presentation.xml.rels'] = prels.encode('utf-8')
    parts['docProps/app.xml'] = app.encode('utf-8')

    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(dst, 'w', zipfile.ZIP_DEFLATED) as zout:
        written = set()
        for item in zin.infolist():
            zout.writestr(item, parts.get(item.filename, zin.read(item.filename))); written.add(item.filename)
        for pn, data in parts.items():
            if pn not in written:
                zout.writestr(pn, data)
    return npages
