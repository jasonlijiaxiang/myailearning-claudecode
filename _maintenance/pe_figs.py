"""Prompt Engineering 讲义信息图（第一批 4 张）。复用 kb_draw + agent_figs helper。"""
import sys
sys.path.insert(0, "/Users/lijiaxiang/project/myAILearning/_maintenance")
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, BODY_F, TITLE_F)
from agent_figs import line

MOD = "Prompt Engineering · 提示词工程"

# ---------- p14 三种角色 ----------
def draw_pe_roles(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 消息结构  ROLES")
    title(slide, "三种角色：system 定规矩、user 提请求、assistant 是回复")
    label(slide, 0.55, 1.58, 12.2, "和大模型对话不是一段纯文本，而是带「角色」的消息序列——分清角色是写好提示词的第一步", SUB, 12, PP_ALIGN.LEFT, True)
    roles = [(0.9, "system", "系统", "人设、规则、边界", "整段对话生效 · 剧本设定", TEAL, DTEAL, CARD),
             (5.02, "user", "用户", "本轮请求与数据", "这场戏的台词", ORANGE, ODARK, CORAL),
             (9.14, "assistant", "助手", "模型的回复", "也是历史上下文 · 演员表演", TEAL, DTEAL, CARD)]
    for x, en, zh, what, ana, edge, tc, fill in roles:
        band(slide, x, 2.5, 3.98, 2.65, fill=fill, edge=edge)
        tf = tb(slide, x+0.3, 2.72, 3.5, 2.3)
        para(tf, en, BODY_F, 15, tc, bold=True, first=True)
        para(tf, zh, BODY_F, 11, edge)
        para(tf, what, BODY_F, 12, INK)
        para(tf, ana, BODY_F, 10.5, SUB, italic=True)
    arrow(slide, 4.9, 3.82, 5.0, 3.82, color=SUB, w=1.5)
    arrow(slide, 9.02, 3.82, 9.12, 3.82, color=SUB, w=1.5)
    band(slide, 0.7, 5.45, 11.93, 1.4, fill=CARD, edge=LINE)
    tf2 = tb(slide, 1.0, 5.58, 11.4, 1.25)
    para(tf2, "像拍一部戏：system 是剧本设定与导演须知(角色是谁、什么调性、什么不能碰)，user 是台词，assistant 是表演——设定错了，台词再好也走样。", BODY_F, 11.5, DTEAL, bold=True, first=True)
    para(tf2, "案例：客服机器人 system 写「你是 XX 银行客服，只答业务、不谈投资、语气礼貌」——人设放 system，一次设定、通篇约束。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 2 章 · 消息结构")

# ---------- p27 思维链 CoT ----------
def draw_pe_cot(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 思维链  CHAIN-OF-THOUGHT")
    title(slide, "思维链：让模型先想再答，给它一张「草稿纸」")
    label(slide, 0.55, 1.58, 12.2, "引导模型把中间推理一步步写出来再给答案——一句「让我们一步步思考」就能触发", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.9, 2.4, 5.4, 2.55, fill=CARD, edge=SUB)
    label(slide, 1.15, 2.55, 5.0, "直接问 · 张口就来", SUB, 12.5, PP_ALIGN.LEFT, False, True)
    tf = tb(slide, 1.2, 3.02, 4.9, 1.85)
    para(tf, "「原价 120，先打八折再减 20，多少？」", BODY_F, 11.5, INK, first=True)
    para(tf, "→ 一步到位，容易跳步算错", BODY_F, 12, ODARK, bold=True)
    para(tf, "复杂多步题，直接答不稳。", BODY_F, 10.5, SUB, italic=True)
    band(slide, 7.0, 2.4, 5.4, 2.55, fill=CORAL, edge=ORANGE)
    label(slide, 7.25, 2.55, 5.0, "加「一步步算」 · 先想再答", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    tf2 = tb(slide, 7.3, 3.02, 4.9, 1.85)
    para(tf2, "① 120 × 0.8 = 96", BODY_F, 12.5, INK, first=True)
    para(tf2, "② 96 − 20 = 76", BODY_F, 12.5, INK)
    para(tf2, "→ 76（写出过程，跳步和低级错大幅减少）", BODY_F, 11, DTEAL, bold=True)
    arrow(slide, 6.35, 3.68, 6.98, 3.68, color=TEAL, w=2.5)
    band(slide, 0.7, 5.25, 11.93, 1.6, fill=CARD, edge=LINE)
    tf3 = tb(slide, 1.0, 5.4, 11.4, 1.45)
    para(tf3, "像考试让学生「写出解题过程」而不是只填答案：过程写出来，答案更可靠——对多步计算、逻辑推理提升最明显。", BODY_F, 11.5, DTEAL, bold=True, first=True)
    para(tf3, "注：第 4 章会讲——推理模型(o 系、Claude 思考模式)已把 CoT 内建，反而不需要你再手写「一步步思考」。", BODY_F, 11, SUB, italic=True)
    footer(slide, MOD, "第 3 章 · 思维链")

# ---------- p36 自洽性 ----------
def draw_pe_consistency(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · 自洽性  SELF-CONSISTENCY")
    title(slide, "自洽性：同一题问几遍，多数表决压掉随机错误")
    label(slide, 0.55, 1.58, 12.2, "对同一提示词多次采样(带随机性生成多条思维链)，取「出现最多」的答案", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 0.7, 3.35, 1.9, 1.05, "同一个问题", "多次采样", "", fill=CARD, edge=SUB, zh_c=INK, en_c=SUB, zh_size=13, en_size=9.5, desc_size=9)
    chains = [(2.65, "推理链 A", "→ 76", True), (3.72, "推理链 B", "→ 76", True), (4.79, "推理链 C", "→ 74", False)]
    for cy, zh, ans, ok in chains:
        node(slide, 3.4, cy, 2.6, 0.82, zh, "思维链", ans, fill=CARD, edge=TEAL, zh_c=DTEAL, en_c=TEAL, sub_c=(DTEAL if ok else SUB), zh_size=12, en_size=9, desc_size=10)
        line(slide, 2.62, 3.87, 3.38, cy+0.41, color=TEAL, w=1.2)
        line(slide, 6.02, cy+0.41, 7.18, 3.87, color=(TEAL if ok else LINE), w=1.3)
    node(slide, 7.2, 3.35, 2.5, 1.05, "投票", "多数表决", "76 ×2 胜 74 ×1", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=WHITE, zh_size=14, en_size=9.5, desc_size=10)
    node(slide, 10.3, 3.35, 2.3, 1.05, "最终答案", "76", "更可信", fill=ORANGE, edge=ODARK, zh_c=OINK, sub_c=OINK, en_c=OINK, zh_size=14, en_size=13, desc_size=10)
    arrow(slide, 9.72, 3.87, 10.28, 3.87, color=ORANGE, w=1.75)
    band(slide, 0.7, 5.7, 11.93, 1.15, fill=CARD, edge=LINE)
    tf = tb(slide, 1.0, 5.84, 11.4, 1.0)
    para(tf, "像重要决策请几位专家各自独立判断，再看多数意见——单个可能失手，多数一致的结论更可信。", BODY_F, 11.5, DTEAL, bold=True, first=True)
    para(tf, "代价是多花几倍 token / 时间。只在「答错代价高」的高价值推理题用，别对所有请求都开。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 4 章 · 自洽性")

# ---------- p38 提示词链 ----------
def draw_pe_chaining(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · 提示词链  PROMPT CHAINING")
    title(slide, "提示词链：拆成几步，每步只干一件清楚的事")
    label(slide, 0.55, 1.58, 12.2, "与其一条巨型提示词做完所有事，不如拆成有序多次调用，上一步输出喂给下一步", SUB, 12, PP_ALIGN.LEFT, True)
    steps = [("① 读长文档", "抽取要点", "专注抽要点"),
             ("② 写摘要", "凝练", "专注凝练"),
             ("③ 写邮件", "措辞", "专注措辞")]
    W = 3.3; y = 2.95; h = 1.4; xs = [0.9, 4.9, 8.9]
    for k, (zh, en, d) in enumerate(steps):
        node(slide, xs[k], y, W, h, zh, en, d + " · 可单独调试 / 替换", zh_size=15, en_size=10.5, desc_size=10)
    for k in range(2):
        arrow(slide, xs[k]+W+0.02, y+h/2, xs[k+1]-0.02, y+h/2, w=2.0)
        label(slide, xs[k]+W-0.2, y+h/2-0.42, 1.4, "输出 → 输入", SUB, 9.5, PP_ALIGN.CENTER, True)
    band(slide, 0.7, 5.25, 11.93, 1.6, fill=CARD, edge=LINE)
    tf = tb(slide, 1.0, 5.4, 11.4, 1.45)
    para(tf, "像流水线：不指望一个工位造完整台车，拆成「冲压→焊接→喷漆→总装」，每站只干一件事，出问题也知道卡在哪站。", BODY_F, 11.5, DTEAL, bold=True, first=True)
    para(tf, "好处：每步更可控、可单独调试、可插入检查。坏处：多次调用、延迟和成本上升。复杂业务流几乎都值得拆链。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 4 章 · 提示词链")

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in [draw_pe_roles, draw_pe_cot, draw_pe_consistency, draw_pe_chaining]:
        fn(prs.slides.add_slide(B))
    out = "/Users/lijiaxiang/project/myAILearning/_maintenance/_pe_preview.pptx"
    prs.save(out); print("saved 4 ->", out)


# ========== 第二轮补图 ==========
from kb_draw import GBG as _GBG
from pptx.dml.color import RGBColor as _RGB
_GINK = _RGB(0x27, 0x50, 0x0A)

def draw_pe_cheapest_lever(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 能力杠杆  THE CHEAPEST LEVER")
    title(slide, "三种手段，成本差着数量级")
    label(slide, 0.55, 1.58, 12.2, "提示词工程几乎零成本、分钟级见效——任何优化都应该「先把提示词试到头」，再考虑更贵的手段", SUB, 12, PP_ALIGN.LEFT, True)
    tiers = [("提示词工程", "分钟级 · 几乎零成本", "只改「输入的话」，不动权重、不搭系统，改完立刻生效", "×1", 0.12, TEAL, DTEAL, CARD),
             ("RAG / 微调", "天级~周级 · 几万~几十万", "搭外挂知识库、或用自己数据继续训练——要工程、要数据", "×100", 0.45, ORANGE, ODARK, CORAL),
             ("从头预训练", "月级 · 数百万+", "重训基座——极少数场景才值得", "×10000", 1.0, _RGB(0x99,0x3C,0x1D), _RGB(0x5A,0x24,0x12), _RGB(0xF3,0xE2,0xD8))]
    x0 = 0.7; y = 2.55; h = 0.95; gap = 0.14
    for i, (name, when, what, mult, frac, bc, tc, fill) in enumerate(tiers):
        yy = y + i*(h+gap)
        band(slide, x0, yy, 3.15, h, fill=fill, edge=tc)
        tf = tb(slide, x0+0.2, yy+0.11, 2.85, 0.75); p = tf.paragraphs[0]
        sr(p.add_run(), name, BODY_F, 13.5, tc, bold=True)
        para(tf, when, BODY_F, 10, SUB)
        wt = tb(slide, x0+3.35, yy+0.13, 6.0, 0.72); sr(wt.paragraphs[0].add_run(), what, BODY_F, 11, INK)
        band(slide, 9.7, yy+0.2, 2.4*frac, h-0.4, fill=bc, edge=bc)
        vt = tb(slide, 9.7+2.4*frac+0.08, yy+0.24, 1.2, 0.5); sr(vt.paragraphs[0].add_run(), mult, BODY_F, 12, bc, bold=True)
    label(slide, 9.7, y-0.28, 3.0, "相对成本 →", SUB, 10, PP_ALIGN.LEFT, True)
    label(slide, x0, y-0.28, 6.0, "每往下一档，贵一个数量级", SUB, 10, PP_ALIGN.LEFT, True)
    band(slide, 0.7, 5.82, 11.93, 1.0, fill=_GBG, edge=GREEN)
    tf = tb(slide, 1.0, 5.92, 11.4, 0.85); p = tf.paragraphs[0]
    sr(p.add_run(), "售前记住这句：", BODY_F, 12, _GINK, bold=True)
    sr(p.add_run(), "约 80% 的「模型效果不好」其实是提示词没写好，不是模型不行——先怀疑提示词。", BODY_F, 12, INK)
    para(tf, "案例：客户一上来就想「微调一个自己的模型」，多数情况下先把提示词 / few-shot / 输出格式调好就够用，省下几十万训练成本和几周工期。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 1 章 · 什么是提示词工程")

def draw_pe_zero_few_shot(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 零样本 vs 少样本  ZERO vs FEW-SHOT")
    title(slide, "零样本 vs 少样本：给不给例子")
    label(slide, 0.55, 1.58, 12.2, "最基础、最常用的一对技巧——对「格式固定、风格特定、有细微规则」的任务，few-shot 往往一给就准", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 2.95, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.52, 5.6, "零样本 Zero-shot —— 直接口头派活", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 3.1, 5.3, 0.56, "不给例子，直接下指令", "", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, zh_size=12)
    tf = tb(slide, 0.9, 3.85, 5.5, 1.4)
    para(tf, "· 简单、通用任务首选", BODY_F, 11.5, INK, bold=True, first=True)
    para(tf, "· 提示词短、上手快", BODY_F, 11, SUB)
    para(tf, "· 风险：口径不明时，中性易被判成负面", BODY_F, 11, ODARK)
    band(slide, 6.75, 2.4, 6.0, 2.95, fill=CORAL, edge=ORANGE)
    label(slide, 7.0, 2.52, 5.6, "少样本 Few-shot —— 给两份样板照做", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.1, 5.3, 0.56, "放几个「输入 → 输出」范例", "", "", fill=ORANGE, edge=ODARK, zh_c=OINK, zh_size=12)
    tf2 = tb(slide, 7.05, 3.85, 5.5, 1.4)
    para(tf2, "· 数量通常 2–5 个，不是越多越好", BODY_F, 11.5, INK, bold=True, first=True)
    para(tf2, "· 格式固定 / 风格特定 / 有细微规则 → 一给就准", BODY_F, 11, SUB)
    para(tf2, "· 例子要覆盖边界（含一个中性样本）", BODY_F, 11, ODARK, bold=True)
    band(slide, 0.6, 5.5, 12.03, 1.32, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.62, 11.5, 1.15); p = tf3.paragraphs[0]
    sr(p.add_run(), "类比：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "zero-shot 像直接口头派活；few-shot 像「给你两份做好的样板，照这个样子做」——对着样板做，比只听描述更不容易跑偏。", BODY_F, 11.5, INK)
    para(tf3, "案例：情感分类 zero-shot 有时把中性判成负面；给 3 个标注好的例子（含一个中性）后，模型立刻对齐你的口径。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 3 章 · 核心技巧")
