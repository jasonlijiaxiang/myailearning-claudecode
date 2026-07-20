# Agent 讲义 · 第 10 章「多智能体 / Sub-agent 编排」新页绘制(12 页,大画布 13.333×7.5)
# 复用 kb_draw 设计系统 helper;页面模式对齐本册既有章(章扉/学习目标/动手做/对练/小结按
# slide 94/95/100/101/102 逐坐标仿制,深度页用 kb_draw 惯用版式)。
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from kb_draw import (bg, tb, sr, para, eyebrow, title, node, arrow, label, band,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL,
                     LINE, WHITE, NEARW, PALE, GREEN, GBG, TITLE_F, BODY_F)
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

MOD = "Agent 讲义 · 智能体与编排"
CH = "第 10 章 · Sub-agent 编排"
LIGHTCH = RGBColor(0xD6, 0xE1, 0xE8)
EAF = RGBColor(0xEA, 0xF3, 0xF5)
BRIGHT = RGBColor(0x1C, 0xA3, 0xB8)
DBOX = RGBColor(0x0C, 0x3A, 0x4C)
MONO = "Courier New"

def dbg(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    s.fill.solid(); s.fill.fore_color.rgb = NAVY
    s.line.color.rgb = NAVY; s.line.width = Pt(0.75); s.shadow.inherit = False
    return s

def deco_circle(slide, x=10.20, y=4.20, d=5.0):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = DTEAL
    s.line.fill.background(); s.shadow.inherit = False

def foot(slide, right):
    tf = tb(slide, 0.70, 7.08, 6.0, 0.30)
    sr(tf.paragraphs[0].add_run(), MOD, BODY_F, 9, RGBColor(0x9F, 0xB3, 0xC0))
    tf2 = tb(slide, 8.63, 7.08, 4.0, 0.30)
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    sr(p.add_run(), right, BODY_F, 9, RGBColor(0x9F, 0xB3, 0xC0))

def src_note(slide, text, y=6.78):
    tf = tb(slide, 0.55, y, 12.2, 0.26)
    sr(tf.paragraphs[0].add_run(), text, BODY_F, 9.5, SUB, italic=True)

# ---------- p1 章扉(深色) ----------
def p01_cover(slide):
    dbg(slide); deco_circle(slide)
    tf = tb(slide, 0.70, 1.60, 4.0, 0.50)
    sr(tf.paragraphs[0].add_run(), "第 10 章", BODY_F, 16, ORANGE, bold=True)
    tf = tb(slide, 0.70, 2.15, 11.0, 1.10)
    sr(tf.paragraphs[0].add_run(), "多智能体与 Sub-agent 编排", TITLE_F, 42, WHITE, bold=True)
    tf = tb(slide, 0.70, 3.35, 11.0, 0.50)
    sr(tf.paragraphs[0].add_run(), "Sub-agents · 一个 agent 雇佣另一个 agent", BODY_F, 20, LIGHTCH)
    tf = tb(slide, 0.70, 4.20, 8.8, 2.40)
    lines = [
        "派遣是运行时的能力：模型只是发起一个工具调用，新循环由 harness 启动",
        "三个性质定一切：独立上下文 / 窄带通信 / 深度上限",
        "一本账定去留：15× token 换 90.2% 提升——要结果外包，要过程留下",
    ]
    for i, t in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        sr(p.add_run(), t, BODY_F, 15, EAF)
    foot(slide, "第 10 章 · Sub-agent 编排")

# ---------- p2 学习目标 ----------
def p02_goals(slide):
    bg(slide)
    tf = tb(slide, 0.70, 0.55, 11.0, 0.35)
    sr(tf.paragraphs[0].add_run(), "第 10 章 · Sub-agent 编排 · 学习目标", BODY_F, 12, TEAL, bold=True)
    tf = tb(slide, 0.70, 0.90, 11.93, 0.95)
    sr(tf.paragraphs[0].add_run(), "学完这一章，你能对客户讲清楚", TITLE_F, 30, NAVY, bold=True)
    goals = [
        "用「模型 / 运行时 / 产品」三层说清 sub-agent 是谁的能力——客户问的是平台能力，不是模型能力",
        "讲清三个关键性质：独立上下文（任务书自包含）、窄带通信（只回结论）、深度上限（Claude Code 5 层 / Codex 默认 1 层）",
        "亲手配出一个自定义 sub-agent（markdown / TOML 各一），并用四层触发机制答好「怎么确保被调用」",
        "算清决策账：多 agent ≈ 15× token 换 90.2% 提升——一句「要结果还是要过程」当场定去留",
    ]
    for k, g in enumerate(goals):
        c = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.70), Inches(2.19 + 1.25 * k), Inches(0.50), Inches(0.50))
        c.fill.solid(); c.fill.fore_color.rgb = BRIGHT
        c.line.fill.background(); c.shadow.inherit = False
        ctf = c.text_frame; ctf.word_wrap = True
        pp = ctf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        sr(pp.add_run(), str(k + 1), BODY_F, 11, WHITE, bold=True)
        tf = tb(slide, 1.50, 2.15 + 1.25 * k, 11.03, 1.15)
        sr(tf.paragraphs[0].add_run(), g, BODY_F, 15.5, INK)
    foot(slide, "Sub-agent 编排")

# ---------- p3 深度① 关系 ----------
def p03_relation(slide):
    bg(slide)
    eyebrow(slide, CH)
    title(slide, "Sub-agent 就是 Agent：递归一层，不是新物种")
    rows = [
        ("概念", TEAL, DTEAL, CARD,
         "Sub-agent 是一个普通 agent，只是雇主换成了另一个 agent。主 agent 的工具箱里有一个「派遣工具」，调用它＝现场启动一个全新的 agent 循环——独立上下文、自己的工具、跑完交回一段结论。层级是运行时的，机制是同一套，递归了一层。"),
        ("类比", TEAL, DTEAL, CARD,
         "项目负责人与临时外包：外包没参加过你们之前的会（任务书必须自包含）；干完交一份报告就走人（过程不占你的脑子）；活能拆就同时雇几个（并行）。"),
        ("三层框架", ORANGE, ODARK, CORAL,
         "模型层只会输出「我要调这个工具」；运行时层（harness）负责真的把新循环跑起来——派遣能力住在这里；产品层把它包装成配置、路由与面板。普通 chatbot 没有 sub-agent，是 harness 没实现，不是模型不会。"),
        ("小结", GREEN, GREEN, GBG,
         "给客户一句话：sub-agent 不是更聪明的模型，是更聪明的组织方式——把一个人的长对话，改组成一支队伍的分工。"),
    ]
    y = 1.78
    heights = [1.30, 1.06, 1.30, 0.90]
    for (tag, edge, tagc, fillc, text), h in zip(rows, heights):
        band(slide, 0.55, y, 12.23, h, fill=fillc, edge=fillc)
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(y + 0.18), Inches(1.30), Inches(0.44))
        try: chip.adjustments[0] = 0.5
        except Exception: pass
        chip.fill.solid(); chip.fill.fore_color.rgb = WHITE
        chip.line.color.rgb = edge; chip.line.width = Pt(1.25); chip.shadow.inherit = False
        ctf = chip.text_frame; ctf.word_wrap = True
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ctf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        sr(pp.add_run(), tag, BODY_F, 12.5, tagc, bold=True)
        tf = tb(slide, 2.30, y + 0.10, 10.30, h - 0.20)
        sr(tf.paragraphs[0].add_run(), text, BODY_F, 12.5, INK)
        y += h + 0.12
    foot(slide, "Sub-agent 编排")

# ---------- p4 深度① 图示:扇出/扇入 ----------
def p04_fanout(slide):
    bg(slide)
    eyebrow(slide, CH + "  FAN-OUT / FAN-IN")
    title(slide, "扇出与扇入（Fan-out / Fan-in）：一份任务书进去，一段结论回来", size=26)
    node(slide, 5.17, 1.55, 3.0, 0.90, "主 Agent", "Main Agent", "上下文 = 与你的整段对话",
         fill=DTEAL, edge=NAVY, zh_c=WHITE, sub_c=PALE, en_c=WHITE, en_size=11, desc_size=11)
    subs = [("Sub-agent A", "搜代码库", 1.60), ("Sub-agent B", "查官方文档", 5.17), ("Sub-agent C", "跑安全扫描", 8.74)]
    for name, job, x in subs:
        node(slide, x, 3.05, 3.0, 1.02, name, job, "独立上下文 · 自己的工具循环", en_size=11, desc_size=11)
    arrow(slide, 6.67, 2.45, 3.10, 3.03)
    arrow(slide, 6.67, 2.45, 6.67, 3.03)
    arrow(slide, 6.67, 2.45, 10.24, 3.03)
    node(slide, 4.17, 4.75, 5.0, 0.90, "主对话只收三段结论", "Fan-in",
         "过程废料随 sub-agent 一起销毁", fill=GBG, edge=GREEN, zh_c=GREEN, sub_c=SUB, en_c=GREEN, en_size=11, desc_size=11)
    arrow(slide, 3.10, 4.07, 5.30, 4.73, color=GREEN)
    arrow(slide, 6.67, 4.07, 6.67, 4.73, color=GREEN)
    arrow(slide, 10.24, 4.07, 8.04, 4.73, color=GREEN)
    label(slide, 1.55, 2.42, 2.6, "任务书 · 自包含 prompt", TEAL, 11, PP_ALIGN.LEFT, True)
    label(slide, 9.95, 4.42, 2.9, "窄带通信：只交回一段报告", GREEN, 11, PP_ALIGN.LEFT, True)
    band(slide, 0.55, 6.02, 12.23, 0.60, fill=PALE, edge=PALE)
    tf = tb(slide, 0.85, 6.12, 11.7, 0.42)
    sr(tf.paragraphs[0].add_run(),
       "实测：复杂查询典型并行 3–5 个 sub-agent，并行把研究时长最多压缩 90%——放大的账单见本章「派与不派」页。",
       BODY_F, 12.5, INK, bold=True)
    src_note(slide, "来源：Anthropic Engineering《How we built our multi-agent research system》（2025-06；访问 2026-07-17）")
    foot(slide, "Sub-agent 编排")

# ---------- p5 深度② 三个关键性质 ----------
def p05_props(slide):
    bg(slide)
    eyebrow(slide, CH)
    title(slide, "三个关键性质：一切用法与局限的根源")
    cards = [
        ("独立上下文", "Own Context",
         ["看不到主对话历史，只看到主 agent 写给它的那份「任务书」——连你们刚达成的共识它也不知道。",
          "纪律：任务书必须自包含。目标、范围、约束、产出格式，一样都不能省。"]),
        ("窄带通信", "Narrow Return",
         ["只交回一段最终文字；读过的一百个文件、试错的十轮，随它的上下文一起销毁。",
          "主对话干净了；代价是过程不可见、中途难纠偏——multi-agent 难调试，难就难在交接处。"]),
        ("深度上限", "Depth Limit",
         ["可以嵌套再派，但有硬顶：Claude Code 5 层封顶（固定不可配）；Codex max_depth 默认 1。",
          "为什么要顶：没有上限的扇出（fan-out），就是没有上限的账单。"]),
    ]
    for i, (zh, en, paras) in enumerate(cards):
        x = 0.55 + i * 4.115
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.80), Inches(3.90), Inches(3.75))
        try: card.adjustments[0] = 0.04
        except Exception: pass
        card.fill.solid(); card.fill.fore_color.rgb = CARD
        card.line.color.rgb = TEAL; card.line.width = Pt(1.25); card.shadow.inherit = False
        tf = card.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18); tf.margin_top = Inches(0.20)
        p = tf.paragraphs[0]
        sr(p.add_run(), f"{i+1} · {zh}", BODY_F, 15, DTEAL, bold=True)
        pe = tf.add_paragraph(); pe.space_after = Pt(8)
        sr(pe.add_run(), en, BODY_F, 11, TEAL)
        for t in paras:
            pp = tf.add_paragraph(); pp.space_after = Pt(7)
            sr(pp.add_run(), t, BODY_F, 12, INK)
    band(slide, 0.55, 5.80, 12.23, 0.62, fill=GBG, edge=GBG)
    tf = tb(slide, 0.85, 5.92, 11.7, 0.40)
    sr(tf.paragraphs[0].add_run(),
       "三个性质一句话：「隔离创造价值，隔离也制造成本」——什么时候划算，本章最后用一本账算清。",
       BODY_F, 12.5, GREEN, bold=True)
    src_note(slide, "来源：code.claude.com/docs（sub-agents）；developers.openai.com/codex（subagents）（访问 2026-07-17）")
    foot(slide, "Sub-agent 编排")

# ---------- p6 深度③ Claude Code 实操 ----------
def p06_cc(slide):
    bg(slide)
    eyebrow(slide, CH)
    title(slide, "Claude Code 实操：一份 markdown，一个分身")
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.78), Inches(6.35), Inches(4.35))
    try: card.adjustments[0] = 0.03
    except Exception: pass
    card.fill.solid(); card.fill.fore_color.rgb = CARD
    card.line.color.rgb = LINE; card.line.width = Pt(1.0); card.shadow.inherit = False
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.20); tf.margin_right = Inches(0.15); tf.margin_top = Inches(0.14)
    p = tf.paragraphs[0]; p.space_after = Pt(6)
    sr(p.add_run(), ".claude/agents/code-reviewer.md（项目级；用户级放 ~/.claude/agents/）", BODY_F, 12, DTEAL, bold=True)
    code = [
        ("---", INK),
        ("name: code-reviewer", INK),
        ("description: 代码安全审查。涉及鉴权、加密、", INK),
        ("  外部输入的改动，主动用我复核", INK),
        ("tools: Read, Grep, Glob    # 只读白名单", DTEAL),
        ("model: haiku               # 可换更便宜的模型", DTEAL),
        ("---", INK),
        ("你是一名安全审查员，只找漏洞，不改代码……", INK),
        ("（正文＝这个 sub-agent 的系统提示词）", SUB),
    ]
    for t, c in code:
        pp = tf.add_paragraph(); pp.space_after = Pt(3)
        sr(pp.add_run(), t, MONO, 11.5, c)
    pp = tf.add_paragraph(); pp.space_before = Pt(8)
    sr(pp.add_run(), "description 是写给主 agent 看的路由依据——想被自动调用，先把「什么时候用我」写具体。",
       BODY_F, 11.5, ODARK, bold=True)
    band(slide, 7.10, 1.78, 5.68, 1.30, fill=PALE, edge=PALE)
    tf = tb(slide, 7.35, 1.90, 5.2, 1.10)
    p = tf.paragraphs[0]; p.space_after = Pt(4)
    sr(p.add_run(), "内置三件套", BODY_F, 12.5, DTEAL, bold=True)
    pp = tf.add_paragraph()
    sr(pp.add_run(), "Explore 只读探路 / Plan 规划期调研 / general-purpose 通吃；v2.1.198 起默认后台运行，面板可见并行树。",
       BODY_F, 11.5, INK)
    band(slide, 7.10, 3.24, 5.68, 2.89, fill=CARD, edge=TEAL)
    tf = tb(slide, 7.35, 3.36, 5.2, 2.70)
    p = tf.paragraphs[0]; p.space_after = Pt(5)
    sr(p.add_run(), "触发的四层机制（从软到硬）", BODY_F, 12.5, DTEAL, bold=True)
    for t, hard in [
        ("① description 自动路由——写得越具体命中越高", False),
        ("② 对话点名——「用 code-reviewer 审一下」，基本必中", False),
        ("③ CLAUDE.md 常备指令——放倾向，别放硬流程", False),
        ("④ skill 绑 context: fork——用 /名字 显式触发，必进 sub-agent，硬保证", True),
    ]:
        pp = tf.add_paragraph(); pp.space_after = Pt(5)
        sr(pp.add_run(), t, BODY_F, 11.5, ODARK if hard else INK, bold=hard)
    src_note(slide, "来源：code.claude.com/docs（sub-agents、skills；访问 2026-07-17）", y=6.40)
    foot(slide, "Sub-agent 编排")

# ---------- p7 深度④ Codex 对照 ----------
def p07_codex(slide):
    bg(slide)
    eyebrow(slide, CH)
    title(slide, "Codex 对照：同一件事的另一种拼法")
    cols = [(0.55, 2.00), (2.65, 5.05), (7.80, 4.98)]
    hdr = ["维度", "Claude Code", "Codex"]
    y0 = 1.72
    for (x, w), t in zip(cols, hdr):
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y0), Inches(w), Inches(0.42))
        b.fill.solid(); b.fill.fore_color.rgb = DTEAL
        b.line.color.rgb = DTEAL; b.shadow.inherit = False
        tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        sr(pp.add_run(), t, BODY_F, 12.5, WHITE, bold=True)
    rows = [
        ("定义文件", ".claude/agents/*.md（Markdown + frontmatter）", ".codex/agents/*.toml（TOML）"),
        ("必填字段", "name、description", "name、description、developer_instructions"),
        ("内置角色", "Explore / Plan / general-purpose", "explorer / worker / default"),
        ("并发与运行", "默认后台运行（v2.1.198 起）", "max_threads 默认 6 个并发线程"),
        ("嵌套深度", "5 层封顶，固定不可配", "max_depth 默认 1"),
        ("批量特技", "agent teams：多会话协作（独立特性）", "spawn_agents_on_csv：CSV 一行一工人（实验）"),
    ]
    y = 2.18
    for k, (lab, cc, cx) in enumerate(rows):
        if k % 2 == 0:
            band(slide, 0.55, y - 0.03, 12.23, 0.66, fill=CARD, edge=CARD)
        tf = tb(slide, 0.66, y + 0.06, 1.85, 0.52)
        sr(tf.paragraphs[0].add_run(), lab, BODY_F, 12, DTEAL, bold=True)
        tf = tb(slide, 2.75, y + 0.06, 4.85, 0.56)
        sr(tf.paragraphs[0].add_run(), cc, BODY_F, 11.5, INK)
        tf = tb(slide, 7.90, y + 0.06, 4.80, 0.56)
        sr(tf.paragraphs[0].add_run(), cx, BODY_F, 11.5, INK)
        y += 0.66
    band(slide, 0.55, 6.22, 12.23, 0.52, fill=GBG, edge=GBG)
    tf = tb(slide, 0.85, 6.30, 11.7, 0.38)
    sr(tf.paragraphs[0].add_run(),
       "概念完全同构：配置文件＋内置角色＋并发 / 深度限额——学会一家，另一家只是方言。两家迭代都快，引用前先核日期。",
       BODY_F, 12, GREEN, bold=True)
    src_note(slide, "来源：developers.openai.com/codex/subagents；code.claude.com/docs/sub-agents（访问 2026-07-17）", y=6.80)
    foot(slide, "Sub-agent 编排")

# ---------- p8 深度⑤ 决策账 ----------
def p08_ledger(slide):
    bg(slide)
    eyebrow(slide, CH)
    title(slide, "派与不派：先背口诀，再算账")
    b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.72), Inches(12.23), Inches(0.78))
    try: b.adjustments[0] = 0.10
    except Exception: pass
    b.fill.solid(); b.fill.fore_color.rgb = NAVY
    b.line.color.rgb = NAVY; b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    sr(pp.add_run(), "口诀：「我要的是结果，还是过程？」——只要结果，外包给 sub-agent；过程影响后续决策，留在主对话。",
       BODY_F, 15.5, WHITE, bold=True)
    costs = [
        ("冷启动", "不带主对话记忆，进场先重新摸索一遍"),
        ("窄带通信", "任务书写不清就跑偏，中途难纠偏"),
        ("token 放大", "队伍越大账单越陡——见下面的账"),
    ]
    for i, (zh, d) in enumerate(costs):
        x = 0.55 + i * 4.115
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.68), Inches(3.90), Inches(0.92))
        try: c.adjustments[0] = 0.08
        except Exception: pass
        c.fill.solid(); c.fill.fore_color.rgb = CORAL
        c.line.color.rgb = ORANGE; c.line.width = Pt(1.25); c.shadow.inherit = False
        tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.12)
        pp = tf.paragraphs[0]
        sr(pp.add_run(), "成本 · " + zh, BODY_F, 12.5, ODARK, bold=True)
        p2 = tf.add_paragraph()
        sr(p2.add_run(), d, BODY_F, 11.5, INK)
    band(slide, 0.55, 3.82, 12.23, 1.66, fill=PALE, edge=PALE)
    tf = tb(slide, 0.85, 3.94, 11.7, 1.46)
    p = tf.paragraphs[0]; p.space_after = Pt(5)
    sr(p.add_run(), "手算一本账（基线：一次 10K token 的 chat 任务）", BODY_F, 13, DTEAL, bold=True)
    for t in [
        "单 agent 循环 ≈ 4× → 40K；multi-agent ≈ 15× → 150K——放大一个数量级。",
        "买到什么：内部评测比单 agent Claude Opus 4 高 90.2%；token 用量单因素解释 80% 的性能差异。",
        "结论：性能就是用 token 买的——任务价值配得上，这本账才划算。",
    ]:
        pp = tf.add_paragraph(); pp.space_after = Pt(4)
        sr(pp.add_run(), t, BODY_F, 12.5, INK)
    band(slide, 0.55, 5.64, 12.23, 1.00, fill=CARD, edge=CARD)
    tf = tb(slide, 0.85, 5.74, 11.7, 0.84)
    for i, t in enumerate([
        "不派：两三步的活 / 步骤强顺序依赖 / 依赖主对话已积累的理解。",
        "派 1 个：过程废料多、结论短（全库搜索、调研）。 派 N 个：能切成互不依赖的块（并行扇出）。",
    ]):
        pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        pp.space_after = Pt(3)
        sr(pp.add_run(), t, BODY_F, 12, INK, bold=(i == 0))
    src_note(slide, "来源：Anthropic Engineering（2025-06）：multi-agent 研究系统内部评测口径（BrowseComp）")
    foot(slide, "Sub-agent 编排")

# ---------- p9 上云怎么落地 ----------
def p09_cloud(slide):
    bg(slide)
    eyebrow(slide, CH + " · 上云怎么落地")
    title(slide, "上云怎么落地：放大后的账单谁来接")
    band(slide, 0.55, 1.78, 12.23, 0.70, fill=CORAL, edge=CORAL)
    tf = tb(slide, 0.85, 1.90, 11.7, 0.50)
    sr(tf.paragraphs[0].add_run(),
       "multi-agent 把 token 消耗放大一个数量级、并发 sub-agent 同时打满 API——限流、预算、观测从「可选」变「必须」。",
       BODY_F, 12.5, ODARK, bold=True)
    rows = [
        ("并发与限流", "sub-agent 并发上限（max_threads / 后台任务数）要对齐模型网关的并发配额与排队策略——AI Gateway 模块的老朋友"),
        ("预算护栏", "按任务价值分层设 token 上限与超额熔断；主流模型平台（方舟 / 百炼 / Bedrock 等）都有配额与费用告警"),
        ("全链路观测", "每个 sub-agent 的输入输出与交接都留 trace（LangSmith / Langfuse 类）——multi-agent 的错误常藏在交接里（第 6 章）"),
    ]
    y = 2.68
    for lab, d in rows:
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(y), Inches(2.55), Inches(0.78))
        try: b.adjustments[0] = 0.10
        except Exception: pass
        b.fill.solid(); b.fill.fore_color.rgb = PALE
        b.line.color.rgb = TEAL; b.line.width = Pt(1.0); b.shadow.inherit = False
        tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
        pp = tf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        sr(pp.add_run(), lab, BODY_F, 13, DTEAL, bold=True)
        tf = tb(slide, 3.35, y + 0.08, 9.4, 0.66)
        sr(tf.paragraphs[0].add_run(), d, BODY_F, 12, INK)
        y += 0.94
    band(slide, 0.55, 5.62, 12.23, 0.66, fill=GBG, edge=GBG)
    tf = tb(slide, 0.85, 5.74, 11.7, 0.44)
    sr(tf.paragraphs[0].add_run(),
       "给客户的顺序：先上 trace 和预算，再开 multi-agent——顺序反了，第一张账单就是事故复盘会。",
       BODY_F, 12.5, GREEN, bold=True)
    foot(slide, "Sub-agent 编排")

# ---------- p10 动手做(深色) ----------
def p10_handson(slide):
    dbg(slide)
    tf = tb(slide, 0.70, 0.60, 9.0, 0.40)
    sr(tf.paragraphs[0].add_run(), "动手做 · HANDS-ON", BODY_F, 12, ORANGE, bold=True)
    tf = tb(slide, 0.70, 1.00, 11.0, 0.90)
    sr(tf.paragraphs[0].add_run(), "第 10 章 · Sub-agent 编排", TITLE_F, 28, WHITE, bold=True)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.70), Inches(2.10), Inches(11.93), Inches(3.00))
    try: box.adjustments[0] = 0.04
    except Exception: pass
    box.fill.solid(); box.fill.fore_color.rgb = DBOX
    box.line.color.rgb = DBOX; box.shadow.inherit = False
    tf = tb(slide, 1.10, 2.35, 11.13, 2.60)
    tasks = [
        "在自己项目建 .claude/agents/code-reviewer.md（只读 tools ＋ 审查提示词），点名触发一次——确认主对话只收到结论、面板出现它的运行条。",
        "让主 agent「用三个 sub-agent 并行」调研三个目录，再串行做一遍——记录两轮的墙钟时间与 token 消耗，亲手算出并行收益与放大成本。",
        "给一个 skill 加 context: fork 与 disable-model-invocation: true，用 /名字 触发——验证「确保调用」的硬保证链路。",
    ]
    for i, t in enumerate(tasks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        sr(p.add_run(), f"{i+1} · ", BODY_F, 15, ORANGE, bold=True)
        r = p.add_run(); sr(r, t, BODY_F, 15, EAF, bold=True)
    tf = tb(slide, 0.70, 5.35, 11.93, 1.40)
    p = tf.paragraphs[0]
    sr(p.add_run(), "提示：", BODY_F, 13.5, ORANGE, bold=True)
    sr(p.add_run(), "第 2 题量出来的数字就是你的售前弹药——给客户讲决策账时，自己项目里实测的放大倍数，比任何博客都有说服力。",
       BODY_F, 13.5, EAF, bold=True)
    foot(slide, "Sub-agent 编排 · 动手做")

# ---------- p11 对练页 ----------
def p11_qa(slide):
    bg(slide)
    tf = tb(slide, 0.70, 0.55, 11.5, 0.35)
    sr(tf.paragraphs[0].add_run(), "第 10 章 · Sub-agent 编排 · 客户会怎么问（面试题）", BODY_F, 12, TEAL, bold=True)
    tf = tb(slide, 0.70, 0.90, 11.93, 0.95)
    sr(tf.paragraphs[0].add_run(), "现场最可能被追问的问题", TITLE_F, 30, NAVY, bold=True)
    qas = [
        ("你们第 3 章讲的 multi-agent 和这章的 sub-agent，是一回事吗？",
         "一体两面：orchestrator-workers 是模式（蓝图），sub-agent 是 coding agent 产品里的机制（车间）。给客户讲蓝图，落地时看车间。"),
        ("为什么我的任务有时开好几个 sub-agent，有时一个不开？",
         "主 agent 在算账：能切成独立块才并行；强顺序依赖、或依赖主对话已有理解的活，自己干更快更准。口诀：要结果外包，要过程留下。"),
        ("sub-agent 会不会把费用跑爆？",
         "会放大——研究类任务约 15× token。治理三件：按任务价值分层开、并发限额（Codex 默认 6 线程）、token 预算与超额熔断。"),
        ("怎么保证某个 sub-agent 每次都被调用？",
         "分层：日常写好 description 让它自动路由；要硬保证就用 skill 绑 context: fork——/名字 显式触发，必然进 sub-agent。"),
        ("sub-agent 把活干歪了怎么发现？",
         "窄带通信＝过程不可见，所以每个 sub-agent 的输入输出都要留 trace、可回放（第 6 章）；任务书写清产出格式与边界是第一道防线。"),
    ]
    for k, (q, a) in enumerate(qas):
        y = 2.10 + 0.97 * k
        band(slide, 0.70, y, 11.93, 0.81, fill=RGBColor(0xF3, 0xF8, 0xFA), edge=RGBColor(0xF3, 0xF8, 0xFA))
        tf = tb(slide, 0.98, y + 0.08, 11.37, 0.30)
        p = tf.paragraphs[0]
        sr(p.add_run(), "Q  ", BODY_F, 14, TEAL, bold=True)
        sr(p.add_run(), q, BODY_F, 13.5, INK, bold=True)
        tf = tb(slide, 0.98, y + 0.36, 11.37, 0.44)
        p = tf.paragraphs[0]
        sr(p.add_run(), "A  ", BODY_F, 13, ORANGE, bold=True)
        sr(p.add_run(), a, BODY_F, 12, INK)
    foot(slide, "Sub-agent 编排")

# ---------- p12 本章小结(深色) ----------
def p12_recap(slide):
    dbg(slide)
    tf = tb(slide, 0.70, 0.70, 9.0, 0.40)
    sr(tf.paragraphs[0].add_run(), "第 10 章 · 本章小结", BODY_F, 13, ORANGE, bold=True)
    tf = tb(slide, 0.70, 1.15, 11.0, 0.80)
    sr(tf.paragraphs[0].add_run(), "多智能体 / Sub-agent 编排", TITLE_F, 30, WHITE, bold=True)
    tf = tb(slide, 0.70, 2.20, 11.93, 3.40)
    for i, t in enumerate([
        "Sub-agent 就是 agent：同一机制递归一层；派遣能力住在运行时层——平台能力，不是模型独门绝技。",
        "三个性质定一切：独立上下文（任务书自包含）、窄带通信（只回结论）、深度上限（Claude Code 5 层 / Codex 默认 1）。",
        "触发四层从软到硬：description 路由 → 对话点名 → CLAUDE.md → skill 绑 context: fork（硬保证）。",
        "决策一本账：约 15× token 换 90.2% 提升——要结果的活外包，要过程的活留在主对话。",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        sr(p.add_run(), t, BODY_F, 15.5, EAF)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.70), Inches(5.75), Inches(11.93), Inches(1.00))
    try: box.adjustments[0] = 0.06
    except Exception: pass
    box.fill.solid(); box.fill.fore_color.rgb = DBOX
    box.line.color.rgb = DBOX; box.shadow.inherit = False
    tf = tb(slide, 1.00, 5.85, 11.33, 0.80)
    p = tf.paragraphs[0]
    sr(p.add_run(), "承上启下：", BODY_F, 13, ORANGE, bold=True)
    sr(p.add_run(), "编排从蓝图（第 3 章）讲到了车间（本章）——最后一章把全书折叠成三页速查，现场 15 秒一条用。",
       BODY_F, 13, EAF)
    foot(slide, "Sub-agent 编排")

PAGES = [
    (p01_cover, "多智能体与 Sub-agent 编排"),
    (p02_goals, "学完这一章，你能对客户讲清楚"),
    (p03_relation, "Sub-agent 就是 Agent：递归一层，不是新物种"),
    (p04_fanout, "扇出与扇入（Fan-out / Fan-in）：一份任务书进去，一段结论回来"),
    (p05_props, "三个关键性质：一切用法与局限的根源"),
    (p06_cc, "Claude Code 实操：一份 markdown，一个分身"),
    (p07_codex, "Codex 对照：同一件事的另一种拼法"),
    (p08_ledger, "派与不派：先背口诀，再算账"),
    (p09_cloud, "上云怎么落地：放大后的账单谁来接"),
    (p10_handson, "第 10 章 · Sub-agent 编排"),
    (p11_qa, "现场最可能被追问的问题"),
    (p12_recap, "多智能体 / Sub-agent 编排"),
]
