import os
"""Data-Engineering 讲义信息图（5 张）。复用 kb_draw + agent_figs helper。
嫁接：p5 后(五处证据)、p19 后(连接器五件事)、p27 后(向量库演进)、p34 后(坏答案回流)、p48 后(治理四输入)。
DE 无页码，图系统页脚匹配。独立预览：python3 de_figs.py"""
import sys, math
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "数据工程讲义 · Data-Engineering"
GINK = RGBColor(0x27, 0x50, 0x0A)

# ---------- p5 五处证据 hub ----------
def draw_de_five_evidence(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 立论  ONE PIT, FIVE VOICES")
    title(slide, "五处证据：全库都在喊同一个坑")
    label(slide, 0.55, 1.58, 12.2, "五句话都是果，因在同一处——没人把数据搞就绪；本模块就是那个「接住」的动作", SUB, 12, PP_ALIGN.LEFT, True)
    # 中心
    node(slide, 5.35, 3.7, 2.65, 1.05, "数据就绪度", "是第一风险", "把隐性成本变显性工程件", fill=DTEAL, edge=NAVY, zh_c=WHITE, sub_c=PALE, en_c=PALE, zh_size=15, en_size=10.5, desc_size=9)
    cx, cy = 6.68, 4.22
    spokes = [("SP 第 2 章 · POC 三要素", "数据坑=第一风险：脏、缺、散", 0.7, 2.5),
              ("SP 第 4 章 · 企业搜索", "连接器是报价大头：每个源都是小项目", 9.35, 2.5),
              ("SP 第 5 章 · 内容生成", "产品库是唯一事实源——前提是干净的", 9.7, 4.15),
              ("AI-Ops 第 3 章", "垃圾进垃圾出：坏答案大半源于脏数据", 9.35, 5.7),
              ("RAG 第 3 章", "解析质量决定切分质量，这关没过全白搭", 0.7, 5.7)]
    for name, txt, bx, by in spokes:
        b = band(slide, bx, by, 3.28, 1.0, fill=CARD, edge=TEAL)
        tf = tb(slide, bx+0.18, by+0.1, 2.95, 0.85)
        para(tf, name, BODY_F, 11, DTEAL, bold=True, first=True)
        para(tf, txt, BODY_F, 9.5, SUB)
        # 箭头指向中心
        ax = bx + (3.28 if bx < 5 else 0); ay = by + 0.5
        line(slide, ax, ay, cx + (-1.35 if bx < 5 else 1.35), cy + (by-cy)*0.28, color=RGBColor(0x9F,0xB3,0xC0), w=1.2)
    footer(slide, MOD, "第 1 章 · 数据就绪度")

# ---------- p19 连接器五件事 ----------
def draw_de_connector(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 工程清单  ONE CONNECTOR = 5 THINGS")
    title(slide, "一个连接器 = 五件事")
    label(slide, 0.55, 1.58, 12.2, "报价口径：按「源系统 × 五件事」估工作量——自研系统才是真正的大头", SUB, 12, PP_ALIGN.LEFT, True)
    steps = [("① 认证", "OAuth / 服务账号"), ("② 拉取", "API 限速与分页"), ("③ 解析", "格式转换（第 2 章）"),
             ("④ ACL 映射", "权限同步（第 7 章）"), ("⑤ 增量", "轮询 / webhook / 日志")]
    W = 2.28; gap = 0.14; x0 = 0.6; y = 2.5; h = 1.1; xs = []
    for i, (zh, d) in enumerate(steps):
        x = x0 + i*(W+gap); xs.append(x)
        node(slide, x, y, W, h, zh, "", d, fill=CARD, edge=TEAL, zh_c=DTEAL, sub_c=SUB, zh_size=13, desc_size=9.5)
        if i < 4:
            arrow(slide, x+W+0.01, y+h/2, x+W+gap-0.01, y+h/2, color=TEAL, w=1.5)
    label(slide, 0.6, 3.95, 6.0, "第 ⑤ 件事：增量同步的三种模式取舍", DTEAL, 12, PP_ALIGN.LEFT, False, True)
    modes = [("轮询", "实现最简单；新鲜度差、浪费配额——低频源够用", TEAL, DTEAL, CARD),
             ("Webhook", "准实时；要源系统支持回调、要处理丢消息补偿", TEAL, DTEAL, CARD),
             ("变更日志 CDC", "最可靠；只有数据库类源才有，工程最重", ORANGE, ODARK, CORAL)]
    mx = [0.7, 4.68, 8.66]; my = 4.4; mw = 3.9; mh = 1.15
    for x, (h1, h2, tc, ec, fill) in zip(mx, modes):
        band(slide, x, my, mw, mh, fill=fill, edge=ec)
        tf = tb(slide, x+0.22, my+0.12, mw-0.44, 0.95)
        para(tf, h1, BODY_F, 12, tc, bold=True, first=True)
        para(tf, h2, BODY_F, 10.5, SUB)
    band(slide, 0.7, 5.75, 11.93, 1.07, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 5.86, 11.4, 0.9); p = tf.paragraphs[0]
    sr(p.add_run(), "报价口径：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "SharePoint / Confluence 这类有成熟连接器（Airbyte 模式可参照）；自研系统每个都是一个小项目——那才是大头。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 3 章 · 连接器与同步")

# ---------- p27 向量库按规模演进 ----------
def draw_de_vectordb_scaling(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · 决策与迁移  SCALE-DRIVEN")
    title(slide, "向量库按规模演进：起步、触顶、迁移")
    label(slide, 0.55, 1.58, 12.2, "「按规模演进」的决策，不是「一步到位」的决策——起步简单，把迁移路径想清楚", SUB, 12, PP_ALIGN.LEFT, True)
    steps = [("起步", "pgvector / ES 就地", TEAL, CARD), ("增长", "监控召回与延迟", TEAL, CARD),
             ("触顶", "50–100M / 重建变慢", ORANGE, CORAL), ("迁移", "Qdrant / Milvus", ORANGE, CORAL),
             ("超大规模", "分布式 + 分租户", DTEAL, PALE)]
    W = 2.28; gap = 0.14; x0 = 0.6; y = 2.5; h = 0.95; xs = []
    for i, (zh, d, ec, fill) in enumerate(steps):
        x = x0 + i*(W+gap); xs.append(x)
        node(slide, x, y, W, h, zh, "", d, fill=fill, edge=ec, zh_c=(WHITE if fill == PALE else ec if ec != CORAL else ODARK), sub_c=SUB, zh_size=13, desc_size=9.5)
        if i < 4:
            arrow(slide, x+W+0.01, y+h/2, x+W+gap-0.01, y+h/2, color=SUB, w=1.5)
    band(slide, 0.6, 3.85, 5.95, 1.65, fill=CORAL, edge=ORANGE)
    tf = tb(slide, 0.85, 3.97, 5.5, 1.45)
    para(tf, "触顶信号", BODY_F, 12, ODARK, bold=True, first=True)
    para(tf, "· HNSW 索引重建时间不可接受", BODY_F, 10.5, SUB)
    para(tf, "· P99 延迟随库量线性恶化", BODY_F, 10.5, SUB)
    para(tf, "· 过滤条件多导致召回骤降", BODY_F, 10.5, SUB)
    band(slide, 6.68, 3.85, 5.95, 1.65, fill=CARD, edge=TEAL)
    tf2 = tb(slide, 6.93, 3.97, 5.5, 1.45)
    para(tf2, "迁移纪律", BODY_F, 12, DTEAL, bold=True, first=True)
    para(tf2, "· 嵌入模型不换则向量可搬；换模型必须全量重算", BODY_F, 10.5, SUB)
    para(tf2, "· 双写灰度：新旧库并行跑评估集对比", BODY_F, 10.5, SUB)
    para(tf2, "· 别为「将来可能的十亿级」提前买复杂度", BODY_F, 10.5, SUB)
    band(slide, 0.6, 5.7, 12.03, 1.12, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.82, 11.5, 0.95); p = tf3.paragraphs[0]
    sr(p.add_run(), "对客一句话：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "起步简单、把迁移路径想清楚，比压注哪家都稳——向量库是演进决策，不是一锤定音。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 4 章 · 向量库选型")

# ---------- p34 坏答案回流闭环 ----------
def draw_de_bad_answer_loop(slide):
    bg(slide)
    eyebrow(slide, "第 5 章 · 回流闭环  FIX DATA, NOT PROMPT")
    title(slide, "坏答案回流：修数据，而不是改提示词")
    label(slide, 0.55, 1.58, 12.2, "八成问题在检索侧，检索问题的大半又在数据侧——坏答案要一路追到数据层", SUB, 12, PP_ALIGN.LEFT, True)
    steps = [("① 在线评估", "AI-Ops 抓坏答案"), ("② 归因", "检索? 解析? 缺文档?"),
             ("③ 修数据", "补 / 换 / 下架"), ("④ 回归验证", "评估集重跑")]
    cx, cy = 4.0, 4.15; rx, ry = 2.75, 1.35; W, H = 2.15, 0.85
    cen = []
    for i in range(4):
        ang = math.radians(-90 + i*90)
        cen.append((cx + rx*math.cos(ang), cy + ry*math.sin(ang)))
    for i in range(4):
        x1, y1 = cen[i]; x2, y2 = cen[(i+1) % 4]
        dx, dy = x2-x1, y2-y1; d = math.hypot(dx, dy); ux, uy = dx/d, dy/d
        arrow(slide, x1+ux*1.15, y1+uy*0.55, x2-ux*1.15, y2-uy*0.55, color=TEAL, w=1.6)
    for i, (zh, d) in enumerate(steps):
        x, y = cen[i]
        solid = (i == 2)
        node(slide, x-W/2, y-H/2, W, H, zh, "", d, fill=(ORANGE if solid else CARD), edge=(ODARK if solid else TEAL),
             zh_c=(OINK if solid else DTEAL), sub_c=(OINK if solid else SUB), zh_size=12.5, desc_size=9)
    label(slide, cx-1.4, cy-0.16, 2.8, "改提示词治不了数据病", ODARK, 11, PP_ALIGN.CENTER, False, True)
    band(slide, 7.6, 2.5, 5.05, 3.05, fill=CARD, edge=LINE)
    tf = tb(slide, 7.85, 2.63, 4.55, 2.85)
    para(tf, "覆盖率仪表盘的售前故事", BODY_F, 12.5, DTEAL, bold=True, first=True)
    para(tf, "高频无答案问题清单 = 客户该补哪些文档的优先级表——把「AI 答不上」变成知识治理的抓手。", BODY_F, 11, SUB)
    para(tf, "质量运营的节奏", BODY_F, 12.5, ODARK, bold=True)
    para(tf, "日看新鲜度告警、周看坏答案归因、月出四指标报告——和 AI-Ops 运营包同一节奏合并交付。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 5 章 · 数据质量")

# ---------- p48 治理四类输入 ----------
def draw_de_access_decision(slide):
    bg(slide)
    eyebrow(slide, "第 7 章 · 治理衔接 · 深潜  ACCESS DECISION")
    title(slide, "每次访问 = 四类输入 + 一个策略决策")
    label(slide, 0.55, 1.58, 12.2, "执行点管「在哪拦」，四输入模型管「怎么判」——授权在资源服务端强制执行", SUB, 12, PP_ALIGN.LEFT, True)
    ins = [("主体 + 资源", "谁在请求（用户/应用/agent/租户）· 访问什么（表/列/文档/索引）", "客服 agent 替张三查——主体是「张三经由 agent」，不是 agent 自己"),
           ("动作 + 上下文", "要做什么（读/检索/训练/分享/删除）· 什么条件与目的", "同一份文档：客服场景可检索，营销场景不可用于训练")]
    y = 2.4
    for i, (h1, h2, ex) in enumerate(ins):
        yy = y + i*1.15
        band(slide, 0.6, yy, 7.4, 1.02, fill=CARD, edge=TEAL)
        tf = tb(slide, 0.82, yy+0.1, 7.0, 0.85); p = tf.paragraphs[0]
        sr(p.add_run(), f"输入 {i+1} · " + h1 + "　", BODY_F, 12, DTEAL, bold=True)
        sr(p.add_run(), h2, BODY_F, 10, SUB)
        para(tf, "例：" + ex, BODY_F, 10, ODARK)
    arrow(slide, 8.05, 3.5, 8.75, 3.5, color=SUB, w=2.0)
    node(slide, 8.8, 2.95, 3.8, 1.1, "策略引擎", "决策大脑", "allow / deny / 附加义务", fill=DTEAL, edge=NAVY, zh_c=WHITE, sub_c=PALE, en_c=PALE, zh_size=15, en_size=10, desc_size=10)
    band(slide, 8.8, 4.25, 3.8, 1.25, fill=CORAL, edge=ORANGE)
    tf2 = tb(slide, 9.0, 4.36, 3.45, 1.1)
    para(tf2, "附加义务", BODY_F, 11, ODARK, bold=True, first=True)
    para(tf2, "脱敏 · 行列过滤 · 水印 · 短保留 · 人工审批", BODY_F, 10, SUB)
    band(slide, 0.6, 5.7, 12.03, 1.12, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.82, 11.5, 0.95); p = tf3.paragraphs[0]
    sr(p.add_run(), "拼齐才叫治理闭环：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "执行点管「在哪拦」、四输入管「怎么判」、审计记主体与理由、血缘定位派生影响——模型 / 查询 / 客户端都不能自行裁决。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 7 章 · 治理衔接")

# ---------- p45 评估集 vs 训练集 ----------
def draw_de_two_lines(slide):
    bg(slide)
    eyebrow(slide, "第 6 章 · 双线运营  ONE BAD CASE, TWO PATHS")
    title(slide, "评估集与训练集：一份坏例，两条去向")
    label(slide, 0.55, 1.58, 12.2, "坏例先问「是不知道还是不听话」——分流错了，钱就白花", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 5.35, 2.35, 2.65, 0.72, "一份坏例", "分流", "", fill=NAVY, edge=NAVY, zh_c=WHITE, sub_c=PALE, zh_size=14, en_size=9.5)
    label(slide, 3.4, 3.35, 2.6, "「不知道」知识缺", DTEAL, 11, PP_ALIGN.CENTER, False, True)
    label(slide, 7.5, 3.35, 2.6, "「不听话」行为偏", ODARK, 11, PP_ALIGN.CENTER, False, True)
    arrow(slide, 5.9, 3.07, 4.4, 3.6, color=TEAL, w=1.7)
    arrow(slide, 7.4, 3.07, 8.9, 3.6, color=ORANGE, w=1.7)
    band(slide, 0.7, 3.7, 5.75, 1.95, fill=CARD, edge=TEAL)
    label(slide, 0.95, 3.82, 5.3, "评估集线", DTEAL, 13, PP_ALIGN.LEFT, False, True)
    tf = tb(slide, 0.98, 4.28, 5.3, 1.3)
    para(tf, "· 管：验收与回归（Evaluation 第 5 章方法）", BODY_F, 11, INK, bold=True, first=True)
    para(tf, "· 坏例打标签入集；留 10–20% 保留集防应试", BODY_F, 11, SUB)
    para(tf, "· 修数据 / RAG 更新（知识缺 → 补库）", BODY_F, 11, DTEAL, bold=True)
    band(slide, 6.85, 3.7, 5.75, 1.95, fill=CORAL, edge=ORANGE)
    label(slide, 7.1, 3.82, 5.3, "训练集线", ODARK, 13, PP_ALIGN.LEFT, False, True)
    tf2 = tb(slide, 7.13, 4.28, 5.3, 1.3)
    para(tf2, "· 管：微调数据（Fine-tuning 第 3 章配方）", BODY_F, 11, INK, bold=True, first=True)
    para(tf2, "· 只收「行为类」坏例", BODY_F, 11, SUB)
    para(tf2, "· 知识类交给 RAG 更新，别硬塞进训练集", BODY_F, 11, ODARK, bold=True)
    band(slide, 0.6, 5.85, 12.03, 0.97, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.95, 11.5, 0.82); p = tf3.paragraphs[0]
    sr(p.add_run(), "运营纪律：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "业务专家的时间是最稀缺资源，只用在「定标准、写种子、审边界样本」三件事上——每周两小时能撑起整个标注运营。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 6 章 · 标注与合成")

FIGS = [draw_de_five_evidence, draw_de_connector, draw_de_vectordb_scaling,
        draw_de_bad_answer_loop, draw_de_access_decision, draw_de_two_lines]
NEW = [draw_de_two_lines]

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in FIGS:
        fn(prs.slides.add_slide(B))
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_de_preview.pptx")
    prs.save(out); print(f"saved {len(FIGS)} ->", out)
