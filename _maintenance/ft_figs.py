"""Fine-tuning 讲义信息图（7 张）。复用 kb_draw + agent_figs helper。
嫁接：p7 后(定制光谱)、p8 后(知识vs行为)、p18 后(LoRA)、p22 后(三方对比)、
      p43 后(loss曲线)、p52 后(SFT→DPO→RFT)、p72 后(部署两形态)。
FT 页脚无页码，图系统页脚(MOD/章名)正好匹配。独立预览：python3 ft_figs.py"""
import sys
sys.path.insert(0, "/Users/lijiaxiang/project/myAILearning/_maintenance")
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "Fine-tuning 讲义 · 微调工程实践"
GINK = RGBColor(0x27, 0x50, 0x0A)
REDBG = RGBColor(0xFB, 0xE9, 0xE7)
REDINK = RGBColor(0x9B, 0x2C, 0x2C)

# ---------- p7 定制光谱（阶梯）----------
def draw_ft_spectrum(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 定制光谱  CUSTOMIZATION LADDER")
    title(slide, "定制手段光谱：一级一级往上爬")
    label(slide, 0.55, 1.58, 12.2, "梯子原则——从左往右逐级升级，左边一档没榨干之前不动右边", SUB, 12, PP_ALIGN.LEFT, True)
    rungs = [("提示词工程", "分钟级 · 零训练", "改说法、给示例，立即生效；上限是上下文窗口", TEAL, DTEAL, CARD),
             ("RAG", "小时级 · 建索引", "知识外挂：可溯源、随时更新；不改模型行为", TEAL, DTEAL, CARD),
             ("微调", "天级 · 要数据", "行为内化：语气、格式、领域习惯；改的是权重", ORANGE, ODARK, CORAL),
             ("继续预训练", "月级 · 重投入", "大规模灌注领域语料，改知识底盘；极少数场景", ODARK, RGBColor(0x5A,0x24,0x12), RGBColor(0xF3,0xE2,0xD8))]
    x0 = 0.7; W = 2.82; gap = 0.16; ybase = 4.68
    xs = []
    for i, (zh, when, desc, ec, tc, fill) in enumerate(rungs):
        x = x0 + i*(W+gap); y = ybase - i*0.6; xs.append((x, y))
        band(slide, x, y, W, 1.4, fill=fill, edge=ec)
        tf = tb(slide, x+0.2, y+0.12, W-0.4, 1.22); p = tf.paragraphs[0]
        sr(p.add_run(), f"{i+1}. {zh}", BODY_F, 14, tc, bold=True)
        para(tf, when, BODY_F, 11, ec if fill != CARD else DTEAL, bold=True)
        para(tf, desc, BODY_F, 11, SUB)
    for i in range(3):
        (x1, y1), (x2, y2) = xs[i], xs[i+1]
        arrow(slide, x1+W+0.02, y1+0.2, x2-0.02, y2+0.95, color=SUB, w=1.6)
    label(slide, 0.55, 2.3, 3.6, "▲ 成本 / 周期 / 维护 逐级升", ODARK, 10.5, PP_ALIGN.LEFT, False, True)
    label(slide, 9.3, 6.55, 3.6, "生效速度逐级降 ▼", SUB, 10.5, PP_ALIGN.LEFT, False, True)
    band(slide, 0.6, 6.15, 8.4, 0.72, fill=PALE, edge=LINE)
    tf = tb(slide, 0.85, 6.25, 8.0, 0.6); p = tf.paragraphs[0]
    sr(p.add_run(), "售前口径：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "客户说「模型不好用」先问卡在哪一档——九成的「要微调」，追问两句会退回提示词或 RAG。", BODY_F, 11, INK)
    footer(slide, MOD, "第 1 章 · 什么时候该微调")

# ---------- p8 知识 vs 行为 ----------
def draw_ft_knowledge_vs_behavior(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 知识 vs 行为  THE KEY FORK")
    title(slide, "最重要的分界线：知识问题，还是行为问题")
    label(slide, 0.55, 1.58, 12.2, "先分清是「不知道」还是「不听话」——分错了，微调也白搭", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 5.35, 2.35, 2.65, 0.8, "模型的问题是？", "", "", fill=NAVY, edge=NAVY, zh_c=WHITE, zh_size=14)
    # 左 知识→RAG
    band(slide, 0.7, 3.55, 5.75, 2.05, fill=CARD, edge=TEAL)
    label(slide, 0.95, 3.67, 5.3, "「不知道」/ 过期 → 知识问题", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 4.18, 1.7, 0.66, "上 RAG", "", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, zh_size=13)
    tf = tb(slide, 2.8, 4.14, 3.5, 1.35)
    para(tf, "改数据不改模型 · 可溯源 · 便宜", BODY_F, 11, INK, bold=True, first=True)
    para(tf, "类比：发最新版「手册」，内容天天可换新。", BODY_F, 10.5, SUB)
    para(tf, "例：「答错了产品价格」→ 知识过期。", BODY_F, 10.5, ODARK)
    # 右 行为→微调
    band(slide, 6.85, 3.55, 5.75, 2.05, fill=CORAL, edge=ORANGE)
    label(slide, 7.1, 3.67, 5.3, "「不听话」/ 风格不对 → 行为问题", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 4.18, 1.7, 0.66, "才上微调", "", "", fill=ORANGE, edge=ODARK, zh_c=OINK, zh_size=13)
    tf2 = tb(slide, 8.95, 4.14, 3.5, 1.35)
    para(tf2, "改权重 · 内化语气/格式/习惯", BODY_F, 11, INK, bold=True, first=True)
    para(tf2, "类比：送员工去「培训」，改做事习惯。", BODY_F, 10.5, SUB)
    para(tf2, "例：「老不按工单 JSON 格式」→ 行为不合规。", BODY_F, 10.5, ODARK)
    arrow(slide, 6.0, 3.15, 3.5, 3.53, color=TEAL, w=1.7)
    arrow(slide, 7.35, 3.15, 9.7, 3.53, color=ORANGE, w=1.7)
    band(slide, 0.7, 5.85, 11.93, 0.98, fill=GBG, edge=GREEN)
    tf3 = tb(slide, 1.0, 5.96, 11.4, 0.85); p = tf3.paragraphs[0]
    sr(p.add_run(), "生产中常组合：", BODY_F, 11.5, GINK, bold=True)
    sr(p.add_run(), "微调管「风格与格式」，RAG 管「事实与知识」——与 RAG 模块互为镜像，一道选择题两侧收口。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 1 章 · 什么时候该微调")

# ---------- p18 LoRA 便利贴原理 ----------
def draw_ft_lora(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · LoRA  HOW IT WORKS")
    title(slide, "LoRA：给权重贴便利贴")
    label(slide, 0.55, 1.58, 12.2, "冻结原权重 W，旁边加一对低秩矩阵 B×A，只训这条旁路——可训参数通常不到 1%", SUB, 12, PP_ALIGN.LEFT, True)
    # 机制图（W 冻结 + BA 训练 → 汇入生效）
    band(slide, 0.6, 2.4, 7.45, 3.0, fill=RGBColor(0xF7,0xFB,0xFC), edge=LINE)
    node(slide, 0.9, 2.72, 2.4, 0.98, "原权重 W", "冻结 frozen", "一个字不动", fill=RGBColor(0xEC,0xEC,0xEA), edge=SUB, zh_c=RGBColor(0x55,0x55,0x52), sub_c=SUB, en_c=SUB, zh_size=13.5, en_size=10, desc_size=9.5)
    node(slide, 0.9, 4.08, 2.4, 0.98, "低秩旁路 B × A", "秩 r = 8–64", "只训这条 · <1% 参数", fill=ORANGE, edge=ODARK, zh_c=OINK, sub_c=OINK, en_c=OINK, zh_size=12.5, en_size=9.5, desc_size=9.5)
    node(slide, 5.55, 3.36, 1.85, 1.1, "生效", "W + BA", "推理时相加", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=WHITE, zh_size=14, en_size=11, desc_size=9.5)
    label(slide, 3.55, 3.62, 1.4, "＋", INK, 22, PP_ALIGN.CENTER, False, True)
    arrow(slide, 3.32, 3.28, 5.5, 3.72, color=RGBColor(0x9F,0xB3,0xC0), w=1.5, dashed=True)
    arrow(slide, 3.32, 4.5, 5.5, 4.06, color=ODARK, w=1.6)
    label(slide, 3.75, 3.0, 2.0, "冻结·直接并入", SUB, 9, PP_ALIGN.LEFT, True)
    label(slide, 3.75, 4.5, 2.0, "训练后相加", ODARK, 9, PP_ALIGN.LEFT, True)
    # 右侧类比+案例
    band(slide, 8.2, 2.4, 4.43, 3.0, fill=CARD, edge=TEAL)
    tf = tb(slide, 8.45, 2.53, 3.95, 2.8)
    para(tf, "类比 · 便利贴", BODY_F, 12.5, DTEAL, bold=True, first=True)
    para(tf, "不重写整本书，只在关键页贴便利贴——原书一字不动，读时把便利贴一起读进去；一套便利贴对应一个任务，可揭下换一套。", BODY_F, 11, SUB)
    para(tf, "案例：8B 挂 LoRA（r=16）", BODY_F, 11.5, ODARK, bold=True)
    para(tf, "可训参数约 0.3% · 单卡几小时 · 产物只是几十 MB 的 adapter 文件。", BODY_F, 11, SUB)
    band(slide, 0.6, 5.62, 12.03, 1.2, fill=GBG, edge=GREEN)
    tf2 = tb(slide, 0.9, 5.74, 11.5, 1.05); p = tf2.paragraphs[0]
    sr(p.add_run(), "默认起点：", BODY_F, 12, GINK, bold=True)
    sr(p.add_run(), "生产微调的默认选择——adapter 小、好存好换好回滚，还能多套并存（第 7 章多租户部署就靠它）。", BODY_F, 12, INK)
    footer(slide, MOD, "第 2 章 · 微调方法谱系")

# ---------- p22 全参 / LoRA / QLoRA 三方对比 ----------
def draw_ft_three_methods(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 微调方法谱系  FULL / LoRA / QLoRA")
    title(slide, "三种微调对比：显存、效果、遗忘怎么权衡")
    label(slide, 0.55, 1.58, 12.2, "客户问「要几张卡」，看这张——从整本书重编，到贴便利贴，到压缩影印再贴", SUB, 12, PP_ALIGN.LEFT, True)
    cards = [("全参微调", "Full", "整本书重编", 1.0, "≈ 参数 × 16 字节", "7B ≈ 112GB · 多卡集群", "效果上限最高，遗忘风险也高", REDINK, ODARK, CORAL),
             ("LoRA", "低秩旁路", "贴便利贴", 0.42, "16-bit 底座 + 小旁路", "单卡几小时 · adapter 几十 MB", "学得少也忘得少 · 默认起点", DTEAL, TEAL, CARD),
             ("QLoRA", "4-bit 底座", "压缩影印再贴", 0.16, "4-bit 量化 + 分页优化器", "7B 只需 6–8GB · 单卡 4090", "轻微质量损失 · 没集群第一选择", GINK, GREEN, GBG)]
    W = 3.95; gap = 0.14; x0 = 0.6; y = 2.4; h = 2.9
    for i, (zh, en, ana, memfrac, memrule, memcase, note, tc, ec, fill) in enumerate(cards):
        x = x0 + i*(W+gap)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, y+0.15, W-0.5, 0.72); p = tf.paragraphs[0]
        sr(p.add_run(), f"{i+1}. {zh}", BODY_F, 14.5, tc, bold=True)
        sr(p.add_run(), "　" + ana, BODY_F, 10.5, SUB, italic=True)
        # 显存条
        label(slide, x+0.28, y+0.86, W-0.56, "显存", SUB, 9.5, PP_ALIGN.LEFT, False, True)
        band(slide, x+0.9, y+0.9, (W-1.2)*memfrac, 0.24, fill=ec, edge=ec)
        tf2 = tb(slide, x+0.28, y+1.2, W-0.56, 0.4); sr(tf2.paragraphs[0].add_run(), memrule, BODY_F, 11, INK, bold=True)
        tf3 = tb(slide, x+0.28, y+1.62, W-0.56, 0.4); sr(tf3.paragraphs[0].add_run(), memcase, BODY_F, 11, SUB)
        tf4 = tb(slide, x+0.28, y+2.12, W-0.56, 0.65); sr(tf4.paragraphs[0].add_run(), note, BODY_F, 11, tc, bold=True)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.67, 11.5, 1.1); p = tf.paragraphs[0]
    sr(p.add_run(), "决策口径：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "行为定制、多任务共存 → LoRA 先行；要把大量新知识「灌」进模型 → 才考虑全参 + 防遗忘；没有大集群 → QLoRA。", BODY_F, 12, INK)
    para(tf, "心算规则给客户先讲量级（全参×16 字节、QLoRA×0.5 字节），精确数留给 POC。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 2 章 · 微调方法谱系")

# ---------- p43 loss 曲线 ----------
def draw_ft_loss_curves(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · 训练监控  LOSS CURVES")
    title(slide, "盯 eval 不盯 train：loss 曲线怎么看")
    label(slide, 0.55, 1.58, 12.2, "train loss 下降只说明「在学」；eval loss 同步下降才说明「学到能泛化的东西」", SUB, 12, PP_ALIGN.LEFT, True)
    # 坐标系
    ox, oy, pw, ph = 1.5, 5.2, 5.6, 2.5   # 原点(左下)，宽高
    line(slide, ox, oy-ph, ox, oy, color=SUB, w=1.4)      # y 轴
    line(slide, ox, oy, ox+pw, oy, color=SUB, w=1.4)      # x 轴
    label(slide, 0.5, oy-ph-0.05, 1.0, "loss", SUB, 10, PP_ALIGN.LEFT, False, True)
    label(slide, ox+pw-2.5, oy+0.1, 2.6, "训练步 / epoch →", SUB, 10, PP_ALIGN.RIGHT, True)
    def P(fx, fv):  # fx,fv in 0..1 (fv: 0=low loss 底, 1=high loss 顶)
        return (ox + fx*pw, oy - fv*ph)
    train = [(0.02, 0.92), (0.18, 0.68), (0.36, 0.46), (0.54, 0.30), (0.72, 0.20), (0.9, 0.14), (0.99, 0.11)]
    ev = [(0.02, 0.95), (0.18, 0.72), (0.36, 0.52), (0.5, 0.44), (0.58, 0.43), (0.72, 0.5), (0.86, 0.62), (0.99, 0.74)]
    for series, col in [(train, TEAL), (ev, ORANGE)]:
        pts = [P(a, b) for a, b in series]
        for j in range(len(pts)-1):
            line(slide, pts[j][0], pts[j][1], pts[j+1][0], pts[j+1][1], color=col, w=2.4)
    # 过拟合点
    mx, my = P(0.56, 0.43)
    dot(slide, mx, my, 0.08, ODARK)
    label(slide, mx-0.4, my-0.5, 2.9, "eval 掉头 = 过拟合起点", ODARK, 10, PP_ALIGN.LEFT, False, True)
    label(slide, mx-0.2, my-0.26, 2.9, "取这个 checkpoint 交付", ODARK, 9.5, PP_ALIGN.LEFT, True)
    label(slide, ox+pw*0.72, oy-ph*0.16, 1.6, "train ↓", TEAL, 10.5, PP_ALIGN.LEFT, False, True)
    label(slide, ox+pw*0.78, oy-ph*0.68, 1.8, "eval ↑", ODARK, 10.5, PP_ALIGN.LEFT, False, True)
    # 右侧
    band(slide, 7.5, 2.4, 5.13, 2.8, fill=CARD, edge=TEAL)
    tf = tb(slide, 7.75, 2.53, 4.65, 2.6)
    para(tf, "类比 · 刷题分 vs 模拟考", BODY_F, 12.5, DTEAL, bold=True, first=True)
    para(tf, "train loss 是刷题正确率——刷多了自然高，可能只是背了题；eval loss 是没见过的模拟考，这个分才算数。", BODY_F, 11, SUB)
    para(tf, "标准操作", BODY_F, 12, ODARK, bold=True)
    para(tf, "小数据训 3 个 epoch，eval 在第 2 个 epoch 后上翘 → 取第 2 个 epoch 的 checkpoint，而不是「再训久一点」。", BODY_F, 11, SUB)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=GBG, edge=GREEN)
    tf2 = tb(slide, 0.9, 5.68, 11.5, 1.1); p = tf2.paragraphs[0]
    sr(p.add_run(), "铁律：", BODY_F, 12, GINK, bold=True)
    sr(p.add_run(), "永远留一份不参与训练的验证集；永远按 eval 曲线（加上任务指标）选 checkpoint——别被 train loss 骗了。", BODY_F, 12, INK)
    footer(slide, MOD, "第 4 章 · 训练实操与框架图鉴")

# ---------- p52 SFT → DPO → RFT ----------
def draw_ft_alignment_ladder(slide):
    bg(slide)
    eyebrow(slide, "第 5 章 · 偏好对齐  SFT → DPO → RFT")
    title(slide, "对齐三级：照做、分寸、对着分数练")
    label(slide, 0.55, 1.58, 12.2, "示范教不会「哪个更好」的分寸——按症状逐级升级，多数项目止步 SFT + 少量 DPO", SUB, 12, PP_ALIGN.LEFT, True)
    steps = [("SFT", "示范照做", "范文 · 教格式 / 语气", "天花板：教不会「两个都通顺、哪个更好」的分寸", TEAL, DTEAL, CARD),
             ("DPO", "好坏对纠偏", "批改 · chosen / rejected 对", "SFT 后「味道不对」的首选；成本≈再跑一次 SFT", ORANGE, ODARK, CORAL),
             ("RFT", "对着分数练", "判卷 · 打分器 grader", "结果可机器判分的窄任务（代码跑测试）；成败在 grader", GINK, GREEN, GBG)]
    W = 3.7; gap = 0.42; x0 = 0.7; y = 2.75; h = 2.15
    xs = []
    for i, (en, zh, way, note, tc, ec, fill) in enumerate(steps):
        x = x0 + i*(W+gap); xs.append(x)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, y+0.16, W-0.5, 0.75); p = tf.paragraphs[0]
        sr(p.add_run(), f"{i+1}. {en}", BODY_F, 16, tc, bold=True)
        sr(p.add_run(), "　" + zh, BODY_F, 12, INK, bold=True)
        node(slide, x+0.3, y+0.86, W-0.6, 0.52, way, "", "", fill=WHITE, edge=ec, zh_c=tc, zh_size=11)
        tf2 = tb(slide, x+0.28, y+1.48, W-0.56, 0.6); sr(tf2.paragraphs[0].add_run(), note, BODY_F, 10.5, SUB)
    for i in range(2):
        arrow(slide, xs[i]+W+0.02, y+h/2, xs[i+1]-0.02, y+h/2, color=SUB, w=1.9)
        label(slide, xs[i]+W-0.05, y+h/2-0.42, gap+0.1, "不够再升", SUB, 9, PP_ALIGN.CENTER, True)
    band(slide, 0.7, 5.5, 11.93, 1.32, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 5.62, 11.4, 1.15); p = tf.paragraphs[0]
    sr(p.add_run(), "现实比例：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "多数企业项目止步于 SFT（+ 少量 DPO 纠偏）；RFT 是尖刀任务的利器，不是标配。", BODY_F, 12, INK)
    para(tf, "别让客户以为每个项目都要走全流程——判分规则含糊、能被钻空子（reward hacking），就退回 DPO / SFT。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 5 章 · 偏好对齐落地")

# ---------- p72 部署两形态 ----------
def draw_ft_deploy(slide):
    bg(slide)
    eyebrow(slide, "第 7 章 · 部署形态  MERGE vs ADAPTER")
    title(slide, "部署两形态：合并权重 vs 挂 adapter")
    label(slide, 0.55, 1.58, 12.2, "单一稳定任务合并权重换运维简单；多任务 / 多租户挂 adapter 是标准答案", SUB, 12, PP_ALIGN.LEFT, True)
    # 左 合并
    band(slide, 0.6, 2.4, 6.0, 3.0, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.52, 5.6, "合并权重 —— W + BA 合成新模型", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 1.1, 3.15, 1.6, 0.72, "W + BA", "合并", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, en_c=PALE, zh_size=12.5, en_size=9)
    arrow(slide, 2.72, 3.51, 3.15, 3.51, color=TEAL, w=1.5)
    node(slide, 3.2, 3.15, 2.9, 0.72, "一个普通模型", "和平常无异", "", fill=WHITE, edge=TEAL, zh_c=DTEAL, en_c=SUB, zh_size=12.5, en_size=9)
    tf = tb(slide, 0.9, 4.15, 5.5, 1.15)
    para(tf, "· 部署简单，和普通模型完全一样", BODY_F, 11, SUB, first=True)
    para(tf, "· 一次一个任务，换任务要重发整模型", BODY_F, 11, SUB)
    para(tf, "· 适合：单一、稳定的任务，换运维简单", BODY_F, 11, DTEAL, bold=True)
    # 右 adapter
    band(slide, 6.75, 2.4, 6.0, 3.0, fill=CORAL, edge=ORANGE)
    label(slide, 7.0, 2.52, 5.6, "挂 adapter —— 底座共享，运行时加载", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 7.25, 3.12, 2.0, 1.5, "共享底座", "one base", "", fill=ORANGE, edge=ODARK, zh_c=OINK, en_c=OINK, zh_size=12.5, en_size=9)
    for j, name in enumerate(["部门A", "部门B", "部门C…"]):
        yy = 3.12 + j*0.52
        node(slide, 9.7, yy, 2.7, 0.44, name + " LoRA", "", "", fill=WHITE, edge=ODARK, zh_c=ODARK, zh_size=11)
        line(slide, 9.27, 3.87, 9.68, yy+0.22, color=ODARK, w=1.1)
    tf2 = tb(slide, 7.0, 4.72, 5.6, 0.7)
    para(tf2, "可热插拔、多套并存：5 部门 = 1 底座 + 5 LoRA，vLLM 多 LoRA 按请求路由，省一个数量级显存。", BODY_F, 11, SUB, first=True)
    band(slide, 0.6, 5.6, 12.03, 1.22, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.72, 11.5, 1.05); p = tf3.paragraphs[0]
    sr(p.add_run(), "默认姿势：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "多任务 / 多租户 → adapter 化部署是标准答案；单一稳定任务 → 合并权重换运维简单。", BODY_F, 12, INK)
    para(tf3, "（动态加载 adapter 有官方安全警告：仅限隔离受信环境；上线还要过 LLM-Inference 的量化 / 批处理 / SLO 关。）", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 7 章 · 评估与部署")

FIGS = [draw_ft_spectrum, draw_ft_knowledge_vs_behavior, draw_ft_lora, draw_ft_three_methods,
        draw_ft_loss_curves, draw_ft_alignment_ladder, draw_ft_deploy]

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in FIGS:
        fn(prs.slides.add_slide(B))
    out = "/Users/lijiaxiang/project/myAILearning/_maintenance/_ft_preview.pptx"
    prs.save(out); print(f"saved {len(FIGS)} ->", out)
