"""Solution-Patterns 讲义信息图（5 张）。复用 kb_draw + agent_figs helper。
嫁接：p5 后(两条轴)、p29 后(解决率bar)、p57 后(AI Coding双层)、p76 后(ChatBI bar)、p94 后(分诊树)。
SP 无页码，图系统页脚匹配。独立预览：python3 sp_figs.py"""
import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "场景方案讲义 · Solution-Patterns"
GINK = RGBColor(0x27, 0x50, 0x0A)
REDBG = RGBColor(0xFB, 0xE9, 0xE7)
REDINK = RGBColor(0x9B, 0x2C, 0x2C)

# ---------- p5 两条轴 ----------
def draw_sp_two_axes(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 两条轴  TECH vs SCENARIO")
    title(slide, "技术轴管深度，场景轴管对客")
    label(slide, 0.55, 1.58, 12.2, "同一批技术，两种组织方式——解决方案人日常活在场景轴", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 5.0, 2.35, 3.35, 0.72, "方案 = 场景 × 积木", "", "", fill=NAVY, edge=NAVY, zh_c=WHITE, zh_size=15)
    # 左 技术轴
    band(slide, 0.7, 3.35, 5.75, 2.0, fill=CARD, edge=TEAL)
    label(slide, 0.95, 3.47, 5.3, "技术轴 · 你已有的 15 个模块", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    bl = tb(slide, 0.98, 3.95, 5.3, 1.3)
    para(bl, "· 按技术组织：RAG 是什么、Agent 怎么编排、推理怎么提速", BODY_F, 11, SUB, first=True)
    para(bl, "· 回答「这个技术怎么用好」", BODY_F, 11, INK, bold=True)
    para(bl, "· 学习与内功的视角，越深越好", BODY_F, 11, SUB)
    # 右 场景轴
    band(slide, 6.85, 3.35, 5.75, 2.0, fill=CORAL, edge=ORANGE)
    label(slide, 7.1, 3.47, 5.3, "场景轴 · 本模块", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    br = tb(slide, 7.13, 3.95, 5.3, 1.3)
    para(br, "· 按客户场景：客服、搜索、营销、Coding、数字人", BODY_F, 11, SUB, first=True)
    para(br, "· 回答「这个需求怎么解」", BODY_F, 11, INK, bold=True)
    para(br, "· 对客与方案的视角，越快越好", BODY_F, 11, SUB)
    arrow(slide, 5.7, 3.09, 3.6, 3.33, color=TEAL, w=1.6)
    arrow(slide, 7.65, 3.09, 9.7, 3.33, color=ORANGE, w=1.6)
    band(slide, 0.7, 5.55, 11.93, 1.28, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 5.66, 11.4, 1.12); p = tf.paragraphs[0]
    sr(p.add_run(), "类比：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "15 个模块是乐高零件盒，本模块是拼装说明书——零件认得全，不等于拼得出船。", BODY_F, 11.5, INK)
    para(tf, "案例：客户不会问「你们 RAG 用什么重排序」，他问「我的客服能不能少雇十个人」——场景轴负责把后者翻译成前者。", BODY_F, 11, SUB)
    footer(slide, MOD, "从技术轴到场景轴")

def _barrow(slide, x, y, w, h, frac, lab, val, sub, barc, edgec, valc, labbold=True):
    band(slide, x, y, w, h, fill=CARD, edge=LINE)
    tf = tb(slide, x+0.15, y+0.06, 2.85, h-0.1); p = tf.paragraphs[0]
    sr(p.add_run(), lab, BODY_F, 11.5, INK, bold=labbold)
    if sub:
        para(tf, sub, BODY_F, 9, SUB)
    bx = x+3.05; bw = (w-3.05-1.5)*frac
    band(slide, bx, y+0.14, max(bw, 0.12), h-0.28, fill=barc, edge=edgec)
    vt = tb(slide, bx+max(bw, 0.12)+0.08, y+0.06, 1.5, h-0.1); q = vt.paragraphs[0]
    sr(q.add_run(), val, BODY_F, 13, valc, bold=True)

# ---------- p29 解决率真相 ----------
def draw_sp_resolution_truth(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 口径战场  SELF-REPORTED vs MEASURED")
    title(slide, "解决率的真相：自报与实测差多远")
    label(slide, 0.55, 1.58, 12.2, "厂商自报 51–80%，独立口径中位只有 41.2%——差距全在「口径」", SUB, 12, PP_ALIGN.LEFT, True)
    x0 = 0.7; W = 11.0; y = 2.35; h = 0.62; gap = 0.14
    rows = [(0.67, "Fin（Intercom）自报", "51–67%", "7000+ 客户两种口径", TEAL, DTEAL),
            (0.70, "Sierra 自报", "~70%", "WeightWatchers 案例", TEAL, DTEAL),
            (0.80, "Decagon 自报", "80%", "注意：是 deflection", ORANGE, ODARK),
            (0.412, "独立口径中位", "41.2%", "Zendesk 企业全量数据", REDINK, REDINK)]
    for i, (frac, lab, val, sub, bc, ec) in enumerate(rows):
        yy = y + i*(h+gap)
        _barrow(slide, x0, yy, W, h, frac, lab, val, sub, bc, ec, (REDINK if i == 3 else INK))
    label(slide, 7.8, y+3*(h+gap)+0.12, 4.8, "← 独立全量数据，几乎腰斩", REDINK, 10.5, PP_ALIGN.LEFT, False, True)
    band(slide, 0.7, 5.6, 11.93, 1.24, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 5.71, 11.4, 1.08); p = tf.paragraphs[0]
    sr(p.add_run(), "差距来自口径：", BODY_F, 11.5, ODARK, bold=True)
    sr(p.add_run(), "deflection（挡走）≠ resolution（真解决）· 演示数据 ≠ 生产数据 · 头部案例 ≠ 全量中位。", BODY_F, 11.5, INK)
    para(tf, "售前答法：「用您的历史工单抽样建评估集，POC 测出您自己的基线，验收线双方签字」——比任何数字都有说服力。", BODY_F, 11, DTEAL, bold=True)
    footer(slide, MOD, "智能客服")

# ---------- p57 AI Coding 双层 ----------
def draw_sp_ai_coding_layers(slide):
    bg(slide)
    eyebrow(slide, "第 6 章 · 双层格局  TWO LAYERS")
    title(slide, "补全管单人效率，agentic 管整任务")
    label(slide, 0.55, 1.58, 12.2, "「三选一」是伪命题——给客户配组合：IDE 层选一个 + agent 层选一个", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.35, 6.0, 2.95, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.47, 5.6, "IDE 补全 / 编辑层", DTEAL, 13, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 3.0, 5.3, 0.56, "逐行补全 · 局部改写", "Copilot / Cursor", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, en_c=PALE, zh_size=12, en_size=9.5)
    tf = tb(slide, 0.9, 3.72, 5.5, 1.4)
    para(tf, "· 价值：单人打字与查找效率", BODY_F, 11, INK, bold=True, first=True)
    para(tf, "· 按人头月费，最容易起量", BODY_F, 11, SUB)
    para(tf, "· 适合：普惠铺开、先让全员用起来", BODY_F, 11, SUB)
    band(slide, 6.75, 2.35, 6.0, 2.95, fill=CORAL, edge=ORANGE)
    label(slide, 7.0, 2.47, 5.6, "agentic 编码层", ODARK, 13, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.0, 5.3, 0.56, "整任务委托 · 端到端修 bug", "Claude Code 类终端 agent", "", fill=ORANGE, edge=ODARK, zh_c=OINK, en_c=OINK, zh_size=12, en_size=9.5)
    tf2 = tb(slide, 7.05, 3.72, 5.5, 1.4)
    para(tf2, "· 价值：把「活」交出去，不是「补字」", BODY_F, 11, INK, bold=True, first=True)
    para(tf2, "· 按用量计费", BODY_F, 11, SUB)
    para(tf2, "· 适合：跨文件改造、重复性工程任务", BODY_F, 11, SUB)
    band(slide, 0.6, 5.5, 12.03, 1.33, fill=GBG, edge=GREEN)
    tf3 = tb(slide, 0.9, 5.62, 11.5, 1.15); p = tf3.paragraphs[0]
    sr(p.add_run(), "数据锚点：", BODY_F, 11.5, GINK, bold=True)
    sr(p.add_run(), "70% 工程师同时用 2–4 个工具；主流组合 = IDE 层（Cursor / Copilot）+ agent 层（Claude Code 类）。", BODY_F, 11.5, INK)
    para(tf3, "要点：给客户配组合而不是站队——网关统一管预算，工具迭代不伤架构。", BODY_F, 11, SUB)
    footer(slide, MOD, "AI Coding")

# ---------- p76 ChatBI 准确率 ----------
def draw_sp_chatbi_accuracy(slide):
    bg(slide)
    eyebrow(slide, "第 8 章 · 口径战场  SEMANTIC LAYER JUMP")
    title(slide, "准确率的真相：差距不在模型，在语义层")
    label(slide, 0.55, 1.58, 12.2, "把「你们准确率多少」转化为「语义层建到多厚」——口径治理问题，就赢了", SUB, 12, PP_ALIGN.LEFT, True)
    steps = [(0.16, "裸 Text-to-SQL", "10–21%", "真实企业库实测（乱 schema、黑话口径）", RGBColor(0x88,0x87,0x80), RGBColor(0x66,0x66,0x60)),
             (0.72, "加语义模型后", "51% → 90%+", "Snowflake 内部 150 题基准", TEAL, DTEAL),
             (0.99, "语义层基准上限", "98–100%", "dbt 语义层基准口径", GREEN, GINK)]
    x0 = 0.7; W = 11.0; y = 2.45; h = 0.82; gap = 0.24
    for i, (frac, lab, val, sub, bc, ec) in enumerate(steps):
        yy = y + i*(h+gap)
        _barrow(slide, x0, yy, W, h, frac, lab, val, sub, bc, ec, ec)
        if i < 2:
            arrow(slide, x0+0.5, yy+h+0.02, x0+0.5, yy+h+gap-0.02, color=ORANGE, w=1.8)
    label(slide, 1.35, y+h+0.02, 4.5, "加语义层 → 跃迁", ODARK, 10, PP_ALIGN.LEFT, False, True)
    label(slide, 1.35, y+2*(h+gap)-0.22, 4.5, "再加治理 → 逼近上限", GINK, 10, PP_ALIGN.LEFT, False, True)
    band(slide, 0.7, 5.72, 11.93, 1.1, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 5.82, 11.4, 0.95); p = tf.paragraphs[0]
    sr(p.add_run(), "这页是本章的灵魂：", BODY_F, 11.5, ODARK, bold=True)
    sr(p.add_run(), "学术基准（Spider 类）裸 SQL 能到五六成，但真实企业库只有一两成——给客户一定讲后者，再用语义层的跃迁数据卖治理。", BODY_F, 11.5, INK)
    footer(slide, MOD, "ChatBI 与数据分析")

# ---------- p94 分诊决策树 ----------
def draw_sp_triage_tree(slide):
    bg(slide)
    eyebrow(slide, "第 10 章 · 分诊决策树  TRIAGE")
    title(slide, "客户开口第一句，先分诊")
    label(slide, 0.55, 1.55, 12.2, "分诊错了后面全错——宁可多问三个问题再分诊，也别顺着第一句话直接进方案", SUB, 12, PP_ALIGN.LEFT, True)
    rows = [("「想让 AI 答问题 / 少雇客服」", "对外 → 第 3 章客服 ｜ 对内 → 第 4 章知识库"),
            ("「想让 AI 做内容 / 做视频」", "第 5 章内容生成；要「人形」→ 第 7 章数字人"),
            ("「研发想提效」", "第 6 章 AI Coding"),
            ("「想用自然语言问数据 / 智能报表」", "第 8 章 ChatBI"),
            ("「会议纪要 / 办公提效」", "第 9 章会议助手"),
            ("「都想要，先做哪个」", "回第 1 章：按数据就绪度 × 容错成本排矩阵"),
            ("「你们有没有平台」", "Agent 模块第 7 章低代码平台选型")]
    y0 = 2.28; h = 0.5; gap = 0.058
    label(slide, 0.7, y0-0.3, 4.6, "客户说的话", ODARK, 11, PP_ALIGN.LEFT, False, True)
    label(slide, 5.75, y0-0.3, 6.9, "分诊去向", DTEAL, 11, PP_ALIGN.LEFT, False, True)
    for i, (say, to) in enumerate(rows):
        yy = y0 + i*(h+gap)
        band(slide, 0.7, yy, 4.85, h, fill=CORAL, edge=ORANGE)
        tf = tb(slide, 0.9, yy+0.09, 4.5, 0.4); sr(tf.paragraphs[0].add_run(), say, BODY_F, 11, ODARK, bold=True)
        arrow(slide, 5.58, yy+h/2, 5.98, yy+h/2, color=SUB, w=1.5)
        band(slide, 6.02, yy, 6.6, h, fill=CARD, edge=TEAL)
        tf2 = tb(slide, 6.24, yy+0.09, 6.3, 0.4); sr(tf2.paragraphs[0].add_run(), to, BODY_F, 11, DTEAL, bold=True)
    footer(slide, MOD, "售前速查")

import math
# ---------- p8 解决方案人五步循环 ----------
def draw_sp_work_loop(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 工作循环  THE LOOP")
    title(slide, "解决方案人的五步循环")
    label(slide, 0.55, 1.55, 12.2, "循环不是瀑布——POC 失败回到需求发现重挑场景，是常态不是事故", SUB, 12, PP_ALIGN.LEFT, True)
    steps = [("① 需求发现", "听痛点、挖场景"), ("② 参考架构", "拼积木、画图"), ("③ POC", "小切口验证"),
             ("④ 报价方案书", "三本账"), ("⑤ 落地运营", "评估、迭代")]
    cx, cy = 4.05, 4.2; rx, ry = 2.85, 1.5; W, H = 2.15, 0.82
    cen = []
    for i in range(5):
        ang = math.radians(-90 + i*72)
        cen.append((cx + rx*math.cos(ang), cy + ry*math.sin(ang)))
    for i in range(5):
        x1, y1 = cen[i]; x2, y2 = cen[(i+1) % 5]
        dx, dy = x2-x1, y2-y1; d = math.hypot(dx, dy); ux, uy = dx/d, dy/d
        col = ORANGE if i == 4 else TEAL
        arrow(slide, x1+ux*1.15, y1+uy*0.52, x2-ux*1.15, y2-uy*0.52, color=col, w=1.5)
    for i, (zh, d) in enumerate(steps):
        x, y = cen[i]
        node(slide, x-W/2, y-H/2, W, H, zh, "", d, fill=(TEAL if i == 2 else CARD), edge=(DTEAL if i == 2 else TEAL),
             zh_c=(WHITE if i == 2 else DTEAL), sub_c=(PALE if i == 2 else SUB), zh_size=12, desc_size=9)
    label(slide, cx-1.5, cy-0.32, 3.0, "回路", ODARK, 12, PP_ALIGN.CENTER, False, True)
    label(slide, cx-1.7, cy+0.05, 3.4, "POC 失败 → 回需求发现", ODARK, 9.5, PP_ALIGN.CENTER, True)
    band(slide, 8.15, 2.5, 4.5, 3.05, fill=CARD, edge=LINE)
    tf = tb(slide, 8.4, 2.63, 4.05, 2.85)
    para(tf, "角色 · 像全科医生", BODY_F, 12.5, DTEAL, bold=True, first=True)
    para(tf, "· 先分诊、开检查单（POC 指标）", BODY_F, 11, SUB)
    para(tf, "· 深度问题转回 15 个技术模块（转专科）", BODY_F, 11, SUB)
    para(tf, "· 价值在「翻译」：业务语言 ↔ 技术语言", BODY_F, 11, DTEAL, bold=True)
    para(tf, "每步对应章节", BODY_F, 12, ODARK, bold=True)
    para(tf, "需求发现 → 第 1、8 章 · 架构与三本账 → 第 2 章 · POC 验收 → 第 2 章 + Evaluation 模块", BODY_F, 11, SUB)
    footer(slide, MOD, "从技术轴到场景轴")

# ---------- p42 权限感知检索 ----------
def draw_sp_permission_search(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · 权限命门  PERMISSION-AWARE")
    title(slide, "权限感知检索：存在性也是机密")
    label(slide, 0.55, 1.58, 12.2, "搜索越好、泄密越快——权限做不对，AI 搜索就是最高效的泄密器", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 3.0, fill=GBG, edge=GREEN)
    label(slide, 0.85, 2.52, 5.6, "怎么做对", GINK, 12.5, PP_ALIGN.LEFT, False, True)
    tf = tb(slide, 0.9, 3.05, 5.5, 2.2)
    para(tf, "· ACL 从源系统继承，不另建一套权限", BODY_F, 11.5, INK, bold=True, first=True)
    para(tf, "· 查询时实时过滤，不是索引时过滤一次", BODY_F, 11.5, INK, bold=True)
    para(tf, "· 无权限时，连「有这份文档」都不暴露", BODY_F, 11.5, GINK, bold=True)
    para(tf, "类比：图书馆内部书库——馆员只把你有借阅证的书拿给你，且不会告诉你「有一本你不能看的书」。", BODY_F, 11, SUB)
    band(slide, 6.75, 2.4, 6.0, 3.0, fill=REDBG, edge=REDINK)
    label(slide, 7.0, 2.52, 5.6, "做错的后果", REDINK, 12.5, PP_ALIGN.LEFT, False, True)
    tf2 = tb(slide, 7.05, 3.05, 5.5, 2.2)
    para(tf2, "· 员工搜到薪酬表 / 裁员计划 → 安全事故", BODY_F, 11.5, REDINK, bold=True, first=True)
    para(tf2, "· 权限缓存过期 → 离职员工还能搜", BODY_F, 11.5, INK)
    para(tf2, "· 搜索越好，泄密越快 → 放大器效应", BODY_F, 11.5, REDINK, bold=True)
    para(tf2, "根因：索引时过滤一次、或另建权限，都会漏——权限必须跟着「查询这一刻的身份」走。", BODY_F, 11, SUB)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.67, 11.5, 1.1); p = tf3.paragraphs[0]
    sr(p.add_run(), "POC 验收必须包含「越权测试集」：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "拿无权限账号问 20 条敏感问题，一条都不能漏——这是安全团队签字放行的前提。", BODY_F, 12, INK)
    para(tf3, "对客一句话：权限不是「搜索的一个功能」，是「搜索能不能上线」的开关。", BODY_F, 11, SUB)
    footer(slide, MOD, "企业知识库与 AI 搜索")

FIGS = [draw_sp_two_axes, draw_sp_resolution_truth, draw_sp_ai_coding_layers,
        draw_sp_chatbi_accuracy, draw_sp_triage_tree,
        draw_sp_work_loop, draw_sp_permission_search]
NEW = [draw_sp_work_loop, draw_sp_permission_search]

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in FIGS:
        fn(prs.slides.add_slide(B))
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_sp_preview.pptx")
    prs.save(out); print(f"saved {len(FIGS)} ->", out)
