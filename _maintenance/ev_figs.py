import os
"""Evaluation 讲义信息图（4 张）。复用 kb_draw + agent_figs helper。
嫁接：p8 后（评估三层）、p38 后（三种判法）、p60 后（RAG 三角）、p78 后（质量飞轮）。
注：Evaluation 是页码册，图沿用图系统统一页脚（无页码），与首批 PE/LLM 一致。
独立预览：python3 ev_figs.py"""
import sys, math
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "Evaluation 讲义 · 模型评估"
GINK = RGBColor(0x27, 0x50, 0x0A)

# ---------- p8 评估三层 ----------
def draw_ev_three_layers(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 三层分工  THREE LAYERS")
    title(slide, "评估三层：基准 / 应用评估 / 生产监控")
    label(slide, 0.55, 1.58, 12.2, "三层各管一段、互相不能替代——像招人不会只看高考分", SUB, 12, PP_ALIGN.LEFT, True)
    rows = [("基准测试", "Benchmark", "高考", "量模型通用能力：一套固定考卷横向比模型", TEAL, DTEAL, CARD),
            ("应用评估", "Application Evals", "面试", "量「你的场景 + 你的数据」：岗位到底匹不匹配", ORANGE, ODARK, CORAL),
            ("生产监控", "Production Monitoring", "试用期", "量真实流量下的持续质量：上线后还稳不稳", TEAL, DTEAL, CARD)]
    y = 2.35; h = 0.98; gap = 0.16
    for i, (zh, en, ana, desc, ec, tc, fill) in enumerate(rows):
        yy = y + i*(h+gap)
        band(slide, 0.7, yy, 11.93, h, fill=fill, edge=ec)
        node(slide, 0.95, yy+0.16, 1.75, 0.66, ana, "", "", fill=ec, edge=tc, zh_c=WHITE, zh_size=15)
        tf = tb(slide, 3.0, yy+0.12, 9.4, 0.78); p = tf.paragraphs[0]
        sr(p.add_run(), zh + "　", BODY_F, 15, tc, bold=True)
        sr(p.add_run(), en, BODY_F, 11, SUB, italic=True)
        para(tf, desc, BODY_F, 11.5, INK)
    band(slide, 0.7, 5.75, 11.93, 1.1, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 5.87, 11.4, 0.95); p = tf.paragraphs[0]
    sr(p.add_run(), "售前口径：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "榜单筛入围 → 应用评估定选型 → 生产监控保长期。", BODY_F, 12, INK)
    para(tf, "案例：某客户按榜单选第一名，POC 一跑发现自家行业术语不如第二名——榜单是「高考分」，客户场景是「岗位面试」，两码事。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 1 章 · 为什么评估这么难")

# ---------- p38 三种判法 ----------
def draw_ev_three_judging(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · LLM-as-a-Judge  THREE WAYS")
    title(slide, "三种判法：打分、比较、按细则批改")
    label(slide, 0.55, 1.58, 12.2, "LLM 更擅长「判别」——比较和逐条核对，远比裸的绝对打分稳定", SUB, 12, PP_ALIGN.LEFT, True)
    cards = [("单点打分", "Pointwise", "给单个输出打 1–5 分", "最直觉，但绝对分漂移大、难对齐", CARD, TEAL, DTEAL, "输出 → ★ 3/5"),
             ("成对比较", "Pairwise", "两个输出选更好", "最稳：谁都答得准「哪个更好」", GBG, GREEN, GINK, "A  vs  B → 选 A"),
             ("评分细则", "Rubric", "拆成可核对条款逐条判", "最可控：把「好」拆成 ✓ 清单", CORAL, ORANGE, ODARK, "✓ 准确 ✓ 完整 ✓ 安全")]
    W = 3.95; gap = 0.14; x0 = 0.6; y = 2.4; h = 2.65
    for i, (zh, en, what, note, fill, ec, tc, demo) in enumerate(cards):
        x = x0 + i*(W+gap)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, y+0.16, W-0.5, 0.66); p = tf.paragraphs[0]
        sr(p.add_run(), f"{i+1}. {zh}　", BODY_F, 14.5, tc, bold=True)
        sr(p.add_run(), en, BODY_F, 10.5, SUB, italic=True)
        node(slide, x+0.35, y+0.86, W-0.7, 0.56, demo, "", "", fill=WHITE, edge=ec, zh_c=tc, zh_size=12)
        tf2 = tb(slide, x+0.28, y+1.56, W-0.56, 1.0)
        para(tf2, what, BODY_F, 11.5, INK, bold=True, first=True)
        para(tf2, note, BODY_F, 11, SUB)
    band(slide, 0.6, 5.35, 12.03, 1.5, fill=RGBColor(0xF7, 0xFB, 0xFC), edge=LINE)
    tf = tb(slide, 0.9, 5.48, 11.5, 1.35); p = tf.paragraphs[0]
    sr(p.add_run(), "类比：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "问「这碗面打几分」很难答稳，问「这两碗哪碗好吃」谁都答得准——比较比打分容易，机器也一样。", BODY_F, 12, INK)
    para(tf, "选题型：回归监控用细则逐条判 · 版本对比用成对比较 · 绝对分数只作趋势参考。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 4 章 · LLM-as-a-Judge")

# ---------- p60 RAG 评估三角 ----------
def draw_ev_rag_triangle(slide):
    bg(slide)
    eyebrow(slide, "第 6 章 · 场景验收  RAG TRIANGLE")
    title(slide, "RAG 评估三角：开卷考试怎么判")
    label(slide, 0.55, 1.58, 12.2, "Ragas 定义的三角——检索与生成分开评，才知道「谁错了」", SUB, 12, PP_ALIGN.LEFT, True)
    ax, ay = 3.5, 2.5    # 忠实度（顶·生成端）
    bx, by = 1.4, 5.3    # 召回率（左下·检索端）
    cx, cy = 5.6, 5.3    # 精确率（右下·检索端）
    line(slide, ax, ay+0.62, bx+0.55, by-0.1, color=TEAL, w=2.0)
    line(slide, ax, ay+0.62, cx-0.55, cy-0.1, color=TEAL, w=2.0)
    line(slide, bx+0.7, by+0.25, cx-0.7, cy+0.25, color=TEAL, w=2.0)
    node(slide, 2.55, 2.5, 2.0, 0.72, "忠实度", "Faithfulness", "", fill=ORANGE, edge=ODARK, zh_c=OINK, en_c=OINK, zh_size=13, en_size=9)
    node(slide, 0.6, 5.28, 2.0, 0.72, "召回率", "Context Recall", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, en_c=PALE, zh_size=13, en_size=8.5)
    node(slide, 4.55, 5.28, 2.0, 0.72, "精确率", "Context Precision", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, en_c=PALE, zh_size=13, en_size=8.5)
    label(slide, 2.55, 3.35, 2.0, "生成端", ODARK, 10, PP_ALIGN.CENTER, True, True)
    label(slide, 0.6, 6.06, 2.0, "检索端", DTEAL, 10, PP_ALIGN.CENTER, True, True)
    label(slide, 4.55, 6.06, 2.0, "检索端", DTEAL, 10, PP_ALIGN.CENTER, True, True)
    label(slide, 2.35, 4.35, 2.5, "三关各自把守", SUB, 11, PP_ALIGN.CENTER, True)
    # 右侧：类比 + 诊断
    band(slide, 7.05, 2.4, 5.85, 1.75, fill=CARD, edge=LINE)
    tf = tb(slide, 7.3, 2.52, 5.4, 1.55)
    para(tf, "开卷考试三件事：", BODY_F, 12, DTEAL, bold=True, first=True)
    para(tf, "翻到对的页（召回）· 别抄不相干段落（精确率）· 答案只依据书上写的（忠实度）。", BODY_F, 11, SUB)
    band(slide, 7.05, 4.35, 5.85, 2.0, fill=CORAL, edge=ORANGE)
    tf2 = tb(slide, 7.3, 4.47, 5.4, 1.8)
    para(tf2, "用三角定位问题：", BODY_F, 12, ODARK, bold=True, first=True)
    para(tf2, "忠实度低 → 生成端在编造：改 prompt / 换模型", BODY_F, 11, INK)
    para(tf2, "召回率低 → 检索端没找到：修 chunking / embedding / 混合检索", BODY_F, 11, INK)
    footer(slide, MOD, "第 6 章 · 场景验收")

# ---------- p78 质量飞轮 ----------
def draw_ev_flywheel(slide):
    bg(slide)
    eyebrow(slide, "第 7 章 · 工具链与闭环  QUALITY FLYWHEEL")
    title(slide, "质量飞轮：评估驱动的迭代闭环")
    label(slide, 0.55, 1.55, 12.2, "每转一圈，评估集厚一分、系统稳一分——上线不是结束，是飞轮起点", SUB, 12, PP_ALIGN.LEFT, True)
    steps = ["线上坏例", "错误分析归类", "进评估集", "针对性修\nprompt/检索/微调", "过回归门禁", "上线", "监控采样"]
    cx, cy = 4.05, 4.15; rx, ry = 3.15, 1.62
    W, H = 1.85, 0.72
    cen = []
    for i, s in enumerate(steps):
        ang = math.radians(-90 + i*360/len(steps))
        x = cx + rx*math.cos(ang); y = cy + ry*math.sin(ang)
        cen.append((x, y))
    # 箭头（沿环，顺时针）
    for i in range(len(steps)):
        x1, y1 = cen[i]; x2, y2 = cen[(i+1) % len(steps)]
        dx, dy = x2-x1, y2-y1; d = math.hypot(dx, dy); ux, uy = dx/d, dy/d
        col = ORANGE if i == len(steps)-1 else TEAL
        arrow(slide, x1+ux*1.02, y1+uy*0.62, x2-ux*1.02, y2-uy*0.62, color=col, w=1.5)
    # 节点
    for i, s in enumerate(steps):
        x, y = cen[i]
        two = "\n" in s
        zh = s.split("\n")[0]
        node(slide, x-W/2, y-H/2, W, H, zh, (s.split("\n")[1] if two else ""), "",
             fill=(TEAL if i in (0,) else CARD), edge=(DTEAL if i == 0 else TEAL),
             zh_c=(WHITE if i == 0 else DTEAL), en_c=(PALE if i == 0 else SUB), zh_size=11.5, en_size=8.5)
    # 中心
    label(slide, cx-1.4, cy-0.34, 2.8, "质量飞轮", DTEAL, 16, PP_ALIGN.CENTER, False, True)
    label(slide, cx-1.6, cy+0.08, 3.2, "转起来越转越省力", SUB, 10.5, PP_ALIGN.CENTER, True)
    # 右侧收束（分块定位，不用小字空行）
    band(slide, 8.5, 2.5, 4.4, 3.5, fill=CARD, edge=LINE)
    t1 = tb(slide, 8.75, 2.66, 3.95, 1.15)
    para(t1, "类比 · 飞轮", BODY_F, 12.5, DTEAL, bold=True, first=True)
    para(t1, "起初推得费劲（建集、校判官），转起来后每圈更省力——数据飞轮讲数据，这个是「质量飞轮」。", BODY_F, 11, SUB)
    t2 = tb(slide, 8.75, 3.92, 3.95, 1.15)
    para(t2, "对客运营口径", BODY_F, 12.5, ODARK, bold=True, first=True)
    para(t2, "建议客户设每周半天「评估例会」一起看坏例——整个体系里回报率最高的投入。", BODY_F, 11, SUB)
    t3 = tb(slide, 8.75, 5.2, 3.95, 0.7)
    para(t3, "客户买的不是一次性交付，是一套会自我改进的质量体系。", BODY_F, 11, GINK, bold=True, first=True)
    footer(slide, MOD, "第 7 章 · 工具链与闭环")

# ---------- p6 非确定性：传统 vs LLM ----------
def draw_ev_nondeterminism(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 非确定性  WHY TESTING FAILS")
    title(slide, "非确定性：传统软件测试为什么失灵")
    label(slide, 0.55, 1.58, 12.2, "LLM 评估天生是「改作文」，不是「对答案」——拿单元测试思路硬套注定失败", SUB, 12, PP_ALIGN.LEFT, True)
    # 左 传统软件
    band(slide, 0.6, 2.35, 6.0, 2.5, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.47, 5.6, "传统软件 —— 数学题阅卷", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    for i, (zh, x) in enumerate([("同一输入", 0.95), ("同一输出", 3.55)]):
        node(slide, x, 3.0, 2.2, 0.62, zh, "", "", fill=WHITE, edge=TEAL, zh_c=DTEAL, zh_size=12.5)
        if i == 0:
            arrow(slide, 3.18, 3.31, 3.5, 3.31, color=TEAL, w=1.6)
    tf = tb(slide, 0.9, 3.85, 5.5, 0.9)
    para(tf, "写个 assert 就能测：对答案就行。", BODY_F, 11.5, INK, bold=True, first=True)
    para(tf, "确定性：同一问题永远同一答案，有唯一正确解。", BODY_F, 11, SUB)
    # 右 LLM
    band(slide, 6.75, 2.35, 6.0, 2.5, fill=CORAL, edge=ORANGE)
    label(slide, 7.0, 2.47, 5.6, "LLM —— 作文阅卷", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.0, 2.0, 0.62, "同一问题", "", "", fill=WHITE, edge=ORANGE, zh_c=ODARK, zh_size=12.5)
    node(slide, 9.65, 2.72, 2.8, 0.52, "答案 A：措辞甲", "", "", fill=WHITE, edge=ORANGE, zh_c=ODARK, zh_size=11)
    node(slide, 9.65, 3.4, 2.8, 0.52, "答案 B：措辞乙", "", "", fill=WHITE, edge=ORANGE, zh_c=ODARK, zh_size=11)
    arrow(slide, 9.12, 3.18, 9.6, 2.98, color=ODARK, w=1.4)
    arrow(slide, 9.12, 3.34, 9.6, 3.6, color=ODARK, w=1.4)
    label(slide, 7.1, 3.78, 2.4, "概率采样 Sampling", SUB, 9.5, PP_ALIGN.LEFT, True)
    tf2 = tb(slide, 7.0, 4.05, 5.6, 0.7)
    para(tf2, "多数任务没有唯一正确答案：只能按维度（切题 / 结构 / 文采）打分。", BODY_F, 11, SUB, first=True)
    # 隐性回归 + 收束
    band(slide, 0.6, 5.05, 12.03, 1.8, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.18, 11.5, 1.6); p = tf3.paragraphs[0]
    sr(p.add_run(), "最常见的事故 · 隐性回归：", BODY_F, 12, ODARK, bold=True)
    sr(p.add_run(), "改一版 prompt，盯着的问题答好了，原本正常的三个场景却悄悄答坏——没有回归评估根本发现不了。", BODY_F, 12, INK)
    para(tf3, "所以：LLM 评估量的不是「对不对」，而是「好不好 + 稳不稳」——需要一套专门的评估工程。", BODY_F, 11.5, DTEAL, bold=True)
    footer(slide, MOD, "第 1 章 · 为什么评估这么难")

# ---------- p42 判官三大偏差 ----------
def draw_ev_judge_biases(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · LLM-as-a-Judge 深潜  THREE BIASES")
    title(slide, "判官的三大偏差：都不是故意，但都真实存在")
    label(slide, 0.55, 1.58, 12.2, "评审客户判官方案的第一批检查点——裁判也会被带偏", SUB, 12, PP_ALIGN.LEFT, True)
    cards = [("位置偏差", "Position Bias", "先说的占便宜", "换序判决漂移 > 10%", "A-B 与 B-A 各跑一次，一致才计数（换序复核）"),
             ("冗长偏差", "Verbosity Bias", "长的占便宜", "改长不加信息，90%+ 偏好长版", "细则写明简洁条款 + 要求引用证据"),
             ("自我偏好", "Self-Preference", "像自己的占便宜", "给「更像自己写的」打高分", "判官与被评模型错开厂家、多判官交叉")]
    W = 3.95; gap = 0.14; x0 = 0.6; y = 2.4; h = 2.95
    for i, (zh, en, one, evid, fix) in enumerate(cards):
        x = x0 + i*(W+gap)
        band(slide, x, y, W, h, fill=CORAL, edge=ORANGE)
        tf = tb(slide, x+0.25, y+0.16, W-0.5, 0.7); p = tf.paragraphs[0]
        sr(p.add_run(), f"{i+1}. {zh}", BODY_F, 14.5, ODARK, bold=True)
        para(tf, en, BODY_F, 10, SUB, italic=True)
        node(slide, x+0.3, y+0.92, W-0.6, 0.56, one, "", "", fill=WHITE, edge=ORANGE, zh_c=ODARK, zh_size=12)
        eb = tb(slide, x+0.28, y+1.58, W-0.56, 0.5); q = eb.paragraphs[0]
        sr(q.add_run(), "证据：", BODY_F, 10.5, SUB, bold=True); sr(q.add_run(), evid, BODY_F, 10.5, SUB)
        fbx = band(slide, x+0.28, y+2.12, W-0.56, 0.68, fill=GBG, edge=GREEN)
        fb = tb(slide, x+0.4, y+2.2, W-0.78, 0.55); r = fb.paragraphs[0]
        sr(r.add_run(), "缓解：", BODY_F, 10.5, GINK, bold=True); sr(r.add_run(), fix, BODY_F, 10.5, RGBColor(0x33, 0x44, 0x28))
    band(slide, 0.6, 5.62, 12.03, 1.2, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.74, 11.5, 1.05); p = tf.paragraphs[0]
    sr(p.add_run(), "一句话：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "「用哪家模型当判官」本身就是评估设计决策——裁判和运动员同门，是客户方案里最常见的坑。", BODY_F, 12, INK)
    footer(slide, MOD, "第 4 章 · LLM-as-a-Judge")

# ---------- p83 四道门 ----------
def draw_ev_four_gates(slide):
    bg(slide)
    eyebrow(slide, "第 7 章 · 发布门禁  FOUR GATES")
    title(slide, "四道门：把评测结果变成发布决策")
    label(slide, 0.55, 1.58, 12.2, "像机场安检四道闸——上一道过了不等于下一道免检，更不能用「平均分」混过", SUB, 12, PP_ALIGN.LEFT, True)
    gates = [("① 契约门", "Contract", "schema / 工具 / 权限", "零关键失败，挂了阻断构建", TEAL, DTEAL, CARD),
             ("② 质量门", "Quality", "数据集 + 评分器", "基线与分片达标", TEAL, DTEAL, CARD),
             ("③ 风险门", "Risk", "安全 / 红队 / 隐私", "高风险清零或例外签批", ORANGE, ODARK, CORAL),
             ("④ 生产门", "Production", "金丝雀 SLI / 业务结果", "预算内持续达标", GREEN, GINK, GBG)]
    W = 2.86; gap = 0.16; x0 = 0.6; y = 2.5; h = 2.5
    xs = []
    for i, (zh, en, what, cond, ec, tc, fill) in enumerate(gates):
        x = x0 + i*(W+gap); xs.append(x)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.22, y+0.18, W-0.44, 0.7); p = tf.paragraphs[0]
        sr(p.add_run(), zh, BODY_F, 15, tc, bold=True)
        para(tf, en, BODY_F, 10, SUB, italic=True)
        node(slide, x+0.26, y+0.98, W-0.52, 0.56, what, "", "", fill=WHITE, edge=ec, zh_c=tc, zh_size=10.5)
        cf = tb(slide, x+0.26, y+1.66, W-0.52, 0.75); q = cf.paragraphs[0]
        sr(q.add_run(), "通过条件：", BODY_F, 10, tc, bold=True)
        para(cf, cond, BODY_F, 10.5, INK)
        if i < 3:
            arrow(slide, x+W+0.01, y+h/2, x+W+gap-0.01, y+h/2, color=SUB, w=1.8)
    band(slide, 0.6, 5.42, 12.03, 1.4, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.54, 11.5, 1.25); p = tf.paragraphs[0]
    sr(p.add_run(), "两条铁规矩：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "① 不能用下游平均分覆盖上游硬失败（比如高危注入用例）；", BODY_F, 12, INK)
    para(tf, "② 数据集与评分器的版本是门禁的一部分——换了 Judge 必须重新校准、重建可比基线，否则「过门」本身就不可信。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 7 章 · 工具链与闭环")

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    figs = [draw_ev_three_layers, draw_ev_three_judging, draw_ev_rag_triangle, draw_ev_flywheel,
            draw_ev_nondeterminism, draw_ev_judge_biases, draw_ev_four_gates]
    for fn in figs:
        fn(prs.slides.add_slide(B))
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_ev_preview.pptx")
    prs.save(out); print(f"saved {len(figs)} ->", out)
