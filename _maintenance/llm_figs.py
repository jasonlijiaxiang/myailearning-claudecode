"""LLM 讲义信息图（第一批 4 张）。复用 kb_draw + agent_figs 的 helper。"""
import sys
sys.path.insert(0, "/Users/lijiaxiang/project/myAILearning/_maintenance")
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, BODY_F, TITLE_F)
from agent_figs import line

MOD = "LLM 讲义 · 大模型原理"

# ---------- p22 多头注意力 ----------
def draw_llm_mha(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 多头注意力  MULTI-HEAD ATTENTION")
    title(slide, "多头注意力：多位专家各看一面，再拼成会诊结论")
    label(slide, 0.55, 1.58, 12.2, "把向量空间切成 h 份，每份独立做一遍注意力（一个头），结果拼接后再投影融合", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 0.7, 3.15, 1.8, 1.0, "输入向量", "每个 token", "", fill=CARD, edge=SUB, zh_c=INK, en_c=SUB, zh_size=13, en_size=9.5, desc_size=9)
    heads = [("语法头", "盯主谓一致", True), ("指代头", "盯「它」指谁", True), ("语气头", "看情绪走向", True), ("… 更多头", "各看一面", False)]
    hx = 4.1; hw = 2.0; hy = [2.3, 3.25, 4.2, 5.15]
    for k, (zh, d, solid) in enumerate(heads):
        node(slide, hx, hy[k], hw, 0.8, zh, "Head", d, fill=(CARD if solid else NEARW), edge=(TEAL if solid else LINE), zh_c=(DTEAL if solid else SUB), en_c=TEAL, sub_c=SUB, zh_size=12.5, en_size=9, desc_size=9)
        line(slide, 2.52, 3.65, hx-0.02, hy[k]+0.4, color=TEAL, w=1.1)
    node(slide, 7.7, 3.15, 2.2, 1.0, "拼接 + 投影", "concat", "融合成一个结论", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=WHITE, zh_size=13, en_size=9.5, desc_size=9)
    for k in range(4):
        line(slide, hx+hw+0.02, hy[k]+0.4, 7.68, 3.65, color=TEAL, w=1.1)
    node(slide, 10.55, 3.15, 2.05, 1.0, "输出", "融合表示", "", fill=CORAL, edge=ORANGE, zh_c=ODARK, en_c=ODARK, zh_size=13, en_size=9.5, desc_size=9)
    arrow(slide, 9.92, 3.65, 10.53, 3.65, color=ORANGE, w=1.75)
    band(slide, 0.7, 6.15, 11.93, 0.75, fill=NEARW, edge=NEARW)
    tf = tb(slide, 0.75, 6.13, 12.0, 0.72)
    para(tf, "像专家会诊：每人独立写一份小报告，最后拼成会诊结论——一个人看一遍，视角必然单一。GPT-2 用 12 头、Llama-70B 用 64 头。", BODY_F, 11.5, DTEAL, bold=True, first=True)
    footer(slide, MOD, "第 2 章 · 多头注意力")

# ---------- p33 堆叠 N 层 ----------
def draw_llm_layers(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 堆叠 N 层  STACKING")
    title(slide, "堆叠 N 层：同一个块过 N 遍，理解逐层加深")
    label(slide, 0.55, 1.58, 12.2, "一个 Transformer 块 = 注意力 + FFN（各带残差）；所谓大模型，就是同一个块堆 N 次", SUB, 12, PP_ALIGN.LEFT, True)
    layers = [("高层", "任务级抽象", "情感、推理线索", "终审看战略取向", ORANGE, ODARK, CORAL),
              ("中层", "语义", "指代消解、实体关系", "复审看逻辑漏洞", TEAL, DTEAL, CARD),
              ("低层", "句法", "词性搭配、谁修饰谁", "初审看格式错字", TEAL, DTEAL, CARD)]
    y0 = 2.4; h = 1.0; gap = 0.2
    for k, (zh, what, detail, ana, edge, tc, fill) in enumerate(layers):
        y = y0 + k*(h+gap)
        band(slide, 2.6, y, 8.3, h, fill=fill, edge=edge)
        tf = tb(slide, 2.9, y+0.16, 7.8, h-0.24)
        p = tf.paragraphs[0]
        sr(p.add_run(), f"{zh}：{what}", BODY_F, 13.5, tc, bold=True)
        sr(p.add_run(), f"　{detail}", BODY_F, 11, INK)
        para(tf, f"像公司审阅：{ana}", BODY_F, 10.5, SUB, italic=True)
    arrow(slide, 1.7, y0+3*h+2*gap-0.1, 1.7, y0+0.1, color=DTEAL, w=2.5)
    label(slide, 0.45, 3.55, 1.15, "越上越抽象", DTEAL, 10.5, PP_ALIGN.CENTER, False, True)
    band(slide, 0.7, 6.15, 11.93, 0.75, fill=NEARW, edge=NEARW)
    tf2 = tb(slide, 0.75, 6.13, 12.0, 0.72)
    para(tf2, "层数给了模型「反复咀嚼」的机会——同一份稿子层层过手，每层看的东西不一样。GPT-2 堆 12 层、Llama-70B 堆 80 层。", BODY_F, 11.5, DTEAL, bold=True, first=True)
    footer(slide, MOD, "第 3 章 · 堆叠 N 层")

# ---------- p35 MoE ----------
def draw_llm_moe(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · MoE  MIXTURE OF EXPERTS")
    title(slide, "MoE：编制一千人，每次只叫两位相关的")
    label(slide, 0.55, 1.58, 12.2, "把每层 FFN 复制成多个「专家」，路由器给每个 token 挑少数几个激活——参数总量大，单 token 计算量不变", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 0.7, 3.5, 1.7, 1.0, "token", "一个词", "", fill=CARD, edge=SUB, zh_c=INK, en_c=SUB, zh_size=13, en_size=9.5, desc_size=9)
    node(slide, 3.0, 3.5, 1.95, 1.0, "路由器", "Router", "挑几个专家", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=WHITE, zh_size=13, en_size=9.5, desc_size=9)
    arrow(slide, 2.42, 4.0, 2.98, 4.0, w=1.75)
    ex = [("专家 1", True), ("专家 2", False), ("专家 3", True), ("专家 4", False), ("专家 5", False), ("专家 6", False)]
    exx = 5.7; ew = 2.0; ey0 = 2.35; eh = 0.62; egap = 0.13
    for k, (zh, active) in enumerate(ex):
        y = ey0 + k*(eh+egap)
        if active:
            node(slide, exx, y, ew, eh, zh, "激活", "", fill=ORANGE, edge=ODARK, zh_c=OINK, en_c=OINK, zh_size=12, en_size=9, desc_size=9)
            line(slide, 4.97, 4.0, exx-0.02, y+eh/2, color=ORANGE, w=1.4)
        else:
            node(slide, exx, y, ew, eh, zh, "休眠", "", fill=NEARW, edge=LINE, zh_c=SUB, en_c=SUB, zh_size=12, en_size=9, desc_size=9)
    label(slide, 8.1, 3.7, 4.4, "总参数很大（全体在编），每 token 只激活少数几个", SUB, 11, PP_ALIGN.LEFT, True)
    band(slide, 0.7, 6.15, 11.93, 0.75, fill=NEARW, edge=NEARW)
    tf = tb(slide, 0.75, 6.13, 12.0, 0.72)
    para(tf, "像三甲医院分诊台：1000 位专家在编（总参数），你挂一次号只见 2 位相关科室的（激活参数）。DeepSeek-V3：6710 亿总参 / 每 token 激活约 370 亿。", BODY_F, 11.5, DTEAL, bold=True, first=True)
    footer(slide, MOD, "第 3 章 · MoE")

# ---------- p43 Prefill 与 Decode ----------
def draw_llm_prefill(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · Prefill 与 Decode  TWO PHASES")
    title(slide, "Prefill 与 Decode：首 token 慢的真相")
    label(slide, 0.55, 1.58, 12.2, "备菜（prefill）慢、上菜（decode）快——催菜没用，要么菜单短一点，要么后厨加人", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 0.7, 3.0, 2.0, 1.3, "Prompt 进入", "整段提示词", "", fill=CARD, edge=SUB, zh_c=INK, en_c=SUB, zh_size=13, en_size=10, desc_size=9)
    node(slide, 3.3, 3.0, 3.0, 1.3, "Prefill", "整段并行计算", "算力密集 · O(n²) · 慢", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=WHITE, zh_size=15, en_size=10.5, desc_size=10)
    node(slide, 6.9, 3.0, 3.0, 1.3, "Decode", "逐 token 循环", "带宽密集 · 流式 · 快", fill=ORANGE, edge=ODARK, zh_c=OINK, sub_c=OINK, en_c=OINK, zh_size=15, en_size=10.5, desc_size=10)
    node(slide, 10.5, 3.0, 2.1, 1.3, "流式输出", "一个个吐字", "", fill=CORAL, edge=ORANGE, zh_c=ODARK, en_c=ODARK, zh_size=13, en_size=10, desc_size=9)
    arrow(slide, 2.72, 3.65, 3.28, 3.65, w=1.75)
    arrow(slide, 6.32, 3.65, 6.88, 3.65, color=TEAL, w=1.75)
    label(slide, 5.95, 4.42, 2.0, "▲ 首 token · TTFT", DTEAL, 10.5, PP_ALIGN.CENTER, False, True)
    arrow(slide, 9.92, 3.65, 10.48, 3.65, color=ORANGE, w=1.75)
    label(slide, 9.45, 4.42, 2.0, "▲ 每 token · TPOT", ODARK, 10.5, PP_ALIGN.CENTER, False, True)
    band(slide, 0.7, 5.35, 11.93, 1.5, fill=CARD, edge=LINE)
    tf = tb(slide, 1.0, 5.5, 11.4, 1.35)
    para(tf, "客户抱怨「贴了 50 页文档进去，半天不出字！」——那是 prefill 在啃 O(n²) 的注意力计算；一旦开始输出就流畅了。", BODY_F, 11.5, DTEAL, bold=True, first=True)
    para(tf, "两个指标分开谈：首 token 延迟 TTFT（Time To First Token）、每 token 时延 TPOT（Time Per Output Token）。要么精简 prompt / 命中缓存，要么堆算力。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 4 章 · Prefill 与 Decode")

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in [draw_llm_mha, draw_llm_layers, draw_llm_moe, draw_llm_prefill]:
        fn(prs.slides.add_slide(B))
    out = "/Users/lijiaxiang/project/myAILearning/_maintenance/_llm_preview.pptx"
    prs.save(out); print("saved 4 ->", out)


# ========== 第二轮补图 ==========
def draw_llm_transformer_overview(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · Transformer 总览  A ROUND TABLE")
    title(slide, "Transformer 总览：一场圆桌会议")
    label(slide, 0.55, 1.58, 12.2, "这张图是全讲义的地基——第 2 章拆「注意力」这一格，第 3 章拆其余每一格", SUB, 12, PP_ALIGN.LEFT, True)
    steps = [("文本", "一句话"), ("切词", "Token"), ("向量化", "Embedding"), ("N 层", "注意力 + FFN"), ("下一词", "概率分布")]
    W = 2.15; gap = 0.28; x0 = 0.6; y = 2.55; h = 1.0; xs = []
    for i, (zh, d) in enumerate(steps):
        x = x0 + i*(W+gap); xs.append(x)
        solid = (i == 3)
        node(slide, x, y, W, h, zh, "", d, fill=(TEAL if solid else CARD), edge=(DTEAL if solid else TEAL),
             zh_c=(WHITE if solid else DTEAL), sub_c=(PALE if solid else SUB), zh_size=13, desc_size=9.5)
        if i < 4:
            arrow(slide, x+W+0.02, y+h/2, x+W+gap-0.02, y+h/2, color=TEAL, w=1.6)
    label(slide, xs[3]-0.3, y+h+0.08, W+0.6, "会开 N 轮（N 层），每轮更深", SUB, 9.5, PP_ALIGN.CENTER, True)
    band(slide, 0.6, 4.35, 6.0, 1.85, fill=CARD, edge=TEAL)
    tf = tb(slide, 0.85, 4.47, 5.5, 1.65)
    para(tf, "类比 · 圆桌会议", BODY_F, 12.5, DTEAL, bold=True, first=True)
    para(tf, "每个词是一位与会者：发言前先环顾全场（注意力）——谁与我相关就多听谁的；然后回座位独立消化整理（前馈网络 FFN）。", BODY_F, 11, SUB)
    band(slide, 6.75, 4.35, 5.88, 1.85, fill=CORAL, edge=ORANGE)
    tf2 = tb(slide, 7.0, 4.47, 5.4, 1.65)
    para(tf2, "要点 · 与「传声筒」的本质区别", BODY_F, 12.5, ODARK, bold=True, first=True)
    para(tf2, "任意两个词直接连线，信息路径长度恒为 1——「第 12 页的甲方」和「第 1 页的定义」之间是直达电梯，不是 12 页的接力。", BODY_F, 11, SUB)
    band(slide, 0.6, 6.32, 12.03, 0.5, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 6.37, 11.5, 0.42); p = tf3.paragraphs[0]
    sr(p.add_run(), "全讲义地基：", BODY_F, 11, DTEAL, bold=True)
    sr(p.add_run(), "第 2 章拆「注意力」，第 3 章拆其余每一格，第 4、5 章讲这张图跑起来要花多少钱。", BODY_F, 11, INK)
    footer(slide, MOD, "第 1 章 · 语言与序列建模")

def draw_llm_qkv(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · QKV 三件套  LIBRARY SEARCH")
    title(slide, "QKV 三件套：图书馆检索类比")
    label(slide, 0.55, 1.58, 12.2, "每个词的向量分别乘三个可学习矩阵，分身出三个角色——都是训练出来的", SUB, 12, PP_ALIGN.LEFT, True)
    roles = [("Q · Query", "查询：我想找什么", "你手里的检索词", TEAL, DTEAL, CARD),
             ("K · Key", "键：我是讲什么的", "每本书的书脊标签", ORANGE, ODARK, CORAL),
             ("V · Value", "值：真正被取走的", "书的正文内容", DTEAL, NAVY, PALE)]
    W = 3.95; gap = 0.14; x0 = 0.6; y = 2.45; h = 1.85
    for i, (en, zh, ana, tc, ec, fill) in enumerate(roles):
        x = x0 + i*(W+gap)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, y+0.18, W-0.5, 0.6); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        sr(p.add_run(), en, BODY_F, 15, tc, bold=True)
        para(tf, zh, BODY_F, 11.5, INK, bold=True, align=PP_ALIGN.CENTER)
        t2 = tb(slide, x+0.28, y+1.15, W-0.56, 0.55); q = t2.paragraphs[0]; q.alignment = PP_ALIGN.CENTER
        sr(q.add_run(), ana, BODY_F, 11, SUB)
    label(slide, 0.6, 4.45, 12.0, "检索词（Q）和书脊（K）匹配度越高，从那本书（V）里取的内容就越多", DTEAL, 11, PP_ALIGN.CENTER, False, True)
    band(slide, 0.6, 4.85, 12.03, 1.97, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 4.97, 11.5, 1.8); p = tf.paragraphs[0]
    sr(p.add_run(), "为什么要三个分身？", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "因为「我在找什么」和「我能提供什么」是两回事——K ≠ V。", BODY_F, 12, INK)
    para(tf, "例：语法上「动词要找它的主语」（Q），但主语那个词能提供的语义远不止「我是主语」（K 挂的标签 ≠ V 给的内容）。", BODY_F, 11, SUB)
    para(tf, "三个矩阵都是训练出来的——模型自己学会「该发什么检索词、该挂什么标签、该给什么内容」。", BODY_F, 11, DTEAL, bold=True)
    footer(slide, MOD, "第 2 章 · 注意力机制")
