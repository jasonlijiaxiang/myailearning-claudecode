"""LLM-Inference 讲义信息图（5 张）。复用 kb_draw + agent_figs helper。
嫁接：p7 后(Prefill/Decode)、p19 后(KV并发)、p27 后(批处理三代)、p28 后(PagedAttention)、p58 后(投机解码)。
LLM-Inference 无页码，图系统页脚匹配。独立预览：python3 inf_figs.py"""
import sys
sys.path.insert(0, "/Users/lijiaxiang/project/myAILearning/_maintenance")
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "LLM-Inference 讲义 · 推理服务"
GINK = RGBColor(0x27, 0x50, 0x0A)

# ---------- p7 Prefill vs Decode ----------
def draw_inf_prefill_decode(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 两阶段  PREFILL vs DECODE")
    title(slide, "Prefill 与 Decode：性格相反的两个阶段")
    label(slide, 0.55, 1.58, 12.2, "一个吃算力、一个吃带宽——混在同一张卡上必然互相干扰（第 7 章 P/D 分离的伏笔）", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 3.0, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.52, 5.6, "Prefill 预填充 —— 读题，一目十行", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 3.1, 5.3, 0.62, "所有输入 token 并行算 · 一次大矩阵乘吃满算力", "", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, zh_size=11)
    tf = tb(slide, 0.9, 3.9, 5.5, 1.35)
    para(tf, "· 瓶颈：算力受限（Compute-bound）", BODY_F, 11.5, INK, bold=True, first=True)
    para(tf, "· 决定：首字延迟 TTFT", BODY_F, 11, ODARK, bold=True)
    para(tf, "· 读题快慢看「理解力」（算力）", BODY_F, 11, SUB)
    band(slide, 6.75, 2.4, 6.0, 3.0, fill=CORAL, edge=ORANGE)
    label(slide, 7.0, 2.52, 5.6, "Decode 解码 —— 答题，一字一字写", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.1, 5.3, 0.62, "逐 token 生成 · 每步把整个权重从显存搬一遍", "", "", fill=ORANGE, edge=ODARK, zh_c=OINK, zh_size=11)
    tf2 = tb(slide, 7.05, 3.9, 5.5, 1.35)
    para(tf2, "· 瓶颈：带宽受限（Memory-bound）", BODY_F, 11.5, INK, bold=True, first=True)
    para(tf2, "· 决定：每字延迟 TPOT", BODY_F, 11, ODARK, bold=True)
    para(tf2, "· 「不在算得慢，在搬得慢」——手速（带宽）", BODY_F, 11, SUB)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.67, 11.5, 1.1); p = tf3.paragraphs[0]
    sr(p.add_run(), "类比 · 考试：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "读题可以一目十行（prefill），答题只能一个字一个字写（decode）——读题快慢看理解力（算力），写字快慢看手速（带宽）。", BODY_F, 11.5, INK)
    para(tf3, "所以推理优化的两条主线：prefill 拼算力利用率，decode 拼显存带宽——两阶段的账要分开算。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 1 章 · 推理是怎么跑起来的")

# ---------- p19 KV Cache 座位学 ----------
def draw_inf_kv_concurrency(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 并发天花板  KV = SEATS")
    title(slide, "KV Cache 决定并发：显存里的「座位学」")
    label(slide, 0.55, 1.58, 12.2, "并发上限 = 剩余显存 ÷ 每请求 KV——容量规划第一问是「峰值并发 × 平均上下文 = 多少 KV」", SUB, 12, PP_ALIGN.LEFT, True)
    # 一条显存条：权重（固定） + KV×N + 剩余
    label(slide, 0.7, 2.35, 6.0, "一张卡的显存账", DTEAL, 11.5, PP_ALIGN.LEFT, False, True)
    X0 = 0.7; y = 2.75; BW = 11.9; h = 0.95
    band(slide, X0, y, BW*0.42, h, fill=DTEAL, edge=NAVY)
    label(slide, X0, y+0.3, BW*0.42, "权重 · 固定成本（占掉大半）", WHITE, 12, PP_ALIGN.CENTER, False, True)
    segs = [("KV·请求1", 0.13, ORANGE), ("KV·请求2", 0.11, ORANGE), ("KV·请求3", 0.09, CORAL), ("剩余", 0.25, GBG)]
    x = X0 + BW*0.42
    for name, frac, c in segs:
        w = BW*frac
        band(slide, x, y, w, h, fill=c, edge=(ODARK if c != GBG else GREEN))
        label(slide, x, y+0.32, w, name, (OINK if c == ORANGE else ODARK if c == CORAL else GINK), 10, PP_ALIGN.CENTER, False, True)
        x += w
    label(slide, X0 + BW*0.42, y+h+0.08, BW*0.5, "↑ KV 是变动成本：每接一路请求预留一段、且随对话越聊越大", SUB, 10, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 4.35, 6.0, 1.75, fill=CARD, edge=TEAL)
    tf = tb(slide, 0.85, 4.47, 5.5, 1.55)
    para(tf, "类比 · 餐厅", BODY_F, 12.5, DTEAL, bold=True, first=True)
    para(tf, "厨房设备（权重）占掉大半面积后，剩下的面积决定能摆几张桌子（并发）；客人坐下后桌子还越变越大（对话变长、KV 膨胀）——不清桌就爆满。", BODY_F, 11, SUB)
    band(slide, 6.75, 4.35, 5.88, 1.75, fill=CORAL, edge=ORANGE)
    tf2 = tb(slide, 7.0, 4.47, 5.4, 1.55)
    para(tf2, "案例 · 静态预分配的浪费", BODY_F, 12.5, ODARK, bold=True, first=True)
    para(tf2, "早期按「最大可能长度」给每桌预留——为一个可能聊到 32K 的对话空占 10GB，结果 60–80% 显存空转。解法：下一章 PagedAttention 分页管理。", BODY_F, 11, SUB)
    band(slide, 0.6, 6.28, 12.03, 0.54, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 6.34, 11.5, 0.44); p = tf3.paragraphs[0]
    sr(p.add_run(), "底层原因：", BODY_F, 11, DTEAL, bold=True)
    sr(p.add_run(), "显存耗尽的瞬间，新请求只能排队或被抢占——这就是高峰期 API 变慢的真相。", BODY_F, 11, INK)
    footer(slide, MOD, "第 2 章 · KV Cache 与显存账")

# ---------- p27 批处理三代 ----------
def draw_inf_batching(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 批处理  STATIC → CONTINUOUS")
    title(slide, "从静态到连续：批处理的三代演进")
    label(slide, 0.55, 1.58, 12.2, "连续批解决「时间维度」的空转，PagedAttention 解决「空间维度」的浪费", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 3.0, fill=CORAL, edge=ORANGE)
    label(slide, 0.85, 2.52, 5.6, "静态批 —— 旅游包车", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    for i, (lab, w) in enumerate([("请求 A（短）", 1.6), ("请求 B（长）", 4.6), ("请求 C（短）", 1.6)]):
        yy = 3.1 + i*0.5
        band(slide, 0.95, yy, w, 0.36, fill=ORANGE, edge=ODARK)
        if w < 3:
            band(slide, 0.95+w, yy, 4.6-w, 0.36, fill=WHITE, edge=RGBColor(0xD9,0xC9,0xB8))
            label(slide, 0.95+w, yy+0.02, 4.6-w, "空等", RGBColor(0xB0,0x7A,0x40), 9, PP_ALIGN.CENTER, True)
    label(slide, 0.9, 4.75, 5.5, "攒够 N 个一起跑、全完成才收工——短请求陪长请求空等（木桶效应），新请求等下一班车。", SUB, 10.5, PP_ALIGN.LEFT, True)
    band(slide, 6.75, 2.4, 6.0, 3.0, fill=GBG, edge=GREEN)
    label(slide, 7.0, 2.52, 5.6, "连续批 —— 地铁（迭代级调度）", GINK, 12.5, PP_ALIGN.LEFT, False, True)
    for i, (lab, x0v, w) in enumerate([("A", 7.1, 1.5), ("B", 7.1, 4.6), ("C", 7.1, 1.5), ("D 新", 8.9, 2.8)]):
        yy = 3.1 + i*0.5
        band(slide, x0v, yy, w, 0.36, fill=GREEN, edge=GINK)
        label(slide, x0v, yy+0.02, w, lab, WHITE, 9.5, PP_ALIGN.CENTER, False, True)
    label(slide, 7.0, 4.75, 5.7, "每生成一步都重新组批——完成的立刻下车、新请求下一步就上车（D），GPU 没有「等一班车」的空窗。", SUB, 10.5, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.67, 11.5, 1.1); p = tf.paragraphs[0]
    sr(p.add_run(), "类比：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "从「旅游包车」（人齐才发、全程同进退）进化到「地铁」（到站随时上下、车厢永远在跑）。", BODY_F, 11.5, INK)
    para(tf, "连续批是 vLLM / SGLang / TensorRT-LLM 的默认行为——但要配下一页的 PagedAttention 才真正落地（否则「随时上车」会把显存撕碎）。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 3 章 · 批处理与调度")

# ---------- p28 PagedAttention ----------
def draw_inf_paged_attention(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · PagedAttention  MEMORY PAGING")
    title(slide, "PagedAttention：把显存当虚拟内存管")
    label(slide, 0.55, 1.58, 12.2, "操作系统虚拟内存的翻版——逻辑连续、物理任意，碎片近零，直接换来 2–4 倍吞吐", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 2.95, fill=CORAL, edge=ORANGE)
    label(slide, 0.85, 2.52, 5.6, "旧做法 · 连续预留 → 碎片化", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    for i, (used, x0v) in enumerate([(0.2, 0.95), (0.5, 0.95), (0.3, 0.95)]):
        yy = 3.1 + i*0.5
        band(slide, x0v, yy, 5.1*used, 0.36, fill=ORANGE, edge=ODARK)
        band(slide, x0v+5.1*used, yy, 5.1*(1-used), 0.36, fill=WHITE, edge=RGBColor(0xD9,0xC9,0xB8))
    label(slide, 0.95, 4.75, 5.5, "按「最大可能长度」预留连续块，实际只用 20%——研究测到 60–80% KV 显存被浪费。", SUB, 10.5, PP_ALIGN.LEFT, True)
    band(slide, 6.75, 2.4, 6.0, 2.95, fill=GBG, edge=GREEN)
    label(slide, 7.0, 2.52, 5.6, "分页 · 16 token/页，按需分配", GINK, 12.5, PP_ALIGN.LEFT, False, True)
    for r in range(3):
        for c in range(8):
            filled = ((r*8+c) % 5 != 4)
            band(slide, 7.1+c*0.62, 3.1+r*0.5, 0.5, 0.36, fill=(GREEN if filled else WHITE), edge=GINK)
    label(slide, 7.0, 4.75, 5.7, "切成固定小页、用页表映射：逻辑连续、物理任意，用一格给一格——碎片近零、利用率逼近 100%。", SUB, 10.5, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 5.5, 12.03, 1.32, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.62, 11.5, 1.15); p = tf.paragraphs[0]
    sr(p.add_run(), "类比 · 虚拟内存：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "进程以为自己拿到整段连续内存，物理页其实散落各处——图书馆不再「按人预留一整排书架」，改成「用一格给一格」。", BODY_F, 11.5, INK)
    para(tf, "赠品：同一份前缀的页可被多请求共享（copy-on-write）——前缀缓存由此而来。论文已入库 arXiv-2309.06180。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 3 章 · 批处理与调度")

# ---------- p58 投机解码 ----------
def draw_inf_speculative(slide):
    bg(slide)
    eyebrow(slide, "第 6 章 · 投机解码  SPECULATIVE DECODING")
    title(slide, "投机解码：秘书起草，老板审签")
    label(slide, 0.55, 1.58, 12.2, "用一次并行验证替代 K 次串行生成——把 decode 的串行枷锁，部分变成 prefill 式的并行", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 0.7, 2.9, 2.5, 1.0, "草稿模型", "便宜、小", "一口气猜 K 个 token", fill=CARD, edge=TEAL, zh_c=DTEAL, sub_c=SUB, en_c=TEAL, zh_size=14, en_size=9.5, desc_size=9.5)
    arrow(slide, 3.22, 3.4, 3.95, 3.4, color=TEAL, w=1.8)
    # K 个草稿 token
    toks = [("the", True), ("cat", True), ("sat", True), ("in", True), ("mat", False)]
    for i, (t, ok) in enumerate(toks):
        band(slide, 4.0 + i*1.0, 3.05, 0.9, 0.66, fill=(GBG if ok else RGBColor(0xFB,0xE9,0xE7)), edge=(GREEN if ok else RGBColor(0x9B,0x2C,0x2C)))
        label(slide, 4.0 + i*1.0, 3.14, 0.9, t, (GINK if ok else RGBColor(0x9B,0x2C,0x2C)), 11, PP_ALIGN.CENTER, False, True)
        label(slide, 4.0 + i*1.0, 3.42, 0.9, "接受" if ok else "接管", (GINK if ok else RGBColor(0x9B,0x2C,0x2C)), 8.5, PP_ALIGN.CENTER, True)
    label(slide, 4.0, 2.68, 5.0, "大模型一次前向并行验证这 K 个", SUB, 10, PP_ALIGN.LEFT, True)
    label(slide, 4.0, 3.85, 5.5, "接受猜对的前缀 → 从第一个猜错处，大模型自己接管", ODARK, 10.5, PP_ALIGN.LEFT, False, True)
    node(slide, 9.7, 2.9, 2.9, 1.0, "大模型", "并行验证", "验证一次 ≈ 生成一个 token", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=PALE, zh_size=14, en_size=9.5, desc_size=9)
    band(slide, 0.6, 4.5, 6.0, 1.55, fill=CARD, edge=TEAL)
    tf = tb(slide, 0.85, 4.62, 5.5, 1.35)
    para(tf, "类比 · 老板口述邮件太慢", BODY_F, 12, DTEAL, bold=True, first=True)
    para(tf, "让秘书先起草一段，老板扫一眼：前半句都对（签字通过），第五个词不对（从这里亲自改写）——审批比逐字口述快得多。", BODY_F, 11, SUB)
    band(slide, 6.75, 4.5, 5.88, 1.55, fill=GBG, edge=GREEN)
    tf2 = tb(slide, 7.0, 4.62, 5.4, 1.35)
    para(tf2, "为什么「无损」", BODY_F, 12, GINK, bold=True, first=True)
    para(tf2, "验证按大模型自己的概率分布逐 token 把关（拒绝采样），数学上可证输出分布与大模型独跑一模一样——不是「差不多」，是分布级相等。", BODY_F, 11, SUB)
    label(slide, 0.7, 6.25, 12.0, "小结：猜对几个就白赚几个；「验证」可并行、「生成」只能串行——加速不改变质量，所以敢进生产默认配置。", DTEAL, 11, PP_ALIGN.LEFT, False, True)
    footer(slide, MOD, "第 6 章 · 投机解码与算法加速")

# ---------- p10 三大指标 ----------
def draw_inf_three_metrics(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 三大指标  TTFT · TPOT · THROUGHPUT")
    title(slide, "三个仪表盘：TTFT · TPOT · 吞吐")
    label(slide, 0.55, 1.58, 12.2, "所有推理服务的体验与成本，都能落到这三个可测量的指标上——售前谈 SLA 就谈它们", SUB, 12, PP_ALIGN.LEFT, True)
    cards = [("TTFT", "首字延迟", "Time To First Token", "用户等多久看到第一个字", "由 Prefill 决定", TEAL, DTEAL, CARD),
             ("TPOT", "每字延迟", "Time Per Output Token", "字与字之间的「打字速度」", "由 Decode 决定", ORANGE, ODARK, CORAL),
             ("吞吐", "Throughput", "tokens / 秒", "一张卡每秒总产出、摊成本的分母", "由批处理决定", DTEAL, NAVY, PALE)]
    W = 3.95; gap = 0.14; x0 = 0.6; y = 2.4; h = 2.75
    for i, (en, zh, full, what, phase, tc, ec, fill) in enumerate(cards):
        x = x0 + i*(W+gap)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, y+0.18, W-0.5, 0.95); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        sr(p.add_run(), en, BODY_F, 24, tc, bold=True)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        sr(p2.add_run(), zh + "　" + full, BODY_F, 10.5, SUB)
        node(slide, x+0.35, y+1.35, W-0.7, 0.56, phase, "", "", fill=WHITE, edge=ec, zh_c=tc, zh_size=11.5)
        tf2 = tb(slide, x+0.3, y+2.05, W-0.6, 0.6); q = tf2.paragraphs[0]; q.alignment = PP_ALIGN.CENTER
        sr(q.add_run(), what, BODY_F, 11, INK)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=RGBColor(0xFB,0xE9,0xE7), edge=RGBColor(0x9B,0x2C,0x2C))
    tf = tb(slide, 0.9, 5.67, 11.5, 1.1); p = tf.paragraphs[0]
    sr(p.add_run(), "延迟与吞吐天然打架：", BODY_F, 12, RGBColor(0x9B,0x2C,0x2C), bold=True)
    sr(p.add_run(), "批越大 → 吞吐越高、但单请求越慢（TPOT 上升）。所有调度设计都是在这条曲线上选点（第 3、8 章展开）。", BODY_F, 12, INK)
    para(tf, "售前口径：先问客户「首字快」重要还是「总吞吐 / 成本低」重要——SLA 从这三个数谈起。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 1 章 · 推理是怎么跑起来的")

# ---------- p22 长上下文两头吃钱 ----------
def draw_inf_long_context(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 长上下文的价格  TWO-SIDED COST")
    title(slide, "长上下文为什么贵：两头吃钱")
    label(slide, 0.55, 1.58, 12.2, "上下文是「按面积收费的仓库」——能检索就别硬塞（串联 RAG 模块）", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 2.95, fill=CORAL, edge=ORANGE)
    label(slide, 0.85, 2.52, 5.6, "Prefill 端：计算「平方级」涨", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 3.1, 5.3, 0.62, "输入 4K → 128K（32 倍）", "", "注意力计算量 ∝ 长度²", fill=WHITE, edge=ODARK, zh_c=ODARK, sub_c=SUB, zh_size=12, desc_size=9)
    tf = tb(slide, 0.9, 3.9, 5.5, 1.35)
    para(tf, "· 计算量涨约 1000 倍", BODY_F, 11.5, REDINK if False else RGBColor(0x9B,0x2C,0x2C), bold=True, first=True)
    para(tf, "· TTFT 从零点几秒 → 几十秒", BODY_F, 11, SUB)
    para(tf, "· 还独占算力、挤压别的请求", BODY_F, 11, SUB)
    band(slide, 6.75, 2.4, 6.0, 2.95, fill=CARD, edge=TEAL)
    label(slide, 7.0, 2.52, 5.6, "Decode 端：显存 + 带宽「线性」涨", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.1, 5.3, 0.62, "128K 的 KV ≈ 41GB", "", "不仅占座，每生成一字还扫一遍", fill=WHITE, edge=TEAL, zh_c=DTEAL, sub_c=SUB, zh_size=12, desc_size=9)
    tf2 = tb(slide, 7.05, 3.9, 5.5, 1.35)
    para(tf2, "· TPOT 随之变慢", BODY_F, 11.5, DTEAL, bold=True, first=True)
    para(tf2, "· API 把长输入按 token 计费、有的分档加价", BODY_F, 11, SUB)
    para(tf2, "· 成本结构就是这么来的", BODY_F, 11, SUB)
    band(slide, 0.6, 5.5, 12.03, 1.32, fill=GBG, edge=GREEN)
    tf3 = tb(slide, 0.9, 5.62, 11.5, 1.15); p = tf3.paragraphs[0]
    sr(p.add_run(), "案例：", BODY_F, 11.5, GINK, bold=True)
    sr(p.add_run(), "「把整本手册塞进上下文」每次提问为 10 万 token 付 prefill+KV 的钱；「RAG 检索相关三页」每次只带几千 token。", BODY_F, 11.5, INK)
    para(tf3, "小结：能检索就别硬塞；必须长驻的前缀，用 prefix caching 只付一次——这正是「1M 窗口不是 RAG 终结者」的推理侧论据。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 2 章 · KV Cache 与显存账")

# ---------- p75 P/D 分离 ----------
def draw_inf_pd_split(slide):
    bg(slide)
    eyebrow(slide, "第 7 章 · P/D 分离  DISAGGREGATION")
    title(slide, "P/D 分离：读题的和写字的，别再抢一张桌子")
    label(slide, 0.55, 1.58, 12.2, "第 1 章埋的雷在集群规模爆发——两种负载体质不同，分开各自优化优于混在一起互相迁就", SUB, 12, PP_ALIGN.LEFT, True)
    # 问题
    band(slide, 0.6, 2.35, 12.03, 1.55, fill=RGBColor(0xFB,0xE9,0xE7), edge=RGBColor(0x9B,0x2C,0x2C))
    label(slide, 0.85, 2.45, 8.0, "问题 · P/D 干扰：混在同一张卡上，两个 SLO 互相绑架", RGBColor(0x9B,0x2C,0x2C), 12, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 2.95, 5.3, 0.78, "Prefill 突发大活（算力密集）", "", "长 prefill 一来，decode 集体卡顿（TPOT 尖刺）", fill=WHITE, edge=RGBColor(0x9B,0x2C,0x2C), zh_c=RGBColor(0x9B,0x2C,0x2C), sub_c=SUB, zh_size=12, desc_size=10)
    node(slide, 6.45, 2.95, 5.3, 0.78, "Decode 持续细活（带宽密集）", "", "为保 decode 限制 prefill，TTFT 又恶化", fill=WHITE, edge=RGBColor(0x9B,0x2C,0x2C), zh_c=RGBColor(0x9B,0x2C,0x2C), sub_c=SUB, zh_size=12, desc_size=10)
    label(slide, 6.05, 3.05, 0.55, "vs", RGBColor(0x9B,0x2C,0x2C), 13, PP_ALIGN.CENTER, False, True)
    # 解法
    band(slide, 0.6, 4.05, 12.03, 1.5, fill=GBG, edge=GREEN)
    label(slide, 0.85, 4.14, 8.0, "解法 · 分池部署 + KV 快递（前厅后厨彻底分区）", GINK, 12, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 4.62, 3.3, 0.82, "Prefill 池", "读题、产出 KV", "配算力猛的卡", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=PALE, zh_size=13, en_size=10, desc_size=10)
    arrow(slide, 4.3, 5.02, 5.15, 5.02, color=ODARK, w=2.0); label(slide, 4.18, 4.62, 1.15, "KV 快递", ODARK, 9.5, PP_ALIGN.CENTER, False, True); label(slide, 4.18, 4.84, 1.15, "NIXL", ODARK, 9, PP_ALIGN.CENTER, True)
    node(slide, 5.2, 4.62, 3.3, 0.82, "Decode 池", "写字、批量生成", "配带宽大、显存多的卡", fill=ORANGE, edge=ODARK, zh_c=OINK, sub_c=OINK, en_c=OINK, zh_size=13, en_size=10, desc_size=10)
    tf = tb(slide, 8.75, 4.62, 3.85, 0.85)
    para(tf, "两池可用不同型号 GPU、按各自负载独立扩缩。", BODY_F, 11, SUB, first=True)
    para(tf, "旅游团再多，只加接单员，不动后厨。", BODY_F, 11, GINK, bold=True)
    band(slide, 0.6, 5.7, 12.03, 1.12, fill=PALE, edge=LINE)
    tf2 = tb(slide, 0.9, 5.81, 11.5, 0.95); p = tf2.paragraphs[0]
    sr(p.add_run(), "数据：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "DistServe（OSDI'24）证明分离可数倍提升 goodput；NVIDIA Dynamo 在 Blackwell 上跑 DeepSeek-R1 吞吐最高提升 7 倍（厂商基准，理解量级即可）。", BODY_F, 11.5, INK)
    para(tf2, "代价是 KV 搬运的网络开销——请求短、KV 小时不划算，它是「大集群 + 长上下文」的解药，不是万能默认。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 7 章 · P/D 分离与分布式")

FIGS = [draw_inf_prefill_decode, draw_inf_kv_concurrency, draw_inf_batching,
        draw_inf_paged_attention, draw_inf_speculative,
        draw_inf_three_metrics, draw_inf_long_context, draw_inf_pd_split]
NEW = [draw_inf_three_metrics, draw_inf_long_context, draw_inf_pd_split]

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in FIGS:
        fn(prs.slides.add_slide(B))
    out = "/Users/lijiaxiang/project/myAILearning/_maintenance/_inf_preview.pptx"
    prs.save(out); print(f"saved {len(FIGS)} ->", out)
