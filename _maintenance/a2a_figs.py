"""A2A 讲义信息图（4 张）。复用 kb_draw 设计系统 + agent_figs 的 line/dot helper。
嫁接位置：p7 后（A2A vs MCP）、p17 后（五大对象）、p20 后（状态机）、p47 后（委派拓扑）。
独立预览：python3 a2a_figs.py"""
import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "A2A 讲义 · 智能体间协议"
REDBG = RGBColor(0xFB, 0xE9, 0xE7)
REDINK = RGBColor(0x9B, 0x2C, 0x2C)
GINK = RGBColor(0x27, 0x50, 0x0A)

# ---------- p7 A2A vs MCP：纵向 vs 横向 ----------
def draw_a2a_vs_mcp(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 是什么/为什么  A2A vs MCP")
    title(slide, "A2A vs MCP：一个接工具，一个接 Agent")
    label(slide, 0.55, 1.58, 12.2, "协议层的两兄弟——最常被问「选哪个」，标准答案是「都要」：一纵一横，不同层、互补", SUB, 12, PP_ALIGN.LEFT, True)
    # 左 MCP · 纵向
    band(slide, 0.6, 2.35, 6.0, 3.05, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.47, 5.6, "MCP · 模型上下文协议 —— 纵向", DTEAL, 13, PP_ALIGN.LEFT, False, True)
    node(slide, 2.5, 3.02, 2.2, 0.72, "Agent", "", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, zh_size=14)
    node(slide, 0.95, 4.32, 2.05, 0.72, "工具", "Tools", "", fill=WHITE, edge=TEAL, zh_c=DTEAL, en_c=TEAL, zh_size=12.5, en_size=9)
    node(slide, 4.2, 4.32, 2.05, 0.72, "数据 / 文件", "Data", "", fill=WHITE, edge=TEAL, zh_c=DTEAL, en_c=TEAL, zh_size=12.5, en_size=9)
    arrow(slide, 3.2, 3.76, 2.0, 4.3, color=TEAL, w=1.6)
    arrow(slide, 4.0, 3.76, 5.2, 4.3, color=TEAL, w=1.6)
    label(slide, 0.85, 5.02, 5.6, "类比：给 Agent 装的 USB 接口", DTEAL, 10.5, PP_ALIGN.LEFT, True)
    tfl = tb(slide, 0.85, 5.28, 5.6, 0.35); sr(tfl.paragraphs[0].add_run(), "粒度：一次工具调用 tools/call", BODY_F, 10.5, SUB)
    # 右 A2A · 横向
    band(slide, 6.75, 2.35, 6.0, 3.05, fill=CORAL, edge=ORANGE)
    label(slide, 7.0, 2.47, 5.6, "A2A · 智能体间协议 —— 横向", ODARK, 13, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.5, 2.2, 0.82, "Agent A", "委派方", "", fill=ORANGE, edge=ODARK, zh_c=OINK, sub_c=OINK, en_c=OINK, zh_size=13, en_size=9.5)
    node(slide, 10.2, 3.5, 2.2, 0.82, "Agent B", "执行方", "", fill=ORANGE, edge=ODARK, zh_c=OINK, sub_c=OINK, en_c=OINK, zh_size=13, en_size=9.5)
    arrow(slide, 9.35, 3.72, 10.15, 3.72, color=ODARK, w=1.6)
    arrow(slide, 10.15, 4.12, 9.35, 4.12, color=ODARK, w=1.6)
    label(slide, 9.15, 3.36, 1.5, "委派 Task", ODARK, 9, PP_ALIGN.CENTER, True)
    label(slide, 9.15, 4.16, 1.5, "回传 Artifact", ODARK, 9, PP_ALIGN.CENTER, True)
    label(slide, 7.0, 5.02, 5.6, "类比：Agent 之间的协作 / 外交规范", ODARK, 10.5, PP_ALIGN.LEFT, True)
    tfr = tb(slide, 7.0, 5.28, 5.6, 0.35); sr(tfr.paragraphs[0].add_run(), "粒度：一次任务委派 Task（可长时运行）", BODY_F, 10.5, SUB)
    # 底部收束
    band(slide, 0.7, 5.72, 11.93, 1.13, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 5.84, 11.4, 1.0)
    para(tf, "「MCP 接工具，A2A 接 Agent」——一个系统里两个都要，不是二选一。", BODY_F, 12, DTEAL, bold=True, first=True)
    para(tf, "别把两者对立：A2A 编排的多个 Agent，每个内部往往正用 MCP 调工具——横向协作里嵌着纵向调用。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 1 章 · 是什么/为什么")

# ---------- p17 五大对象怎么套在一起 ----------
def draw_a2a_five_objects(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 协议解剖  FIVE OBJECTS")
    title(slide, "五大对象怎么套在一起：从发现到交付")
    label(slide, 0.55, 1.58, 12.2, "Agent Card 发现 → 开一个 Task 干活 → 用 Message/Part 你来我往 → 交付 Artifact", SUB, 12, PP_ALIGN.LEFT, True)
    # 主干：Agent Card → Task → Artifact
    node(slide, 0.7, 2.5, 3.4, 1.12, "Agent Card", "数字名片", "身份/能力/端点/鉴权 · 发现的起点", zh_size=15, en_size=11, desc_size=11)
    node(slide, 4.95, 2.5, 3.4, 1.12, "Task", "有状态的活儿", "唯一 ID + 生命周期 · 可长时运行", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=WHITE, zh_size=15, en_size=11, desc_size=11)
    node(slide, 9.2, 2.5, 3.4, 1.12, "Artifact", "产出物", "由 Part 组成 · 增量流式返回", fill=CORAL, edge=ORANGE, zh_c=ODARK, sub_c=SUB, en_c=ODARK, zh_size=15, en_size=11, desc_size=11)
    arrow(slide, 4.12, 3.04, 4.92, 3.04, color=TEAL, w=1.9)
    arrow(slide, 8.37, 3.04, 9.17, 3.04, color=ORANGE, w=1.9)
    label(slide, 3.9, 2.66, 1.3, "发现", TEAL, 10, PP_ALIGN.CENTER, True, True)
    label(slide, 8.15, 2.66, 1.3, "交付", ODARK, 10, PP_ALIGN.CENTER, True, True)
    # Task 里的内容怎么装
    arrow(slide, 6.65, 3.6, 6.65, 4.5, color=DTEAL, w=1.7)
    label(slide, 6.75, 3.92, 2.0, "装载内容", DTEAL, 9.5, PP_ALIGN.LEFT, True)
    node(slide, 3.35, 4.6, 3.4, 1.08, "Message", "一轮通信", "带 role（user / agent）", fill=CARD, edge=TEAL, zh_size=14, en_size=11, desc_size=11)
    node(slide, 7.55, 4.6, 3.4, 1.08, "Part", "最小内容单元", "文本 / 文件 / 结构化数据", fill=CARD, edge=TEAL, zh_size=14, en_size=11, desc_size=11)
    arrow(slide, 6.77, 5.14, 7.52, 5.14, color=TEAL, w=1.6)
    label(slide, 6.6, 4.8, 1.4, "含", TEAL, 10.5, PP_ALIGN.CENTER, True, True)
    label(slide, 7.55, 5.72, 3.4, "Artifact 也由 Part 组成", SUB, 10.5, PP_ALIGN.LEFT, True)
    # 底部
    band(slide, 0.7, 6.05, 11.93, 0.82, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 6.16, 11.4, 0.7); p = tf.paragraphs[0]
    sr(p.add_run(), "一句话：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "Card 是「名片」、Task 是「工单」、Message/Part 是「一来一回的内容」、Artifact 是「交付件」——五个对象各就各位，一次协作就跑起来了。", BODY_F, 11.5, SUB)
    footer(slide, MOD, "第 2 章 · 协议解剖")

# ---------- p20 Task 状态机 ----------
def draw_a2a_task_states(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 协议解剖  STATE MACHINE")
    title(slide, "Task 状态机：九种状态、三条路")
    label(slide, 0.55, 1.55, 12.2, "正常线一条道走完，中断线补料后回来，异常线随时可断——终态不可回跳", SUB, 12, PP_ALIGN.LEFT, True)
    # 中断线（上排）
    node(slide, 3.55, 2.05, 2.55, 0.78, "input-required", "待补料", "", fill=CORAL, edge=ORANGE, zh_c=ODARK, sub_c=ODARK, zh_size=12, en_size=9.5)
    node(slide, 6.45, 2.05, 2.55, 0.78, "auth-required", "待授权", "", fill=CORAL, edge=ORANGE, zh_c=ODARK, sub_c=ODARK, zh_size=12, en_size=9.5)
    label(slide, 8.35, 1.72, 4.5, "中断线：补料 / 补授权后回 working 继续", ODARK, 10, PP_ALIGN.LEFT, True)
    # 正常线（中排）
    node(slide, 0.65, 3.35, 2.15, 0.9, "submitted", "已提交", "", fill=CARD, edge=TEAL, zh_size=12.5, en_size=10)
    node(slide, 5.15, 3.32, 2.35, 0.98, "working", "进行中", "状态真相在服务端", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=WHITE, zh_size=14, en_size=10.5, desc_size=9)
    node(slide, 10.45, 3.35, 2.15, 0.9, "completed", "完成", "", fill=GBG, edge=GREEN, zh_c=GINK, sub_c=GINK, en_c=GINK, zh_size=12.5, en_size=10)
    arrow(slide, 2.82, 3.8, 5.13, 3.8, color=TEAL, w=1.9)
    arrow(slide, 7.52, 3.8, 10.43, 3.8, color=GREEN, w=1.9)
    label(slide, 3.1, 3.46, 1.9, "开始", TEAL, 9.5, PP_ALIGN.CENTER, True, True)
    label(slide, 8.0, 3.46, 2.3, "正常线", GINK, 9.5, PP_ALIGN.CENTER, True, True)
    # 中断双向箭头
    arrow(slide, 5.7, 3.3, 4.7, 2.85, color=ORANGE, w=1.3)
    arrow(slide, 4.55, 2.83, 5.55, 3.28, color=ORANGE, w=1.3, dashed=True)
    arrow(slide, 6.6, 3.3, 7.4, 2.85, color=ORANGE, w=1.3)
    arrow(slide, 7.55, 2.83, 6.75, 3.28, color=ORANGE, w=1.3, dashed=True)
    # 异常线（下排）
    exc = [("failed", "失败", 2.6), ("canceled", "已取消", 5.35), ("rejected", "被拒绝", 8.1)]
    for en, zh, x in exc:
        node(slide, x, 5.15, 2.35, 0.82, en, zh, "", fill=REDBG, edge=REDINK, zh_c=REDINK, sub_c=REDINK, zh_size=12, en_size=9.5)
        arrow(slide, 6.15 + (x-5.35)*0.42, 4.32, x+1.17, 5.13, color=REDINK, w=1.2, dashed=True)
    tfe = tb(slide, 10.62, 5.2, 2.05, 0.78)
    para(tfe, "异常线", BODY_F, 10, REDINK, bold=True, italic=True, first=True)
    para(tfe, "随时可进 · 均为终态", BODY_F, 9.5, REDINK, italic=True)
    # 底部
    band(slide, 0.7, 6.12, 11.93, 0.74, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 6.2, 11.4, 0.62); p = tf.paragraphs[0]
    sr(p.add_run(), "终态不可回跳：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "重试 ＝ 创建新 Task（可关联旧任务），绝不改终态历史——审计与回放靠这条纪律；另有兜底 unspecified（未知态）。", BODY_F, 11.5, SUB)
    footer(slide, MOD, "第 2 章 · 协议解剖")

# ---------- p47 三种委派拓扑 ----------
def draw_a2a_delegation(slide):
    bg(slide)
    eyebrow(slide, "第 5 章 · 多智能体协作  TOPOLOGY")
    title(slide, "任务委派的三种协作拓扑")
    label(slide, 0.55, 1.58, 12.2, "跨进程、跨厂商的编排，落到三种基本形状——和单机多 Agent 编排一个道理", SUB, 12, PP_ALIGN.LEFT, True)
    cols = [(0.6, "① 编排者—执行者", "差旅审批：主 Agent 调度订票 / 预算 / 风控"),
            (4.52, "② 链式（流水线）", "文档：翻译 → 校对 → 排版，顺序接力"),
            (8.44, "③ 并行汇聚", "尽调：并行查财务 / 法务 / 舆情，汇成一份")]
    for x0, head, _ in cols:
        band(slide, x0, 2.3, 3.72, 3.35, fill=CARD, edge=LINE)
        label(slide, x0+0.2, 2.42, 3.4, head, DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    # ① hub-spoke
    cx = 0.6
    node(slide, cx+1.16, 3.0, 1.4, 0.62, "主 Agent", "", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, zh_size=11.5)
    sub1 = [(cx+0.35, "订票"), (cx+1.46, "预算"), (cx+2.57, "风控")]
    for sx, zh in sub1:
        node(slide, sx, 4.35, 0.9, 0.56, zh, "", "", fill=WHITE, edge=TEAL, zh_c=DTEAL, zh_size=10.5)
        line(slide, cx+1.86, 3.62, sx+0.45, 4.35, color=TEAL, w=1.2)
    # ② chain
    cx = 4.52
    ch = [(cx+0.2, "翻译"), (cx+1.31, "校对"), (cx+2.42, "排版")]
    for i, (sx, zh) in enumerate(ch):
        node(slide, sx, 3.75, 1.0, 0.66, zh, "", "", fill=(TEAL if i == 1 else WHITE), edge=(DTEAL if i == 1 else TEAL), zh_c=(WHITE if i == 1 else DTEAL), zh_size=11.5)
        if i < 2:
            arrow(slide, sx+1.0, 4.08, ch[i+1][0], 4.08, color=TEAL, w=1.5)
    # ③ fan-in
    cx = 8.44
    node(slide, cx+1.16, 2.95, 1.4, 0.56, "主 Agent", "", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, zh_size=11)
    par = [(cx+0.35, "财务"), (cx+1.46, "法务"), (cx+2.57, "舆情")]
    for sx, zh in par:
        node(slide, sx, 3.95, 0.9, 0.54, zh, "", "", fill=WHITE, edge=TEAL, zh_c=DTEAL, zh_size=10.5)
        line(slide, cx+1.86, 3.51, sx+0.45, 3.95, color=TEAL, w=1.1)
        line(slide, sx+0.45, 4.49, cx+1.86, 5.0, color=ORANGE, w=1.1)
    node(slide, cx+1.06, 5.0, 1.6, 0.56, "汇总", "merge", "", fill=CORAL, edge=ORANGE, zh_c=ODARK, en_c=ODARK, zh_size=11, en_size=8.5)
    # 场景标注
    for x0, _, scene in cols:
        tf = tb(slide, x0+0.16, 5.72, 3.5, 0.5); p = tf.paragraphs[0]
        sr(p.add_run(), scene, BODY_F, 9.5, SUB, italic=True)
    # 底部
    band(slide, 0.7, 6.32, 11.93, 0.6, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 6.38, 11.4, 0.5); p = tf.paragraphs[0]
    sr(p.add_run(), "一句话：", BODY_F, 11, DTEAL, bold=True)
    sr(p.add_run(), "拓扑没变，变的只是「跨进程 / 跨厂商」——A2A 把编排的边，从一个系统之内扩到了系统之外。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 5 章 · 多智能体协作")

# ---------- p32 三种传输绑定 ----------
def draw_a2a_transports(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 发现与传输  TRANSPORTS")
    title(slide, "三种传输绑定：语义相同，按需选")
    label(slide, 0.55, 1.58, 12.2, "三种绑定由 protobuf 统一约束、可互换——A2A 不逼你换技术栈，用顺手的传输即可", SUB, 12, PP_ALIGN.LEFT, True)
    cards = [("JSON-RPC 2.0", "通用默认 · Web 友好", TEAL, DTEAL, CARD,
              ["基于 HTTP POST，人类可读", "生态最广", "preferredTransport 常用它"]),
             ("gRPC", "高性能 · 内网服务间", DTEAL, NAVY, PALE,
              ["HTTP/2 + protobuf，低延迟高吞吐", "强类型", "适合大规模 Agent 集群"]),
             ("HTTP + JSON / REST", "已有 REST 基建", ORANGE, ODARK, CORAL,
              ["贴合传统 REST 习惯", "易接入现有网关 / 鉴权 / 监控", "改造成本最低"])]
    W = 3.95; gap = 0.14; x0 = 0.6; y = 2.4; h = 2.75
    for i, (name, scene, ec, tc, fill, feats) in enumerate(cards):
        x = x0 + i*(W+gap)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, y+0.16, W-0.5, 0.5)
        sr(tf.paragraphs[0].add_run(), name, BODY_F, 14.5, tc, bold=True)
        node(slide, x+0.3, y+0.72, W-0.6, 0.56, scene, "", "", fill=WHITE, edge=ec, zh_c=tc, zh_size=11.5)
        bf = tb(slide, x+0.28, y+1.42, W-0.56, 1.2)
        for j, ft in enumerate(feats):
            para(bf, "· " + ft, BODY_F, 11, INK, first=(j == 0))
    band(slide, 0.6, 5.42, 12.03, 1.4, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.55, 11.5, 1.25); p = tf.paragraphs[0]
    sr(p.add_run(), "要点：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "三种绑定语义由 protobuf 统一约束、可互换——换传输不换语义。", BODY_F, 12, INK)
    para(tf, "对外一句话：A2A 不逼你换技术栈，JSON-RPC 起步、内网大流量上 gRPC、有 REST 网关就走 REST。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 3 章 · 发现与传输")

# ---------- p34 四种交付姿势 ----------
def draw_a2a_delivery(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 发现与传输 · 深潜  DELIVERY")
    title(slide, "拿结果的四种姿势：按运行条件选")
    label(slide, 0.55, 1.58, 12.2, "先判运行条件，再选交付方式——别为了「实时」默认上长连接", SUB, 12, PP_ALIGN.LEFT, True)
    cards = [("Blocking", "同步等", "短活", "连接占用、易超时", TEAL, DTEAL, CARD),
             ("Polling", "轮询", "网络最简单", "空轮询、状态陈旧（退避 + 变更标记）", TEAL, DTEAL, CARD),
             ("Streaming", "SSE 流式", "看进度、低延迟", "断线重连、背压", ORANGE, ODARK, CORAL),
             ("Push", "回调推送", "长任务 + 客户端离线", "防重复投递、回调伪造", ORANGE, ODARK, CORAL)]
    W = 3.0; gap = 0.11; x0 = 0.55; y = 2.4; h = 2.55
    for i, (en, zh, fit, cost, ec, tc, fill) in enumerate(cards):
        x = x0 + i*(W+gap)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.22, y+0.16, W-0.44, 0.7); p = tf.paragraphs[0]
        sr(p.add_run(), en, BODY_F, 14, tc, bold=True)
        para(tf, zh, BODY_F, 11, SUB, italic=True)
        node(slide, x+0.28, y+0.92, W-0.56, 0.56, "适合：" + fit, "", "", fill=WHITE, edge=ec, zh_c=tc, zh_size=11)
        cf = tb(slide, x+0.26, y+1.62, W-0.52, 0.85); q = cf.paragraphs[0]
        sr(q.add_run(), "代价：", BODY_F, 10.5, ODARK, bold=True)
        para(cf, cost, BODY_F, 10.5, SUB)
    band(slide, 0.55, 5.42, 12.08, 1.4, fill=PALE, edge=LINE)
    tf = tb(slide, 0.85, 5.54, 11.5, 1.25); p = tf.paragraphs[0]
    sr(p.add_run(), "选型口诀：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "短活阻塞 · 简单轮询 · 看进度流式 · 长任务推送。", BODY_F, 12, INK)
    para(tf, "同一个 Task 可同时被轮询 / 订阅 / 推送观察——通道只是观察方式，状态真相在服务端；Push 回调必须幂等 + URL 校验 + 签名。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 3 章 · 发现与传输")

# ---------- p71 五个信任边界 ----------
def draw_a2a_trust_boundaries(slide):
    bg(slide)
    eyebrow(slide, "第 7 章 · 安全 · 深潜  TRUST BOUNDARIES")
    title(slide, "五个信任边界，对应五组控制")
    label(slide, 0.55, 1.55, 12.2, "自动协作放大攻击面——每个边界配一组控制，且每组控制都要有可复现的验证证据", SUB, 12, PP_ALIGN.LEFT, True)
    rows = [("发现", "Card 可伪造", "可信域 + 签名 Card"),
            ("身份", "凭据扩散、混淆代理", "带外授权 + audience 限定 + 短期凭证"),
            ("任务", "越权读取 / 取消", "逐操作鉴权（RBAC / ABAC）+ 租户绑定"),
            ("模型", "跨 Agent 提示注入", "指令与数据隔离 + 工具最小权限"),
            ("回调", "伪造回调、重放、洪泛", "URL 校验 + 签名 + 幂等 + 限流")]
    y0 = 2.32; h = 0.6; gap = 0.09
    # 表头
    label(slide, 0.7, y0-0.32, 4.3, "信任边界（风险）", REDINK, 11, PP_ALIGN.LEFT, False, True)
    label(slide, 6.05, y0-0.32, 6.6, "对应控制", DTEAL, 11, PP_ALIGN.LEFT, False, True)
    for i, (b, risk, ctrl) in enumerate(rows):
        yy = y0 + i*(h+gap)
        band(slide, 0.7, yy, 4.55, h, fill=REDBG, edge=REDINK)
        tf = tb(slide, 0.92, yy+0.08, 4.2, 0.48); p = tf.paragraphs[0]
        sr(p.add_run(), b + "　", BODY_F, 12.5, REDINK, bold=True)
        sr(p.add_run(), risk, BODY_F, 11, RGBColor(0x88, 0x3A, 0x2E))
        arrow(slide, 5.28, yy+h/2, 5.98, yy+h/2, color=SUB, w=1.6)
        band(slide, 6.02, yy, 6.6, h, fill=CARD, edge=TEAL)
        tf2 = tb(slide, 6.24, yy+0.11, 6.3, 0.42)
        sr(tf2.paragraphs[0].add_run(), ctrl, BODY_F, 11.5, DTEAL, bold=True)
    band(slide, 0.7, 5.82, 11.93, 1.0, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 5.92, 11.4, 0.85); p = tf.paragraphs[0]
    sr(p.add_run(), "一句话：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "身份验证答「你是谁」，授权答「你能对这个 Task 做什么」——两个问题都不能交给模型猜。", BODY_F, 11.5, INK)
    para(tf, "安全主张要能被测试复现（签名校验日志 / 越权回归 / 注入语料 / 重放测试），不靠口头承诺。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 7 章 · 安全 · 售前速查")

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    figs = [draw_a2a_vs_mcp, draw_a2a_five_objects, draw_a2a_task_states, draw_a2a_delegation,
            draw_a2a_transports, draw_a2a_delivery, draw_a2a_trust_boundaries]
    for fn in figs:
        fn(prs.slides.add_slide(B))
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_a2a_preview.pptx")
    prs.save(out); print(f"saved {len(figs)} ->", out)
