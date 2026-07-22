"""Agent 讲义信息图（第二批 4 张）。复用 kb_draw + agent_figs 的 helper。
python3 agent_figs2.py 生成预览。"""
import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, BODY_F, TITLE_F)
from agent_figs import line, MOD

# ---------- p31 A2A 握手时序 ----------
def draw_a2a_handshake(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 编排模式  A2A")
    title(slide, "A2A：两个 agent 怎么握手——先验名片，再交任务")
    label(slide, 0.55, 1.58, 12.2, "A2A = 跨团队 / 跨厂商的 agent 互相发现与协作的开放协议（2026-03 发布 v1.0）", SUB, 12.5, PP_ALIGN.LEFT, True)
    ax, bx = 2.6, 10.0; ytop = 2.55; ybot = 5.4
    node(slide, ax-1.35, ytop, 2.7, 0.72, "采购 Agent", "发起方", "", zh_size=14, en_size=10, desc_size=9)
    node(slide, bx-1.35, ytop, 2.7, 0.72, "航司 Agent", "被委托方", "", fill=CORAL, edge=ORANGE, zh_c=ODARK, en_c=ODARK, zh_size=14, en_size=10, desc_size=9)
    line(slide, ax, ytop+0.72, ax, ybot, color=LINE, w=1.5, dashed=True)
    line(slide, bx, ytop+0.72, bx, ybot, color=LINE, w=1.5, dashed=True)
    steps = [(ytop+1.15, "① 读签名 Agent Card", "确认对方是谁、会什么（签名防伪装）", TEAL),
             (ytop+2.05, "② 委托任务", "经标准接口把任务交过去", ORANGE),
             (ytop+2.75, "③ 标准报文往返、全程可审计", "", TEAL)]
    for y, t, d, col in steps:
        arrow(slide, ax+0.1, y, bx-0.1, y, color=col, w=1.75)
        label(slide, ax+0.2, y-0.36, bx-ax-0.4, t, DTEAL if col == TEAL else ODARK, 12, PP_ALIGN.CENTER, False, True)
        if d:
            label(slide, ax+0.2, y+0.06, bx-ax-0.4, d, SUB, 10.5, PP_ALIGN.CENTER, True)
    band(slide, 0.7, 5.75, 11.93, 1.15, fill=CARD, edge=LINE)
    tf = tb(slide, 1.0, 5.88, 11.4, 1.0)
    para(tf, "MCP 是「agent 用工具」的 USB-C；A2A 是「agent 找 agent」的名片 + 合同——一个对物、一个对人。", BODY_F, 12, DTEAL, bold=True, first=True)
    para(tf, "现状(2026)：Linux Foundation 治理，150+ 组织采用、五种语言 SDK。组织内部编排用框架就够，跨组织协作看 A2A。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 3 章 · 编排模式")

# ---------- p80 四种记忆分类 ----------
def draw_four_memory(slide):
    bg(slide)
    eyebrow(slide, "第 8 章 · 记忆系统  FOUR TYPES")
    title(slide, "四种记忆：像人一样，分工记不同的东西")
    label(slide, 0.55, 1.58, 12.2, "全塞脑子里会过载，所以分层存取——小核心常驻上下文 + 向量检索按需取 + 显式遗忘策略", SUB, 12, PP_ALIGN.LEFT, True)
    cards = [(0.9, 2.35, "工作记忆", "working", "当前上下文", "像手头便签", TEAL, DTEAL, CARD),
             (7.0, 2.35, "情景记忆", "episodic", "过往发生的事件", "像日记", TEAL, DTEAL, CARD),
             (0.9, 4.3, "语义记忆", "semantic", "事实与用户偏好", "像常识卡片", ORANGE, ODARK, CORAL),
             (7.0, 4.3, "程序记忆", "procedural", "习得的自身指令", "像肌肉记忆", ORANGE, ODARK, CORAL)]
    for x, y, zh, en, desc, ana, edge, tc, fill in cards:
        band(slide, x, y, 5.45, 1.75, fill=fill, edge=edge)
        tf = tb(slide, x+0.32, y+0.2, 4.9, 1.45)
        p = tf.paragraphs[0]
        sr(p.add_run(), zh, BODY_F, 14, tc, bold=True)
        sr(p.add_run(), f"  {en}", BODY_F, 11, edge)
        para(tf, desc, BODY_F, 12, INK)
        para(tf, ana, BODY_F, 11, SUB, italic=True)
    tf2 = tb(slide, 0.7, 6.4, 12.0, 0.5)
    para(tf2, "记忆 ≠ 把历史全塞上下文——记忆工程的本质是「挑什么值得记、取什么值得带」。", BODY_F, 12, DTEAL, bold=True, first=True)
    footer(slide, MOD, "第 8 章 · 记忆系统")

# ---------- p81 四层状态分层 ----------
def draw_four_state(slide):
    bg(slide)
    eyebrow(slide, "第 8 章 · 记忆系统 · 深潜  FOUR LAYERS")
    title(slide, "四层状态：按生命周期和归属分开，别混成一坨聊天记录")
    layers = [("工作状态", "当前目标 / 计划 / 中间产物", "任务结束即归档", "Agent 运行时"),
              ("会话状态", "短期上下文 / 用户选择", "会话期", "Session 存储"),
              ("长期记忆", "稳定偏好 / 经确认的经验", "带 TTL、可更正删除", "记忆服务"),
              ("业务事实", "订单 / 合同 / 权威数据", "长期权威", "业务库")]
    y0 = 2.2; h = 0.8; gap = 0.13
    for i, (zh, cont, life, own) in enumerate(layers):
        y = y0 + i*(h+gap)
        solid = (i == 3)
        band(slide, 0.9, y, 11.5, h, fill=(DTEAL if solid else CARD), edge=(NAVY if solid else TEAL))
        tf = tb(slide, 1.2, y+0.14, 11.0, h-0.22)
        p = tf.paragraphs[0]
        sr(p.add_run(), f"{i+1}. {zh}", BODY_F, 13.5, (WHITE if solid else DTEAL), bold=True)
        sr(p.add_run(), f"　{cont}", BODY_F, 11.5, (PALE if solid else INK))
        sr(p.add_run(), f"　·　{life}　·　归 {own}", BODY_F, 10.5, (PALE if solid else SUB))
    band(slide, 0.7, 5.95, 11.93, 0.95, fill=CORAL, edge=ORANGE)
    tf2 = tb(slide, 1.0, 6.06, 11.4, 0.8)
    para(tf2, "铁律：模型可以从四层读，但写回权威事实库或长期记忆，必须过验证。", BODY_F, 12, ODARK, bold=True, first=True)
    para(tf2, "上一页「四种记忆」是记什么(认知分类)；这一页是存哪、归谁管、怎么删(工程分层)——设计评审两张都要过。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 8 章 · 记忆系统")

# ---------- p89 Computer Use 三路线 ----------
def draw_cu_routes(slide):
    bg(slide)
    eyebrow(slide, "第 9 章 · Computer Use  THREE ROUTES")
    title(slide, "Computer Use 三条路线：任务在哪，就走哪条")
    label(slide, 0.55, 1.58, 12.2, "Computer Use = agent 通过「看屏幕截图 + 操作鼠标键盘」直接使用人的界面", SUB, 12.5, PP_ALIGN.LEFT, True)
    cards = [(0.7, "Claude Computer Use", "截图 + 鼠键通用原语", "跨 VM / 远程桌面，唯一全桌面", "跨桌面应用，只有它能覆盖", TEAL, DTEAL, CARD),
             (5.02, "Gemini Computer Use", "DOM 感知", "浏览器任务更稳", "纯网页任务，更稳更便宜", ORANGE, ODARK, CORAL),
             (9.34, "Codex Background CU", "后台批量并行", "2026-04 发布", "批量并行任务看它", DTEAL, NAVY, PALE)]
    for x, name, feat, note, fit, edge, tc, fill in cards:
        band(slide, x, 2.4, 3.98, 3.05, fill=fill, edge=edge)
        tf = tb(slide, x+0.28, 2.62, 3.5, 2.7)
        para(tf, name, BODY_F, 13, tc, bold=True, first=True)
        para(tf, feat, BODY_F, 11.5, edge)
        para(tf, note, BODY_F, 11, INK)
        para(tf, "适合：" + fit, BODY_F, 11, SUB)
    band(slide, 0.7, 5.7, 11.93, 1.2, fill=CARD, edge=LINE)
    tf2 = tb(slide, 1.0, 5.83, 11.4, 1.05)
    para(tf2, "选路线先问一句：任务在哪？", BODY_F, 12.5, DTEAL, bold=True, first=True)
    para(tf2, "纯网页 → DOM 路线更稳更便宜；跨桌面应用 → 只有截图路线能覆盖；批量并行 → Codex 后台模式。", BODY_F, 11, INK)
    para(tf2, "类比：API 集成像给系统装专用接口，Computer Use 像雇了个会用电脑的实习生——不用改造系统，但得盯着点。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 9 章 · Computer Use")

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in [draw_a2a_handshake, draw_four_memory, draw_four_state, draw_cu_routes]:
        fn(prs.slides.add_slide(B))
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_agent2_preview.pptx")
    prs.save(out); print("saved 4 ->", out)
