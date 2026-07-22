"""Agent 讲义信息图（第一批 4 张）。复用 _maintenance/kb_draw.py 的设计系统 helper。
独立运行生成预览：python3 agent_figs.py"""
import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, BODY_F, TITLE_F)

MOD = "Agent 讲义 · 智能体与编排"

def line(slide, x1, y1, x2, y2, color=ORANGE, w=2.0, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w); c.shadow.inherit = False
    if dashed:
        from pptx.oxml.ns import qn
        ln = c.line._get_or_add_ln(); ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return c

def dot(slide, cx, cy, r, color):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx-r), Inches(cy-r), Inches(2*r), Inches(2*r))
    d.fill.solid(); d.fill.fore_color.rgb = color; d.line.color.rgb = color; d.shadow.inherit = False
    return d

# ---------- p16 ReAct 循环 ----------
def draw_react_loop(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 核心组件  PLANNING")
    title(slide, "ReAct：想一步、做一步、看一步，再回头调")
    label(slide, 0.55, 1.58, 12.2, "规划 = 把大任务拆成小步，并用每一步的结果修正下一步", SUB, 12.5, PP_ALIGN.LEFT, True)
    steps = [("思考", "Reasoning", "先写下推理：这步该查/该做什么", True),
             ("行动", "Acting", "调用工具、执行一步操作", False),
             ("观察", "Observation", "看这步的结果，喂给下一轮", False)]
    W = 3.0; y = 2.55; h = 1.35; xs = [1.1, 4.9, 8.7]
    for i, (zh, en, desc, solid) in enumerate(steps):
        if solid:
            node(slide, xs[i], y, W, h, zh, en, desc, fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=WHITE, zh_size=17, en_size=11.5, desc_size=10.5)
        else:
            node(slide, xs[i], y, W, h, zh, en, desc, zh_size=17, en_size=11.5, desc_size=10.5)
    yc = y + h/2
    arrow(slide, xs[0]+W+0.02, yc, xs[1]-0.02, yc, w=2.0)
    arrow(slide, xs[1]+W+0.02, yc, xs[2]-0.02, yc, w=2.0)
    ix = xs[0]+W/2; ox = xs[2]+W/2; yb = y+h+0.55
    line(slide, ox, y+h, ox, yb, color=ORANGE, w=2.0)
    line(slide, ox, yb, ix, yb, color=ORANGE, w=2.0)
    arrow(slide, ix, yb, ix, y+h, color=ORANGE, w=2.0)
    label(slide, ix+0.3, yb-0.32, 6.0, "带着结果进入下一轮，直到任务完成 —— 这条回边就是 agent 的灵魂", ODARK, 11, PP_ALIGN.LEFT, True, True)
    band(slide, 0.7, 5.5, 11.93, 1.42, fill=CARD, edge=LINE)
    tf = tb(slide, 1.0, 5.64, 11.4, 1.2)
    para(tf, "像老练的维修工：不是拿到工单就闷头拆，而是先判断多半哪坏了、拆一步看一眼、随时修正。", BODY_F, 12, DTEAL, bold=True, first=True)
    para(tf, "案例：查「上季度华东销售额下滑原因」→ 查数据 → 对比往期 → 定位异常品类 → 归因，每步都按上一步结果决定下一步。规划能力决定 agent 能接多大的活。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 2 章 · 核心组件")

# ---------- p17 记忆：桌面与档案柜 ----------
def draw_memory_desk(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 核心组件  MEMORY")
    title(slide, "记忆两层：桌面放当下，档案柜管长远")
    label(slide, 0.55, 1.58, 12.2, "记忆设计的核心，就是决定「什么留在桌面、什么进档案柜」", SUB, 12.5, PP_ALIGN.LEFT, True)
    c1 = band(slide, 0.9, 2.5, 5.1, 2.7, fill=CARD, edge=TEAL)
    tf = tb(slide, 1.2, 2.66, 4.6, 2.4)
    para(tf, "短期记忆 = 上下文窗口", BODY_F, 15, DTEAL, bold=True, first=True)
    para(tf, "Context Window · 「桌面」", BODY_F, 11, TEAL)
    para(tf, "· 随对话增长，容量有限", BODY_F, 12, INK)
    para(tf, "· 每轮都占 token，越堆越贵", BODY_F, 12, INK)
    para(tf, "· 东西摊多了反而找不到笔", BODY_F, 12, SUB)
    c2 = band(slide, 7.35, 2.5, 5.1, 2.7, fill=CORAL, edge=ORANGE)
    tf2 = tb(slide, 7.65, 2.66, 4.6, 2.4)
    para(tf2, "长期记忆 = 外部存储", BODY_F, 15, ODARK, bold=True, first=True)
    para(tf2, "文件 / 数据库 / 向量库 · 「档案柜」", BODY_F, 11, ODARK)
    para(tf2, "· 跨会话保留，用时按需取回", BODY_F, 12, INK)
    para(tf2, "· 便宜、可无限扩", BODY_F, 12, INK)
    para(tf2, "· 只把当前要用的调回桌面", BODY_F, 12, SUB)
    arrow(slide, 6.05, 3.55, 7.32, 3.55, color=SUB, w=1.75)
    label(slide, 5.95, 3.2, 1.5, "归档", SUB, 10, PP_ALIGN.CENTER, True)
    arrow(slide, 7.32, 4.15, 6.05, 4.15, color=SUB, w=1.75)
    label(slide, 5.95, 4.2, 1.5, "取回", SUB, 10, PP_ALIGN.CENTER, True)
    band(slide, 0.7, 5.5, 11.93, 1.42, fill=CARD, edge=LINE)
    tf3 = tb(slide, 1.0, 5.64, 11.4, 1.2)
    para(tf3, "案例：客服 agent 记住「这位客户上月投诉过物流」——", BODY_F, 12, DTEAL, bold=True, first=True)
    para(tf3, "会话结束写入客户档案(长期记忆)，下次开场自动取回，而不是一直占着每轮对话的上下文。记忆 ≠ 把历史全塞进桌面。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 2 章 · 核心组件")

# ---------- p38 M×N → M+N ----------
def draw_mxn(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · 工具接入与 MCP")
    title(slide, "MCP 把 M×N 的网状连接，收成 M+N 的标准口")
    label(slide, 0.55, 1.58, 12.2, "USB-C 之前每台设备一种专用充电器；接口统一后，任何充电器充任何设备", SUB, 12.5, PP_ALIGN.LEFT, True)
    # 左：M×N 网状
    label(slide, 0.6, 2.15, 5.8, "没有标准：M 应用 × N 工具 = M×N 套连接器", ODARK, 12, PP_ALIGN.CENTER, False, True)
    appL = [(1.2, 3.0), (1.2, 4.0), (1.2, 5.0)]
    toolL = [(5.4, 3.0), (5.4, 4.0), (5.4, 5.0)]
    for ax, ay in appL:
        for tx, ty in toolL:
            line(slide, ax+0.35, ay, tx-0.35, ty, color=RGB_LIGHT(), w=1.0)
    for ax, ay in appL:
        n = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ax-0.35), Inches(ay-0.28), Inches(0.7), Inches(0.56))
        n.fill.solid(); n.fill.fore_color.rgb = TEAL; n.line.color.rgb = DTEAL; n.shadow.inherit = False
        _c(n, "应用")
    for tx, ty in toolL:
        n = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(tx-0.35), Inches(ty-0.28), Inches(0.7), Inches(0.56))
        n.fill.solid(); n.fill.fore_color.rgb = CARD; n.line.color.rgb = TEAL; n.shadow.inherit = False
        _c(n, "工具", DTEAL)
    # 右：M+N 星型（MCP 中枢）
    label(slide, 7.0, 2.15, 5.8, "有了 MCP：M + N，接一次处处用", DTEAL, 12, PP_ALIGN.CENTER, False, True)
    hub = (9.7, 4.0)
    appR = [(7.6, 3.0), (7.6, 4.0), (7.6, 5.0)]
    toolR = [(11.8, 3.0), (11.8, 4.0), (11.8, 5.0)]
    for ax, ay in appR:
        line(slide, ax+0.35, ay, hub[0]-0.75, hub[1], color=TEAL, w=1.25)
    for tx, ty in toolR:
        line(slide, hub[0]+0.75, hub[1], tx-0.35, ty, color=ORANGE, w=1.25)
    for ax, ay in appR:
        n = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ax-0.35), Inches(ay-0.28), Inches(0.7), Inches(0.56))
        n.fill.solid(); n.fill.fore_color.rgb = TEAL; n.line.color.rgb = DTEAL; n.shadow.inherit = False
        _c(n, "应用")
    for tx, ty in toolR:
        n = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(tx-0.35), Inches(ty-0.28), Inches(0.7), Inches(0.56))
        n.fill.solid(); n.fill.fore_color.rgb = CORAL; n.line.color.rgb = ORANGE; n.shadow.inherit = False
        _c(n, "工具", ODARK)
    hb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(hub[0]-0.75), Inches(hub[1]-0.5), Inches(1.5), Inches(1.0))
    hb.fill.solid(); hb.fill.fore_color.rgb = DTEAL; hb.line.color.rgb = NAVY; hb.shadow.inherit = False
    tf = hb.text_frame; tf.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    sr(p.add_run(), "MCP", BODY_F, 15, WHITE, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    sr(p2.add_run(), "标准中枢", BODY_F, 10, PALE)
    band(slide, 0.7, 5.9, 11.93, 1.0, fill=CARD, edge=LINE)
    tf3 = tb(slide, 1.0, 6.0, 11.4, 0.8)
    para(tf3, "同一个 GitHub MCP server，Claude Desktop、IDE 插件、自研 agent 平台都能直接接上——集成写一次、处处能跑，工具生态从此沉淀成企业资产。", BODY_F, 11.5, INK, first=True)
    footer(slide, MOD, "第 4 章 · 工具接入与 MCP")

# ---------- p41 六层工具契约 ----------
def draw_six_layer(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · 工具接入与 MCP · 深潜")
    title(slide, "生产工具的六层契约：缺哪层，事故从哪层来")
    layers = [("发现", "Discovery", "名称/用途/反例写清"),
              ("输入", "Input", "严格 schema、租户字段"),
              ("授权", "Auth", "主体/scope/资源级策略"),
              ("执行", "Exec", "超时/幂等/限流/审计"),
              ("输出", "Output", "高信号、分页截断、带源"),
              ("错误", "Error", "可重试/不可恢复/需审批")]
    W = 1.83; gap = 0.2; y = 2.35; h = 1.5; x0 = 0.62
    xs = []
    for i, (zh, en, desc) in enumerate(layers):
        x = x0 + i*(W+gap); xs.append(x)
        solid = i in (1, 5)
        if solid:
            node(slide, x, y, W, h, f"{i+1} {zh}", en, desc, fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=WHITE, zh_size=14, en_size=10, desc_size=9.5)
        else:
            node(slide, x, y, W, h, f"{i+1} {zh}", en, desc, zh_size=14, en_size=10, desc_size=9.5)
    yc = y + h/2
    for i in range(5):
        arrow(slide, xs[i]+W+0.01, yc, xs[i+1]-0.01, yc, w=1.6)
    band(slide, 0.62, 4.2, 12.09, 1.35, fill=RGB_REDBG(), edge=ORANGE)
    tf = tb(slide, 0.9, 4.34, 11.6, 1.15)
    para(tf, "缺一层，漏一类事故：", BODY_F, 12.5, ODARK, bold=True, first=True)
    para(tf, "没有租户字段 = 跨客户越权(缺授权)　·　错误不分类 = 模型对着「500」死重试(缺错误)　·　输出不截断 = 一次调用吃掉半个上下文窗口(缺输出)", BODY_F, 11, INK)
    tf2 = tb(slide, 0.7, 5.85, 12.0, 0.9)
    para(tf2, "投入在工具契约上的一小时，常比调一天提示词更提成功率——工具返回的质量决定模型「看到什么」。", BODY_F, 12.5, DTEAL, bold=True, first=True)
    para(tf2, "评审 agent 方案，先翻它的工具清单。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 4 章 · 工具接入与 MCP")

from pptx.dml.color import RGBColor
def RGB_LIGHT(): return RGBColor(0xC9, 0xD8, 0xE2)
def RGB_REDBG(): return RGBColor(0xFB, 0xF0, 0xE6)
def _c(shape, text, color=None):
    from pptx.enum.text import MSO_ANCHOR
    tf = shape.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02); tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    sr(p.add_run(), text, BODY_F, 9.5, color if color else WHITE, bold=True)

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in [draw_react_loop, draw_memory_desk, draw_mxn, draw_six_layer]:
        fn(prs.slides.add_slide(B))
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_agent_preview.pptx")
    prs.save(out)
    print("saved 4 figs ->", out)
