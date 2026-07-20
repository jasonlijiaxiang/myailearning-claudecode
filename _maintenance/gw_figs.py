"""AI-Gateway 讲义信息图（5 张）。复用 kb_draw + agent_figs helper。
嫁接：p6 后(M×N)、p7 后(六件套)、p26 后(容灾)、p35 后(缓存)、p45 后(护栏)。
GW 是页码册；图沿用图系统页脚（无页码），与首批 PE/LLM 一致（接受页码轻微漂移）。
独立预览：python3 gw_figs.py"""
import sys
sys.path.insert(0, "/Users/lijiaxiang/project/myAILearning/_maintenance")
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "AI 网关讲义 · 统一入口与运营管控"
GINK = RGBColor(0x27, 0x50, 0x0A)
REDBG = RGBColor(0xFB, 0xE9, 0xE7)
REDINK = RGBColor(0x9B, 0x2C, 0x2C)

# ---------- p6 M×N 痛点 → 网关 ----------
def draw_gw_mxn(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 是什么 / 为什么  M×N → HUB")
    title(slide, "痛点：一旦不止一个模型，就撞上「多对多」")
    label(slide, 0.55, 1.58, 12.2, "M 个应用 × N 家模型 = M×N 条私接线——没配电箱，每个电器各拉一根线到变电站", SUB, 12, PP_ALIGN.LEFT, True)
    # 左：没网关 mesh
    band(slide, 0.6, 2.35, 5.5, 3.15, fill=REDBG, edge=REDINK)
    label(slide, 0.85, 2.46, 5.0, "没网关：M×N 意大利面式接线", REDINK, 12, PP_ALIGN.LEFT, False, True)
    apps = [(1.05, 3.15), (1.05, 3.9), (1.05, 4.65)]
    mods = [(5.05, 3.15), (5.05, 3.9), (5.05, 4.65)]
    for ax, ay in apps:
        node(slide, ax, ay-0.2, 1.1, 0.42, "应用", "", "", fill=WHITE, edge=SUB, zh_c=INK, zh_size=10.5)
    for mx, my in mods:
        node(slide, mx-0.05, my-0.2, 1.1, 0.42, "模型", "", "", fill=WHITE, edge=SUB, zh_c=INK, zh_size=10.5)
    for ax, ay in apps:
        for mx, my in mods:
            line(slide, ax+1.1, ay, mx-0.05, my, color=RGBColor(0xC9,0x7A,0x70), w=0.9)
    label(slide, 0.85, 5.02, 5.0, "格式 / 鉴权 / 参数 / 限流各不同，改模型要动几十处代码", SUB, 9.5, PP_ALIGN.LEFT, True)
    arrow(slide, 6.15, 3.95, 6.95, 3.95, color=TEAL, w=2.4)
    # 右：有网关 hub
    band(slide, 7.05, 2.35, 5.55, 3.15, fill=GBG, edge=GREEN)
    label(slide, 7.3, 2.46, 5.0, "有网关：M + N，一处配置", GINK, 12, PP_ALIGN.LEFT, False, True)
    gapps = [(7.45, 3.15), (7.45, 3.9), (7.45, 4.65)]
    gmods = [(11.5, 3.15), (11.5, 3.9), (11.5, 4.65)]
    node(slide, 9.15, 3.35, 1.5, 1.0, "AI 网关", "", "总闸+电表", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, zh_size=13, desc_size=9)
    for ax, ay in gapps:
        node(slide, ax, ay-0.2, 1.1, 0.42, "应用", "", "", fill=WHITE, edge=TEAL, zh_c=DTEAL, zh_size=10.5)
        line(slide, ax+1.1, ay, 9.15, 3.85, color=TEAL, w=1.0)
    for mx, my in gmods:
        node(slide, mx-0.05, my-0.2, 1.1, 0.42, "模型", "", "", fill=WHITE, edge=TEAL, zh_c=DTEAL, zh_size=10.5)
        line(slide, 10.65, 3.85, mx-0.05, my, color=TEAL, w=1.0)
    band(slide, 0.6, 5.68, 12.03, 1.14, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.79, 11.5, 1.0); p = tf.paragraphs[0]
    sr(p.add_run(), "案例：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "团队想把简单请求切国产模型省钱、长文切 Claude——不上网关要在几十个服务里改 SDK、重配密钥；", BODY_F, 11.5, INK)
    para(tf, "上网关只改一处配置。这正是「配电箱」的价值：总闸、电表、空开集中管，电器只管用电。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 1 章 · 是什么 / 为什么")

# ---------- p7 网关六件套 ----------
def draw_gw_definition(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 定义  UNIFIED ENTRY + CONTROL")
    title(slide, "AI 网关 = 统一入口 + 运营管控层")
    label(slide, 0.55, 1.58, 12.2, "把横切关注点从应用抽出来，收到一层统一管——配电箱 + 电表 + 空气开关", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 3.9, 2.35, 5.5, 0.62, "应用（对上给统一接口）", "", "", fill=CARD, edge=SUB, zh_c=INK, zh_size=12.5)
    arrow(slide, 6.65, 2.99, 6.65, 3.32, color=SUB, w=1.7)
    # 网关层 + 六件套
    band(slide, 0.7, 3.35, 11.93, 1.75, fill=PALE, edge=TEAL)
    label(slide, 0.95, 3.46, 6.0, "AI 网关 · LLM 专属六件套", DTEAL, 13, PP_ALIGN.LEFT, False, True)
    six = ["统一接入 / 协议转换", "按 token 限流", "成本归集", "语义缓存", "护栏 / 内容安全", "AI 可观测"]
    for i, s in enumerate(six):
        col = i % 3; row = i // 3
        x = 1.0 + col*3.95; y = 3.9 + row*0.56
        node(slide, x, y, 3.7, 0.48, s, "", "", fill=WHITE, edge=TEAL, zh_c=DTEAL, zh_size=11)
    arrow(slide, 6.65, 5.12, 6.65, 5.45, color=SUB, w=1.7)
    node(slide, 3.4, 5.48, 6.5, 0.62, "各家大模型 / Agent / 工具（对下接所有供应商）", "", "", fill=CARD, edge=SUB, zh_c=INK, zh_size=12)
    band(slide, 0.7, 6.28, 11.93, 0.56, fill=GBG, edge=GREEN)
    tf = tb(slide, 1.0, 6.35, 11.4, 0.46); p = tf.paragraphs[0]
    sr(p.add_run(), "类比：", BODY_F, 11, GINK, bold=True)
    sr(p.add_run(), "统一进线=统一接口，分路=路由，电表=按 token 计费，空开=限流与护栏——电器（应用）只管用电。", BODY_F, 11, INK)
    footer(slide, MOD, "第 1 章 · 是什么 / 为什么")

# ---------- p26 容灾 fallback 链 ----------
def draw_gw_fallback(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 路由 · 容灾  FALLBACK")
    title(slide, "容灾：fallback、重试与「爆炸半径」")
    label(slide, 0.55, 1.58, 12.2, "让「某家挂了」不等于「业务挂了」——像家里的双回路供电 + 空气开关", SUB, 12, PP_ALIGN.LEFT, True)
    # fallback 链
    chain = [("gpt-4o", "主模型", TEAL, DTEAL, CARD), ("azure-gpt-4o", "备用供应商", ORANGE, ODARK, CORAL), ("claude", "再兜底", TEAL, DTEAL, CARD)]
    xs = [0.9, 5.0, 9.1]; y = 2.7; W = 3.0; h = 1.0
    conds = ["限流", "超时"]
    for i, (name, role, tc, ec, fill) in enumerate(chain):
        node(slide, xs[i], y, W, h, name, role, "", fill=fill, edge=ec, zh_c=tc, sub_c=SUB, en_c=tc, zh_size=15, en_size=10.5)
        if i < 2:
            arrow(slide, xs[i]+W+0.02, y+h/2, xs[i+1]-0.02, y+h/2, color=ODARK, w=2.0)
            label(slide, xs[i]+W-0.15, y+h/2-0.42, 1.35, "（" + conds[i] + "）", ODARK, 10, PP_ALIGN.CENTER, True, True)
    label(slide, 0.9, 3.85, 11.5, "对调用方全程无感——只在日志里看到「这次走了备用」；多区域部署再叠一层地域容灾", SUB, 10.5, PP_ALIGN.LEFT, True)
    # 三机制
    mechs = [("Fallback", "主挂了自动改走备用模型 / 供应商", TEAL, DTEAL, CARD),
             ("指数退避重试", "失败后隔一会儿再试，间隔越拉越长", TEAL, DTEAL, CARD),
             ("熔断 = 控爆炸半径", "连续失败的后端隔离一段时间；某路短路，空开只跳那一路", ORANGE, ODARK, CORAL)]
    mx = [0.7, 4.68, 8.66]; my = 4.35; mw = 3.9; mh = 1.15
    for x, (h1, h2, tc, ec, fill) in zip(mx, mechs):
        band(slide, x, my, mw, mh, fill=fill, edge=ec)
        tf = tb(slide, x+0.22, my+0.12, mw-0.44, 0.95)
        para(tf, h1, BODY_F, 12, tc, bold=True, first=True)
        para(tf, h2, BODY_F, 10.5, SUB)
    band(slide, 0.7, 5.72, 11.93, 1.1, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 5.83, 11.4, 0.95); p = tf.paragraphs[0]
    sr(p.add_run(), "一句话：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "主线停电自动切备用线（fallback）；某路短路，空开只跳那一路、不牵连全屋（熔断）——这就是把「爆炸半径」关进笼子。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 3 章 · 路由 · 容灾")

# ---------- p35 缓存 精确 vs 语义 ----------
def draw_gw_cache(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · 成本治理  EXACT vs SEMANTIC CACHE")
    title(slide, "缓存：精确缓存 vs 语义缓存")
    label(slide, 0.55, 1.58, 12.2, "给网关配一个「常见问题抽屉」——意思相近就直接返回，不必再花钱调模型", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 2.9, fill=CARD, edge=SUB)
    label(slide, 0.85, 2.52, 5.6, "精确缓存 —— 一字不差才命中", RGBColor(0x55,0x55,0x52), 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 3.15, 5.3, 0.6, "「怎么退货」==「怎么退货」", "", "", fill=WHITE, edge=SUB, zh_c=INK, zh_size=11.5)
    tf = tb(slide, 0.9, 3.95, 5.5, 1.2)
    para(tf, "· 请求完全一样才命中", BODY_F, 11, SUB, first=True)
    para(tf, "· 换一个说法就落空、又得调模型", BODY_F, 11, ODARK)
    para(tf, "· 命中率低，省得有限", BODY_F, 11, SUB)
    band(slide, 6.75, 2.4, 6.0, 2.9, fill=GBG, edge=GREEN)
    label(slide, 7.0, 2.52, 5.6, "语义缓存 —— 意思相近就命中", GINK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.15, 2.5, 0.6, "「怎么退货」", "", "", fill=WHITE, edge=GREEN, zh_c=GINK, zh_size=11)
    node(slide, 10.15, 3.15, 2.5, 0.6, "「退货流程是啥」", "", "", fill=WHITE, edge=GREEN, zh_c=GINK, zh_size=11)
    line(slide, 9.62, 3.45, 10.12, 3.45, color=GREEN, w=1.4, dashed=True)
    label(slide, 9.15, 3.5, 1.6, "≈ 同一问题", GINK, 9, PP_ALIGN.CENTER, True)
    tf2 = tb(slide, 7.05, 3.95, 5.5, 1.2)
    para(tf2, "· 请求算成向量，比相似度命中", BODY_F, 11, SUB, first=True)
    para(tf2, "· 省钱又降延迟（命中即毫秒级返回）", BODY_F, 11, GINK, bold=True)
    para(tf2, "· 实现：Embeddings 算向量 + 向量库近似匹配", BODY_F, 11, SUB)
    band(slide, 0.6, 5.5, 12.03, 1.32, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.62, 11.5, 1.15); p = tf3.paragraphs[0]
    sr(p.add_run(), "类比 · 客服的「常见问题」抽屉：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "「怎么退货」和「退货流程是啥」其实是同一个问题，客服不必每次重新组织语言，直接翻抽屉。", BODY_F, 11.5, INK)
    para(tf3, "案例：Azure APIM llm-semantic-cache 用 Embeddings + Managed Redis/RediSearch 按相似度阈值命中。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 4 章 · 成本治理")

# ---------- p45 护栏 pre/post ----------
def draw_gw_guardrails(slide):
    bg(slide)
    eyebrow(slide, "第 5 章 · 安全 · 护栏  PRE / POST HOOKS")
    title(slide, "护栏：请求前后各拦一道")
    label(slide, 0.55, 1.58, 12.2, "进门搜身 + 出门过检——两道关，缺一不可；一处挂载，所有走网关的应用都受保护", SUB, 12, PP_ALIGN.LEFT, True)
    # 流水线：应用 → pre → 模型 → post → 应用
    node(slide, 0.7, 3.15, 1.7, 0.9, "应用", "请求", "", fill=CARD, edge=SUB, zh_c=INK, sub_c=SUB, zh_size=13, desc_size=9)
    node(slide, 3.0, 3.05, 2.2, 1.1, "pre 护栏", "进模型前", "搜身", fill=CORAL, edge=ORANGE, zh_c=ODARK, sub_c=ODARK, en_c=ODARK, zh_size=13, en_size=9.5, desc_size=9)
    node(slide, 5.75, 3.15, 1.9, 0.9, "大模型", "", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, zh_size=13)
    node(slide, 8.2, 3.05, 2.2, 1.1, "post 护栏", "出模型后", "过检", fill=CORAL, edge=ORANGE, zh_c=ODARK, sub_c=ODARK, en_c=ODARK, zh_size=13, en_size=9.5, desc_size=9)
    node(slide, 10.95, 3.15, 1.7, 0.9, "应用", "响应", "", fill=CARD, edge=SUB, zh_c=INK, sub_c=SUB, zh_size=13, desc_size=9)
    for x1, x2 in [(2.4, 2.98), (5.22, 5.73), (7.67, 8.18), (10.42, 10.93)]:
        arrow(slide, x1, 3.6, x2, 3.6, color=SUB, w=1.6)
    band(slide, 2.5, 4.4, 3.2, 1.15, fill=RGBColor(0xFF,0xF6,0xEA), edge=ORANGE)
    tf = tb(slide, 2.7, 4.5, 2.85, 1.0)
    para(tf, "pre（进模型前）", BODY_F, 11, ODARK, bold=True, first=True)
    para(tf, "PII 脱敏 · 注入检测 · 敏感话题 · 越权工具拦截", BODY_F, 10.5, SUB)
    band(slide, 7.7, 4.4, 3.2, 1.15, fill=RGBColor(0xFF,0xF6,0xEA), edge=ORANGE)
    tf2 = tb(slide, 7.9, 4.5, 2.85, 1.0)
    para(tf2, "post（出模型后）", BODY_F, 11, ODARK, bold=True, first=True)
    para(tf2, "有害内容过滤 · 机密信息回收 · 格式 / 合规校验", BODY_F, 10.5, SUB)
    label(slide, 4.9, 4.62, 3.6, "命中就拦截 / 改写 / 告警", SUB, 10, PP_ALIGN.CENTER, True)
    band(slide, 0.7, 5.85, 11.93, 0.97, fill=PALE, edge=LINE)
    tf3 = tb(slide, 1.0, 5.95, 11.4, 0.82); p = tf3.paragraphs[0]
    sr(p.add_run(), "一处挂，处处生效：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "Higress 内置注入检测与敏感识别、LiteLLM 有零成本 PII/注入护栏——挂在网关，所有应用都受保护。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 5 章 · 安全 · 护栏")

# ---------- p19 虚拟密钥 ----------
def draw_gw_virtual_key(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 虚拟密钥  VIRTUAL KEY")
    title(slide, "虚拟密钥：一把逻辑钥匙管多把真钥匙")
    label(slide, 0.55, 1.58, 12.2, "真 key 只存在网关里，发给应用的是「房卡」——只能开你这间、到期作废、丢了立刻注销", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 0.7, 3.3, 2.2, 0.95, "应用 / 团队", "拿到虚拟 key", "（房卡）", fill=CARD, edge=TEAL, zh_c=DTEAL, sub_c=SUB, en_c=TEAL, zh_size=13, en_size=10, desc_size=9.5)
    arrow(slide, 2.92, 3.77, 3.75, 3.77, color=TEAL, w=1.8); label(slide, 2.85, 3.35, 1.0, "虚拟 key", TEAL, 9, PP_ALIGN.CENTER, True)
    node(slide, 3.8, 3.15, 3.0, 1.25, "AI 网关", "真 key 保险箱", "真 key 只存在这里", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=PALE, zh_size=15, en_size=10, desc_size=9)
    for j, name in enumerate(["OpenAI", "Anthropic", "国产模型"]):
        yy = 3.0 + j*0.62
        node(slide, 9.5, yy, 2.6, 0.5, name, "", "", fill=WHITE, edge=SUB, zh_c=INK, zh_size=11)
        line(slide, 6.83, 3.77, 9.45, yy+0.25, color=SUB, w=1.1)
    label(slide, 6.9, 3.45, 2.6, "真 key（网关持有）", SUB, 9, PP_ALIGN.CENTER, True)
    band(slide, 3.8, 4.62, 3.0, 0.62, fill=CORAL, edge=ORANGE)
    tf = tb(slide, 3.95, 4.7, 2.75, 0.5); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    sr(p.add_run(), "虚拟 key 上可挂：", BODY_F, 9.5, ODARK, bold=True)
    for j, t in enumerate(["预算", "限速", "可用模型范围", "有效期", "随时吊销"]):
        node(slide, 0.7+j*2.44, 5.5, 2.2, 0.5, t, "", "", fill=CORAL, edge=ORANGE, zh_c=ODARK, zh_size=11)
    label(slide, 0.7, 5.2, 6.0, "虚拟 key 上挂的策略（每把独立）：", DTEAL, 10.5, PP_ALIGN.LEFT, False, True)
    band(slide, 0.6, 6.2, 12.03, 0.62, fill=PALE, edge=LINE)
    tf2 = tb(slide, 0.9, 6.27, 11.5, 0.5); p = tf2.paragraphs[0]
    sr(p.add_run(), "类比 · 酒店房卡 vs 万能钥匙：", BODY_F, 11, DTEAL, bold=True)
    sr(p.add_run(), "前台不给你大门总钥匙，只发一张房卡——换供应商、调预算、封停某个应用，都在网关一处操作，不动真 key。", BODY_F, 11, INK)
    footer(slide, MOD, "第 2 章 · 统一接入")

FIGS = [draw_gw_mxn, draw_gw_definition, draw_gw_fallback, draw_gw_cache, draw_gw_guardrails, draw_gw_virtual_key]
NEW = [draw_gw_virtual_key]

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in FIGS:
        fn(prs.slides.add_slide(B))
    out = "/Users/lijiaxiang/project/myAILearning/_maintenance/_gw_preview.pptx"
    prs.save(out); print(f"saved {len(FIGS)} ->", out)
