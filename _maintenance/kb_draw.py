"""RAG 讲义信息图绘制库：设计系统令牌 + 形状 helper + 各页绘制函数。
所有图共用讲义的深蓝青橙令牌、Cambria/Calibri 字体、16:9 画布。
页脚遵循本讲义惯例：左下模块名，右下章名（本讲义页脚不带页码）。"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

NAVY  = RGBColor(0x0B,0x25,0x40)
INK   = RGBColor(0x1A,0x27,0x33)
SUB   = RGBColor(0x51,0x60,0x6E)
TEAL  = RGBColor(0x12,0x81,0x99)
DTEAL = RGBColor(0x0E,0x4D,0x64)
ORANGE= RGBColor(0xF2,0xA6,0x5A)
ODARK = RGBColor(0x99,0x3C,0x1D)
OINK  = RGBColor(0x4A,0x1B,0x0C)
CARD  = RGBColor(0xF3,0xF8,0xFA)
CORAL = RGBColor(0xFB,0xF0,0xE6)
LINE  = RGBColor(0xD6,0xE1,0xE8)
WHITE = RGBColor(0xFF,0xFF,0xFF)
NEARW = RGBColor(0xFE,0xFE,0xFE)
PALE  = RGBColor(0xEA,0xF3,0xF5)
GREEN = RGBColor(0x2F,0x85,0x5A)
GBG   = RGBColor(0xE6,0xF2,0xEA)
TITLE_F = "Cambria"
BODY_F  = "Calibri"

def bg(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(13.333),Inches(7.5))
    s.fill.solid(); s.fill.fore_color.rgb = NEARW
    s.line.color.rgb = NEARW; s.line.width = Pt(0.75); s.shadow.inherit=False
    return s

def tb(slide,x,y,w,h):
    box = slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = box.text_frame; tf.word_wrap=True
    tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05)
    tf.margin_top=Inches(0.02); tf.margin_bottom=Inches(0.02)
    return tf

def sr(r,text,font=BODY_F,size=12,color=INK,bold=False,italic=False):
    r.text=text; r.font.name=font; r.font.size=Pt(size)
    r.font.color.rgb=color; r.font.bold=bold; r.font.italic=italic
    rPr=r._r.get_or_add_rPr()
    rPr.append(rPr.makeelement(qn('a:ea'),{'typeface':font}))

def para(tf,text,font=BODY_F,size=12,color=INK,bold=False,italic=False,align=PP_ALIGN.LEFT,first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=align; sr(p.add_run(),text,font,size,color,bold,italic); return p

def eyebrow(s,t):
    tf=tb(s,0.55,0.32,11.8,0.4); sr(tf.paragraphs[0].add_run(),t,BODY_F,13,TEAL,bold=True)
def title(s,t,size=27):
    tf=tb(s,0.55,0.66,12.2,0.95); sr(tf.paragraphs[0].add_run(),t,TITLE_F,size,INK,bold=True)
def footer(s,left,right):
    tf=tb(s,0.55,7.08,7,0.3); sr(tf.paragraphs[0].add_run(),left,BODY_F,9,SUB)
    tf2=tb(s,7.9,7.08,4.9,0.3); p=tf2.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT
    sr(p.add_run(),right,BODY_F,9,SUB)

def node(s,x,y,w,h,zh,en,desc,fill=CARD,edge=TEAL,zh_c=DTEAL,sub_c=SUB,en_c=TEAL,zh_size=14,en_size=10.5,desc_size=10):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    try: sp.adjustments[0]=0.08
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    sp.line.color.rgb=edge; sp.line.width=Pt(1.25); sp.shadow.inherit=False
    tf=sp.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    tf.margin_left=Inches(0.08); tf.margin_right=Inches(0.08)
    tf.margin_top=Inches(0.03); tf.margin_bottom=Inches(0.03)
    def _tight(pp):
        pp.alignment=PP_ALIGN.CENTER
        try: pp.line_spacing=1.0; pp.space_before=Pt(0); pp.space_after=Pt(1.5)
        except Exception: pass
    p=tf.paragraphs[0]; _tight(p)
    sr(p.add_run(),zh,BODY_F,zh_size,zh_c,bold=True)
    if en:
        pe=tf.add_paragraph(); _tight(pe)
        sr(pe.add_run(),en,BODY_F,en_size,en_c)
    if desc:
        p2=tf.add_paragraph(); _tight(p2)
        sr(p2.add_run(),desc,BODY_F,desc_size,sub_c)
    return sp

def arrow(s,x1,y1,x2,y2,color=TEAL,w=1.75,dashed=False):
    c=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=color; c.line.width=Pt(w); c.shadow.inherit=False
    ln=c.line._get_or_add_ln()
    if dashed: ln.append(ln.makeelement(qn('a:prstDash'),{'val':'dash'}))
    ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))
    return c

def label(s,x,y,w,t,color=SUB,size=10,align=PP_ALIGN.CENTER,italic=True,bold=False):
    tf=tb(s,x,y,w,0.3); p=tf.paragraphs[0]; p.alignment=align
    sr(p.add_run(),t,BODY_F,size,color,bold=bold,italic=italic)

def band(s,x,y,w,h,fill=CARD,edge=LINE):
    sp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    try: sp.adjustments[0]=0.04
    except Exception: pass
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    sp.line.color.rgb=edge; sp.line.width=Pt(1.0); sp.shadow.inherit=False
    return sp

MOD="RAG 讲义 · 检索增强生成"

# ---------- p57 双泳道架构 ----------
def draw_arch_swimlane(slide):
    bg(slide)
    eyebrow(slide,"第 6 章 · 最小 RAG 管线  ARCHITECTURE")
    title(slide,"RAG 两条线：离线把知识灌进库，在线拿问题去查")
    label(slide,0.55,1.18,12,"① 离线建库 Offline Indexing —— 可以慢、可批量、可贵",DTEAL,12.5,PP_ALIGN.LEFT,False,True)
    off=[("文档","Documents","PDF / 网页 / 工单"),("解析","Parsing · OCR","抽成结构化文本"),
         ("切分","Chunking","按语义分块"),("嵌入","Embedding","算成向量")]
    W=2.2; y=1.60; h=1.0; x0=1.5; step=2.7; ox=[]
    for i,(zh,en,desc) in enumerate(off):
        x=x0+i*step; ox.append(x); node(slide,x,y,W,h,zh,en,desc)
    for i in range(3): arrow(slide,ox[i]+W+0.01,y+h/2,ox[i+1]-0.01,y+h/2)
    vw=2.6; vh=1.05; vx=(13.333-vw)/2; vy=3.18
    node(slide,vx,vy,vw,vh,"向量库","Vector DB","离线写入 · 在线读出",fill=DTEAL,edge=NAVY,zh_c=WHITE,sub_c=PALE,en_c=WHITE)
    arrow(slide,ox[3]+W/2,y+h,vx+vw*0.62,vy,color=TEAL)
    label(slide,vx+vw+0.05,2.55,1.4,"写入 index",TEAL,10,PP_ALIGN.LEFT,True)
    label(slide,0.55,4.55,12,"② 在线检索 Online Serving —— 要快、要稳、要省",ODARK,12.5,PP_ALIGN.LEFT,False,True)
    on=[("提问","Query","用户问题"),("检索","Retrieval","召回候选块"),("重排","Reranking","精排 Top-K"),
        ("生成","Generation","据证据作答"),("答案","Answer","带出处")]
    W2=1.95; y2=5.0; h2=1.0; x0b=1.09; step2=2.3; nx=[]
    for i,(zh,en,desc) in enumerate(on):
        x=x0b+i*step2; nx.append(x)
        if zh=="生成":
            node(slide,x,y2,W2,h2,zh,en,desc,fill=ORANGE,edge=ODARK,zh_c=OINK,sub_c=OINK,en_c=OINK)
        else:
            node(slide,x,y2,W2,h2,zh,en,desc,fill=CORAL,edge=ORANGE,zh_c=ODARK,sub_c=SUB,en_c=ODARK)
    for i in range(4): arrow(slide,nx[i]+W2+0.01,y2+h2/2,nx[i+1]-0.01,y2+h2/2,color=ORANGE)
    arrow(slide,vx+vw*0.42,vy+vh,nx[1]+W2/2,y2,color=DTEAL,dashed=True)
    label(slide,nx[1]+W2+0.05,4.62,1.4,"读取 query",DTEAL,10,PP_ALIGN.LEFT,True)
    tf=tb(slide,0.55,6.55,12.2,0.5)
    sr(tf.paragraphs[0].add_run(),"同一个向量库，两种脾气：离线追求「灌得全」，在线追求「查得快」。",BODY_F,12.5,INK,bold=True)
    footer(slide,MOD,"第 6 章 · 最小 RAG 管线")

# ---------- p7 三步流程 ----------
def draw_flow_3step(slide):
    bg(slide)
    eyebrow(slide,"第 1 章 · 是什么/为什么  HOW IT WORKS")
    title(slide,"RAG 三步：先检索、再增强、后生成")
    label(slide,0.55,1.6,12.2,"把「开卷考试」拆成三步：检索找资料 → 增强喂资料 → 生成照资料答",SUB,12.5,PP_ALIGN.LEFT,True)
    steps=[("① 检索","Retrieval","把用户问题拿去知识库，找出最相关的若干片段"),
           ("② 增强","Augment","把这些片段和问题拼成新提示词，一起交给大模型"),
           ("③ 生成","Generation","模型只依据给定片段作答，可逐句标注来源")]
    W=3.6; y=2.55; h=1.75; xs=[0.77,4.87,8.97]
    for i,(zh,en,desc) in enumerate(steps):
        solid = (i==1)
        if solid:
            node(slide,xs[i],y,W,h,zh,en,desc,fill=TEAL,edge=DTEAL,zh_c=WHITE,sub_c=PALE,en_c=WHITE,zh_size=18,en_size=12,desc_size=11)
        else:
            node(slide,xs[i],y,W,h,zh,en,desc,zh_size=18,en_size=12,desc_size=11)
    yc=y+h/2
    arrow(slide,xs[0]+W+0.02,yc,xs[1]-0.02,yc,w=2.0)
    arrow(slide,xs[1]+W+0.02,yc,xs[2]-0.02,yc,w=2.0)
    band(slide,0.7,4.62,11.93,1.5,fill=CARD,edge=LINE)
    tf=tb(slide,1.0,4.78,11.4,1.2)
    para(tf,"Retrieval → Augmented → Generation：名字本身就是流程。",BODY_F,13,DTEAL,bold=True,first=True)
    para(tf,"最容易被忽略的是中间的「增强 Augment」——正是它把检索到的资料喂进模型。缺了这步，就退回「模型凭记忆答」，也就退回了一本正经编造的老问题。",BODY_F,11.5,SUB)
    footer(slide,MOD,"第 1 章 · 是什么/为什么")

# ---------- p16 语义空间散点（直觉） ----------
def draw_semantic_space(slide):
    from pptx.dml.color import RGBColor as _R
    bg(slide)
    eyebrow(slide,"第 2 章 · 向量检索与 Embedding  INTUITION")
    title(slide,"把语义变成坐标：意思相近，坐标就相近")
    tf=tb(slide,0.6,1.62,4.9,1.0)
    para(tf,"Embedding 模型把一段文字压成一串数字（向量），相当于在一个高维空间里给它分配一个坐标。",BODY_F,12.5,INK,first=True)
    steps=[("一段文字","例：“怎么退货”"),("Embedding 模型","文字 → 向量"),("空间里一个点","[0.12, -0.4, …]")]
    nx=0.9; nw=3.7; nh=0.7; ys=[]
    for i,(zh,en) in enumerate(steps):
        y=2.75+i*1.02; ys.append(y)
        node(slide,nx,y,nw,nh,zh,"",en,zh_size=12.5,desc_size=10)
    for i in range(2): arrow(slide,nx+nw/2,ys[i]+nh,nx+nw/2,ys[i+1],color=TEAL)
    tf2=tb(slide,0.6,5.98,4.9,0.9)
    para(tf2,"检索 = 在这个空间里，找离问题最近的邻居。",BODY_F,12.5,DTEAL,bold=True,first=True)
    fx,fy,fw,fh=5.85,1.62,6.95,4.75
    fr=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(fx),Inches(fy),Inches(fw),Inches(fh))
    try: fr.adjustments[0]=0.03
    except Exception: pass
    fr.fill.solid(); fr.fill.fore_color.rgb=_R(0xF7,0xFB,0xFC)
    fr.line.color.rgb=_R(0x9F,0xB3,0xC0); fr.line.width=Pt(1.0); fr.shadow.inherit=False
    lnf=fr.line._get_or_add_ln(); lnf.append(lnf.makeelement(qn('a:prstDash'),{'val':'dash'}))
    label(slide,fx+0.18,fy+0.08,4.5,"语义空间 semantic space（距离 ≈ 语义相似度）",SUB,11,PP_ALIGN.LEFT,True)
    def dot(cx,cy,zh,en,color,tc):
        d=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(cx-0.09),Inches(cy-0.09),Inches(0.18),Inches(0.18))
        d.fill.solid(); d.fill.fore_color.rgb=color; d.line.color.rgb=color; d.shadow.inherit=False
        t=tb(slide,cx+0.13,cy-0.15,1.7,0.3); p=t.paragraphs[0]
        sr(p.add_run(),zh,BODY_F,11,tc,bold=True)
        if en: sr(p.add_run()," "+en,BODY_F,9,SUB)
    for cx,cy,zh,en in [(7.0,2.6,"狗","dog"),(7.75,2.42,"犬","canine"),(7.2,3.25,"小狗","puppy"),(8.75,2.95,"猫","cat"),(9.35,3.5,"幼猫","kitten")]:
        dot(cx,cy,zh,en,TEAL,DTEAL)
    for cx,cy,zh,en in [(10.7,5.05,"汽车","car"),(11.45,4.8,"轿车","auto"),(10.95,5.6,"卡车","truck")]:
        dot(cx,cy,zh,en,ORANGE,ODARK)
    qx,qy=8.2,3.9
    dq=slide.shapes.add_shape(MSO_SHAPE.OVAL,Inches(qx-0.11),Inches(qy-0.11),Inches(0.22),Inches(0.22))
    dq.fill.solid(); dq.fill.fore_color.rgb=NAVY; dq.line.color.rgb=NAVY; dq.shadow.inherit=False
    arrow(slide,qx+0.02,qy-0.08,8.68,3.05,color=NAVY,w=1.5,dashed=True)
    label(slide,qx-1.75,qy+0.16,2.4,"新问题：宠物怎么养",NAVY,10,PP_ALIGN.LEFT,False,True)
    label(slide,8.5,3.28,2.0,"最近邻 nearest",NAVY,9.5,PP_ALIGN.LEFT,True)
    label(slide,6.55,3.62,2.6,"动物：意思相近，挤成一团",DTEAL,9.5,PP_ALIGN.LEFT,True)
    label(slide,10.2,6.0,2.6,"交通：另一个语义区，离得远",ODARK,9.5,PP_ALIGN.LEFT,True)
    footer(slide,MOD,"第 2 章 · 向量检索与 Embedding")

# ---------- p36 两阶段检索漏斗 ----------
def draw_two_stage_funnel(slide):
    bg(slide)
    eyebrow(slide,"第 4 章 · 重排序 Reranking  TWO-STAGE")
    title(slide,"两阶段检索：先粗召回一大批，再精排一小撮")
    label(slide,0.55,1.6,12.2,"reranking 的价值，全在「两阶段」这个分工里",SUB,12.5,PP_ALIGN.LEFT,True)
    blocks=[(1.0,2.15,2.9,1.5,"候选池","全部文档块","~1,000,000",CARD,TEAL,DTEAL),
            (5.0,2.32,2.5,1.16,"① 召回","Recall","Top-50",CARD,TEAL,DTEAL),
            (8.7,2.48,2.1,0.84,"② 精排","Rerank","Top-5",ORANGE,ODARK,OINK)]
    cx=[]
    for x,y,w,h,zh,en,num,fill,edge,tc in blocks:
        node(slide,x,y,w,h,zh,en,num,fill=fill,edge=edge,zh_c=tc,sub_c=tc,en_c=tc,zh_size=13,en_size=10.5,desc_size=13)
        cx.append((x,y,w,h))
    mid=lambda b:(b[1]+b[3]/2)
    arrow(slide,cx[0][0]+cx[0][2]+0.05,mid(cx[0]),cx[1][0]-0.05,mid(cx[1]),w=2.0)
    arrow(slide,cx[1][0]+cx[1][2]+0.05,mid(cx[1]),cx[2][0]-0.05,mid(cx[2]),w=2.0)
    label(slide,cx[0][0]+cx[0][2]+0.0,mid(cx[0])-0.42,1.1,"召回",TEAL,10,PP_ALIGN.CENTER,True,True)
    label(slide,cx[1][0]+cx[1][2]+0.0,mid(cx[1])-0.42,1.1,"精排",ODARK,10,PP_ALIGN.CENTER,True,True)
    label(slide,cx[2][0]+cx[2][2]+0.15,mid(cx[2])-0.12,2.2,"→ 交给大模型",SUB,11,PP_ALIGN.LEFT,True)
    band(slide,0.7,4.2,5.85,1.85,fill=CARD,edge=TEAL)
    tf=tb(slide,0.98,4.36,5.35,1.55)
    para(tf,"① 召回 Recall — 快、便宜、糙",BODY_F,12.5,DTEAL,bold=True,first=True)
    para(tf,"向量 / BM25 一次捞一大批，宁可多召、不能漏。看的是召回率 recall：该找到的别落下。",BODY_F,11,SUB)
    band(slide,6.78,4.2,5.85,1.85,fill=CORAL,edge=ORANGE)
    tf2=tb(slide,7.06,4.36,5.35,1.55)
    para(tf2,"② 精排 Rerank — 准、贵、慢",BODY_F,12.5,ODARK,bold=True,first=True)
    para(tf2,"交叉编码器逐条细读 query+doc，重排出最相关的少数。看的是精确率 precision：排在前面的要对味。",BODY_F,11,SUB)
    tf3=tb(slide,0.7,6.2,12.2,0.4)
    para(tf3,"宽召回 + 精排，是当下 RAG 提质的标准组合拳。",BODY_F,12.5,INK,bold=True,first=True)
    footer(slide,MOD,"第 4 章 · 重排序 Reranking")

# ---------- 在线检索链路 6 步（样板演示，非入库） ----------
def draw_online_flow(slide):
    bg(slide)
    eyebrow(slide,"RAG 在线检索链路  ONLINE RETRIEVAL")
    title(slide,"从提问到答案，RAG 检索链路要过 6 道关")
    label(slide,0.55,1.60,12.2,"输入：用户的自然语言问题　　→　　输出：带出处、可核查的答案",SUB,12,PP_ALIGN.LEFT,True)
    steps=[("查询改写","Query Rewrite","补全指代、拆解子问题",False),
           ("混合检索","Hybrid Search","向量+关键词并行召回",True),
           ("重排","Reranking","交叉编码器精排 Top-K",False),
           ("上下文组装","Context Build","拼进 Prompt、控 token",False),
           ("生成","Generation","LLM 依据证据作答",True),
           ("引用核对","Citation Check","附出处、拦截无据",False)]
    W=1.8; gap=0.286; y=2.7; h=1.4; xs=[]
    for i,(zh,en,desc,solid) in enumerate(steps):
        x=0.55+i*(W+gap); xs.append(x)
        if solid: node(slide,x,y,W,h,zh,en,desc,fill=TEAL,edge=DTEAL,zh_c=WHITE,sub_c=PALE,en_c=WHITE)
        else: node(slide,x,y,W,h,zh,en,desc)
    yc=y+h/2
    for i in range(5): arrow(slide,xs[i]+W+0.01,yc,xs[i+1]-0.01,yc)
    tf=tb(slide,0.55,4.55,12.2,1.0)
    para(tf,"为什么每一步都不能省：",BODY_F,12,DTEAL,bold=True,first=True)
    para(tf,"「检索」环节是朴素 RAG 的失败高发区——查询改写、混合检索、重排就是把失败一层层压下来。",BODY_F,11.5,SUB)
    footer(slide,MOD,"在线链路总览")

# ---------- p37 双 vs 交叉编码器 ----------
def draw_encoder_compare(slide):
    bg(slide)
    eyebrow(slide,"第 4 章 · 重排序 Reranking  BI vs CROSS")
    title(slide,"双编码器 vs 交叉编码器：准与快的取舍")
    band(slide,0.6,1.62,12.15,1.98,fill=CARD,edge=TEAL)
    label(slide,0.85,1.74,11.5,"双编码器 Bi-Encoder —— 问题、文档各自独立编码，快、可预建索引",DTEAL,12.5,PP_ALIGN.LEFT,False,True)
    node(slide,1.0,2.28,1.7,0.58,"问题","Query","",fill=WHITE,edge=TEAL,zh_size=12.5,en_size=9.5)
    node(slide,1.0,3.02,1.7,0.58,"文档","Document","",fill=WHITE,edge=TEAL,zh_size=12.5,en_size=9.5)
    node(slide,3.3,2.5,2.4,0.9,"各自编码成向量","independent","事先算好、可预建索引",fill=CARD,edge=TEAL,zh_size=12.5,en_size=9.5,desc_size=10)
    node(slide,6.3,2.5,2.3,0.9,"只比向量距离","cosine","查询时才算一次",fill=TEAL,edge=DTEAL,zh_c=WHITE,sub_c=PALE,en_c=WHITE,zh_size=12.5,en_size=9.5,desc_size=10)
    arrow(slide,2.7,2.57,3.28,2.75,color=TEAL,w=1.5)
    arrow(slide,2.7,3.31,3.28,3.15,color=TEAL,w=1.5)
    arrow(slide,5.7,2.95,6.28,2.95,color=TEAL,w=1.75)
    label(slide,8.8,2.7,3.8,"能扫全库、快，但略粗",DTEAL,11,PP_ALIGN.LEFT,True)
    band(slide,0.6,3.9,12.15,1.95,fill=CORAL,edge=ORANGE)
    label(slide,0.85,4.02,11.5,"交叉编码器 Cross-Encoder —— 问题+文档拼一起联合打分，准、只精排少量",ODARK,12.5,PP_ALIGN.LEFT,False,True)
    node(slide,1.0,4.6,2.7,0.92,"问题 ＋ 文档","Query + Doc","拼在一起",fill=WHITE,edge=ORANGE,zh_c=ODARK,zh_size=12.5,en_size=9.5,desc_size=10)
    node(slide,4.5,4.6,2.7,0.92,"模型联合打分","Cross-Encoder","捕捉细粒度关联",fill=ORANGE,edge=ODARK,zh_c=OINK,sub_c=OINK,en_c=OINK,zh_size=12.5,en_size=9.5,desc_size=10)
    node(slide,8.0,4.6,2.2,0.92,"精准相关分","score","",fill=WHITE,edge=ORANGE,zh_c=ODARK,zh_size=12.5,en_size=9.5)
    arrow(slide,3.7,5.06,4.48,5.06,color=ORANGE,w=1.75)
    arrow(slide,7.2,5.06,7.98,5.06,color=ORANGE,w=1.75)
    label(slide,10.4,4.75,2.3,"每候选实时算一次，只排少量，准但慢",ODARK,10,PP_ALIGN.LEFT,True)
    tf=tb(slide,0.6,6.08,12.15,0.7); p=tf.paragraphs[0]
    sr(p.add_run(),"分工：",BODY_F,12,INK,bold=True)
    sr(p.add_run(),"双编码器扫全库粗召回（快），交叉编码器精排少量候选（准）——不是替代，是流水线的上下游。",BODY_F,11.5,SUB)
    footer(slide,MOD,"第 4 章 · 重排序 Reranking")

# ---------- p48 四层证明 ----------
def draw_four_layers(slide):
    bg(slide)
    eyebrow(slide,"第 5 章 · 常见评估方法  FOUR LAYERS")
    title(slide,"一条可信回答需要四层证明")
    label(slide,0.55,1.6,12.2,"「有引用」不等于「有根据」——四层缺一不可，断一环整条作废",SUB,12.5,PP_ALIGN.LEFT,True)
    layers=[("可访问","Accessible","当前用户有权看","身份 / ACL / 租户"),
            ("可定位","Locatable","能指认出处","源文档 / 版本 / 页码"),
            ("可支持","Supported","真被证据支撑","引用准确率 + 覆盖率"),
            ("可呈现","Presentable","当庭可出示","可点摘录 / 冲突提示 / 拒答")]
    W=2.85; gap=0.13; y=2.45; x0=0.7
    for i,(zh,en,q,d) in enumerate(layers):
        x=x0+i*(W+gap)
        node(slide,x,y,W,1.05,f"{i+1}. {zh}",en,q,fill=(CARD if i%2==0 else PALE),edge=TEAL,zh_size=15,en_size=10.5,desc_size=10.5)
        tf=tb(slide,x,y+1.2,W,0.9); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        sr(p.add_run(),d,BODY_F,10.5,SUB)
        if i<3: arrow(slide,x+W+0.005,y+0.52,x+W+gap-0.005,y+0.52,color=TEAL,w=2.0)
    band(slide,0.7,4.95,12.0,0.82,fill=RGBColor(0xF7,0xFB,0xFC),edge=LINE)
    tf=tb(slide,0.95,5.08,11.5,0.6); p=tf.paragraphs[0]
    sr(p.add_run(),"像法庭证据链：",BODY_F,11.5,DTEAL,bold=True)
    sr(p.add_run(),"合法取得 → 指认出处 → 支撑指控 → 当庭出示，断一环，整条作废。",BODY_F,11.5,SUB)
    tf2=tb(slide,0.6,6.05,12.2,0.6); p=tf2.paragraphs[0]
    sr(p.add_run(),"引用链要从「用户身份 + 资产版本」开始建，不是生成完再往后面贴链接——可信是设计出来的。",BODY_F,11.5,INK,bold=True)
    footer(slide,MOD,"第 5 章 · 常见评估方法")

# ---------- p49 六段诊断 ----------
def draw_six_diagnosis(slide):
    bg(slide)
    eyebrow(slide,"第 5 章 · 常见评估方法  SIX-STAGE")
    title(slide,"六段诊断：把总分拆成可修复问题")
    label(slide,0.55,1.6,12.2,"总分掉了不可怕，不知道哪段掉的才可怕——六段各装一个「压力表」",SUB,12.5,PP_ALIGN.LEFT,True)
    stages=[("解析","Parsing","内容完整进库？"),("切分","Chunking","块边界合理？"),
            ("召回","Recall","证据进候选池？"),("重排","Rerank","证据排到前面？"),
            ("上下文","Context","证据密度够？"),("生成引用","Answer","主张被支撑？")]
    metrics=["抽样比对原文","gold chunk 对照","Recall","MRR / Precision","噪声 / 重复率","忠实度 + 引用"]
    W=1.85; gap=0.13; y=2.55; h=1.45; x0=0.6; xs=[]
    for i,(zh,en,q) in enumerate(stages):
        x=x0+i*(W+gap); xs.append(x)
        node(slide,x,y,W,h,zh,en,q,fill=CARD,edge=TEAL,zh_size=13,en_size=9.5,desc_size=10)
        tf=tb(slide,x,y+h+0.06,W,0.5); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        sr(p.add_run(),"查："+metrics[i],BODY_F,9.5,DTEAL)
    for i in range(5): arrow(slide,xs[i]+W+0.005,y+h/2,xs[i+1]-0.005,y+h/2,color=TEAL,w=1.5)
    band(slide,0.7,4.95,12.0,1.1,fill=CARD,edge=ORANGE)
    tf=tb(slide,0.95,5.05,11.6,0.95)
    para(tf,"同样是「答错」，三种病吃三种药：",BODY_F,12,ODARK,bold=True,first=True)
    para(tf,"gold chunk 没进池 → 修检索　｜　进了池没进上下文 → 修重排与预算　｜　进了上下文还答错 → 修提示词或换模型",BODY_F,11,SUB)
    footer(slide,MOD,"第 5 章 · 常见评估方法")

# ---------- p67 RRF 算例 ----------
def draw_rrf(slide):
    bg(slide)
    eyebrow(slide,"第 7 章 · 混合检索  RRF")
    title(slide,"RRF：不比分数，只比名次")
    label(slide,0.55,1.6,12.2,"倒数排名融合 Reciprocal Rank Fusion —— 用名次绕开「分数不可比」",SUB,12.5,PP_ALIGN.LEFT,True)
    band(slide,0.6,2.3,3.6,3.25,fill=CARD,edge=TEAL)
    label(slide,0.82,2.42,3.2,"两份榜单，分数不可比",DTEAL,11.5,PP_ALIGN.LEFT,False,True)
    label(slide,0.82,2.9,3.2,"BM25 榜（分数无上界）",SUB,10.5,PP_ALIGN.LEFT,True)
    label(slide,0.92,3.2,3.1,"①文档A　②…　③文档B",INK,11.5,PP_ALIGN.LEFT,False)
    label(slide,0.82,3.75,3.2,"向量榜（0~1 相似度）",SUB,10.5,PP_ALIGN.LEFT,True)
    label(slide,0.92,4.05,3.1,"①…　③文档B　⑧文档A",INK,11.5,PP_ALIGN.LEFT,False)
    label(slide,0.82,4.65,3.3,"两套分数没法直接加权平均",ODARK,10.5,PP_ALIGN.LEFT,True)
    node(slide,4.55,3.25,2.35,1.35,"只看名次","1/(k+r) 相加","r=名次  k=60",fill=TEAL,edge=DTEAL,zh_c=WHITE,sub_c=PALE,en_c=WHITE,zh_size=15,en_size=11,desc_size=10.5)
    arrow(slide,4.2,3.92,4.53,3.92,color=TEAL,w=1.75)
    arrow(slide,6.9,3.92,7.28,3.92,color=ORANGE,w=1.75)
    band(slide,7.3,2.3,5.4,3.25,fill=RGBColor(0xF7,0xFB,0xFC),edge=ORANGE)
    label(slide,7.52,2.42,5.0,"算一遍就懂：谁更稳谁赢",ODARK,11.5,PP_ALIGN.LEFT,False,True)
    tf=tb(slide,7.56,2.95,5.0,2.4)
    para(tf,"文档A：BM25 第 1、向量第 8",BODY_F,11.5,INK,first=True)
    para(tf,"　→ 1/(60+1) + 1/(60+8) ≈ 0.031",BODY_F,11,SUB)
    para(tf,"文档B：两榜都第 3",BODY_F,11.5,INK)
    para(tf,"　→ 2 × 1/(60+3) ≈ 0.032　✓ 胜",BODY_F,11,GREEN)
    para(tf,"两边都不错的稳健选手，赢过单边冠军。",BODY_F,11,DTEAL,bold=True)
    footer(slide,MOD,"第 7 章 · 混合检索")

# ---------- p86 可观测性 / 归因 ----------
def draw_observability(slide):
    bg(slide)
    eyebrow(slide,"第 9 章 · 生产化与常见坑  OBSERVABILITY")
    title(slide,"可观测性：从「答错了」到「哪一环错了」")
    label(slide,0.55,1.6,12.2,"把每道工序的输入输出记成结构化 Trace，坏答案逐环回放定位（Langfuse / LangSmith）",SUB,12,PP_ALIGN.LEFT,True)
    stages=[("查询改写","Rewrite"),("检索","Retrieval"),("融合","Fusion"),("重排","Rerank"),("生成","Generation")]
    W=2.15; gap=0.2; y=2.5; h=1.0; x0=0.7; xs=[]
    for i,(zh,en) in enumerate(stages):
        x=x0+i*(W+gap); xs.append(x)
        node(slide,x,y,W,h,zh,en,"记 Trace",fill=CARD,edge=TEAL,zh_size=13,en_size=10,desc_size=9.5)
    for i in range(4): arrow(slide,xs[i]+W+0.02,y+h/2,xs[i+1]-0.02,y+h/2,color=TEAL,w=1.5)
    band(slide,0.7,4.05,12.0,1.05,fill=CORAL,edge=ORANGE)
    tf=tb(slide,0.95,4.15,11.6,0.9)
    para(tf,"坏答案回放（Trace 就是 RAG 的检查单）：",BODY_F,12,ODARK,bold=True,first=True)
    para(tf,"正确片段其实在 Top-5、但排第 4、被截断在上下文之外 → 问题在「截断策略」，不在检索。没有 Trace，多半冤枉嵌入模型。",BODY_F,11,SUB)
    tf2=tb(slide,0.6,5.35,12.2,0.9)
    para(tf2,"排障三问固定问：",BODY_F,12,DTEAL,bold=True,first=True)
    para(tf2,"检索到没有？（Recall）　→　排上来没有？（MRR）　→　用上没有？（忠实度）——评估与观测，是同一件事的离线态与在线态。",BODY_F,11,SUB)
    footer(slide,MOD,"第 9 章 · 生产化与常见坑")

# ---------- p87 三层缓存 ----------
def draw_cache(slide):
    bg(slide)
    eyebrow(slide,"第 9 章 · 生产化与常见坑  CACHING")
    title(slide,"缓存与成本：三层缓存，三本账")
    label(slide,0.55,1.6,12.2,"重复劳动是 RAG 最大的隐形成本——三层缓存各砍一件",SUB,12.5,PP_ALIGN.LEFT,True)
    band(slide,0.6,2.35,4.3,3.1,fill=CARD,edge=RGBColor(0x9F,0xB3,0xC0))
    label(slide,0.85,2.47,3.9,"RAG 成本四大件",INK,12,PP_ALIGN.LEFT,False,True)
    costs=[("嵌入调用","建库一次性大头"),("向量库","存储与查询"),("重排调用","精排阶段"),("生成 token","最大头")]
    for i,(c,d) in enumerate(costs):
        tf=tb(slide,0.9,2.98+i*0.6,3.9,0.55); p=tf.paragraphs[0]
        sr(p.add_run(),"• "+c+"　",BODY_F,11.5,INK,bold=True); sr(p.add_run(),d,BODY_F,10,SUB)
    band(slide,5.5,2.35,7.2,3.1,fill=GBG,edge=GREEN)
    label(slide,5.75,2.47,6.8,"三层缓存，各砍一件重复劳动",RGBColor(0x27,0x50,0x0A),12,PP_ALIGN.LEFT,False,True)
    caches=[("嵌入缓存","→ 省嵌入重算"),("语义缓存","→ 相似问题直接命中历史答案（省生成）"),("Prompt 缓存","→ 省重复上下文的输入费")]
    for i,(c,d) in enumerate(caches):
        yy=2.98+i*0.72
        node(slide,5.75,yy,2.35,0.6,c,"","",fill=WHITE,edge=GREEN,zh_c=RGBColor(0x27,0x50,0x0A),zh_size=12)
        tf=tb(slide,8.25,yy+0.13,4.3,0.5); p=tf.paragraphs[0]
        sr(p.add_run(),d,BODY_F,11,SUB)
    tf=tb(slide,0.6,5.6,12.2,1.0)
    para(tf,"案例：客服约四成问题重复 → 接语义缓存，生成成本立砍四成、延迟从秒级到毫秒级（常是上线第一个立竿见影的优化）。",BODY_F,11,INK,bold=True,first=True)
    para(tf,"算账口径：月成本 ≈ 建库一次性 + 增量嵌入 + 存储 +（检索+重排+生成）× 月问量 ×（1 − 缓存命中率）。命中率是最值得投资的变量。",BODY_F,10.5,SUB)
    footer(slide,MOD,"第 9 章 · 生产化与常见坑")

# ---------- p96 GraphRAG 四步 + 双路 ----------
def draw_graphrag(slide):
    bg(slide)
    eyebrow(slide,"第 10 章 · GraphRAG  HOW")
    title(slide,"GraphRAG 怎么干：抽取、建图、分区、摘要")
    label(slide,0.55,1.58,12.2,"索引期多干的活，都是为了回答期能「开地图」",SUB,12.5,PP_ALIGN.LEFT,True)
    label(slide,0.6,2.0,8,"索引期 Indexing —— 贵在这里",DTEAL,12,PP_ALIGN.LEFT,False,True)
    steps=[("① 抽取","Extract","LLM 从每块抽实体与关系"),("② 建图","Graph","织成知识图谱"),
           ("③ 分区","Community","Leiden 切成层层社区"),("④ 摘要","Summary","给每社区预写摘要")]
    W=2.85; gap=0.15; y=2.42; h=1.05; x0=0.7; xs=[]
    for i,(zh,en,d) in enumerate(steps):
        x=x0+i*(W+gap); xs.append(x)
        node(slide,x,y,W,h,zh,en,d,fill=CARD,edge=TEAL,zh_size=13,en_size=10,desc_size=10)
    for i in range(3): arrow(slide,xs[i]+W+0.005,y+h/2,xs[i+1]-0.005,y+h/2,color=TEAL,w=1.5)
    label(slide,0.6,3.95,8,"查询期 Query —— 两条路",ODARK,12,PP_ALIGN.LEFT,False,True)
    node(slide,0.7,4.4,5.85,0.95,"本地查询","Local","沿实体邻居找证据（细节题）",fill=CORAL,edge=ORANGE,zh_c=ODARK,sub_c=SUB,en_c=ODARK,zh_size=13,desc_size=10.5)
    node(slide,6.85,4.4,5.85,0.95,"全局查询","Global","汇总社区摘要、map-reduce（全局题）",fill=CORAL,edge=ORANGE,zh_c=ODARK,sub_c=SUB,en_c=ODARK,zh_size=13,desc_size=10.5)
    tf=tb(slide,0.6,5.6,12.2,0.9)
    para(tf,"类比：出版社给一万本书做「导读体系」——索引卡（实体）→ 关系网（图谱）→ 分区（社区）→ 综述（摘要）。",BODY_F,11,DTEAL,bold=True,first=True)
    para(tf,"贵在索引期：每块过 LLM 抽取、每社区预写摘要——这正是它答得了全局问题的原因，也是成本高的根源。",BODY_F,11,SUB)
    footer(slide,MOD,"第 10 章 · GraphRAG")

# ---------- p116 语义层架构 ----------
def draw_semantic_layer(slide):
    bg(slide)
    eyebrow(slide,"第 12 章 · 结构化数据 RAG  SEMANTIC LAYER")
    title(slide,"语义层：给模型一本口径手册")
    label(slide,0.55,1.6,12.2,"把「猜口径」变成「查口径」，准确率 50–70% → 85–95%",SUB,12.5,PP_ALIGN.LEFT,True)
    node(slide,0.7,3.15,2.1,0.95,"业务问题","Question","「上季度营收多少」",fill=WHITE,edge=RGBColor(0x9F,0xB3,0xC0),zh_c=INK,zh_size=13,en_size=9.5,desc_size=9.5)
    node(slide,3.75,2.3,4.1,0.95,"裸 Text-to-SQL","","模型自己猜口径写 SQL",fill=RGBColor(0xF1,0xEF,0xE8),edge=RGBColor(0x88,0x87,0x80),zh_c=RGBColor(0x44,0x44,0x41),sub_c=SUB,zh_size=13,desc_size=10)
    label(slide,8.0,2.5,4.6,"准确率 50–70%：口径全靠猜",RGBColor(0x99,0x3C,0x1D),11.5,PP_ALIGN.LEFT,False,True)
    node(slide,3.75,3.62,4.1,1.3,"语义层 Semantic Layer","","指标/维度/同义词的确定性定义：营收 = 确认收入 − 退款",fill=TEAL,edge=DTEAL,zh_c=WHITE,sub_c=PALE,zh_size=13,desc_size=10)
    label(slide,8.0,4.0,4.6,"准确率 85–95%：照手册出数，接近全对",GREEN,11.5,PP_ALIGN.LEFT,False,True)
    arrow(slide,2.8,3.45,3.73,2.78,color=RGBColor(0x88,0x87,0x80),w=1.5)
    arrow(slide,2.8,3.72,3.73,4.15,color=TEAL,w=1.75)
    tf=tb(slide,0.6,5.5,12.2,1.0)
    para(tf,"类比：给空降的 SQL 高手发一本《公司指标口径手册》——每个指标怎么算、黑话对应哪个字段全写死，他不用猜。",BODY_F,11,DTEAL,bold=True,first=True)
    para(tf,"生产架构：RAG 做检索层（找口径定义、找相关表），语义层做回答层（确定性出数）；查文 / 查数在入口路由。",BODY_F,11,SUB)
    footer(slide,MOD,"第 12 章 · 结构化数据 RAG")
