"""LLM-Training 讲义信息图（5 张）。复用 kb_draw + agent_figs helper。
嫁接：p5 后(生产流水线)、p24 后(下一词预测)、p44 后(奖励模型)、p45 后(RLHF vs DPO)、p52 后(推理模型)。
LLM-Training 无页码，图系统页脚匹配。独立预览：python3 train_figs.py"""
import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "LLM-Training 讲义 · LLM 训练全景"
GINK = RGBColor(0x27, 0x50, 0x0A)

# ---------- p5 生产流水线 ----------
def draw_train_pipeline(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 全景总览  THE PIPELINE")
    title(slide, "一张流水线地图：一个 LLM 是怎么造出来的")
    label(slide, 0.55, 1.58, 12.2, "顺着流水线走一遍——像培养一名医生：读教材、上医学院、实习、立医德、学会诊", SUB, 12, PP_ALIGN.LEFT, True)
    stages = [("数据", "粮食", "读海量教材", TEAL, DTEAL, CARD),
              ("预训练", "身体", "上医学院打底", TEAL, DTEAL, CARD),
              ("SFT", "听话", "跟着示范实习", ORANGE, ODARK, CORAL),
              ("对齐", "分寸", "立医德、懂边界", ORANGE, ODARK, CORAL),
              ("推理模型", "思考", "会诊：先想后答", DTEAL, NAVY, PALE)]
    W = 2.15; gap = 0.28; x0 = 0.6; y = 2.6; h = 1.5; xs = []
    for i, (zh, role, med, tc, ec, fill) in enumerate(stages):
        x = x0 + i*(W+gap); xs.append(x)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.12, y+0.16, W-0.24, 1.25); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        sr(p.add_run(), f"{i+1}. {zh}", BODY_F, 14, tc, bold=True)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        sr(p2.add_run(), "「" + role + "」", BODY_F, 12, ec if fill != CORAL else ODARK, bold=True)
        p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
        sr(p3.add_run(), med, BODY_F, 10, SUB)
        if i < 4:
            arrow(slide, x+W+0.02, y+h/2, x+W+gap-0.02, y+h/2, color=SUB, w=1.7)
    label(slide, 0.6, 4.3, 6.5, "预训练 / SFT / 对齐 / 推理 = 六道工序里的主线", SUB, 10, PP_ALIGN.LEFT, True)
    # 三种出厂形态
    band(slide, 0.6, 4.75, 12.03, 1.05, fill=RGBColor(0xF7,0xFB,0xFC), edge=LINE)
    label(slide, 0.85, 4.85, 6.0, "三种出厂形态：走到哪一步就是哪种", DTEAL, 11.5, PP_ALIGN.LEFT, False, True)
    forms = [("base", "只预训练 · 会续写不会应答", 0.9), ("instruct", "+ SFT/对齐 · 会听话应答", 5.0), ("reasoning", "+ RLVR · 先想后答", 9.1)]
    for name, d, x in forms:
        node(slide, x, 5.2, 3.7, 0.5, name + "　" + d, "", "", fill=WHITE, edge=TEAL, zh_c=DTEAL, zh_size=10.5)
    band(slide, 0.6, 5.95, 12.03, 0.87, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 6.04, 11.5, 0.72); p = tf.paragraphs[0]
    sr(p.add_run(), "售前口径：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "客户的任何「模型能力」问题，先定位到这条流水线的哪一道工序——瓶颈在哪道，方案就从哪道下手。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 1 章 · 全景总览")

# ---------- p24 下一词预测 ----------
def draw_train_next_token(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 预训练  NEXT-TOKEN PREDICTION")
    title(slide, "下一词预测：简单任务，暴力奇迹")
    label(slide, 0.55, 1.58, 12.2, "把「预测下一词」做到极致，就被迫学会了世界的规律——这是整个大模型时代的地基", SUB, 12, PP_ALIGN.LEFT, True)
    # 循环
    steps = [("看一段文本", "「因为下雨，所以他带了…」"), ("猜下一个 token", "「伞」？"), ("猜错就调权重", "梯度下降"), ("重复数万亿次", "loss 一点点往下爬")]
    W = 2.7; gap = 0.28; x0 = 0.6; y = 2.5; h = 1.0; xs = []
    for i, (zh, d) in enumerate(steps):
        x = x0 + i*(W+gap); xs.append(x)
        solid = (i == 2)
        node(slide, x, y, W, h, zh, "", d, fill=(ORANGE if solid else CARD), edge=(ODARK if solid else TEAL),
             zh_c=(OINK if solid else DTEAL), sub_c=(OINK if solid else SUB), zh_size=12.5, desc_size=9.5)
        if i < 3:
            arrow(slide, x+W+0.02, y+h/2, x+W+gap-0.02, y+h/2, color=TEAL, w=1.6)
    line(slide, xs[3]+W/2, y+h, xs[3]+W/2, y+h+0.3, color=SUB, w=1.4)
    line(slide, xs[3]+W/2, y+h+0.3, xs[0]+W/2, y+h+0.3, color=SUB, w=1.4)
    arrow(slide, xs[0]+W/2, y+h+0.3, xs[0]+W/2, y+h+0.02, color=SUB, w=1.4)
    label(slide, xs[1], y+h+0.16, W*2, "循环：看文本 → 猜 → 调 → 再看，直到会「填空」", SUB, 10, PP_ALIGN.CENTER, True)
    band(slide, 0.6, 4.35, 6.0, 1.85, fill=CARD, edge=TEAL)
    tf = tb(slide, 0.85, 4.47, 5.5, 1.65)
    para(tf, "类比 · 无限量完形填空", BODY_F, 12, DTEAL, bold=True, first=True)
    para(tf, "填对「法国的首都是__」要记事实；填对「因为下雨所以他带了__」要懂因果；补全一段 Python 要懂语法——被迫学会世界的规律。", BODY_F, 11, SUB)
    band(slide, 6.75, 4.35, 5.88, 1.85, fill=CORAL, edge=ORANGE)
    tf2 = tb(slide, 7.0, 4.47, 5.4, 1.65)
    para(tf2, "另一个视角 · 有损压缩", BODY_F, 12, ODARK, bold=True, first=True)
    para(tf2, "把几十 TB 文本压进几百 GB 权重——压缩即理解。解压时偶尔失真，就是幻觉的来源之一。", BODY_F, 11, SUB)
    para(tf2, "知识「融」在权重里（像人脑记忆、非硬盘检索）→ 所以需要 RAG 配「可翻的资料柜」。", BODY_F, 11, DTEAL, bold=True)
    footer(slide, MOD, "第 3 章 · 预训练")

# ---------- p44 奖励模型 ----------
def draw_train_reward_model(slide):
    bg(slide)
    eyebrow(slide, "第 5 章 · 对齐  REWARD MODEL")
    title(slide, "奖励模型：先训一个「懂行的评委」")
    label(slide, 0.55, 1.58, 12.2, "人类说不清「好回答」的定义，但两个答案摆面前挑得出更好的——把品味蒸馏成可无限调用的自动评委", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 0.7, 3.0, 2.55, 1.05, "偏好对", "A 比 B 好", "海量人工对比", fill=CARD, edge=TEAL, zh_c=DTEAL, sub_c=SUB, en_c=TEAL, zh_size=14, en_size=10, desc_size=9.5)
    arrow(slide, 3.27, 3.52, 4.0, 3.52, color=TEAL, w=1.8)
    label(slide, 3.2, 3.1, 1.0, "训练", TEAL, 9.5, PP_ALIGN.CENTER, True)
    node(slide, 4.05, 3.0, 2.85, 1.05, "奖励模型 RM", "自动评委", "输入问答 → 输出打分", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=PALE, zh_size=14, en_size=10, desc_size=9.5)
    arrow(slide, 6.92, 3.52, 7.65, 3.52, color=ORANGE, w=1.8)
    label(slide, 6.85, 3.1, 1.0, "照分练", ODARK, 9.5, PP_ALIGN.CENTER, True)
    node(slide, 7.7, 3.0, 2.85, 1.05, "模型（厨师）", "照评委分数练厨艺", "对齐到大众口味", fill=CORAL, edge=ORANGE, zh_c=ODARK, sub_c=SUB, en_c=ODARK, zh_size=14, en_size=10, desc_size=9.5)
    label(slide, 0.7, 4.25, 11.0, "类比 · 米其林评委：说不出满分菜的配方，但两道菜一尝就知哪道更好；吃过十万次对比，打分就能代表「大众口味」。", SUB, 10.5, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 4.75, 12.03, 1.15, fill=RGBColor(0xFB,0xE9,0xE7), edge=RGBColor(0x9B,0x2C,0x2C))
    tf = tb(slide, 0.9, 4.86, 11.5, 1.0); p = tf.paragraphs[0]
    sr(p.add_run(), "注意 · 评委会被钻空子（Reward Hacking）：", BODY_F, 11.5, RGBColor(0x9B,0x2C,0x2C), bold=True)
    sr(p.add_run(), "模型可能发现「答得长、姿态谦逊、堆免责声明」能骗高分——古德哈特定律：指标一旦成为目标就不再是好指标。", BODY_F, 11.5, INK)
    para(tf, "工程上靠评委「定期换届」（重训 RM）+ 多评委制衡来缓解。", BODY_F, 10.5, SUB)
    band(slide, 0.6, 6.02, 12.03, 0.8, fill=PALE, edge=LINE)
    tf2 = tb(slide, 0.9, 6.11, 11.5, 0.66); p = tf2.paragraphs[0]
    sr(p.add_run(), "一句话：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "RM 是 RLHF 的灵魂也是软肋——它有多准对齐就有多好；理解这点就理解了 DPO 和 RLVR 各自想绕开什么。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 5 章 · 对齐")

# ---------- p45 RLHF vs DPO ----------
def draw_train_rlhf_vs_dpo(slide):
    bg(slide)
    eyebrow(slide, "第 5 章 · 对齐  RLHF vs DPO")
    title(slide, "从 RLHF 到 DPO：三件套与暴力简化")
    label(slide, 0.55, 1.58, 12.2, "先 DPO，不够再上 RL——90% 的企业偏好定制场景，DPO 级别的方案就够用", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.35, 6.0, 3.05, fill=CORAL, edge=ORANGE)
    label(slide, 0.85, 2.46, 5.6, "RLHF —— 请评委，再请教练", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    rl = [("SFT 模型生成回答", 2.9), ("奖励模型打分", 3.62), ("PPO 更新 + KL 约束", 4.34)]
    for zh, yy in rl:
        node(slide, 0.95, yy, 5.1, 0.5, zh, "", "", fill=WHITE, edge=ODARK, zh_c=ODARK, zh_size=11.5)
    arrow(slide, 3.5, 3.42, 3.5, 3.6, color=ODARK, w=1.4)
    arrow(slide, 3.5, 4.14, 3.5, 4.32, color=ODARK, w=1.4)
    label(slide, 0.9, 4.95, 5.6, "同时伺候四个模型副本 · 显存翻倍 · 超参敏感 · 走钢丝", RGBColor(0x99,0x3C,0x1D), 10, PP_ALIGN.LEFT, True)
    band(slide, 6.75, 2.35, 6.0, 3.05, fill=GBG, edge=GREEN)
    label(slide, 7.0, 2.46, 5.6, "DPO —— 直接拿对比示范自学", GINK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.15, 5.3, 0.85, "直接在偏好对上训练", "", "抬高好答案概率、压低差答案概率", fill=GREEN, edge=GINK, zh_c=WHITE, sub_c=GBG, zh_size=13.5, desc_size=9.5)
    tf2 = tb(slide, 7.05, 4.15, 5.4, 1.2)
    para(tf2, "· 不训奖励模型、不跑 RL", BODY_F, 11, INK, bold=True, first=True)
    para(tf2, "· 像做分类一样一步到位——简单、稳定、便宜", BODY_F, 11, SUB)
    para(tf2, "· 开源 / 企业偏好调优的默认起手式", BODY_F, 11, GINK, bold=True)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.67, 11.5, 1.1); p = tf3.paragraphs[0]
    sr(p.add_run(), "一句话给客户：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "RLHF 是「请评委再请教练」，DPO 是「直接拿对比示范自学」——别一上来就给自己上 PPO 的工程难度。", BODY_F, 11.5, INK)
    para(tf3, "选型逻辑：先 DPO，不够再上 RL；前沿实验室两者混用（Tülu 3 配方 = SFT → DPO → RLVR）。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 5 章 · 对齐")

# ---------- p52 推理模型 ----------
def draw_train_reasoning(slide):
    bg(slide)
    eyebrow(slide, "第 6 章 · 推理模型  THINK-THEN-ANSWER")
    title(slide, "推理模型：先想后答的新物种")
    label(slide, 0.55, 1.58, 12.2, "在给出答案前先生成一段思维链——慢思考换高正确率，这是继「堆参数」之后的第二增长曲线", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 2.5, fill=CARD, edge=SUB)
    label(slide, 0.85, 2.52, 5.6, "普通模型 —— 看到题立刻报答案", RGBColor(0x55,0x55,0x52), 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 3.15, 2.0, 0.72, "问题", "", "", fill=WHITE, edge=SUB, zh_c=INK, zh_size=12)
    arrow(slide, 2.98, 3.51, 3.7, 3.51, color=SUB, w=1.6)
    node(slide, 3.75, 3.15, 2.4, 0.72, "答案", "凭直觉", "难题只能蒙", fill=WHITE, edge=SUB, zh_c=INK, sub_c=SUB, zh_size=12, desc_size=9)
    label(slide, 0.95, 4.15, 5.5, "像「心算」：没有草稿纸，难题凭直觉蒙", SUB, 10.5, PP_ALIGN.LEFT, True)
    band(slide, 6.75, 2.4, 6.0, 2.5, fill=CORAL, edge=ORANGE)
    label(slide, 7.0, 2.52, 5.6, "推理模型 —— 先打草稿再作答", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.15, 1.5, 0.72, "问题", "", "", fill=WHITE, edge=ORANGE, zh_c=ODARK, zh_size=12)
    node(slide, 8.75, 3.02, 2.35, 0.98, "思维链 CoT", "拆解·尝试", "自查·回溯·换路", fill=ORANGE, edge=ODARK, zh_c=OINK, sub_c=OINK, en_c=OINK, zh_size=12, en_size=9, desc_size=9)
    node(slide, 11.25, 3.15, 1.35, 0.72, "答案", "更准", "", fill=WHITE, edge=ODARK, zh_c=ODARK, sub_c=SUB, zh_size=12, desc_size=9)
    arrow(slide, 8.62, 3.51, 8.72, 3.51, color=ODARK, w=1.5)
    arrow(slide, 11.12, 3.51, 11.22, 3.51, color=ODARK, w=1.5)
    label(slide, 7.0, 4.15, 5.7, "像「打草稿」：写中间步骤、验算、错了划掉重来", ODARK, 10.5, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 5.1, 12.03, 1.72, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.22, 11.5, 1.55); p = tf.paragraphs[0]
    sr(p.add_run(), "推理时扩展（Inference-time Scaling）：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "训练定终身之外的第二条能力曲线——同一个模型，允许它多想（多生成思考 token），难题正确率显著上升。", BODY_F, 11.5, INK)
    para(tf, "代价同样直白：思考 token 也按量计费，延迟也随之上涨。", BODY_F, 11, SUB)
    para(tf, "给客户定性：预训练撞上数据墙后，行业把算力从「训练时」搬到「推理时」——这就是 2025–2026 旗舰发布会都在讲「思考」的原因。", BODY_F, 11, DTEAL, bold=True)
    footer(slide, MOD, "第 6 章 · 推理模型")

# ---------- p27 Scaling Laws 三代 ----------
def draw_train_scaling_laws(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · Scaling Laws  THREE GENERATIONS")
    title(slide, "Scaling Laws 三代：从「越大越好」到「精打细算」")
    label(slide, 0.55, 1.58, 12.2, "「花多少钱能买到多强能力」的可预测曲线——大厂敢投几亿美金，正因为回报可以事先算出来", SUB, 12, PP_ALIGN.LEFT, True)
    eras = [("Kaplan 2020", "参数优先", "能力随参数/数据/算力按幂律可预测提升", "开启军备竞赛：GPT-3 一路到万亿参数", "「做大就完了」", TEAL, DTEAL, CARD),
            ("Chinchilla 2022", "配平数据", "同算力下参数与数据要配平（约 20 token/参数），此前普遍「训不饱」", "从拼参数转向拼数据；撞上数据墙的伏笔", "「吃饱再长个」", ORANGE, ODARK, CORAL),
            ("推理感知 2024–", "总账最优", "把推理服务成本算进最优解：小而多训（超量 token）总账更划算", "小模型过度训练成主流；MoE 崛起", "「好养活才是真便宜」", DTEAL, NAVY, PALE)]
    W = 3.95; gap = 0.14; x0 = 0.6; y = 2.4; h = 3.05
    for i, (gen, tag, claim, impact, mem, tc, ec, fill) in enumerate(eras):
        x = x0 + i*(W+gap)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, y+0.16, W-0.5, 0.72); p = tf.paragraphs[0]
        sr(p.add_run(), f"{i+1}. {gen}", BODY_F, 13.5, tc, bold=True)
        sr(p.add_run(), "　" + tag, BODY_F, 10.5, SUB)
        b = tb(slide, x+0.28, y+0.86, W-0.56, 1.5)
        para(b, "主张：" + claim, BODY_F, 10.5, INK, first=True)
        para(b, "影响：" + impact, BODY_F, 10.5, SUB)
        node(slide, x+0.35, y+2.4, W-0.7, 0.5, mem, "", "", fill=WHITE, edge=ec, zh_c=tc, zh_size=12)
        if i < 2:
            arrow(slide, x+W-0.02, y+h/2, x+W+gap+0.02, y+h/2, color=SUB, w=1.8)
    band(slide, 0.6, 5.62, 12.03, 1.2, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.74, 11.5, 1.05); p = tf.paragraphs[0]
    sr(p.add_run(), "一条主线：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "从「参数越大越好」→「参数数据要配平」→「连推理成本一起算」——Scaling 的目标函数一直在变，但可预测性没变。", BODY_F, 12, INK)
    footer(slide, MOD, "第 3 章 · 预训练")

# ---------- p58 RLVR ----------
def draw_train_rlvr(slide):
    bg(slide)
    eyebrow(slide, "第 6 章 · RLVR  VERIFIABLE REWARD")
    title(slide, "RLVR：用「对错」替代「品味」当奖励")
    label(slide, 0.55, 1.58, 12.2, "只在答案可机器判定的领域训练——答对给奖励、答错不给，不需要奖励模型来「品」", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.4, 6.0, 2.95, fill=CORAL, edge=ORANGE)
    label(slide, 0.85, 2.52, 5.6, "第 5 章：作文 —— 评委打分", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    tf = tb(slide, 0.9, 3.05, 5.5, 2.1)
    para(tf, "· 靠奖励模型（评委）打分", BODY_F, 11, SUB, first=True)
    para(tf, "· 评委有口味、能被讨好（答长/谦逊/堆免责）", BODY_F, 11, SUB)
    para(tf, "· 第 5 章的全部烦恼：Reward Hacking", BODY_F, 11, ODARK, bold=True)
    band(slide, 6.75, 2.4, 6.0, 2.95, fill=GBG, edge=GREEN)
    label(slide, 7.0, 2.52, 5.6, "RLVR：数学 —— 机器批改", GINK, 12.5, PP_ALIGN.LEFT, False, True)
    tf2 = tb(slide, 7.05, 3.05, 5.5, 2.1)
    para(tf2, "· 数学有标准答案、代码有单元测试、逻辑有唯一解", BODY_F, 11, SUB, first=True)
    para(tf2, "· 答案唯一：机器批改、绝对公正、永不疲劳", BODY_F, 11, GINK, bold=True)
    para(tf2, "· 在批不完的数学卷上刷题，「实力」被逼着真实增长——讨好没用", BODY_F, 11, SUB)
    arrow(slide, 6.05, 3.85, 6.95, 3.85, color=SUB, w=1.8); label(slide, 5.85, 3.42, 1.3, "换掉评委", SUB, 9.5, PP_ALIGN.CENTER, True)
    band(slide, 0.6, 5.5, 12.03, 1.32, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.62, 11.5, 1.15); p = tf3.paragraphs[0]
    sr(p.add_run(), "案例 · DeepSeek-R1 震动行业：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "跳过 SFT、直接对 base 做大规模 RLVR（R1-Zero），推理行为——反思、验算、「等等我再想想」——自发涌现，没人教过。", BODY_F, 11.5, INK)
    para(tf3, "边界：RLVR 靠可验证性，数学/代码突飞猛进；写作/共情仍靠第 5 章偏好对齐——完整配方是 SFT → DPO → RLVR 层层叠加。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 6 章 · 推理模型")

# ---------- p71 6ND 成本心算 ----------
def draw_train_6nd(slide):
    bg(slide)
    eyebrow(slide, "第 7 章 · 算力账  6ND RULE")
    title(slide, "训练成本心算：6ND 法则")
    label(slide, 0.55, 1.58, 12.2, "售前的「镇场技能」——客户说「想自己训个 13B」，你当场三行算出卡数、天数、量级成本", SUB, 12, PP_ALIGN.LEFT, True)
    # 公式
    band(slide, 0.6, 2.4, 12.03, 1.25, fill=PALE, edge=TEAL)
    tf = tb(slide, 0.9, 2.5, 11.5, 0.55); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    sr(p.add_run(), "训练总算力 ≈ 6 × N × D", BODY_F, 20, DTEAL, bold=True)
    tf2 = tb(slide, 0.9, 3.05, 11.5, 0.55); p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    sr(p.add_run(), "N = 参数量　·　D = 训练 token 数　·　6 = 每参数每 token 的前向+反向浮点运算量　　÷（集群峰值 × MFU 30–50%）= 训练天数", BODY_F, 11, SUB)
    band(slide, 0.6, 3.85, 6.0, 1.75, fill=CARD, edge=TEAL)
    tf3 = tb(slide, 0.85, 3.97, 5.5, 1.55)
    para(tf3, "案例 · 70B + 15T token", BODY_F, 12, DTEAL, bold=True, first=True)
    para(tf3, "6 × 70e9 × 15e12 = 6.3e24 FLOPs", BODY_F, 11.5, INK, bold=True)
    para(tf3, "一万张 H100（有效 5e14 FLOP/s）≈ 跑 15 天量级；按市价租金，仅算力就是千万美元级。", BODY_F, 11, SUB)
    band(slide, 6.75, 3.85, 5.88, 1.75, fill=CORAL, edge=ORANGE)
    tf4 = tb(slide, 7.0, 3.97, 5.4, 1.55)
    para(tf4, "注意 · 总成本远不止 GPU 租金", BODY_F, 12, ODARK, bold=True, first=True)
    para(tf4, "数据管线、失败重跑（loss 炸了回滚）、消融实验（试几十个小配方）、人力——", BODY_F, 11, SUB)
    para(tf4, "「正式那一跑」只占总投入的一半以下。报价别只算理想值。", BODY_F, 11, ODARK, bold=True)
    band(slide, 0.6, 5.75, 12.03, 1.07, fill=GBG, edge=GREEN)
    tf5 = tb(slide, 0.9, 5.86, 11.5, 0.9); p = tf5.paragraphs[0]
    sr(p.add_run(), "镇场用法：", BODY_F, 11.5, GINK, bold=True)
    sr(p.add_run(), "三行算出卡数/天数/量级成本，再对比「微调 + RAG」的费用——数字一摆，理性结论自然浮现，这就是「基座模型是重资产游戏」的数学依据。", BODY_F, 11.5, INK)
    footer(slide, MOD, "第 7 章 · 基础设施")

FIGS = [draw_train_pipeline, draw_train_next_token, draw_train_reward_model,
        draw_train_rlhf_vs_dpo, draw_train_reasoning,
        draw_train_scaling_laws, draw_train_rlvr, draw_train_6nd]
NEW = [draw_train_scaling_laws, draw_train_rlvr, draw_train_6nd]

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in FIGS:
        fn(prs.slides.add_slide(B))
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_train_preview.pptx")
    prs.save(out); print(f"saved {len(FIGS)} ->", out)
