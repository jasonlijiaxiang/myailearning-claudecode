"""Multimodal 讲义信息图（4 张）。复用 kb_draw + agent_figs helper。
嫁接：p18 后（CLIP）、p31 后（原生 vs 拼管线）、p49 后（扩散 vs 自回归）、p61 后（铁三角）。
独立预览：python3 mm_figs.py"""
import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "多模态讲义 · Multimodal"
GINK = RGBColor(0x27, 0x50, 0x0A)
GRAY = RGBColor(0x88, 0x87, 0x80)
PURP = RGBColor(0x6B, 0x46, 0x9C)
PURPBG = RGBColor(0xF1, 0xEC, 0xF7)

def bullets(slide, x, y, w, h, lines, size=11, gap=None, color=INK, lead=None):
    tf = tb(slide, x, y, w, h)
    for i, ln in enumerate(lines):
        c = lead[1] if (lead and i == 0) else color
        b = bool(lead) and i == 0
        para(tf, ln, BODY_F, size, c, bold=b, first=(i == 0))
    return tf

# ---------- p18 CLIP ----------
def draw_mm_clip(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · CLIP  CONTRASTIVE")
    title(slide, "CLIP：图和文字，走进同一个语义空间")
    label(slide, 0.55, 1.58, 12.2, "用 4 亿对「图 + 配文」训练：匹配的拉近、不匹配的推远——图文从此可比、可搜、可零样本分类", SUB, 12, PP_ALIGN.LEFT, True)
    # 双编码器
    node(slide, 0.7, 2.55, 1.7, 0.72, "图片", "一张金毛照片", "", fill=WHITE, edge=TEAL, zh_c=DTEAL, en_c=TEAL, zh_size=13, en_size=9)
    node(slide, 0.7, 3.75, 1.7, 0.72, "配文", "golden retriever", "", fill=WHITE, edge=ORANGE, zh_c=ODARK, en_c=ODARK, zh_size=13, en_size=9)
    node(slide, 2.75, 2.55, 2.5, 0.72, "图像编码器", "Image Encoder", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, en_c=PALE, zh_size=12.5, en_size=9.5)
    node(slide, 2.75, 3.75, 2.5, 0.72, "文本编码器", "Text Encoder", "", fill=ORANGE, edge=ODARK, zh_c=OINK, en_c=OINK, zh_size=12.5, en_size=9.5)
    arrow(slide, 2.42, 2.91, 2.73, 2.91, color=TEAL, w=1.5)
    arrow(slide, 2.42, 4.11, 2.73, 4.11, color=ODARK, w=1.5)
    # 同一语义空间
    fr = band(slide, 5.75, 2.5, 6.9, 3.05, fill=RGBColor(0xF7, 0xFB, 0xFC), edge=RGBColor(0x9F, 0xB3, 0xC0))
    lnf = fr.line._get_or_add_ln(); lnf.append(lnf.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    label(slide, 5.95, 2.62, 6.0, "同一语义空间 —— 意思相同的自动站到一起", DTEAL, 11.5, PP_ALIGN.LEFT, False, True)
    arrow(slide, 5.28, 2.91, 5.72, 3.2, color=TEAL, w=1.6)
    arrow(slide, 5.28, 4.11, 5.72, 3.6, color=ODARK, w=1.6)
    # 匹配对：拉近
    dot(slide, 7.1, 3.7, 0.11, TEAL); label(slide, 7.25, 3.55, 1.6, "金毛图", DTEAL, 9.5, PP_ALIGN.LEFT, False)
    dot(slide, 7.55, 3.95, 0.11, ORANGE); label(slide, 7.7, 3.88, 2.2, "golden retriever", ODARK, 9, PP_ALIGN.LEFT, False)
    label(slide, 6.9, 4.28, 2.2, "匹配 → 拉近", GINK, 10.5, PP_ALIGN.LEFT, False, True)
    # 不匹配：推远
    dot(slide, 11.4, 3.0, 0.11, GRAY); label(slide, 10.55, 2.86, 2.0, "汽车图", GRAY, 9.5, PP_ALIGN.RIGHT, False)
    label(slide, 10.2, 4.6, 2.3, "不匹配 → 推远", GRAY, 10.5, PP_ALIGN.LEFT, False, True)
    line(slide, 7.6, 3.95, 11.3, 3.05, color=LINE, w=1.0, dashed=True)
    label(slide, 8.7, 3.28, 2.6, "距离大 = 语义远", SUB, 9, PP_ALIGN.CENTER, True)
    # 底部：能力 + 类比
    band(slide, 0.7, 5.72, 11.93, 1.13, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 5.83, 11.4, 1.0)
    para(tf, "得到三样东西：图文进同一空间（可比可搜）· 零样本分类（给类别名就能认新类）· 成为多模态大模型的「眼睛」底座。", BODY_F, 11.5, DTEAL, bold=True, first=True)
    para(tf, "类比：把所有图和所有句子扔进同一个大操场，让「意思相同的」自动站到一起——之后给一句话，就能找到最像它的图。", BODY_F, 11, SUB)
    footer(slide, MOD, "机器怎么看")

# ---------- p31 原生 vs 拼管线 ----------
def draw_mm_native_vs_pipeline(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 灵魂页  NATIVE vs PIPELINE")
    title(slide, "原生多模态 vs 拼管线：不是二选一")
    label(slide, 0.55, 1.58, 12.2, "客户最纠结、你最能体现价值的一页——先看清两条路各自的强项与代价", SUB, 12, PP_ALIGN.LEFT, True)
    # 左：原生
    band(slide, 0.6, 2.35, 6.0, 3.0, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.46, 5.6, "原生多模态模型 —— 一个模型端到端", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    for i, (zh, x) in enumerate([("看图", 0.95), ("理解", 2.55), ("回答", 4.15)]):
        node(slide, x, 2.95, 1.35, 0.52, zh, "", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, zh_size=11.5)
        if i < 2:
            arrow(slide, x+1.35, 3.21, x+1.6, 3.21, color=TEAL, w=1.4)
    bullets(slide, 0.9, 3.62, 5.5, 1.6, [
        "强项：图表推理、复杂版式、开放问答、图文关联",
        "代价：贵、慢、可能幻觉，能力是「黑盒」",
        "适合：理解密集、版式多变、要开放问答"], size=11, lead=(True, DTEAL))
    # 右：拼管线
    band(slide, 6.75, 2.35, 6.0, 3.0, fill=CORAL, edge=ORANGE)
    label(slide, 7.0, 2.46, 5.6, "拼管线 —— OCR/ASR + 文本 LLM 接力", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    for i, (zh, x) in enumerate([("OCR/ASR", 7.1), ("文本", 8.9), ("文本 LLM", 10.35)]):
        node(slide, x, 2.95, 1.6 if i != 1 else 1.2, 0.52, zh, "", "", fill=(WHITE if i == 1 else ORANGE), edge=ODARK, zh_c=(ODARK if i == 1 else OINK), zh_size=11)
        if i < 2:
            arrow(slide, x+(1.6 if i != 1 else 1.2), 3.21, x+(1.8 if i != 1 else 1.75), 3.21, color=ODARK, w=1.4)
    bullets(slide, 7.05, 3.62, 5.5, 1.6, [
        "强项：字段稳定、可控、便宜、每步可替换 / 审计",
        "代价：跨模态弱、误差累积、拼装维护成本",
        "适合：版式固定、字段稳定、对成本敏感"], size=11, lead=(True, ODARK))
    # 底部：混合
    band(slide, 0.7, 5.55, 11.93, 1.3, fill=GBG, edge=GREEN)
    tf = tb(slide, 1.0, 5.68, 11.4, 1.15); p = tf.paragraphs[0]
    sr(p.add_run(), "划重点 · 很多真实系统是「混合」：", BODY_F, 12, GINK, bold=True)
    sr(p.add_run(), "稳定部分走 OCR 管线省钱，疑难部分兜底给原生多模态。", BODY_F, 12, INK)
    para(tf, "别把它当二选一——「主管线 + 原生兜底」往往是成本与效果的最优解。", BODY_F, 11, SUB)
    footer(slide, MOD, "模态怎么拼")

# ---------- p49 扩散 vs 自回归 ----------
def draw_mm_diffusion_vs_ar(slide):
    bg(slide)
    eyebrow(slide, "第 5 章 · 图像生成两大流派  DIFFUSION vs AR")
    title(slide, "扩散 vs 自回归：美术生和逻辑生")
    label(slide, 0.55, 1.58, 12.2, "两种画图的方式，脾气完全不同——按「要艺术感」还是「要按指令」来选", SUB, 12, PP_ALIGN.LEFT, True)
    # 左：扩散
    band(slide, 0.6, 2.35, 6.0, 3.05, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.46, 5.6, "扩散 Diffusion —— 像「美术生」", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    for i, (zh, x) in enumerate([("随机噪声", 0.95), ("逐步去噪", 2.75), ("成图", 4.55)]):
        node(slide, x, 2.95, 1.55 if i < 2 else 1.2, 0.52, zh, "", "", fill=(TEAL if i == 2 else WHITE), edge=(DTEAL if i == 2 else TEAL), zh_c=(WHITE if i == 2 else DTEAL), zh_size=11)
        if i < 2:
            arrow(slide, x+1.55, 3.21, x+1.8, 3.21, color=TEAL, w=1.4)
    bullets(slide, 0.9, 3.62, 5.5, 1.65, [
        "代表：Stable Diffusion、Midjourney、FLUX.2",
        "强项：光影、质感、写实、艺术感",
        "弱项：空间关系（猫在狗左边常画反）、文字渲染"], size=11, lead=(True, DTEAL))
    band(slide, 0.9, 4.82, 5.4, 0.44, fill=PALE, edge=LINE)
    tfd = tb(slide, 1.05, 4.87, 5.2, 0.36); sr(tfd.paragraphs[0].add_run(), "一句话：画面美，但逻辑约束差", BODY_F, 11, DTEAL, bold=True)
    # 右：自回归
    band(slide, 6.75, 2.35, 6.0, 3.05, fill=PURPBG, edge=PURP)
    label(slide, 7.0, 2.46, 5.6, "自回归 Autoregressive —— 像「逻辑生」", PURP, 12.5, PP_ALIGN.LEFT, False, True)
    for i, (zh, x) in enumerate([("tok", 7.15), ("tok", 7.95), ("tok", 8.75)]):
        node(slide, x, 2.95, 0.68, 0.52, zh, "", "", fill=WHITE, edge=PURP, zh_c=PURP, zh_size=10.5)
        if i < 2:
            arrow(slide, x+0.68, 3.21, x+0.8, 3.21, color=PURP, w=1.3)
    arrow(slide, 9.55, 3.21, 9.95, 3.21, color=PURP, w=1.4)
    node(slide, 10.0, 2.95, 1.2, 0.52, "成图", "", "", fill=PURP, edge=PURP, zh_c=WHITE, zh_size=11)
    label(slide, 7.15, 3.52, 4.5, "像写句子一样逐 token「拼」出图", SUB, 9.5, PP_ALIGN.LEFT, True)
    bullets(slide, 7.05, 3.78, 5.5, 1.5, [
        "代表：GPT Image 2、Nano Banana 2、Luma",
        "强项：空间关系、多约束指令遵循、文字渲染",
        "弱项：生成略慢（在快速改善）"], size=11, lead=(True, PURP))
    band(slide, 7.05, 4.82, 5.4, 0.44, fill=PURPBG, edge=PURP)
    tfa = tb(slide, 7.2, 4.87, 5.2, 0.36); sr(tfa.paragraphs[0].add_run(), "一句话：听得懂复杂要求、字写得准", BODY_F, 11, PURP, bold=True)
    # 底部
    band(slide, 0.7, 5.72, 11.93, 1.13, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 5.85, 11.4, 1.0); p = tf.paragraphs[0]
    sr(p.add_run(), "对客一句话：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "要「好看的艺术图」偏扩散；要「按复杂指令、带准确文字的图」（海报、UI、信息图）偏自回归。", BODY_F, 12, INK)
    para(tf, "趋势：自回归正和 LLM「合体」，指令遵循与文字渲染的优势会越来越关键。", BODY_F, 11, SUB)
    footer(slide, MOD, "生成侧盘点")

# ---------- p61 成本/延迟/精度 铁三角 ----------
def draw_mm_tradeoff_triangle(slide):
    bg(slide)
    eyebrow(slide, "第 6 章 · 决策框架  IRON TRIANGLE")
    title(slide, "成本 / 延迟 / 精度：三角里只能选两个")
    label(slide, 0.55, 1.58, 12.2, "像项目管理的铁三角，三者互相拉扯——对客开场先问「哪个最不能妥协」，方案立刻收敛", SUB, 12, PP_ALIGN.LEFT, True)
    # 三角形三顶点
    ax, ay = 3.55, 2.55   # 精度（顶）
    bx, by = 1.5, 5.35    # 成本（左下）
    cx, cy = 5.6, 5.35    # 延迟（右下）
    line(slide, ax, ay+0.35, bx+0.2, by-0.15, color=TEAL, w=2.0)
    line(slide, ax, ay+0.35, cx-0.2, cy-0.15, color=TEAL, w=2.0)
    line(slide, bx+0.35, by, cx-0.35, cy, color=TEAL, w=2.0)
    label(slide, 2.3, 4.0, 2.6, "选两个", DTEAL, 15, PP_ALIGN.CENTER, False, True)
    label(slide, 2.3, 4.42, 2.6, "就得让一个", SUB, 10.5, PP_ALIGN.CENTER, True)
    node(slide, 2.7, 2.5, 1.7, 0.62, "精度", "Accuracy", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, en_c=PALE, zh_size=13, en_size=9)
    node(slide, 0.65, 5.3, 1.7, 0.62, "成本", "Cost", "", fill=ORANGE, edge=ODARK, zh_c=OINK, en_c=OINK, zh_size=13, en_size=9)
    node(slide, 4.75, 5.3, 1.7, 0.62, "延迟", "Latency", "", fill=ORANGE, edge=ODARK, zh_c=OINK, en_c=OINK, zh_size=13, en_size=9)
    # 右侧三策略卡
    cards = [(2.6, "精度优先", TEAL, DTEAL, "上最强旗舰、高分辨率、开放问答；代价是贵和慢。适合医疗 / 金融等容错低场景。"),
             (4.02, "成本优先", ORANGE, ODARK, "拼管线打底、开源 VLM 自部署、低分辨率够用即可。适合海量、字段稳定场景。"),
             (5.44, "延迟优先", ORANGE, ODARK, "小模型 / 缓存 / 预处理降 token；牺牲部分精度换实时。适合实时交互、语音对话。")]
    for y, head, ec, tc, body in cards:
        band(slide, 7.05, y, 5.85, 1.28, fill=CARD, edge=LINE)
        tf = tb(slide, 7.3, y+0.13, 5.4, 1.05); p = tf.paragraphs[0]
        sr(p.add_run(), head + "　", BODY_F, 12.5, tc, bold=True)
        para(tf, body, BODY_F, 11, SUB)
    footer(slide, MOD, "选型与动手做")

# ---------- p29 后 三条融合路线对比 ----------
def draw_mm_fusion_routes(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 模态怎么拼  THREE FUSION ROUTES")
    title(slide, "三条融合路线：把视觉接到大脑的三种接法")
    label(slide, 0.55, 1.58, 12.2, "都在解决同一个问题——视觉 token 怎么进 LLM；差别在 token 开销与工程复杂度", SUB, 12, PP_ALIGN.LEFT, True)
    cards = [("投影层", "Projection · LLaVA", TEAL, DTEAL, CARD,
              "视觉向量 → 投影层 MLP → 拼进文字序列当 token",
              ["token 开销：大", "结构简单、几乎不动 LLM", "开源主流、生态最繁荣"]),
             ("交叉注意力", "Cross-Attn · Flamingo", ORANGE, ODARK, CORAL,
              "LLM 内部插交叉注意力层，生成时按需「回看」图",
              ["token 开销：小（序列不被撑长）", "工程重、改动大、难复刻", "多见于大厂自研旗舰"]),
             ("Q-Former", "BLIP-2", PURP, PURP, PURPBG,
              "查询向量从上千 token 萃取几十个精华再喂 LLM",
              ["token 开销：最省", "编码器 / LLM 都可冻结", "可能丢细节、近年被投影层部分取代"])]
    W = 3.95; gap = 0.14; x0 = 0.6; y = 2.4; h = 2.95
    for i, (zh, en, ec, tc, fill, mech, feats) in enumerate(cards):
        x = x0 + i*(W+gap)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, y+0.16, W-0.5, 0.72); p = tf.paragraphs[0]
        sr(p.add_run(), f"{i+1}. {zh}", BODY_F, 14.5, tc, bold=True)
        para(tf, en, BODY_F, 10, SUB, italic=True)
        mb = tb(slide, x+0.28, y+0.92, W-0.56, 0.85)
        para(mb, mech, BODY_F, 11, INK, bold=True, first=True)
        fb = tb(slide, x+0.28, y+1.82, W-0.56, 1.0)
        for j, ft in enumerate(feats):
            para(fb, "· " + ft, BODY_F, 10.5, SUB, first=(j == 0))
    band(slide, 0.6, 5.62, 12.03, 1.2, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.74, 11.5, 1.05); p = tf.paragraphs[0]
    sr(p.add_run(), "一句话选型：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "开源起步、图不大 → 投影层；要长上下文 / 图文交错、有工程力 → 交叉注意力；视觉 token 太贵要压 → Q-Former。", BODY_F, 11.5, INK)
    footer(slide, MOD, "模态怎么拼")

# ---------- p87 后 级联时延预算 ----------
def draw_mm_latency_budget(slide):
    bg(slide)
    eyebrow(slide, "第 8 章 · 级联管线  LATENCY BUDGET")
    title(slide, "级联时延预算：首音时延花在哪")
    label(slide, 0.55, 1.58, 12.2, "衡量口径永远是「首音时延」不是「说完时延」——瓶颈在 LLM 首 token", SUB, 12, PP_ALIGN.LEFT, True)
    # 比例条
    segs = [("ASR 转写", "100–300ms", 200, TEAL, DTEAL, WHITE),
            ("LLM 首 token", "350–1000ms", 675, ORANGE, ODARK, OINK),
            ("TTS 首音", "90–200ms", 145, TEAL, DTEAL, WHITE),
            ("网络", "50–200ms", 125, RGBColor(0x9F, 0xB3, 0xC0), SUB, WHITE)]
    total = sum(s[2] for s in segs); X0 = 0.7; BW = 11.9; y = 2.55; bh = 0.85
    x = X0
    for name, rng, val, fill, ec, tcol in segs:
        w = BW*val/total
        b = band(slide, x, y, w, bh, fill=fill, edge=ec)
        tf = tb(slide, x, y+0.14, w, 0.6); p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        sr(p.add_run(), name, BODY_F, 11.5 if w > 1.5 else 10, tcol, bold=True)
        t2 = tb(slide, x, y+bh+0.04, w, 0.32); q = t2.paragraphs[0]; q.alignment = PP_ALIGN.CENTER
        sr(q.add_run(), rng, BODY_F, 10, SUB)
        x += w
    label(slide, 0.7, y-0.32, 11.9, "用户说完（VAD 判定） ───────────────► 用户听到（首音）", SUB, 10.5, PP_ALIGN.LEFT, True)
    # 两个开关
    band(slide, 0.6, 3.95, 5.95, 2.0, fill=CARD, edge=TEAL)
    tf = tb(slide, 0.85, 4.08, 5.5, 1.8)
    para(tf, "开关一 · 全链路流式", BODY_F, 12.5, DTEAL, bold=True, first=True)
    para(tf, "· ASR 边听边出字、LLM 边想边吐、TTS 边收边说", BODY_F, 11, SUB)
    para(tf, "· 首音取决于「首 token」，不是「全文生成完」", BODY_F, 11, SUB)
    para(tf, "· 任何一段等全量输出，延迟立刻爆表", BODY_F, 11, SUB)
    band(slide, 6.68, 3.95, 5.95, 2.0, fill=CORAL, edge=ORANGE)
    tf2 = tb(slide, 6.93, 4.08, 5.5, 1.8)
    para(tf2, "开关二 · 瓶颈在 LLM 首 token", BODY_F, 12.5, ODARK, bold=True, first=True)
    para(tf2, "· 占大头：模型档位选小、提示词减脂", BODY_F, 11, SUB)
    para(tf2, "· KV 缓存复用系统提示（见 LLM-Inference）", BODY_F, 11, SUB)
    para(tf2, "· 工具调用最伤延迟——能预取的别现查", BODY_F, 11, SUB)
    band(slide, 0.6, 6.1, 12.03, 0.75, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 6.2, 11.5, 0.62); p = tf3.paragraphs[0]
    sr(p.add_run(), "要点：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "级联的真功夫不在挑最好的三个组件，而在把三段「焊」成流水线——口径永远是「首音时延」。", BODY_F, 11.5, INK)
    footer(slide, MOD, "语音与实时交互")

# ---------- p89 后 级联 vs 端到端 ----------
def draw_mm_cascade_vs_e2e(slide):
    bg(slide)
    eyebrow(slide, "第 8 章 · 级联 vs 端到端  TRADE-OFF")
    title(slide, "级联 vs 端到端：六个维度的取舍")
    label(slide, 0.55, 1.58, 12.2, "要「像 ChatGPT 语音一样快」→ 端到端；要留痕可控（电话客服 / 金融）→ 级联", SUB, 12, PP_ALIGN.LEFT, True)
    rows = [("延迟", "P50 1–1.5 秒（流式后）", "< 1 秒，最快 ~0.3 秒"),
            ("可调试性", "三段可单测、有文字中间态", "黑盒，问题难定位"),
            ("供应商", "ASR / TTS / LLM 各自可换", "仅少数几家（OpenAI / Google 等）"),
            ("合规审计", "全程文字留痕、护栏成熟", "需语音级审计、护栏重做"),
            ("成本", "$0.01–0.17/分钟，可预测", "两端差百倍（$0.002–0.30）"),
            ("2026 现状", "企业生产主流", "消费级助手与 demo 先用")]
    # 表头
    dx, cx, ex = 0.7, 3.2, 8.0
    dw, cw, ew = 2.4, 4.7, 4.63
    yh = 2.28; rh = 0.475; gap = 0.052
    band(slide, cx, yh, cw, rh, fill=TEAL, edge=DTEAL)
    label(slide, cx, yh+0.11, cw, "级联管线 Cascaded", WHITE, 12, PP_ALIGN.CENTER, False, True)
    band(slide, ex, yh, ew, rh, fill=PURP, edge=PURP)
    label(slide, ex, yh+0.11, ew, "端到端 Speech-to-Speech", WHITE, 12, PP_ALIGN.CENTER, False, True)
    for i, (dim, casc, e2e) in enumerate(rows):
        yy = yh + (i+1)*(rh+gap)
        band(slide, dx, yy, dw, rh, fill=NAVY, edge=NAVY)
        label(slide, dx, yy+0.1, dw, dim, WHITE, 11.5, PP_ALIGN.CENTER, False, True)
        band(slide, cx, yy, cw, rh, fill=CARD, edge=LINE)
        tf = tb(slide, cx+0.2, yy+0.09, cw-0.35, 0.4); sr(tf.paragraphs[0].add_run(), casc, BODY_F, 10.5, INK)
        band(slide, ex, yy, ew, rh, fill=PURPBG, edge=PURP)
        tf2 = tb(slide, ex+0.2, yy+0.09, ew-0.35, 0.4); sr(tf2.paragraphs[0].add_run(), e2e, BODY_F, 10.5, INK)
    band(slide, 0.7, 6.12, 11.93, 0.73, fill=GBG, edge=GREEN)
    tf = tb(slide, 1.0, 6.22, 11.4, 0.6); p = tf.paragraphs[0]
    sr(p.add_run(), "售前答法：", BODY_F, 11.5, GINK, bold=True)
    sr(p.add_run(), "理由要说得出口——可审计、可换供应商、成本可预测 → 级联；极致低延迟、消费级体验 → 端到端。", BODY_F, 11.5, INK)
    footer(slide, MOD, "语音与实时交互")

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    figs = [draw_mm_clip, draw_mm_native_vs_pipeline, draw_mm_diffusion_vs_ar, draw_mm_tradeoff_triangle,
            draw_mm_fusion_routes, draw_mm_latency_budget, draw_mm_cascade_vs_e2e]
    for fn in figs:
        fn(prs.slides.add_slide(B))
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_mm_preview.pptx")
    prs.save(out); print(f"saved {len(figs)} ->", out)
