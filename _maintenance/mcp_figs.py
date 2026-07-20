"""MCP 讲义信息图（第一批 4 张）。复用 kb_draw + agent_figs 的 helper。"""
import sys
sys.path.insert(0, "/Users/lijiaxiang/project/myAILearning/_maintenance")
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, BODY_F, TITLE_F)
from agent_figs import line

MOD = "MCP 讲义 · 模型上下文协议"

# ---------- p16 三角色架构 ----------
def draw_mcp_roles(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 协议解剖  ROLES")
    title(slide, "三个角色：Host 管权限、Client 管连接、Server 管能力")
    label(slide, 0.55, 1.58, 12.2, "一个 Host 里可挂多个 Client，每个 Client 一对一连一个 Server", SUB, 12.5, PP_ALIGN.LEFT, True)
    band(slide, 0.9, 2.5, 5.6, 3.35, fill=CARD, edge=TEAL)
    label(slide, 1.15, 2.62, 5.2, "Host（主机）= 用户面对的 AI 应用", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    label(slide, 1.15, 2.98, 5.2, "管权限与用户交互 · 「老板」", SUB, 10.5, PP_ALIGN.LEFT, True)
    cy = [3.55, 4.7]
    srv = [("filesystem", "文件系统"), ("GitHub", "代码仓库")]
    for k, y in enumerate(cy):
        node(slide, 1.35, y, 2.0, 0.95, f"Client {k+1}", "连接器", "一对一维护会话", zh_size=12.5, en_size=9.5, desc_size=9)
        node(slide, 9.4, y, 2.9, 0.95, f"Server {k+1}", srv[k][0], srv[k][1] + " · 提供能力", fill=CORAL, edge=ORANGE, zh_c=ODARK, en_c=ODARK, zh_size=12.5, en_size=9.5, desc_size=9)
        arrow(slide, 3.4, y+0.47, 9.35, y+0.47, color=TEAL, w=1.75)
    label(slide, 4.6, 3.2, 4.6, "各起一个 client、专线连一个 server，互不干扰", SUB, 10, PP_ALIGN.CENTER, True)
    band(slide, 0.7, 6.05, 11.93, 0.85, fill=NEARW, edge=NEARW)
    tf = tb(slide, 0.75, 6.05, 12.0, 0.8)
    para(tf, "像公司里：Host 是「老板」(决定做不做、担责任)，Client 是「对接专员」(每个供应商配一个、专线沟通)，Server 是「供应商」(提供具体能力)。", BODY_F, 11.5, DTEAL, bold=True, first=True)
    footer(slide, MOD, "第 2 章 · 协议解剖")

# ---------- p18 连接生命周期时序 ----------
def draw_mcp_lifecycle(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 协议解剖  LIFECYCLE")
    title(slide, "一次连接怎么开始：先握手、再发现、后调用")
    label(slide, 0.55, 1.58, 12.2, "像新员工入职：先对表(你会什么、我要什么)，再领「能干哪些活」的清单，然后才派活干活", SUB, 12, PP_ALIGN.LEFT, True)
    ax, bx = 3.0, 10.2; ytop = 2.55; ybot = 5.3
    node(slide, ax-1.35, ytop, 2.7, 0.72, "Client", "客户端", "", zh_size=14, en_size=10, desc_size=9)
    node(slide, bx-1.35, ytop, 2.7, 0.72, "Server", "服务端", "", fill=CORAL, edge=ORANGE, zh_c=ODARK, en_c=ODARK, zh_size=14, en_size=10, desc_size=9)
    line(slide, ax, ytop+0.72, ax, ybot, color=LINE, w=1.5, dashed=True)
    line(slide, bx, ytop+0.72, bx, ybot, color=LINE, w=1.5, dashed=True)
    steps = [(ytop+1.05, "① initialize 握手", "双方交换协议版本与能力清单"),
             (ytop+1.85, "② 能力发现 tools/list", "拉取 tools / resources / prompts 清单"),
             (ytop+2.6, "③ tools/call 正常调用", "按需调用，能力协商避免瞎调")]
    for y, t, d in steps:
        arrow(slide, ax+0.1, y, bx-0.1, y, color=TEAL, w=1.75)
        label(slide, ax+0.2, y-0.34, bx-ax-0.4, t, DTEAL, 12, PP_ALIGN.CENTER, False, True)
        label(slide, ax+0.2, y+0.05, bx-ax-0.4, d, SUB, 10, PP_ALIGN.CENTER, True)
    band(slide, 0.7, 5.7, 11.93, 1.15, fill=CARD, edge=LINE)
    tf = tb(slide, 1.0, 5.83, 11.4, 1.0)
    para(tf, "先握手对齐能力、再发现、后调用——能力协商让 client 只去拉 server 支持的那几类。", BODY_F, 12, DTEAL, bold=True, first=True)
    para(tf, "注：2026-07-28 新版把这套握手改成「无状态」(下一张细讲)，但理解经典三步仍是基础。", BODY_F, 10.5, ODARK, italic=True)
    footer(slide, MOD, "第 2 章 · 协议解剖")

# ---------- p31 有状态 → 无状态 ----------
def draw_mcp_stateless(slide):
    bg(slide)
    eyebrow(slide, "第 3 章 · 传输与演进  STATELESS")
    title(slide, "无状态核心：2026-07-28 最大的一次改动")
    label(slide, 0.55, 1.58, 12.2, "把协议核心改成「无状态」：删掉 initialize 握手和会话 ID，改由每个请求自带全部上下文", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.9, 2.35, 5.4, 3.0, fill=CARD, edge=SUB)
    label(slide, 1.15, 2.5, 5.0, "旧 · 有状态", SUB, 13, PP_ALIGN.LEFT, False, True)
    tf = tb(slide, 1.2, 2.98, 4.9, 2.25)
    para(tf, "· initialize 握手 + 会话 ID", BODY_F, 11.5, INK, first=True)
    para(tf, "· 先拨专线、整通话认一个座席", BODY_F, 11.5, INK)
    para(tf, "· 要 sticky session / 共享会话存储", BODY_F, 11.5, INK)
    para(tf, "· 「会话归哪台机器」得一直记着", BODY_F, 11, SUB)
    band(slide, 7.0, 2.35, 5.4, 3.0, fill=CORAL, edge=ORANGE)
    label(slide, 7.25, 2.5, 5.0, "新 · 无状态（2026-07-28）", ODARK, 13, PP_ALIGN.LEFT, False, True)
    tf2 = tb(slide, 7.3, 2.98, 4.9, 2.25)
    para(tf2, "· 每个请求自带版本 + 客户端信息", BODY_F, 11.5, INK, first=True)
    para(tf2, "· 每封信写全地址，谁收都能处理", BODY_F, 11.5, INK)
    para(tf2, "· 前面挂个普通轮询负载均衡就行", BODY_F, 11.5, INK)
    para(tf2, "· 任意请求打任意实例，扩容大简化", BODY_F, 11, ODARK)
    arrow(slide, 6.35, 3.85, 6.98, 3.85, color=TEAL, w=2.5)
    band(slide, 0.7, 5.6, 11.93, 1.25, fill=CARD, edge=LINE)
    tf3 = tb(slide, 1.0, 5.74, 11.4, 1.1)
    para(tf3, "给客户一句话：新版让 MCP server 像普通无状态 Web 服务一样好扩容——这是「向企业规模化倾斜」最实的一步。", BODY_F, 12, DTEAL, bold=True, first=True)
    para(tf3, "（时效性重点：2026-07-28 生效，保鲜巡检优先复查此项。）", BODY_F, 10.5, SUB, italic=True)
    footer(slide, MOD, "第 3 章 · 传输与演进")

# ---------- p42 REST → MCP 转换 ----------
def draw_mcp_wrap(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · 动手写 server  WRAP")
    title(slide, "把现有 REST API 包成 MCP：最常见的落地姿势")
    label(slide, 0.55, 1.58, 12.2, "像给老电器配一个 USB-C 转接头：后端不用改，套一层标准口就能插进所有 MCP 客户端", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 0.7, 2.95, 3.0, 1.55, "现有 REST API", "已有后端", "不用改、不重写业务", fill=CARD, edge=SUB, zh_c=INK, sub_c=SUB, en_c=SUB, zh_size=14.5, en_size=10.5, desc_size=10)
    node(slide, 5.1, 2.95, 3.1, 1.55, "MCP 包装层", "@mcp.tool()", "描述 + 转调 + 整理返回 + 鉴权", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=WHITE, zh_size=14.5, en_size=10.5, desc_size=10)
    node(slide, 9.6, 2.95, 3.0, 1.55, "MCP 客户端", "处处可接", "Claude / IDE / 自研平台", fill=CORAL, edge=ORANGE, zh_c=ODARK, sub_c=SUB, en_c=ODARK, zh_size=14.5, en_size=10.5, desc_size=10)
    arrow(slide, 3.72, 3.72, 5.08, 3.72, w=2.0)
    arrow(slide, 8.22, 3.72, 9.58, 3.72, color=ORANGE, w=2.0)
    label(slide, 3.7, 4.55, 1.5, "套一层", TEAL, 10, PP_ALIGN.CENTER, True)
    band(slide, 0.7, 5.25, 11.93, 1.6, fill=CARD, edge=LINE)
    tf = tb(slide, 1.0, 5.4, 11.4, 1.45)
    para(tf, "案例：已有「查订单」REST 接口 → 写一个 query_order 工具，内部 requests.get 调它、把 JSON 挑关键字段返回。中等复杂度 server 通常几天上手。", BODY_F, 11.5, DTEAL, bold=True, first=True)
    para(tf, "给客户的工作量口径：有现成 API 时，MCP 化主要是「描述 + 转调 + 整理 + 鉴权」，不是重写业务——门槛比想象低。", BODY_F, 11, SUB)
    footer(slide, MOD, "第 4 章 · 动手写 server")

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in [draw_mcp_roles, draw_mcp_lifecycle, draw_mcp_stateless, draw_mcp_wrap]:
        fn(prs.slides.add_slide(B))
    out = "/Users/lijiaxiang/project/myAILearning/_maintenance/_mcp_preview.pptx"
    prs.save(out); print("saved 4 ->", out)


# ========== 第二轮补图 ==========
from kb_draw import GBG
from pptx.dml.color import RGBColor as _RGB2
GINK = _RGB2(0x27, 0x50, 0x0A)

def draw_mcp_what_usb_c(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 是什么 / 为什么  MCP = USB-C")
    title(slide, "MCP 是什么：给大模型装一个「USB-C」")
    label(slide, 0.55, 1.58, 12.2, "一套开放标准——规定 AI 应用与外部工具「怎么发现能力、怎么调用、怎么回传结果」，用同一种语言对话", SUB, 12, PP_ALIGN.LEFT, True)
    # 没标准 vs USB-C
    band(slide, 0.6, 2.4, 5.5, 3.0, fill=CORAL, edge=ORANGE)
    label(slide, 0.85, 2.5, 5.1, "USB-C 之前：每个设备一根专用线", ODARK, 12, PP_ALIGN.LEFT, False, True)
    for i, t in enumerate(["充电线", "数据线", "视频线"]):
        node(slide, 0.95, 3.0+i*0.6, 1.9, 0.48, t, "", "", fill=WHITE, edge=ODARK, zh_c=ODARK, zh_size=11)
        node(slide, 3.6, 3.0+i*0.6, 1.9, 0.48, "各不通用", "", "", fill=WHITE, edge=ODARK, zh_c=SUB, zh_size=11)
    label(slide, 0.9, 4.88, 5.3, "AI 世界的老问题：每个工具为每个应用各写一遍对接", SUB, 10.5, PP_ALIGN.LEFT, True)
    band(slide, 6.75, 2.4, 6.0, 3.0, fill=GBG, edge=GREEN)
    label(slide, 7.0, 2.5, 5.6, "MCP：一个口全搞定", GINK, 12, PP_ALIGN.LEFT, False, True)
    node(slide, 9.0, 3.35, 2.3, 0.9, "MCP 标准口", "USB-C", "工具做一次", fill=GREEN, edge=GINK, zh_c=WHITE, sub_c=GBG, en_c=GBG, zh_size=13, en_size=9.5, desc_size=9)
    for i, t in enumerate(["Claude", "Cursor", "各类 IDE"]):
        node(slide, 6.95, 3.0+i*0.62, 1.7, 0.5, t, "", "", fill=WHITE, edge=GREEN, zh_c=GINK, zh_size=10.5)
        line(slide, 8.65, 3.25+i*0.62, 8.98, 3.8, color=GINK, w=1.1)
    label(slide, 11.4, 3.75, 1.4, "任何支持 MCP 的应用插上就用", GINK, 9, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 5.55, 12.03, 1.27, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.67, 11.5, 1.1); p = tf.paragraphs[0]
    sr(p.add_run(), "案例：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "同一个 GitHub MCP server 做好后，Claude、Cursor、各类 IDE 都能直接接入读 issue、提 PR——工具方不必为每个 AI 应用各写一遍。", BODY_F, 11.5, INK)
    para(tf, "一句话给客户：MCP 不是又一个大模型能力，而是一层「接线标准」——让工具和 AI 应用解耦，各做各的、随处可插。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "第 1 章 · 是什么 / 为什么")

def draw_mcp_primitives(slide):
    bg(slide)
    eyebrow(slide, "第 2 章 · 三大原语  TOOLS · RESOURCES · PROMPTS")
    title(slide, "三大原语：能做什么、能看什么、怎么用得好")
    label(slide, 0.55, 1.58, 12.2, "三者由「谁触发」不同——这是最常被客户问混的点", SUB, 12, PP_ALIGN.LEFT, True)
    cards = [("Tools 工具", "能做什么", "模型可调用的「动作」，会执行、可能改状态", "模型自主决定调用", "查数据库 · 发邮件 · 创建 issue", "动词 · 有副作用", TEAL, DTEAL, CARD),
             ("Resources 资源", "能看什么", "可读取的「数据 / 上下文」，只读、无副作用", "应用 / 用户挂载引用", "一个文件 · 一段记录 · 一份文档", "名词 · 只读", ORANGE, ODARK, CORAL),
             ("Prompts 提示词", "怎么用得好", "预置的「对话模板」，把最佳用法固化下来", "用户显式选用", "「代码审查」「总结会议纪要」模板", "模板 · 固化最佳实践", DTEAL, NAVY, PALE)]
    W = 3.95; gap = 0.14; x0 = 0.6; y = 2.4; h = 3.05
    for i, (zh, one, what, who, ex, tag, tc, ec, fill) in enumerate(cards):
        x = x0 + i*(W+gap)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, y+0.16, W-0.5, 0.72); p = tf.paragraphs[0]
        sr(p.add_run(), f"{i+1}. {zh}", BODY_F, 14, tc, bold=True)
        node(slide, x+0.32, y+0.85, W-0.64, 0.5, one, "", "", fill=WHITE, edge=ec, zh_c=tc, zh_size=12)
        b = tb(slide, x+0.28, y+1.48, W-0.56, 1.5)
        para(b, what, BODY_F, 11, INK, first=True)
        para(b, "谁触发：" + who, BODY_F, 11, tc, bold=True)
        para(b, "例：" + ex, BODY_F, 10.5, SUB)
        para(b, "记：" + tag, BODY_F, 10.5, SUB, italic=True)
    band(slide, 0.6, 5.62, 12.03, 1.2, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.74, 11.5, 1.0); p = tf.paragraphs[0]
    sr(p.add_run(), "一句话记法：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "Tools 是「能做什么」（动词、有副作用）· Resources 是「能看什么」（名词、只读）· Prompts 是「怎么用得好」（模板）。", BODY_F, 12, INK)
    footer(slide, MOD, "第 2 章 · 协议解剖")
