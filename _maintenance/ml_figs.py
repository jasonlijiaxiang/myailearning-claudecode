"""Model-Landscape 讲义信息图（5 张）。复用 kb_draw + agent_figs helper。
嫁接：p5 后(三大阵营)、p6 后(两市场)、p43 后(open weight vs source)、p51 后(价格光谱)、p70 后(三层路由)。
ML 无页码，图系统页脚匹配。独立预览：python3 ml_figs.py"""
import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from kb_draw import (bg, tb, sr, para, eyebrow, title, footer, node, arrow, band, label,
                     NAVY, INK, SUB, TEAL, DTEAL, ORANGE, ODARK, OINK, CARD, CORAL, LINE,
                     WHITE, NEARW, PALE, GREEN, GBG, BODY_F, TITLE_F)
from agent_figs import line, dot

MOD = "模型格局讲义 · Model-Landscape"
GINK = RGBColor(0x27, 0x50, 0x0A)
PURP = RGBColor(0x6B, 0x46, 0x9C)
PURPBG = RGBColor(0xF1, 0xEC, 0xF7)

# ---------- p5 三大阵营 ----------
def draw_ml_three_camps(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 三大阵营  THE MAP")
    title(slide, "一张地图：三大阵营各自在卖什么")
    label(slide, 0.55, 1.58, 12.2, "不是竞争关系，是分工关系——一个企业方案常三者并用", SUB, 12, PP_ALIGN.LEFT, True)
    camps = [("闭源西方旗舰", "卖「最强能力」· API 订阅", TEAL, DTEAL, CARD,
              ["OpenAI / Anthropic / Google / xAI（+Meta 转入）", "前沿基准仍由它们把持", "打最难的仗 → 第 2 章逐家名片"]),
             ("开放权重（中国主导）", "卖「可控与便宜」· 权重可下载", ORANGE, ODARK, CORAL,
              ["DeepSeek / Qwen / Kimi / GLM 占榜首", "西方线：Llama 收尾、Mistral 转宽松", "管可控与走量 → 第 3–4 章"]),
             ("端侧小模型", "下沉到端 · 蒸馏是功臣", PURP, PURP, PURPBG,
              ["Phi-4 / Gemma 4 / SmolLM-3", "sub-10B 已常规超越 2024 版 GPT-4", "4GB 内存可跑（Gemma 4 E2B）"])]
    W = 3.95; gap = 0.14; x0 = 0.6; y = 2.4; h = 3.0
    for i, (zh, sell, ec, tc, fill, items) in enumerate(camps):
        x = x0 + i*(W+gap)
        band(slide, x, y, W, h, fill=fill, edge=ec)
        tf = tb(slide, x+0.25, y+0.16, W-0.5, 0.75); p = tf.paragraphs[0]
        sr(p.add_run(), f"{i+1}. {zh}", BODY_F, 14, tc, bold=True)
        node(slide, x+0.3, y+0.86, W-0.6, 0.5, sell, "", "", fill=WHITE, edge=ec, zh_c=tc, zh_size=10.5)
        ib = tb(slide, x+0.28, y+1.52, W-0.56, 1.4)
        for j, it in enumerate(items):
            para(ib, "· " + it, BODY_F, 11, SUB, first=(j == 0))
    band(slide, 0.6, 5.62, 12.03, 1.2, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.74, 11.5, 1.05); p = tf.paragraphs[0]
    sr(p.add_run(), "要点：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "旗舰打最难的仗、开放权重管可控与走量、小模型下沉到端侧——三者分工，一个企业方案常常三者并用。", BODY_F, 12, INK)
    footer(slide, MOD, "全景地图")

# ---------- p6 两个市场 ----------
def draw_ml_two_markets(slide):
    bg(slide)
    eyebrow(slide, "第 1 章 · 两个市场  ENTERPRISE vs CONSUMER")
    title(slide, "企业与消费：同一批模型，两张榜单")
    label(slide, 0.55, 1.58, 12.2, "对客引用任何份额数字前，先说清是哪个市场的口径——两张榜单已经分裂", SUB, 12, PP_ALIGN.LEFT, True)
    # 左 企业
    band(slide, 0.6, 2.35, 7.15, 3.05, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.47, 6.6, "企业市场（Menlo）—— 看能力、信任、合同", DTEAL, 12.5, PP_ALIGN.LEFT, False, True)
    ent = [("Anthropic", 0.40, "40%", "企业 LLM 支出第一", TEAL), ("OpenAI", 0.27, "27%", "从 2023 年 50% 一路下滑", SUB), ("Google", 0.21, "21%", "三年三倍（7%→21%）", ORANGE)]
    yb = 3.05
    for i, (name, frac, val, sub, c) in enumerate(ent):
        yy = yb + i*0.72
        tf = tb(slide, 0.85, yy, 1.7, 0.6); sr(tf.paragraphs[0].add_run(), name, BODY_F, 11.5, INK, bold=True)
        band(slide, 2.55, yy+0.05, 3.7*frac, 0.42, fill=c, edge=c)
        vt = tb(slide, 2.6+3.7*frac, yy-0.02, 3.6, 0.6); q = vt.paragraphs[0]
        sr(q.add_run(), val + "  ", BODY_F, 13, INK, bold=True); sr(q.add_run(), sub, BODY_F, 9.5, SUB)
    # 右 消费
    band(slide, 7.95, 2.35, 4.68, 3.05, fill=CORAL, edge=ORANGE)
    label(slide, 8.2, 2.47, 4.2, "消费市场 —— 看品牌、生态", ODARK, 12.5, PP_ALIGN.LEFT, False, True)
    node(slide, 8.5, 3.2, 3.6, 1.35, "ChatGPT 74%", "", "日均 25 亿+ prompts", fill=ORANGE, edge=ODARK, zh_c=OINK, sub_c=OINK, zh_size=22, desc_size=10.5)
    label(slide, 8.2, 4.75, 4.2, "消费碾压，与企业榜完全两样", ODARK, 10.5, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 5.6, 12.03, 1.23, fill=GBG, edge=GREEN)
    tf = tb(slide, 0.9, 5.71, 11.5, 1.08); p = tf.paragraphs[0]
    sr(p.add_run(), "数据锚点：", BODY_F, 11.5, GINK, bold=True)
    sr(p.add_run(), "企业 LLM 支出 2026 已达 $8.4B、年底看 $15B——是生产预算不是试验预算；Claude Code 单产品年化 $1B 是标志性事件。", BODY_F, 11.5, INK)
    footer(slide, MOD, "全景地图")

# ---------- p43 open weight vs open source ----------
def draw_ml_weight_vs_source(slide):
    bg(slide)
    eyebrow(slide, "第 5 章 · 两个概念  WEIGHT vs SOURCE")
    title(slide, "open source ≠ open weight：先把词用对")
    label(slide, 0.55, 1.58, 12.2, "对客用词严谨是专业度——法务在场时，这个区别值钱", SUB, 12, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 2.35, 6.0, 3.0, fill=CARD, edge=TEAL)
    label(slide, 0.85, 2.47, 5.6, "open weight · 开放权重", DTEAL, 13, PP_ALIGN.LEFT, False, True)
    node(slide, 0.95, 3.0, 5.3, 0.56, "送你一桌做好的菜", "能吃、能加热、能分给家人", "", fill=TEAL, edge=DTEAL, zh_c=WHITE, en_c=PALE, zh_size=12, en_size=9.5)
    tf = tb(slide, 0.9, 3.72, 5.5, 1.4)
    para(tf, "· 给训练好的权重文件，可下载可部署", BODY_F, 11, INK, bold=True, first=True)
    para(tf, "· 不给训练数据、训练代码、完整配方", BODY_F, 11, SUB)
    para(tf, "· Llama / Gemma / Qwen 都属此类", BODY_F, 11, SUB)
    band(slide, 6.75, 2.35, 6.0, 3.0, fill=PURPBG, edge=PURP)
    label(slide, 7.0, 2.47, 5.6, "open source · 开源（OSAID 定义）", PURP, 13, PP_ALIGN.LEFT, False, True)
    node(slide, 7.1, 3.0, 5.3, 0.56, "连菜谱和食材清单一起给", "你能复现这桌菜", "", fill=PURP, edge=PURP, zh_c=WHITE, en_c=PALE, zh_size=12, en_size=9.5)
    tf2 = tb(slide, 7.05, 3.72, 5.5, 1.4)
    para(tf2, "· OSI 2024-10 发布 OSAID v1.0", BODY_F, 11, INK, bold=True, first=True)
    para(tf2, "· 要求开放训练数据信息 + 架构 + 训练代码", BODY_F, 11, SUB)
    para(tf2, "· 按此标准，绝大多数「开源模型」不合格", BODY_F, 11, PURP, bold=True)
    band(slide, 0.6, 5.5, 12.03, 1.33, fill=PALE, edge=LINE)
    tf3 = tb(slide, 0.9, 5.62, 11.5, 1.15); p = tf3.paragraphs[0]
    sr(p.add_run(), "要点：", BODY_F, 12, DTEAL, bold=True)
    sr(p.add_run(), "市面上「送菜」的多，「给菜谱」的极少——说「开放权重模型」而不是笼统的「开源」，专业度就出来了。", BODY_F, 12, INK)
    footer(slide, MOD, "许可证")

# ---------- p51 价格光谱 ----------
def draw_ml_price_spectrum(slide):
    bg(slide)
    eyebrow(slide, "第 6 章 · 价格光谱  PRICE TIERS")
    title(slide, "价格光谱：四档接住全预算段")
    label(slide, 0.55, 1.58, 12.2, "$/百万 token（输入/输出，2026-07）——报价先定档位，再挑档内型号", SUB, 12, PP_ALIGN.LEFT, True)
    tiers = [("地板档 ~$0.1", 0.16, "DeepSeek V4 Flash $0.14/$0.28 · Gemini Flash-Lite $0.10/$0.40", "海量批处理、分类、抽取", GREEN, GINK, GBG),
             ("走量档 $1–3", 0.42, "Luna $1/$6 · Grok 4.3 $1.25/$2.50 · Sonnet 5 尝鲜 $2/$10 · Terra $2.5/$15", "生产主力：客服、RAG、Agent", TEAL, DTEAL, CARD),
             ("旗舰档 ~$5", 0.72, "GPT-5.6 Sol $5/$30 · Gemini 3.1 Pro（旗舰最低价）", "难任务、深推理", ORANGE, ODARK, CORAL),
             ("超旗舰 $30", 1.0, "GPT-5.4 Pro $30/$180", "极少数场景，先证明再用", RGBColor(0x99,0x3C,0x1D), RGBColor(0x5A,0x24,0x12), RGBColor(0xF3,0xE2,0xD8))]
    x0 = 0.7; W = 11.93; y = 2.4; h = 0.72; gap = 0.12
    for i, (tier, frac, reps, use, bc, tc, fill) in enumerate(tiers):
        yy = y + i*(h+gap)
        band(slide, x0, yy, W, h, fill=fill, edge=tc)
        lt = tb(slide, x0+0.18, yy+0.06, 2.5, h-0.1); p = lt.paragraphs[0]
        sr(p.add_run(), tier, BODY_F, 12.5, tc, bold=True)
        band(slide, x0+2.75, yy+0.13, 1.5*frac+0.25, h-0.26, fill=bc, edge=bc)  # 价格量级条
        rt = tb(slide, x0+4.7, yy+0.04, 7.0, h-0.06)
        para(rt, reps, BODY_F, 10.5, INK, first=True)
        para(rt, "适用：" + use, BODY_F, 10, SUB)
    band(slide, 0.7, 5.9, 11.93, 0.92, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 6.0, 11.4, 0.78); p = tf.paragraphs[0]
    sr(p.add_run(), "人民币参照：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "豆包 2.1 Pro ¥6/¥30（约 $0.85/$4.2）在走量档很有攻击性——先定档位收敛范围，再在档内比型号。", BODY_F, 11.5, INK)
    footer(slide, MOD, "价格带")

# ---------- p70 三层路由 ----------
def draw_ml_three_layer_routing(slide):
    bg(slide)
    eyebrow(slide, "第 8 章 · 三层路由  ROUTING")
    title(slide, "三层路由：任务分流的通用框架")
    label(slide, 0.55, 1.58, 12.2, "比例通常是 1 : 8 : 1——80% 的量走便宜档，是成本工程的核心", SUB, 12, PP_ALIGN.LEFT, True)
    layers = [("最难推理层", "~10%", "前沿闭源：Fable 5 / GPT-5.6 Sol / Gemini 3.1 Pro", "复杂规划、深度分析、关键决策", ORANGE, ODARK, CORAL, 1.4),
              ("高量生产层", "~80%", "走量档：Luna / Flash / Sonnet 5 / 豆包 Turbo / DeepSeek", "客服、RAG 问答、批量抽取", TEAL, DTEAL, CARD, 4.9),
              ("敏感数据层", "~10%", "自托管开放权重：Qwen / DeepSeek / GLM（MIT/Apache）", "数据不出域、强监管场景", PURP, PURP, PURPBG, 1.4)]
    x0 = 0.7; y = 2.4; gap = 0.14; totalh = 3.1; unit = totalh/(1.4+4.9+1.4+2*gap*0)
    yy = y
    for name, ratio, use, task, tc, ec, fill, wgt in layers:
        hh = totalh * wgt / (1.4+4.9+1.4) - gap*0.4
        band(slide, x0, yy, 11.93, hh, fill=fill, edge=ec)
        lt = tb(slide, x0+0.2, yy+hh/2-0.3, 2.6, 0.6); p = lt.paragraphs[0]
        sr(p.add_run(), name, BODY_F, 13.5, tc, bold=True)
        para(lt, "占比 " + ratio, BODY_F, 11, tc, bold=True)
        rt = tb(slide, x0+3.1, yy+hh/2-0.32, 8.5, 0.7)
        para(rt, use, BODY_F, 11, INK, bold=True, first=True)
        para(rt, "→ " + task, BODY_F, 10.5, SUB)
        yy += hh + gap
    band(slide, 0.7, 6.1, 11.93, 0.72, fill=PALE, edge=LINE)
    tf = tb(slide, 1.0, 6.19, 11.4, 0.6); p = tf.paragraphs[0]
    sr(p.add_run(), "要点：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "1:8:1 的比例把成本压在便宜档；路由的落地件（分流、计量、灰度）在 AI-Gateway 模块。", BODY_F, 11.5, INK)
    footer(slide, MOD, "选型方法论")

# ---------- p36 豆包家族树 ----------
def draw_ml_doubao_family(slide):
    bg(slide)
    eyebrow(slide, "第 4 章 · 豆包家族  FAMILY TREE")
    title(slide, "豆包家族树：一条主线加两条专线")
    label(slide, 0.55, 1.58, 12.2, "看豆包家族要看「更新节奏」——它赌的是 Coding / Agent 市场唯快不破", SUB, 12, PP_ALIGN.LEFT, True)
    node(slide, 5.35, 2.4, 2.65, 0.8, "豆包家族", "字节 / 火山", "", fill=NAVY, edge=NAVY, zh_c=WHITE, sub_c=PALE, en_c=PALE, zh_size=15, en_size=10)
    # 主线
    node(slide, 0.9, 3.7, 3.3, 1.0, "① Doubao 2.1 Pro", "主线旗舰", "Coding/Agent/VLM「质变点」· 2026-06-23", fill=TEAL, edge=DTEAL, zh_c=WHITE, sub_c=PALE, en_c=PALE, zh_size=13, en_size=10, desc_size=9)
    node(slide, 4.5, 3.7, 3.3, 1.0, "② Doubao 2.1 Turbo", "高频走量档", "价格为 Pro 一半", fill=CARD, edge=TEAL, zh_c=DTEAL, sub_c=SUB, en_c=TEAL, zh_size=13, en_size=10, desc_size=9)
    arrow(slide, 4.22, 4.2, 4.48, 4.2, color=TEAL, w=1.6)
    # 两专线
    node(slide, 8.6, 3.35, 4.0, 0.85, "③ Seed-Evolving", "专线 · 月更 2–4 次", "Coding/Agent 最激进迭代", fill=CORAL, edge=ORANGE, zh_c=ODARK, sub_c=ODARK, en_c=ODARK, zh_size=12.5, en_size=9.5, desc_size=9)
    node(slide, 8.6, 4.35, 4.0, 0.85, "④ Seed-2.0-lite", "专线 · 全模态", "家族首个原生统一（视频/图/音/文）· 2026-05", fill=CORAL, edge=ORANGE, zh_c=ODARK, sub_c=ODARK, en_c=ODARK, zh_size=12.5, en_size=9.5, desc_size=9)
    line(slide, 6.68, 3.2, 2.55, 3.7, color=TEAL, w=1.3)
    line(slide, 6.68, 3.2, 6.15, 3.7, color=TEAL, w=1.3)
    line(slide, 6.68, 3.2, 8.55, 3.77, color=ORANGE, w=1.3)
    line(slide, 6.68, 3.2, 8.55, 4.77, color=ORANGE, w=1.3)
    label(slide, 0.9, 4.85, 6.0, "主线：旗舰 + 走量，一个家族接住主力预算", SUB, 10, PP_ALIGN.LEFT, True)
    label(slide, 8.6, 5.3, 4.0, "专线：Seed-Evolving 月更是行业最激进", ODARK, 10, PP_ALIGN.LEFT, True)
    band(slide, 0.6, 5.65, 12.03, 1.17, fill=PALE, edge=LINE)
    tf = tb(slide, 0.9, 5.77, 11.5, 1.0); p = tf.paragraphs[0]
    sr(p.add_run(), "数据锚点：", BODY_F, 11.5, DTEAL, bold=True)
    sr(p.add_run(), "豆包系累计 tokens 调用量 180 万亿+（火山引擎口径）——国内调用量第一梯队的证据。", BODY_F, 11.5, INK)
    para(tf, "看点是「更新节奏」：Seed-Evolving 月更 2–4 次，赌的是 Coding / Agent 市场唯快不破。", BODY_F, 10.5, SUB)
    footer(slide, MOD, "中国格局")

FIGS = [draw_ml_three_camps, draw_ml_two_markets, draw_ml_weight_vs_source,
        draw_ml_price_spectrum, draw_ml_three_layer_routing, draw_ml_doubao_family]
NEW = [draw_ml_doubao_family]

if __name__ == '__main__':
    from pptx import Presentation
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    B = prs.slide_layouts[6]
    for fn in FIGS:
        fn(prs.slides.add_slide(B))
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "_ml_preview.pptx")
    prs.save(out); print(f"saved {len(FIGS)} ->", out)
