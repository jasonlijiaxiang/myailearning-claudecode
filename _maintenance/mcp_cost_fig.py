# -*- coding: utf-8 -*-
"""MCP 讲义第 5 章增页：运行期成本（缺口级回流）。大画布直接绘制，正文 ≥11pt。"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import kb_draw as K

MOD = "MCP 讲义 · 模型上下文协议"


def draw_cost(s):
    K.bg(s)
    K.eyebrow(s, "第 5 章 · 生产落地   RUNTIME COST")
    K.title(s, "运行期成本：贵在上下文，不在调用")

    # 引子
    tf = K.tb(s, 0.55, 1.52, 12.23, 0.75)
    K.para(tf, "客户问「这东西贵不贵」，答案常不在他以为的地方。MCP 的运行期开销主要不是每次"
               "调用的网络往返，而是工具定义常驻上下文——而且它随每一条消息进上下文，不是每会话"
               "一次，所以成本是按轮次乘的。", size=12.5, color=K.INK, first=True)

    # 三张量级卡
    cards = [
        ("17.6k", "tokens",
         "GitHub 官方 MCP server 暴露 94 个工具，仅工具定义就占约 17,600 tokens。",
         "边界：Atlassian 工程博客 2026-03-29，未公布模型与计数方法。这是单个大 server 的量级，"
         "不能外推成「每个 server 都这么贵」——工具少、描述简的只有几百到几千。"),
        ("~3.9k → ~0.5k", "压缩后",
         "同样 94 个工具，描述压缩后分别降到约 3,900（低压缩）与约 500（最激进）。",
         "边界：同源同场景。压缩是牺牲描述细节换上下文，而描述质量直接影响模型会不会误调"
         "（见第 4 章），不是免费。"),
        ("1–2ms", "/ 会话",
         "MCP 流量经网关一跳的额外延迟（加密迭代调到约 100 次之后）。",
         "边界：Envoy AI Gateway 官方基准 2025-12-08，M1 八核、echo 工具、按会话计。默认 10 万次"
         "KDF 迭代下为数十毫秒。单机基准，不能当生产 SLA。"),
    ]
    x0, gap, w = 0.55, 0.3, (12.23 - 0.6) / 3
    y, h = 2.4, 2.62
    for i, (big, unit, desc, edge) in enumerate(cards):
        x = x0 + i * (w + gap)
        K.band(s, x, y, w, h, fill=K.CARD, edge=K.LINE)
        tf = K.tb(s, x + 0.18, y + 0.16, w - 0.36, 0.62)
        p = tf.paragraphs[0]
        K.sr(p.add_run(), big, font=K.TITLE_F, size=25, color=K.TEAL, bold=True)
        K.sr(p.add_run(), "  " + unit, font=K.BODY_F, size=12, color=K.SUB)
        tf2 = K.tb(s, x + 0.18, y + 0.86, w - 0.36, 0.9)
        K.para(tf2, desc, size=11.5, color=K.INK, first=True)
        tf3 = K.tb(s, x + 0.18, y + 1.78, w - 0.36, 0.78)
        K.para(tf3, edge, size=11, color=K.SUB, first=True)

    # 底部关键结论
    yb = 5.28
    K.band(s, 0.55, yb, 12.23, 1.28, fill=K.PALE, edge=K.TEAL)
    tf = K.tb(s, 0.78, yb + 0.14, 11.77, 1.0)
    K.para(tf, "这三个数改变的判断：", size=12, color=K.DTEAL, bold=True, first=True)
    K.para(tf, "① 别把工具堆进一个 server——工具数直接乘进每轮上下文，90+ 工具光定义就吃掉 200k "
               "窗口近 9%，接三个就可能过半。② 网关那一跳通常不是瓶颈，客户拿「多一层代理会不会慢」"
               "反对时，真正该谈的是加密配置。③ 选型先问「暴露多少工具」——比「支持哪些功能」更能"
               "预测实际成本。", size=11.5, color=K.INK)

    K.footer(s, MOD, "第 5 章 · 生产落地")


if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    draw_cost(s)
    out = sys.argv[1] if len(sys.argv) > 1 else "/tmp/mcp_cost_preview.pptx"
    prs.save(out)
    print("saved", out)
